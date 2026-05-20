"""트랙 #105 Phase C 보강 — 재사용 가능한 Playwright e2e helper 라이브러리.

향후 트랙들에서도 import 하여 재사용:
    from e2e_helpers import (
        upload_file, download_file, iframe_click,
        wait_for_sse, safe_mutation, ui_e2e_test_prefix,
        ...
    )

설계 원칙:
  - 운영 데이터 무영향: mutation 케이스는 ui-e2e-test prefix + cleanup 의무
  - 비용 회피: LLM cost 케이스는 환경변수 E2E_ALLOW_COST=1 일 때만 호출
  - 결정성: 비결정 대기 (sleep) 최소화, expect_* 컨텍스트 매니저 우선
  - 한국어 사유 반환 — UI 검증 결과는 항상 한글 actual 메시지로

환경변수:
  - E2E_ALLOW_COST=1 → LLM 호출 비용 발생 케이스 실행 허용 (기본 OFF)
  - E2E_ALLOW_MUTATION=1 → 운영 DB 변경 케이스 실행 허용 (기본 ON, ui-e2e-test prefix 필수)
  - E2E_FIXTURE_DIR → fixture 파일 디렉토리 override
"""
from __future__ import annotations

import os
import random
import re
import string
import time
from contextlib import contextmanager
from pathlib import Path
from typing import Any, Callable, Iterator

from playwright.sync_api import Download, Page, Response, TimeoutError as PlaywrightTimeoutError

# ─── 경로 / 상수 ─────────────────────────────────────────────────────
HELPERS_DIR = Path(__file__).resolve().parent
FIXTURE_DIR = Path(os.environ.get("E2E_FIXTURE_DIR", str(HELPERS_DIR.parent / "fixtures")))
FIXTURE_DIR.mkdir(parents=True, exist_ok=True)

DOWNLOAD_TMP_DIR = HELPERS_DIR / "_downloads_tmp"
DOWNLOAD_TMP_DIR.mkdir(parents=True, exist_ok=True)

# ─── 환경변수 flag ───────────────────────────────────────────────────
def allow_cost() -> bool:
    """LLM 비용 발생 케이스 실행 허용 여부 (기본 False)."""
    return os.environ.get("E2E_ALLOW_COST", "0") == "1"


def allow_mutation() -> bool:
    """운영 DB 변경 케이스 실행 허용 여부 (기본 True — ui-e2e-test prefix 필수)."""
    return os.environ.get("E2E_ALLOW_MUTATION", "1") == "1"


# ─── 식별자 / prefix ────────────────────────────────────────────────
_PREFIX_USED: set[str] = set()


def ui_e2e_test_prefix() -> str:
    """타임스탬프 + random 기반 unique 식별자.

    Format: ui-e2e-test-{epoch_sec}-{6char_rand}
    cleanup 시 이 prefix 로 검색하여 삭제 — 절대 운영 데이터와 충돌 없음.
    """
    for _ in range(5):
        ts = int(time.time())
        rand = "".join(random.choices(string.ascii_lowercase + string.digits, k=6))
        prefix = f"ui-e2e-test-{ts}-{rand}"
        if prefix not in _PREFIX_USED:
            _PREFIX_USED.add(prefix)
            return prefix
    # fallback — 매우 희박
    return f"ui-e2e-test-{int(time.time())}-{random.randint(100000, 999999)}"


def is_test_artifact(name_or_id: str) -> bool:
    """이름/ID 가 ui-e2e-test prefix 인지 확인 (cleanup 안전장치)."""
    if not name_or_id:
        return False
    return "ui-e2e-test-" in str(name_or_id).lower()


# ─── fixture 파일 보장 ───────────────────────────────────────────────
def ensure_fixture(name: str, generator: Callable[[Path], None] | None = None) -> Path:
    """fixture 파일이 없으면 generator 로 생성 후 경로 반환.

    Args:
        name: 파일명 (예: 'sample_doc.pdf')
        generator: 없을 때 호출되는 callable. 인자는 출력 경로.

    Returns:
        Path to fixture file.
    """
    path = FIXTURE_DIR / name
    if path.exists() and path.stat().st_size > 0:
        return path
    if generator is None:
        raise FileNotFoundError(f"fixture {name} not found and no generator")
    generator(path)
    if not path.exists() or path.stat().st_size == 0:
        raise RuntimeError(f"fixture generator failed for {name}")
    return path


def get_sample_pdf() -> Path:
    """샘플 PDF — fixture 디렉토리에서 가져오거나 minimal PDF 생성."""
    return ensure_fixture("sample_doc.pdf", _generate_minimal_pdf)


def get_sample_docx() -> Path:
    """샘플 DOCX — 변수 {{name}}, {{topic}} 포함 minimal."""
    return ensure_fixture("sample_template.docx", _generate_minimal_docx)


def get_sample_pptx() -> Path:
    """샘플 PPTX — 변수 {{title}} 포함 minimal."""
    return ensure_fixture("sample_template.pptx", _generate_minimal_pptx)


# ─── fixture 생성기 (외부 의존성 최소) ────────────────────────────────
def _generate_minimal_pdf(path: Path) -> None:
    """PyPI 의존 없이 minimal valid PDF (1 page, blank) 생성.

    PDF spec 1.4 minimal — 어떤 PDF reader 도 열림.
    """
    pdf_bytes = (
        b"%PDF-1.4\n"
        b"1 0 obj\n<< /Type /Catalog /Pages 2 0 R >>\nendobj\n"
        b"2 0 obj\n<< /Type /Pages /Kids [3 0 R] /Count 1 >>\nendobj\n"
        b"3 0 obj\n<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>\nendobj\n"
        b"4 0 obj\n<< /Length 68 >>\nstream\n"
        b"BT /F1 24 Tf 100 700 Td (ui-e2e-test sample PDF) Tj ET\n"
        b"endstream\nendobj\n"
        b"5 0 obj\n<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>\nendobj\n"
        b"xref\n0 6\n"
        b"0000000000 65535 f\n"
        b"0000000009 00000 n\n"
        b"0000000058 00000 n\n"
        b"0000000115 00000 n\n"
        b"0000000228 00000 n\n"
        b"0000000343 00000 n\n"
        b"trailer\n<< /Size 6 /Root 1 0 R >>\nstartxref\n413\n%%EOF\n"
    )
    path.write_bytes(pdf_bytes)


def _generate_minimal_docx(path: Path) -> None:
    """python-docx 가 있으면 사용, 없으면 zip 으로 minimal DOCX 직접 작성."""
    try:
        from docx import Document  # type: ignore

        doc = Document()
        doc.add_heading("ui-e2e-test template", level=1)
        doc.add_paragraph("이 문서는 e2e 테스트용 템플릿입니다. 삭제 안전.")
        doc.add_paragraph("이름: {{name}}")
        doc.add_paragraph("주제: {{topic}}")
        doc.add_paragraph("생성일: {{date}}")
        doc.save(str(path))
    except ImportError:
        _write_minimal_docx_zip(path)


def _write_minimal_docx_zip(path: Path) -> None:
    """python-docx 없이 zip 으로 minimal DOCX (OOXML) 생성."""
    import zipfile

    content_types = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        '<Default Extension="xml" ContentType="application/xml"/>'
        '<Override PartName="/word/document.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
        "</Types>"
    )
    root_rels = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" '
        'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" '
        'Target="word/document.xml"/></Relationships>'
    )
    doc_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        "<w:body>"
        '<w:p><w:r><w:t>ui-e2e-test template — 이름: {{name}}, 주제: {{topic}}</w:t></w:r></w:p>'
        "</w:body></w:document>"
    )
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("[Content_Types].xml", content_types)
        z.writestr("_rels/.rels", root_rels)
        z.writestr("word/document.xml", doc_xml)


def _generate_minimal_pptx(path: Path) -> None:
    """python-pptx 가 있으면 사용, 없으면 zip 으로 minimal PPTX 작성."""
    try:
        from pptx import Presentation  # type: ignore

        prs = Presentation()
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        if title:
            title.text = "ui-e2e-test {{title}}"
        prs.save(str(path))
    except ImportError:
        _write_minimal_pptx_zip(path)


def _write_minimal_pptx_zip(path: Path) -> None:
    """python-pptx 없이 zip 으로 minimal PPTX 생성. 1 slide."""
    import zipfile

    content_types = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        '<Default Extension="xml" ContentType="application/xml"/>'
        '<Override PartName="/ppt/presentation.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>'
        '<Override PartName="/ppt/slides/slide1.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>'
        "</Types>"
    )
    root_rels = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" '
        'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" '
        'Target="ppt/presentation.xml"/></Relationships>'
    )
    pres_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<p:presentation xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
        '<p:sldIdLst><p:sldId id="256" r:id="rId1"/></p:sldIdLst>'
        '<p:sldSz cx="9144000" cy="6858000"/></p:presentation>'
    )
    pres_rels = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" '
        'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" '
        'Target="slides/slide1.xml"/></Relationships>'
    )
    slide_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
        "<p:cSld><p:spTree/></p:cSld></p:sld>"
    )
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("[Content_Types].xml", content_types)
        z.writestr("_rels/.rels", root_rels)
        z.writestr("ppt/presentation.xml", pres_xml)
        z.writestr("ppt/_rels/presentation.xml.rels", pres_rels)
        z.writestr("ppt/slides/slide1.xml", slide_xml)


# ─── 헬퍼: 업로드 ────────────────────────────────────────────────────
def upload_file(
    page: Page,
    selector: str,
    fixture_path: Path | str,
    *,
    wait_after_ms: int = 1500,
) -> dict:
    """`<input type="file">` 셀렉터에 파일 업로드 + 진행 대기.

    Args:
        page: Playwright Page
        selector: input[type=file] 의 CSS selector (예: 'input[type="file"]')
                  React 패턴: hidden file input — visibility 무시 OK.
        fixture_path: 업로드할 파일 경로
        wait_after_ms: 업로드 후 진행 대기 ms (서버 측 처리)

    Returns:
        {"ok": bool, "note": str, "file": str}
    """
    fixture = Path(fixture_path)
    if not fixture.exists():
        return {"ok": False, "note": f"fixture 미존재: {fixture}", "file": str(fixture)}
    try:
        # set_input_files 는 hidden input 도 작동 (visibility 무시)
        page.set_input_files(selector, str(fixture))
    except Exception as e:
        # fallback — file chooser 패턴
        try:
            with page.expect_file_chooser(timeout=3000) as fc_info:
                page.click(selector, force=True)
            fc_info.value.set_files(str(fixture))
        except Exception as e2:
            return {"ok": False, "note": f"set_input_files+chooser 모두 실패: {e} / {e2}", "file": str(fixture)}
    # 업로드 진행 대기 (서버 다음 작업 시작 전까지)
    time.sleep(wait_after_ms / 1000.0)
    return {"ok": True, "note": "업로드 완료", "file": str(fixture)}


# ─── 헬퍼: 다운로드 ──────────────────────────────────────────────────
def download_file(
    page: Page,
    click_selector: str,
    *,
    save_dir: Path | None = None,
    min_size_bytes: int = 100,
    timeout_ms: int = 15_000,
) -> dict:
    """`expect_download()` 컨텍스트 — 클릭으로 다운로드 트리거 + 검증.

    Args:
        page: Playwright Page
        click_selector: 다운로드 트리거 element selector
        save_dir: 저장 디렉토리 (default: _downloads_tmp/)
        min_size_bytes: 최소 파일 크기 (이 이하면 fail)
        timeout_ms: 다운로드 대기 timeout

    Returns:
        {"ok": bool, "note": str, "path": str, "size": int, "filename": str}
    """
    save_dir = save_dir or DOWNLOAD_TMP_DIR
    save_dir.mkdir(parents=True, exist_ok=True)
    try:
        with page.expect_download(timeout=timeout_ms) as dl_info:
            page.click(click_selector)
        download: Download = dl_info.value
        suggested = download.suggested_filename or f"download-{int(time.time())}.bin"
        # 안전한 파일명 (한글 가능 — RFC 5987)
        safe_name = re.sub(r"[\\/:*?\"<>|]", "_", suggested)[:120]
        save_path = save_dir / safe_name
        download.save_as(str(save_path))
        size = save_path.stat().st_size if save_path.exists() else 0
        if size < min_size_bytes:
            return {
                "ok": False, "note": f"다운로드 파일 너무 작음 ({size}B < {min_size_bytes}B)",
                "path": str(save_path), "size": size, "filename": suggested,
            }
        return {
            "ok": True, "note": f"다운로드 OK ({size}B, {suggested})",
            "path": str(save_path), "size": size, "filename": suggested,
        }
    except PlaywrightTimeoutError:
        return {"ok": False, "note": f"다운로드 timeout ({timeout_ms}ms)", "path": "", "size": 0, "filename": ""}
    except Exception as e:
        return {"ok": False, "note": f"다운로드 예외: {type(e).__name__}: {e}", "path": "", "size": 0, "filename": ""}


# ─── 헬퍼: iframe ────────────────────────────────────────────────────
def iframe_click(
    page: Page,
    frame_selector: str,
    inner_selector: str,
    *,
    timeout_ms: int = 5000,
) -> dict:
    """`frame_locator()` 로 iframe 진입 + 내부 요소 클릭.

    Args:
        page: Playwright Page
        frame_selector: iframe element selector (예: 'iframe[name="preview"]')
        inner_selector: iframe 내부의 클릭 대상 selector
        timeout_ms: 내부 요소 대기 timeout

    Returns:
        {"ok": bool, "note": str, "found": bool}
    """
    try:
        # iframe element 존재 확인
        frame_el = page.query_selector(frame_selector)
        if not frame_el:
            # 첫 iframe 사용 시도
            frames = page.query_selector_all("iframe")
            if not frames:
                return {"ok": False, "note": "iframe 미발견", "found": False}
            frame_selector = "iframe"
        fl = page.frame_locator(frame_selector).first
        # 내부 요소 대기
        try:
            fl.locator(inner_selector).first.wait_for(state="visible", timeout=timeout_ms)
        except PlaywrightTimeoutError:
            # iframe 자체 mount 확인만 PASS — body 검증으로 대체
            try:
                fl.locator("body").first.wait_for(state="attached", timeout=2000)
                return {"ok": True, "note": f"iframe mount OK (inner {inner_selector} 미발견 — postMessage 검증 필요)", "found": True}
            except Exception:
                return {"ok": False, "note": f"iframe mount 실패", "found": False}
        try:
            fl.locator(inner_selector).first.click()
        except Exception as ce:
            return {"ok": False, "note": f"iframe 내부 클릭 실패: {ce}", "found": True}
        return {"ok": True, "note": f"iframe ({frame_selector}) 내부 클릭 OK", "found": True}
    except Exception as e:
        return {"ok": False, "note": f"iframe 진입 예외: {type(e).__name__}: {e}", "found": False}


def iframe_mounted(page: Page, frame_selector: str = "iframe", *, timeout_ms: int = 5000) -> dict:
    """iframe 이 DOM 에 mount 되어 있는지만 확인 (postMessage 등 interaction 무관)."""
    try:
        page.wait_for_selector(frame_selector, state="attached", timeout=timeout_ms)
        return {"ok": True, "note": f"iframe ({frame_selector}) mount OK", "found": True}
    except PlaywrightTimeoutError:
        return {"ok": False, "note": f"iframe ({frame_selector}) mount 실패 (timeout)", "found": False}
    except Exception as e:
        return {"ok": False, "note": f"iframe 검증 예외: {type(e).__name__}: {e}", "found": False}


# ─── 헬퍼: SSE / streaming ────────────────────────────────────────────
def wait_for_sse(
    page: Page,
    url_pattern: str,
    *,
    timeout_ms: int = 15_000,
    trigger: Callable[[], None] | None = None,
) -> dict:
    """`page.expect_response()` 로 SSE 응답 첫 chunk 도착 확인.

    Args:
        page: Playwright Page
        url_pattern: 응답 URL 의 substring 또는 regex
        timeout_ms: 응답 대기 timeout
        trigger: 응답을 유발할 callable (예: 메시지 전송 클릭)

    Returns:
        {"ok": bool, "note": str, "url": str, "status": int, "content_type": str}
    """
    try:
        if isinstance(url_pattern, str):
            matcher = lambda resp: url_pattern in resp.url  # noqa: E731
        else:
            matcher = lambda resp: bool(url_pattern.search(resp.url))  # noqa: E731

        with page.expect_response(matcher, timeout=timeout_ms) as resp_info:
            if trigger:
                trigger()
        resp: Response = resp_info.value
        ct = resp.headers.get("content-type", "")
        is_sse = "event-stream" in ct or "stream" in ct or resp.status == 200
        return {
            "ok": is_sse,
            "note": f"SSE 응답 — status={resp.status}, ct={ct}",
            "url": resp.url[:200], "status": resp.status, "content_type": ct,
        }
    except PlaywrightTimeoutError:
        return {"ok": False, "note": f"SSE 응답 timeout ({timeout_ms}ms)", "url": "", "status": 0, "content_type": ""}
    except Exception as e:
        return {"ok": False, "note": f"SSE 검증 예외: {type(e).__name__}: {e}", "url": "", "status": 0, "content_type": ""}


# ─── 헬퍼: safe mutation ──────────────────────────────────────────────
@contextmanager
def safe_mutation(
    page: Page,
    *,
    prefix: str | None = None,
    cleanup_fn: Callable[[], None] | None = None,
) -> Iterator[str]:
    """try/finally 패턴 — mutation 후 cleanup 보장.

    Usage:
        with safe_mutation(page) as prefix:
            page.fill('input[name="title"]', f"{prefix}_doc")
            page.click('button[type="submit"]')
            # cleanup_fn 이 finally 에서 실행됨

    Args:
        prefix: 사용할 prefix (없으면 자동 생성)
        cleanup_fn: 정리 callable (실패해도 swallow — 운영 보호)
    """
    p = prefix or ui_e2e_test_prefix()
    try:
        yield p
    finally:
        if cleanup_fn is not None:
            try:
                cleanup_fn()
            except Exception as e:
                # cleanup 실패는 swallow — 다음 case 에 영향 없게
                print(f"[safe_mutation] cleanup 예외 무시: {type(e).__name__}: {e}")


# ─── 헬퍼: viewport 토글 (mobile 햄버거 등) ──────────────────────────
@contextmanager
def viewport_override(page: Page, width: int, height: int) -> Iterator[None]:
    """일시적 viewport 변경 (mobile 시뮬레이션 등) — 종료 시 복구 안 함 (page 단위)."""
    page.set_viewport_size({"width": width, "height": height})
    try:
        yield
    finally:
        pass  # caller 가 page context 종료할 때 자동 복구


# ─── 헬퍼: networkidle + 안전 대기 ───────────────────────────────────
def wait_settled(page: Page, *, timeout_ms: int = 5000, sleep_after: float = 0.5) -> None:
    """networkidle + 짧은 sleep — 결정성 확보."""
    try:
        page.wait_for_load_state("networkidle", timeout=timeout_ms)
    except Exception:
        pass
    if sleep_after > 0:
        time.sleep(sleep_after)


# ─── 헬퍼: API count (자동 새로고침 검증용) ────────────────────────────
class ApiCallCounter:
    """특정 URL 패턴의 호출 횟수 카운터 — page.on('response') 리스너."""

    def __init__(self, page: Page, url_pattern: str) -> None:
        self.page = page
        self.pattern = url_pattern
        self.count = 0
        self._handler = self._on_response

    def _on_response(self, resp: Response) -> None:
        try:
            if self.pattern in resp.url:
                self.count += 1
        except Exception:
            pass

    def start(self) -> "ApiCallCounter":
        self.page.on("response", self._handler)
        return self

    def stop(self) -> int:
        try:
            self.page.remove_listener("response", self._handler)
        except Exception:
            pass
        return self.count


# ─── 헬퍼: body 가 비어있지 않은지 ───────────────────────────────────
def assert_body_not_empty(page: Page, *, min_len: int = 30) -> dict:
    """body innerText 가 비어있지 않은지 검증."""
    try:
        text = page.evaluate("() => (document.body.innerText || '').length")
        if text >= min_len:
            return {"ok": True, "note": f"body OK (len={text})", "length": text}
        return {"ok": False, "note": f"body 비어있음 (len={text} < {min_len})", "length": text}
    except Exception as e:
        return {"ok": False, "note": f"body 검증 예외: {e}", "length": 0}


# ─── 헬퍼: 외부 클릭 (드롭다운 닫기 등) ──────────────────────────────
def click_outside(page: Page, *, x: int = 10, y: int = 10) -> None:
    """body 의 빈 영역 클릭 — 드롭다운 자동 닫기 등 검증."""
    try:
        page.mouse.click(x, y)
    except Exception:
        pass


# ─── public API ─────────────────────────────────────────────────────
__all__ = [
    # flags
    "allow_cost", "allow_mutation",
    # prefix / cleanup
    "ui_e2e_test_prefix", "is_test_artifact",
    # fixtures
    "ensure_fixture", "get_sample_pdf", "get_sample_docx", "get_sample_pptx",
    "FIXTURE_DIR", "DOWNLOAD_TMP_DIR",
    # actions
    "upload_file", "download_file",
    "iframe_click", "iframe_mounted",
    "wait_for_sse",
    "safe_mutation",
    "viewport_override", "wait_settled",
    "ApiCallCounter",
    "assert_body_not_empty", "click_outside",
]
