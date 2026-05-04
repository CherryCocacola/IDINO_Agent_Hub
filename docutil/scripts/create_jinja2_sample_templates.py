#!/usr/bin/env python3
"""
Jinja2 샘플 템플릿 생성기
=========================
Jinja2 반복문({% for %})과 변수({{ }})를 포함한 DOCX 3종 + PPTX 2종 = 총 5개의
전문적 디자인 템플릿을 생성합니다.

생성되는 템플릿 목록:
    1. weekly_report.docx     — 주간업무보고서
    2. business_proposal.pptx — 사업제안서
    3. meeting_minutes.docx   — 회의록
    4. tech_review.docx       — 기술검토보고서
    5. project_status.pptx    — 프로젝트현황보고

사용법:
    # 템플릿 파일만 생성
    python scripts/create_jinja2_sample_templates.py

    # 생성 후 API 서버에 자동 등록까지 수행
    python scripts/create_jinja2_sample_templates.py --register

의존성:
    pip install python-pptx python-docx requests
"""

import os
import sys
import argparse
from pathlib import Path

# ─── python-docx: Word 문서(.docx)를 생성하기 위한 라이브러리 ────────────────
from docx import Document
from docx.shared import Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn, nsdecls
from docx.oxml import parse_xml

# ─── python-pptx: 파워포인트(.pptx)를 생성하기 위한 라이브러리 ──────────────
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt, Emu as PptxEmu
from pptx.dml.color import RGBColor as PptxRGB
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ═══════════════════════════════════════════════════════════════════════════════
#  IDINO 디자인 시스템 — 색상, 폰트, 크기 상수
#  (create_idino_templates.py와 동일한 디자인 상수를 사용합니다)
# ═══════════════════════════════════════════════════════════════════════════════

HEADER_BAR    = "34495E"   # 진한 남색 — 헤더/푸터 배경, 타이틀 배경
ACCENT_BLUE   = "2980B9"   # 밝은 파랑 — 강조선, 라벨, 링크
HIGHLIGHT_BG  = "EBF5FB"   # 밝은 파랑 배경 — 콘텐츠 박스 배경
DARK_ACCENT   = "1A5276"   # 진한 파랑 — 부제, 강조 텍스트
BODY_TEXT      = "2C3E50"   # 진한 회색 — 본문 텍스트
SUBTITLE_GRAY = "D5D8DC"   # 밝은 회색 — 부제, 날짜 텍스트
WHITE          = "FFFFFF"   # 흰색 — 어두운 배경 위의 텍스트

FONT_NAME = "Malgun Gothic"  # 맑은 고딕 — 전체 문서에 사용하는 기본 폰트

# 출력 디렉토리: 프로젝트 루트의 sample_templates/jinja2/ 폴더
OUTPUT_DIR = Path(__file__).resolve().parent.parent / "sample_templates" / "jinja2"


# ═══════════════════════════════════════════════════════════════════════════════
#  PPTX 헬퍼 함수들 — 슬라이드 생성에 반복적으로 사용되는 유틸리티
# ═══════════════════════════════════════════════════════════════════════════════

# 슬라이드 크기: 16:9 비율 (13.333인치 x 7.5인치)
SLIDE_W = PptxEmu(12192000)  # 13.333인치를 EMU(English Metric Units)로 변환한 값
SLIDE_H = PptxEmu(6858000)   # 7.5인치를 EMU로 변환한 값

# 헤더바/푸터바 높이
HEADER_BAR_H = PptxInches(0.42)                 # 상단 헤더바 높이 (0.42인치)
FOOTER_BAR_H = PptxInches(0.18)                 # 하단 푸터바 높이 (0.18인치)
FOOTER_TOP   = SLIDE_H - FOOTER_BAR_H           # 푸터바 Y 좌표 계산


def _new_pptx():
    """
    새로운 프레젠테이션(PPTX) 객체를 생성합니다.
    슬라이드 크기를 16:9 비율로 설정합니다.

    반환값:
        Presentation 객체
    """
    prs = Presentation()
    prs.slide_width = SLIDE_W    # 슬라이드 너비 설정
    prs.slide_height = SLIDE_H   # 슬라이드 높이 설정
    return prs


def _add_shape_rect(slide, left, top, width, height, fill_hex):
    """
    슬라이드에 색이 채워진 직사각형 도형을 추가합니다.
    배경 영역이나 강조 줄을 만들 때 사용합니다.

    매개변수:
        slide: 도형을 추가할 슬라이드 객체
        left: 왼쪽 여백 (EMU 단위)
        top: 위쪽 여백 (EMU 단위)
        width: 사각형 너비
        height: 사각형 높이
        fill_hex: 채우기 색상 코드 (예: "34495E")

    반환값:
        생성된 도형(Shape) 객체
    """
    # MSO_SHAPE.RECTANGLE = 1 (직사각형 도형 타입 번호)
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()                                    # 단색 채우기 모드 설정
    shape.fill.fore_color.rgb = PptxRGB.from_string(fill_hex)  # 색상 적용
    shape.line.fill.background()                          # 테두리 없음 (투명)
    return shape


def _add_text_box(slide, left, top, width, height, text,
                  size=12, bold=False, color=BODY_TEXT,
                  alignment=PP_ALIGN.LEFT, font=FONT_NAME,
                  anchor=MSO_ANCHOR.TOP):
    """
    슬라이드에 텍스트 상자를 추가합니다.
    제목, 본문, 변수 플레이스홀더 등을 표시할 때 사용합니다.

    매개변수:
        slide: 텍스트를 추가할 슬라이드 객체
        left, top: 텍스트 상자의 왼쪽 위 좌표
        width, height: 텍스트 상자 크기
        text: 표시할 텍스트 내용
        size: 글자 크기 (pt 단위, 기본값 12)
        bold: 굵게 표시 여부 (기본값 False)
        color: 글자 색상 (기본값 본문 회색)
        alignment: 텍스트 정렬 (왼쪽/가운데/오른쪽)
        font: 사용할 폰트 이름
        anchor: 세로 정렬 (상단/중앙/하단)

    반환값:
        생성된 텍스트 상자 객체
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True       # 텍스트가 상자를 넘으면 자동 줄바꿈
    tf.auto_size = None       # 자동 크기 조절 비활성화 (고정 크기 유지)

    # 첫 번째 문단에 텍스트와 서식을 설정
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = PptxPt(size)
    p.font.bold = bold
    p.font.color.rgb = PptxRGB.from_string(color)
    p.font.name = font
    p.alignment = alignment
    return txBox


def _add_multi_line_text_box(slide, left, top, width, height, lines,
                             size=11, bold=False, color=BODY_TEXT,
                             alignment=PP_ALIGN.LEFT, font=FONT_NAME):
    """
    슬라이드에 여러 줄의 텍스트를 가진 텍스트 상자를 추가합니다.
    Jinja2 for 루프가 포함된 블록처럼 여러 줄이 필요한 경우에 사용합니다.

    매개변수:
        slide: 텍스트를 추가할 슬라이드 객체
        left, top: 텍스트 상자의 왼쪽 위 좌표
        width, height: 텍스트 상자 크기
        lines: 표시할 텍스트 줄 목록 (문자열 리스트)
        size: 글자 크기 (pt)
        bold: 굵게 여부
        color: 글자 색상
        alignment: 텍스트 정렬
        font: 폰트 이름

    반환값:
        생성된 텍스트 상자 객체
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    # 첫 번째 줄은 기존 문단에 설정
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            # 두 번째 줄부터는 새 문단 추가
            p = tf.add_paragraph()
        p.text = line
        p.font.size = PptxPt(size)
        p.font.bold = bold
        p.font.color.rgb = PptxRGB.from_string(color)
        p.font.name = font
        p.alignment = alignment
    return txBox


def _pptx_add_header_bar(slide, section_title):
    """
    콘텐츠 슬라이드 상단에 헤더바를 추가합니다.
    진한 남색 배경에 흰색 섹션 제목이 표시됩니다.

    매개변수:
        slide: 헤더를 추가할 슬라이드
        section_title: 헤더에 표시할 섹션 제목 (예: "1. 사업 개요")
    """
    # 전체 너비의 진한 남색 배경 사각형
    _add_shape_rect(slide, 0, 0, SLIDE_W, HEADER_BAR_H, HEADER_BAR)
    # 그 위에 흰색 섹션 제목 텍스트
    _add_text_box(
        slide, PptxInches(0.5), 0,
        PptxInches(10), HEADER_BAR_H,
        section_title,
        size=13, bold=True, color=WHITE,
        alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE
    )


def _pptx_add_footer_bar(slide):
    """
    콘텐츠 슬라이드 하단에 IDINO 로고가 있는 푸터바를 추가합니다.
    진한 남색 배경에 "IDINO" 텍스트가 오른쪽 정렬로 표시됩니다.
    """
    # 하단 진한 남색 배경 사각형
    _add_shape_rect(slide, 0, FOOTER_TOP, SLIDE_W, FOOTER_BAR_H, HEADER_BAR)
    # 오른쪽에 "IDINO" 로고 텍스트
    _add_text_box(
        slide, SLIDE_W - PptxInches(1.2), FOOTER_TOP,
        PptxInches(1.0), FOOTER_BAR_H,
        "IDINO",
        size=7, bold=True, color=WHITE,
        alignment=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE
    )


def _pptx_add_title_slide(prs, title, subtitle, date_text="{{작성일}}"):
    """
    프레젠테이션의 첫 번째 슬라이드(표지)를 추가합니다.
    전체 배경이 진한 남색이고, 중앙에 제목/부제/날짜가 배치됩니다.

    구성:
        - 전체 배경: 진한 남색 (#34495E)
        - 상단 강조선: 파란색 가느다란 줄
        - 제목: 32pt, 굵게, 흰색
        - 부제: 14pt, 밝은 회색
        - 날짜: 11pt, 밝은 회색
        - 하단 IDINO 로고: 11pt, 파란색

    매개변수:
        prs: Presentation 객체
        title: 메인 제목 (예: "{{사업명}}")
        subtitle: 부제 (예: "사업제안서 | {{회사명}}")
        date_text: 날짜 텍스트 (기본값: "{{작성일}}")

    반환값:
        생성된 슬라이드 객체
    """
    slide_layout = prs.slide_layouts[6]       # 빈(blank) 레이아웃 사용
    slide = prs.slides.add_slide(slide_layout)

    # 1) 전체 배경을 진한 남색으로 채움
    _add_shape_rect(slide, 0, 0, SLIDE_W, SLIDE_H, HEADER_BAR)

    # 2) 상단 강조선 — 파란색 가느다란 가로 줄 (전체 너비, 높이 0.03인치)
    _add_shape_rect(slide, 0, 0, SLIDE_W, PptxInches(0.03), ACCENT_BLUE)

    # 3) 메인 제목 — 32pt, 굵게, 흰색, 수평 중앙 정렬, Y=2.5인치
    _add_text_box(
        slide, PptxInches(1), PptxInches(2.5),
        PptxInches(11.333), PptxInches(1.0),
        title,
        size=32, bold=True, color=WHITE,
        alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP
    )

    # 4) 부제 — 14pt, 밝은 회색, 수평 중앙 정렬, Y=3.5인치
    _add_text_box(
        slide, PptxInches(1), PptxInches(3.5),
        PptxInches(11.333), PptxInches(0.5),
        subtitle,
        size=14, color=SUBTITLE_GRAY,
        alignment=PP_ALIGN.CENTER
    )

    # 5) 날짜 — 11pt, 밝은 회색, 수평 중앙 정렬, Y=4.5인치
    _add_text_box(
        slide, PptxInches(1), PptxInches(4.5),
        PptxInches(11.333), PptxInches(0.4),
        date_text,
        size=11, color=SUBTITLE_GRAY,
        alignment=PP_ALIGN.CENTER
    )

    # 6) 하단 IDINO 로고 — 11pt, 파란색, 수평 중앙, Y=6.5인치
    _add_text_box(
        slide, PptxInches(5), PptxInches(6.5),
        PptxInches(3.333), PptxInches(0.4),
        "IDINO",
        size=11, bold=True, color=ACCENT_BLUE,
        alignment=PP_ALIGN.CENTER
    )

    return slide


def _pptx_add_content_slide(prs, section_title, items):
    """
    콘텐츠 슬라이드를 추가합니다.
    상단 헤더바 + 콘텐츠 항목들 + 하단 푸터바로 구성됩니다.

    각 항목(item)은 다음과 같이 렌더링됩니다:
        - "■ 라벨" (파란색, 굵게)
        - 밝은 파란 배경 박스 + 좌측 파란 테두리 안에 변수 텍스트

    매개변수:
        prs: Presentation 객체
        section_title: 헤더바에 표시될 제목 (예: "1. 사업 개요")
        items: [(라벨, 변수_텍스트)] 리스트
               예: [("사업 개요", "{{사업개요}}")]
               변수_텍스트가 리스트이면 여러 줄로 표시됨

    반환값:
        생성된 슬라이드 객체
    """
    slide_layout = prs.slide_layouts[6]       # 빈 레이아웃
    slide = prs.slides.add_slide(slide_layout)

    # 상단 헤더바와 하단 푸터바 추가
    _pptx_add_header_bar(slide, section_title)
    _pptx_add_footer_bar(slide)

    # ── 콘텐츠 영역 레이아웃 계산 ──
    content_left = PptxInches(0.5)    # 콘텐츠 좌측 시작 위치
    content_width = PptxInches(12.3)  # 콘텐츠 너비 (슬라이드 폭 - 양쪽 여백)
    y_start = PptxInches(0.7)         # 콘텐츠 시작 Y 좌표
    y_end = PptxInches(6.6)           # 콘텐츠 끝 Y 좌표

    n = len(items)
    if n == 0:
        return slide

    # 각 항목에 할당되는 높이를 균등 분배
    total_height_emu = y_end - y_start
    item_height_emu = total_height_emu // n     # 항목 하나의 전체 높이
    label_h = PptxInches(0.3)                   # "■ 라벨" 영역 높이 (고정)
    gap = PptxInches(0.1)                       # 항목 간 간격 (고정)
    box_h = item_height_emu - label_h - gap     # 나머지가 콘텐츠 박스 높이

    # 박스 높이가 너무 작거나 크지 않도록 제한
    min_box = PptxInches(0.4)
    max_box = PptxInches(1.5)
    if box_h < min_box:
        box_h = min_box
    elif box_h > max_box:
        box_h = max_box

    y = y_start  # 현재 Y 위치 (위에서부터 아래로 이동)
    for label, variable in items:
        # ── "■ 라벨" 텍스트 (파란색, 굵게) ──
        _add_text_box(
            slide, content_left, y,
            content_width, label_h,
            f"\u25a0 {label}",       # ■ 유니코드 문자 + 라벨명
            size=12, bold=True, color=ACCENT_BLUE
        )
        y += label_h

        # ── 콘텐츠 배경 박스 (밝은 파란 배경) ──
        _add_shape_rect(
            slide, content_left, y,
            content_width, box_h, HIGHLIGHT_BG
        )

        # ── 좌측 파란 테두리 (가느다란 사각형으로 표현) ──
        _add_shape_rect(
            slide, content_left, y,
            PptxInches(0.08), box_h, ACCENT_BLUE
        )

        # ── 변수 텍스트 (배경 박스 안에 표시) ──
        # variable이 리스트이면 여러 줄로, 문자열이면 한 줄로 표시
        if isinstance(variable, list):
            _add_multi_line_text_box(
                slide,
                content_left + PptxInches(0.2), y + PptxInches(0.05),
                content_width - PptxInches(0.35), box_h - PptxInches(0.1),
                variable,
                size=11, color=BODY_TEXT
            )
        else:
            _add_text_box(
                slide,
                content_left + PptxInches(0.2), y + PptxInches(0.05),
                content_width - PptxInches(0.35), box_h - PptxInches(0.1),
                variable,
                size=11, color=BODY_TEXT
            )

        y += box_h + gap  # 다음 항목 위치로 이동

    return slide


def _pptx_add_closing_slide(prs):
    """
    프레젠테이션의 마지막 슬라이드(감사합니다)를 추가합니다.

    구성:
        - 전체 배경: 진한 남색
        - "감사합니다" 텍스트: 36pt, 굵게, 흰색, 중앙
        - 가운데 파란 강조선
        - 연락처 안내문
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 전체 배경
    _add_shape_rect(slide, 0, 0, SLIDE_W, SLIDE_H, HEADER_BAR)

    # "감사합니다" 메인 텍스트
    _add_text_box(
        slide, PptxInches(1), PptxInches(2.5),
        PptxInches(11.333), PptxInches(1.2),
        "감사합니다",
        size=36, bold=True, color=WHITE,
        alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE
    )

    # 가운데 파란 강조선 (너비 2인치, 수평 중앙 배치)
    line_w = PptxInches(2)
    line_left = (SLIDE_W - line_w) // 2
    _add_shape_rect(
        slide, line_left, PptxInches(3.8),
        line_w, PptxInches(0.03), ACCENT_BLUE
    )

    # 하단 연락처 안내문
    _add_text_box(
        slide, PptxInches(2), PptxInches(4.2),
        PptxInches(9.333), PptxInches(0.5),
        "IDINO | 문의사항은 담당자에게 연락 바랍니다",
        size=12, color=SUBTITLE_GRAY,
        alignment=PP_ALIGN.CENTER
    )

    return slide


# ═══════════════════════════════════════════════════════════════════════════════
#  DOCX 헬퍼 함수들 — Word 문서 생성에 반복적으로 사용되는 유틸리티
# ═══════════════════════════════════════════════════════════════════════════════

def _set_cell_shading(cell, color_hex):
    """
    표 셀의 배경색을 설정합니다.
    XML을 직접 다루어 python-docx가 지원하지 않는 셀 음영(shading)을 적용합니다.

    매개변수:
        cell: python-docx 테이블 셀 객체
        color_hex: 배경 색상 코드 (예: "34495E")
    """
    shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{color_hex}"/>')
    cell._tc.get_or_add_tcPr().append(shading)


def _set_cell_border(cell, **kwargs):
    """
    표 셀의 테두리를 설정합니다.
    원하는 변(상하좌우)에만 선택적으로 테두리를 추가할 수 있습니다.

    매개변수:
        cell: python-docx 테이블 셀 객체
        kwargs: 테두리를 설정할 변과 (두께, 색상) 쌍
                예: left=("12", "2980B9") → 왼쪽에 두께 12, 파란색 테두리
    """
    tc_pr = cell._tc.get_or_add_tcPr()
    borders = parse_xml(f'<w:tcBorders {nsdecls("w")}/>')
    for edge, (sz, color) in kwargs.items():
        el = parse_xml(
            f'<w:{edge} {nsdecls("w")} w:val="single" w:sz="{sz}" '
            f'w:space="0" w:color="{color}"/>'
        )
        borders.append(el)
    tc_pr.append(borders)


def _docx_run(paragraph, text, size=11, bold=False, color=BODY_TEXT, font=FONT_NAME):
    """
    문단에 서식이 적용된 텍스트 조각(run)을 추가합니다.
    python-docx에서 하나의 문단 안에 서로 다른 서식을 가진
    텍스트 조각들을 조합할 때 사용합니다.

    매개변수:
        paragraph: 텍스트를 추가할 문단 객체
        text: 추가할 텍스트 내용
        size: 글자 크기 (pt 단위)
        bold: 굵게 여부
        color: 글자 색상 코드
        font: 폰트 이름

    반환값:
        생성된 run(텍스트 조각) 객체
    """
    r = paragraph.add_run(text)
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = RGBColor.from_string(color)
    r.font.name = font
    # 동아시아(한글) 폰트도 동일하게 설정 — 한글이 깨지지 않도록 필수
    r._element.rPr.rFonts.set(qn("w:eastAsia"), font)
    return r


def _docx_add_thin_line(doc, color=ACCENT_BLUE):
    """
    문단 하단 테두리 방식으로 가느다란 수평 구분선을 추가합니다.
    섹션 제목 아래에 사용하여 시각적 구분을 제공합니다.

    매개변수:
        doc: Document 객체
        color: 구분선 색상 (기본값: 파란색)

    반환값:
        구분선이 포함된 문단 객체
    """
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(2)
    p.paragraph_format.space_after = Pt(6)
    pPr = p._element.get_or_add_pPr()
    pBdr = parse_xml(
        f'<w:pBdr {nsdecls("w")}>'
        f'  <w:bottom w:val="single" w:sz="6" w:space="1" w:color="{color}"/>'
        f'</w:pBdr>'
    )
    pPr.append(pBdr)
    return p


def _docx_setup():
    """
    새 DOCX 문서를 생성하고 기본 스타일(용지, 폰트, 여백)을 설정합니다.

    설정 내용:
        - 용지 크기: A4 (21cm x 29.7cm)
        - 여백: 상하좌우 2.5cm
        - 기본 폰트: 맑은 고딕 11pt, 행간 1.15
        - Heading 1: 16pt, 굵게, 진한 남색
        - Heading 2: 14pt, 굵게, 파란색
        - Heading 3: 12pt, 굵게, 파란색

    반환값:
        설정이 완료된 Document 객체
    """
    doc = Document()

    # ── 페이지 크기 및 여백 설정 ──
    section = doc.sections[0]
    section.page_width = Cm(21)         # A4 가로 크기
    section.page_height = Cm(29.7)      # A4 세로 크기
    section.top_margin = Cm(2.5)        # 상단 여백
    section.bottom_margin = Cm(2.5)     # 하단 여백
    section.left_margin = Cm(2.5)       # 좌측 여백
    section.right_margin = Cm(2.5)      # 우측 여백

    # ── 기본(Normal) 스타일 설정 ──
    style = doc.styles["Normal"]
    style.font.name = FONT_NAME         # 기본 폰트를 맑은 고딕으로
    style.font.size = Pt(11)            # 기본 글자 크기 11pt
    style.font.color.rgb = RGBColor.from_string(BODY_TEXT)  # 본문 색상
    style.paragraph_format.line_spacing = 1.15              # 행간 1.15배
    # 한글 폰트 설정 (동아시아 글꼴)
    rpr = style._element.rPr
    rpr.rFonts.set(qn("w:eastAsia"), FONT_NAME)

    # ── Heading 스타일 커스터마이징 ──
    # 레벨별 글자 크기와 색상을 순서대로 설정
    for level, sz, clr in [
        (1, 16, HEADER_BAR),     # Heading 1: 16pt, 진한 남색
        (2, 14, ACCENT_BLUE),    # Heading 2: 14pt, 파란색
        (3, 12, ACCENT_BLUE),    # Heading 3: 12pt, 파란색
    ]:
        hs = doc.styles[f"Heading {level}"]
        hs.font.name = FONT_NAME
        hs.font.size = Pt(sz)
        hs.font.bold = True
        hs.font.color.rgb = RGBColor.from_string(clr)
        hs._element.rPr.rFonts.set(qn("w:eastAsia"), FONT_NAME)

    return doc


def _docx_add_cover_page(doc, title, subtitle_lines):
    """
    문서 첫 페이지에 커버(표지)를 추가합니다.
    표(1행 1열)의 셀 배경색을 이용하여 진한 남색 배너를 만들고,
    그 안에 제목과 부제 정보를 배치합니다.

    매개변수:
        doc: Document 객체
        title: 문서 제목 (예: "주간업무보고서")
        subtitle_lines: 부제 줄 리스트
                        예: ["부서: {{부서명}}", "작성자: {{작성자}}"]
    """
    # ── 배너를 1x1 표로 구현 (셀 배경색 = 진한 남색) ──
    tbl = doc.add_table(rows=1, cols=1)
    tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
    cell = tbl.rows[0].cells[0]
    _set_cell_shading(cell, HEADER_BAR)    # 진한 남색 배경
    cell.width = Cm(16)

    # 셀 높이를 약 1.8인치(3600 twips)로 설정
    tr = tbl.rows[0]._tr
    trPr = tr.get_or_add_trPr()
    trHeight = parse_xml(
        f'<w:trHeight {nsdecls("w")} w:val="3600" w:hRule="atLeast"/>'
    )
    trPr.append(trHeight)

    # ── 제목: 28pt, 굵게, 흰색, 가운데 정렬 ──
    p = cell.paragraphs[0]
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(60)  # 상단 여백으로 세로 중앙 효과
    _docx_run(p, title, size=28, bold=True, color=WHITE)

    # ── 부제 줄들: 13pt, 밝은 회색, 가운데 정렬 ──
    for line in subtitle_lines:
        p2 = cell.add_paragraph()
        p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p2.paragraph_format.space_before = Pt(6)
        _docx_run(p2, line, size=13, color=SUBTITLE_GRAY)

    # ── 파란색 구분선 (━ 문자 30개 반복) ──
    p3 = cell.add_paragraph()
    p3.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p3.paragraph_format.space_before = Pt(20)
    _docx_run(p3, "\u2501" * 30, size=11, color=ACCENT_BLUE)

    # 배너 다음에 빈 줄 + 페이지 넘김 → 본문은 두 번째 페이지부터 시작
    doc.add_paragraph()
    doc.add_page_break()


def _docx_add_header(doc, title_text):
    """
    문서 상단에 반복되는 헤더(머리글)를 추가합니다.
    좌측에 문서 유형명이 표시되고, 하단에 가느다란 구분선이 있습니다.

    매개변수:
        doc: Document 객체
        title_text: 헤더에 표시할 문서 유형명 (예: "주간업무보고서")
    """
    section = doc.sections[0]
    header = section.header
    header.is_linked_to_previous = False  # 이전 섹션 헤더와 연결 해제

    p = header.paragraphs[0] if header.paragraphs else header.add_paragraph()
    _docx_run(p, title_text, size=9, bold=True, color=HEADER_BAR)

    # 헤더 하단에 구분선 추가
    pPr = p._element.get_or_add_pPr()
    pBdr = parse_xml(
        f'<w:pBdr {nsdecls("w")}>'
        f'  <w:bottom w:val="single" w:sz="4" w:space="4" w:color="{HEADER_BAR}"/>'
        f'</w:pBdr>'
    )
    pPr.append(pBdr)


def _docx_add_footer(doc):
    """
    문서 하단에 페이지 번호가 포함된 푸터(바닥글)를 추가합니다.
    "— 1 —" 형식으로 가운데 정렬되며, 상단에 구분선이 있습니다.
    """
    section = doc.sections[0]
    footer = section.footer
    footer.is_linked_to_previous = False

    p = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # "— " 텍스트
    _docx_run(p, "\u2014 ", size=9, color=HEADER_BAR)

    # 페이지 번호 필드 (Word의 PAGE 필드 코드를 XML로 삽입)
    fld_xml = (
        f'<w:fldSimple {nsdecls("w")} w:instr=" PAGE \\* MERGEFORMAT ">'
        f'  <w:r><w:rPr><w:sz w:val="18"/><w:color w:val="{HEADER_BAR}"/></w:rPr>'
        f'  <w:t>1</w:t></w:r></w:fldSimple>'
    )
    p._element.append(parse_xml(fld_xml))

    # " —" 텍스트
    _docx_run(p, " \u2014", size=9, color=HEADER_BAR)

    # 푸터 상단 구분선
    pPr = p._element.get_or_add_pPr()
    pBdr = parse_xml(
        f'<w:pBdr {nsdecls("w")}>'
        f'  <w:top w:val="single" w:sz="4" w:space="4" w:color="{HEADER_BAR}"/>'
        f'</w:pBdr>'
    )
    pPr.append(pBdr)


def _docx_add_section_heading(doc, text):
    """
    섹션 제목(Heading 1)과 그 아래 파란 구분선을 추가합니다.

    매개변수:
        doc: Document 객체
        text: 섹션 제목 (예: "1. 완료 업무")
    """
    doc.add_heading(text, level=1)
    _docx_add_thin_line(doc, ACCENT_BLUE)


def _docx_add_content_block(doc, variable):
    """
    콘텐츠 블록을 추가합니다.
    밝은 파란 배경 + 좌측 파란 테두리를 가진 박스 안에
    Jinja2 변수 텍스트(예: {{특이사항}})가 표시됩니다.

    표(1행 1열)를 이용하여 배경색과 테두리 효과를 구현합니다.

    매개변수:
        doc: Document 객체
        variable: 변수 텍스트 (문자열 또는 줄 리스트)
    """
    tbl = doc.add_table(rows=1, cols=1)
    cell = tbl.rows[0].cells[0]
    _set_cell_shading(cell, HIGHLIGHT_BG)  # 밝은 파란 배경

    # 변수가 리스트이면 여러 줄로, 문자열이면 한 줄로 표시
    if isinstance(variable, list):
        p = cell.paragraphs[0]
        p.paragraph_format.space_before = Pt(8)
        p.paragraph_format.space_after = Pt(4)
        _docx_run(p, variable[0], size=11, color=BODY_TEXT)
        for line in variable[1:]:
            p2 = cell.add_paragraph()
            p2.paragraph_format.space_before = Pt(2)
            p2.paragraph_format.space_after = Pt(2)
            _docx_run(p2, line, size=11, color=BODY_TEXT)
    else:
        p = cell.paragraphs[0]
        p.paragraph_format.space_before = Pt(8)
        p.paragraph_format.space_after = Pt(8)
        _docx_run(p, variable, size=11, color=BODY_TEXT)

    # 좌측 파란 테두리 추가 (두께 12)
    _set_cell_border(cell, left=("12", ACCENT_BLUE))

    # 블록 뒤에 빈 줄 (간격 확보)
    doc.add_paragraph()


def _docx_add_info_table(doc, rows_data):
    """
    정보 테이블(라벨-값 쌍)을 추가합니다.
    라벨 셀은 진한 남색 배경에 흰색 텍스트,
    값 셀은 기본 배경에 본문색 텍스트로 표시됩니다.

    매개변수:
        doc: Document 객체
        rows_data: [(라벨, 값)] 리스트
                   예: [("회의 제목", "{{회의제목}}"), ("일시", "{{일시}}")]
    """
    tbl = doc.add_table(rows=len(rows_data), cols=2)
    tbl.alignment = WD_TABLE_ALIGNMENT.CENTER

    for i, (label, value) in enumerate(rows_data):
        # 라벨 셀 (왼쪽 열): 진한 남색 배경, 흰색 글씨
        lc = tbl.rows[i].cells[0]
        lc.width = Cm(3.5)
        _set_cell_shading(lc, HEADER_BAR)
        p = lc.paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        _docx_run(p, label, size=10, bold=True, color=WHITE)

        # 값 셀 (오른쪽 열): 기본 배경, 본문색 글씨
        vc = tbl.rows[i].cells[1]
        vc.width = Cm(12.5)
        p2 = vc.paragraphs[0]
        _docx_run(p2, value, size=10, color=BODY_TEXT)
        p2.paragraph_format.left_indent = Cm(0.3)  # 약간의 왼쪽 들여쓰기

    # 테이블 뒤 간격
    doc.add_paragraph()


def _docx_add_jinja_table(doc, headers, row_template_lines, loop_var, item_var):
    """
    Jinja2 for 루프를 사용하는 표를 추가합니다.
    표의 헤더 행은 진한 남색 배경에 흰색 텍스트이고,
    데이터 행에는 Jinja2 루프 문법과 변수가 포함됩니다.

    예시 (3열 표):
        ┌───────────┬───────────┬───────────┐
        │  업무명   │  담당자   │  진행률   │  ← 헤더 (진한 남색 배경)
        ├───────────┼───────────┼───────────┤
        │ {% for task in 완료업무 %}         │  ← Jinja2 for 시작
        │{{task.업무명}}│{{task.담당자}}│{{task.진행률}}│
        │ {% endfor %}                       │  ← Jinja2 for 끝
        └───────────┴───────────┴───────────┘

    매개변수:
        doc: Document 객체
        headers: 열 제목 리스트 (예: ["업무명", "담당자", "진행률"])
        row_template_lines: 각 열의 Jinja2 변수 리스트
                            (예: ["{{task.업무명}}", "{{task.담당자}}", "{{task.진행률}}"])
        loop_var: for 루프의 반복 변수명 (예: "task")
        item_var: for 루프의 컬렉션 변수명 (예: "완료업무")
    """
    num_cols = len(headers)

    # 표 생성: 헤더 1행 + for 시작 1행 + 데이터 1행 + for 끝 1행 = 총 4행
    tbl = doc.add_table(rows=4, cols=num_cols)
    tbl.alignment = WD_TABLE_ALIGNMENT.CENTER

    # ── 1행: 헤더 (진한 남색 배경, 흰색 텍스트) ──
    for j, header_text in enumerate(headers):
        cell = tbl.rows[0].cells[j]
        _set_cell_shading(cell, HEADER_BAR)
        p = cell.paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        _docx_run(p, header_text, size=10, bold=True, color=WHITE)

    # ── 2행: Jinja2 for 루프 시작 태그 ──
    # 첫 번째 셀에 for 구문 삽입, 나머지 셀은 비움
    for_start = f"{{% for {loop_var} in {item_var} %}}"
    p = tbl.rows[1].cells[0].paragraphs[0]
    _docx_run(p, for_start, size=9, color=ACCENT_BLUE)
    # 나머지 셀에 밝은 파란 배경 적용 (시각적 구분)
    for j in range(num_cols):
        _set_cell_shading(tbl.rows[1].cells[j], HIGHLIGHT_BG)

    # ── 3행: 데이터 행 (각 열에 Jinja2 변수) ──
    for j, tmpl in enumerate(row_template_lines):
        cell = tbl.rows[2].cells[j]
        _set_cell_shading(cell, HIGHLIGHT_BG)
        p = cell.paragraphs[0]
        _docx_run(p, tmpl, size=10, color=BODY_TEXT)

    # ── 4행: Jinja2 for 루프 끝 태그 ──
    for_end = "{% endfor %}"
    p = tbl.rows[3].cells[0].paragraphs[0]
    _docx_run(p, for_end, size=9, color=ACCENT_BLUE)
    for j in range(num_cols):
        _set_cell_shading(tbl.rows[3].cells[j], HIGHLIGHT_BG)

    # 표 뒤 간격
    doc.add_paragraph()


# ═══════════════════════════════════════════════════════════════════════════════
#  템플릿 1: 주간업무보고서 (DOCX)
# ═══════════════════════════════════════════════════════════════════════════════

def create_weekly_report():
    """
    주간업무보고서 DOCX 템플릿을 생성합니다.

    구조:
        [커버 페이지]
            - 제목: "주간업무보고서"
            - 부서명, 보고기간, 작성자 (Jinja2 변수)

        [본문]
            - 섹션 1: 완료 업무 — Jinja2 for 루프로 표 생성
            - 섹션 2: 진행 업무 — Jinja2 for 루프로 표 생성
            - 섹션 3: 특이사항 — 단일 텍스트 블록

    Jinja2 변수:
        - {{부서명}}, {{보고기간}}, {{작성자}}
        - {% for task in 완료업무 %} → {{task.업무명}}, {{task.담당자}}, {{task.진행률}}
        - {% for task in 진행업무 %} → {{task.업무명}}, {{task.담당자}}, {{task.진행률}}
        - {{특이사항}}
    """
    print("  [1/5] 주간업무보고서 (DOCX) 생성 중...")

    doc = _docx_setup()

    # ── 커버 페이지 ──
    _docx_add_cover_page(
        doc,
        title="주간업무보고서",
        subtitle_lines=[
            "부서: {{부서명}}",
            "보고 기간: {{보고기간}}",
            "작성자: {{작성자}}",
        ],
    )

    # ── 헤더와 푸터 ──
    _docx_add_header(doc, "주간업무보고서")
    _docx_add_footer(doc)

    # ── 섹션 1: 완료 업무 ──
    _docx_add_section_heading(doc, "1. 완료 업무")
    # Jinja2 for 루프를 포함한 업무 표 생성
    _docx_add_jinja_table(
        doc,
        headers=["업무명", "담당자", "진행률"],
        row_template_lines=["{{task.업무명}}", "{{task.담당자}}", "{{task.진행률}}"],
        loop_var="task",
        item_var="완료업무",
    )

    # ── 섹션 2: 진행 업무 ──
    _docx_add_section_heading(doc, "2. 진행 업무")
    _docx_add_jinja_table(
        doc,
        headers=["업무명", "담당자", "진행률"],
        row_template_lines=["{{task.업무명}}", "{{task.담당자}}", "{{task.진행률}}"],
        loop_var="task",
        item_var="진행업무",
    )

    # ── 섹션 3: 특이사항 ──
    _docx_add_section_heading(doc, "3. 특이사항")
    _docx_add_content_block(doc, "{{특이사항}}")

    # ── 저장 ──
    path = OUTPUT_DIR / "weekly_report.docx"
    doc.save(str(path))
    print(f"       → {path.name} 저장 완료 (커버 + 3섹션)")


# ═══════════════════════════════════════════════════════════════════════════════
#  템플릿 2: 사업제안서 (PPTX)
# ═══════════════════════════════════════════════════════════════════════════════

def create_business_proposal():
    """
    사업제안서 PPTX 템플릿을 생성합니다.

    슬라이드 구성:
        1. 표지: {{사업명}}, {{제안일}}, {{회사명}}
        2. 목차 슬라이드
        3. 사업 개요: {{사업개요}}
        4. 추진 전략: {% for item in 추진전략 %} ... {% endfor %}
        5. 기대 효과 & 예산: {{기대효과}}, {{총예산}}
        6. 감사합니다 (마침 슬라이드)

    Jinja2 변수:
        - {{사업명}}, {{제안일}}, {{회사명}}, {{사업개요}}
        - {% for item in 추진전략 %} → {{item}}
        - {{기대효과}}, {{총예산}}
    """
    print("  [2/5] 사업제안서 (PPTX) 생성 중...")

    prs = _new_pptx()

    # ── 슬라이드 1: 표지 ──
    _pptx_add_title_slide(
        prs,
        title="{{사업명}}",
        subtitle="사업제안서  |  {{회사명}}",
        date_text="{{제안일}}",
    )

    # ── 슬라이드 2: 목차 ──
    _pptx_add_content_slide(prs, "목차", [
        ("목차", [
            "1. 사업 개요",
            "2. 추진 전략",
            "3. 기대 효과",
            "4. 예산",
        ]),
    ])

    # ── 슬라이드 3: 사업 개요 ──
    _pptx_add_content_slide(prs, "1. 사업 개요", [
        ("사업 개요", "{{사업개요}}"),
    ])

    # ── 슬라이드 4: 추진 전략 (Jinja2 for 루프) ──
    _pptx_add_content_slide(prs, "2. 추진 전략", [
        ("추진 전략", [
            "{% for item in 추진전략 %}",
            "  - {{item}}",
            "{% endfor %}",
        ]),
    ])

    # ── 슬라이드 5: 기대 효과 & 예산 ──
    _pptx_add_content_slide(prs, "3. 기대 효과 & 예산", [
        ("기대 효과", "{{기대효과}}"),
        ("총 예산", "{{총예산}}"),
    ])

    # ── 슬라이드 6: 마침 (감사합니다) ──
    _pptx_add_closing_slide(prs)

    # ── 저장 ──
    path = OUTPUT_DIR / "business_proposal.pptx"
    prs.save(str(path))
    print(f"       → {path.name} 저장 완료 ({len(prs.slides)}슬라이드)")


# ═══════════════════════════════════════════════════════════════════════════════
#  템플릿 3: 회의록 (DOCX)
# ═══════════════════════════════════════════════════════════════════════════════

def create_meeting_minutes():
    """
    회의록 DOCX 템플릿을 생성합니다.
    커버 페이지 대신 상단에 회의 정보 테이블을 배치하는 특수 구조입니다.

    구조:
        [회의 정보]
            - "회 의 록" 대제목
            - 정보 테이블: 회의제목, 일시, 장소
            - 참석자: Jinja2 for 루프

        [본문]
            - 안건: Jinja2 for 루프 (제목 + 내용)
            - 결정사항: Jinja2 for 루프
            - 차기 회의: 단일 텍스트 블록

    Jinja2 변수:
        - {{회의제목}}, {{일시}}, {{장소}}
        - {% for person in 참석자 %} → {{person}}
        - {% for item in 안건 %} → {{item.제목}}, {{item.내용}}
        - {% for item in 결정사항 %} → {{item}}
        - {{차기회의}}
    """
    print("  [3/5] 회의록 (DOCX) 생성 중...")

    doc = _docx_setup()

    # ── 헤더와 푸터 ──
    _docx_add_header(doc, "회의록")
    _docx_add_footer(doc)

    # ── "회 의 록" 대제목 (가운데 정렬, 큰 글씨) ──
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    _docx_run(p, "회 의 록", size=22, bold=True, color=HEADER_BAR)
    _docx_add_thin_line(doc, HEADER_BAR)  # 제목 아래 구분선

    # ── 회의 기본 정보 테이블 ──
    _docx_add_info_table(doc, [
        ("회의 제목", "{{회의제목}}"),
        ("일    시", "{{일시}}"),
        ("장    소", "{{장소}}"),
    ])

    # ── 참석자 섹션 (Jinja2 for 루프) ──
    _docx_add_section_heading(doc, "참석자")
    _docx_add_content_block(doc, [
        "{% for person in 참석자 %}",
        "  - {{person}}",
        "{% endfor %}",
    ])

    # ── 안건 섹션 (Jinja2 for 루프 — 제목과 내용 쌍) ──
    _docx_add_section_heading(doc, "안건")
    _docx_add_content_block(doc, [
        "{% for item in 안건 %}",
        "■ {{item.제목}}",
        "  {{item.내용}}",
        "",
        "{% endfor %}",
    ])

    # ── 결정사항 섹션 (Jinja2 for 루프) ──
    _docx_add_section_heading(doc, "결정사항")
    _docx_add_content_block(doc, [
        "{% for item in 결정사항 %}",
        "  - {{item}}",
        "{% endfor %}",
    ])

    # ── 차기 회의 섹션 ──
    _docx_add_section_heading(doc, "차기 회의")
    _docx_add_content_block(doc, "{{차기회의}}")

    # ── 저장 ──
    path = OUTPUT_DIR / "meeting_minutes.docx"
    doc.save(str(path))
    print(f"       → {path.name} 저장 완료 (회의록, 4섹션)")


# ═══════════════════════════════════════════════════════════════════════════════
#  템플릿 4: 기술검토보고서 (DOCX)
# ═══════════════════════════════════════════════════════════════════════════════

def create_tech_review():
    """
    기술검토보고서 DOCX 템플릿을 생성합니다.

    구조:
        [커버 페이지]
            - 제목: "기술검토보고서"
            - 프로젝트명, 검토일, 검토자 (Jinja2 변수)

        [본문]
            - 섹션 1: 기술 개요 — 단일 텍스트 블록
            - 섹션 2: 검토 항목 — Jinja2 for 루프 표 (항목명, 결과, 비고)
            - 섹션 3: 결론 — 단일 텍스트 블록
            - 섹션 4: 권고사항 — 단일 텍스트 블록

    Jinja2 변수:
        - {{프로젝트명}}, {{검토일}}, {{검토자}}
        - {{기술개요}}
        - {% for item in 검토항목 %} → {{item.항목명}}, {{item.결과}}, {{item.비고}}
        - {{결론}}, {{권고사항}}
    """
    print("  [4/5] 기술검토보고서 (DOCX) 생성 중...")

    doc = _docx_setup()

    # ── 커버 페이지 ──
    _docx_add_cover_page(
        doc,
        title="기술검토보고서",
        subtitle_lines=[
            "프로젝트: {{프로젝트명}}",
            "검토일: {{검토일}}",
            "검토자: {{검토자}}",
        ],
    )

    # ── 헤더와 푸터 ──
    _docx_add_header(doc, "기술검토보고서")
    _docx_add_footer(doc)

    # ── 섹션 1: 기술 개요 ──
    _docx_add_section_heading(doc, "1. 기술 개요")
    _docx_add_content_block(doc, "{{기술개요}}")

    # ── 섹션 2: 검토 항목 (Jinja2 for 루프 표) ──
    _docx_add_section_heading(doc, "2. 검토 항목")
    _docx_add_jinja_table(
        doc,
        headers=["항목명", "결과", "비고"],
        row_template_lines=["{{item.항목명}}", "{{item.결과}}", "{{item.비고}}"],
        loop_var="item",
        item_var="검토항목",
    )

    # ── 섹션 3: 결론 ──
    _docx_add_section_heading(doc, "3. 결론")
    _docx_add_content_block(doc, "{{결론}}")

    # ── 섹션 4: 권고사항 ──
    _docx_add_section_heading(doc, "4. 권고사항")
    _docx_add_content_block(doc, "{{권고사항}}")

    # ── 저장 ──
    path = OUTPUT_DIR / "tech_review.docx"
    doc.save(str(path))
    print(f"       → {path.name} 저장 완료 (커버 + 4섹션)")


# ═══════════════════════════════════════════════════════════════════════════════
#  템플릿 5: 프로젝트현황보고 (PPTX)
# ═══════════════════════════════════════════════════════════════════════════════

def create_project_status():
    """
    프로젝트현황보고 PPTX 템플릿을 생성합니다.

    슬라이드 구성:
        1. 표지: {{프로젝트명}}, {{보고일}}, {{PM}}
        2. 진행 현황: {{진행률}} 텍스트
        3. 마일스톤: {% for milestone in 마일스톤 %} ... {% endfor %}
        4. 위험요소: {% for risk in 위험요소 %} ... {% endfor %}
        5. 감사합니다 (마침 슬라이드)

    Jinja2 변수:
        - {{프로젝트명}}, {{보고일}}, {{PM}}, {{진행률}}
        - {% for milestone in 마일스톤 %} → {{milestone}}
        - {% for risk in 위험요소 %} → {{risk}}
    """
    print("  [5/5] 프로젝트현황보고 (PPTX) 생성 중...")

    prs = _new_pptx()

    # ── 슬라이드 1: 표지 ──
    _pptx_add_title_slide(
        prs,
        title="{{프로젝트명}}",
        subtitle="프로젝트현황보고  |  PM: {{PM}}",
        date_text="{{보고일}}",
    )

    # ── 슬라이드 2: 진행 현황 ──
    _pptx_add_content_slide(prs, "1. 진행 현황", [
        ("전체 진행률", "{{진행률}}"),
    ])

    # ── 슬라이드 3: 마일스톤 (Jinja2 for 루프) ──
    _pptx_add_content_slide(prs, "2. 마일스톤", [
        ("마일스톤 목록", [
            "{% for milestone in 마일스톤 %}",
            "  - {{milestone}}",
            "{% endfor %}",
        ]),
    ])

    # ── 슬라이드 4: 위험요소 (Jinja2 for 루프) ──
    _pptx_add_content_slide(prs, "3. 위험요소", [
        ("위험요소 목록", [
            "{% for risk in 위험요소 %}",
            "  - {{risk}}",
            "{% endfor %}",
        ]),
    ])

    # ── 슬라이드 5: 마침 (감사합니다) ──
    _pptx_add_closing_slide(prs)

    # ── 저장 ──
    path = OUTPUT_DIR / "project_status.pptx"
    prs.save(str(path))
    print(f"       → {path.name} 저장 완료 ({len(prs.slides)}슬라이드)")


# ═══════════════════════════════════════════════════════════════════════════════
#  API 등록 함수 — 생성된 템플릿을 서버에 업로드
# ═══════════════════════════════════════════════════════════════════════════════

# 등록할 템플릿 메타데이터 목록
# 각 항목: (파일명, 템플릿유형, 출력형식, 이름, 설명)
TEMPLATE_METADATA = [
    (
        "weekly_report.docx",
        "주간업무보고서",       # template_type: 템플릿 유형
        "docx",                # output_format: 출력 파일 형식
        "주간업무보고서",       # name: 템플릿 이름
        "부서별 주간 업무 현황을 정리하는 보고서 템플릿입니다. "
        "완료 업무, 진행 업무, 특이사항을 표 형태로 작성합니다.",
    ),
    (
        "business_proposal.pptx",
        "사업제안서",
        "pptx",
        "사업제안서",
        "신규 사업을 제안할 때 사용하는 프레젠테이션 템플릿입니다. "
        "사업 개요, 추진 전략, 기대 효과, 예산을 포함합니다.",
    ),
    (
        "meeting_minutes.docx",
        "회의록",
        "docx",
        "회의록",
        "회의 내용을 기록하는 템플릿입니다. "
        "참석자, 안건, 결정사항, 차기 회의 일정을 정리합니다.",
    ),
    (
        "tech_review.docx",
        "기술검토보고서",
        "docx",
        "기술검토보고서",
        "기술 검토 결과를 보고하는 템플릿입니다. "
        "검토 항목별 결과를 표로 정리하고, 결론과 권고사항을 기술합니다.",
    ),
    (
        "project_status.pptx",
        "프로젝트현황보고",
        "pptx",
        "프로젝트현황보고",
        "프로젝트 진행 현황을 보고하는 프레젠테이션 템플릿입니다. "
        "진행률, 마일스톤, 위험요소를 포함합니다.",
    ),
]


def register_templates_via_api(base_url="http://localhost:8040/api/v1"):
    """
    생성된 Jinja2 템플릿 파일들을 API 서버에 업로드하여 등록합니다.

    처리 순서:
        1. admin 계정으로 로그인하여 JWT 토큰을 획득합니다.
        2. 각 템플릿 파일을 POST /templates/upload 엔드포인트로 업로드합니다.
        3. 업로드 결과(성공/실패)를 출력합니다.

    매개변수:
        base_url: API 서버의 기본 URL (기본값: "http://localhost:8040/api/v1")
    """
    # requests 라이브러리가 설치되어 있는지 확인
    try:
        import requests
    except ImportError:
        print("\n[오류] requests 라이브러리가 필요합니다.")
        print("       설치: pip install requests")
        return

    print(f"\n{'=' * 60}")
    print("API 서버에 템플릿 등록 시작")
    print(f"서버 URL: {base_url}")
    print(f"{'=' * 60}")

    # ── 1단계: 로그인하여 JWT 토큰 획득 ──
    print("\n[1/2] 로그인 중...")
    login_url = f"{base_url}/auth/login"
    login_data = {
        "username": "admin",       # 관리자 계정 (CLAUDE.md 참조)
        "password": "admin123!",   # 관리자 비밀번호
    }

    try:
        resp = requests.post(login_url, json=login_data, timeout=10)
        resp.raise_for_status()   # HTTP 오류가 있으면 예외 발생
    except requests.RequestException as e:
        print(f"       로그인 실패: {e}")
        print("       API 서버가 실행 중인지 확인해주세요.")
        return

    # 응답에서 액세스 토큰 추출
    token_data = resp.json()
    access_token = token_data.get("access_token")
    if not access_token:
        print(f"       로그인 응답에 access_token이 없습니다: {token_data}")
        return

    print("       로그인 성공!")

    # 이후 요청에 사용할 인증 헤더
    headers = {"Authorization": f"Bearer {access_token}"}

    # ── 2단계: 각 템플릿 파일 업로드 ──
    print("\n[2/2] 템플릿 업로드 중...")
    upload_url = f"{base_url}/templates/upload"

    success_count = 0
    fail_count = 0

    for filename, template_type, output_format, name, description in TEMPLATE_METADATA:
        file_path = OUTPUT_DIR / filename

        # 파일이 존재하는지 확인
        if not file_path.exists():
            print(f"  [건너뜀] {filename} — 파일이 존재하지 않습니다")
            fail_count += 1
            continue

        # multipart/form-data로 파일 업로드
        try:
            with open(file_path, "rb") as f:
                files = {"file": (filename, f)}
                params = {
                    "template_type": template_type,
                    "tone": "formal",          # 공식 문서 어조
                    "output_format": output_format,
                    "name": name,
                    "description": description,
                }
                resp = requests.post(
                    upload_url,
                    headers=headers,
                    files=files,
                    params=params,
                    timeout=30,
                )

            if resp.status_code == 201:
                # 성공 시 응답에서 추출된 변수 수 표시
                result = resp.json()
                var_count = len(result.get("variables", []))
                print(f"  [성공] {filename} — 변수 {var_count}개 자동 추출됨")
                success_count += 1
            else:
                print(f"  [실패] {filename} — HTTP {resp.status_code}: {resp.text[:200]}")
                fail_count += 1

        except requests.RequestException as e:
            print(f"  [실패] {filename} — {e}")
            fail_count += 1

    # ── 결과 요약 ──
    print(f"\n{'=' * 60}")
    print(f"등록 완료: 성공 {success_count}개, 실패 {fail_count}개")
    print(f"{'=' * 60}")


# ═══════════════════════════════════════════════════════════════════════════════
#  메인 실행 블록
# ═══════════════════════════════════════════════════════════════════════════════

def create_all_templates():
    """
    5개의 Jinja2 샘플 템플릿을 모두 생성합니다.
    출력 디렉토리(sample_templates/jinja2/)가 없으면 자동으로 생성합니다.
    """
    print(f"\n{'=' * 60}")
    print("Jinja2 샘플 템플릿 생성기")
    print(f"출력 디렉토리: {OUTPUT_DIR}")
    print(f"{'=' * 60}\n")

    # 출력 디렉토리가 없으면 생성
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    # ── 5개 템플릿 순서대로 생성 ──
    create_weekly_report()           # 1. 주간업무보고서 (DOCX)
    create_business_proposal()       # 2. 사업제안서 (PPTX)
    create_meeting_minutes()         # 3. 회의록 (DOCX)
    create_tech_review()             # 4. 기술검토보고서 (DOCX)
    create_project_status()          # 5. 프로젝트현황보고 (PPTX)

    print(f"\n{'=' * 60}")
    print(f"총 5개 템플릿 생성 완료!")
    print(f"출력 위치: {OUTPUT_DIR}")
    print(f"{'=' * 60}")


if __name__ == "__main__":
    # ── 명령줄 인자 파서 설정 ──
    parser = argparse.ArgumentParser(
        description="Jinja2 샘플 템플릿을 생성하고, 선택적으로 API 서버에 등록합니다."
    )
    parser.add_argument(
        "--register",
        action="store_true",
        help="템플릿 생성 후 API 서버에 자동 등록합니다",
    )
    parser.add_argument(
        "--base-url",
        default="http://localhost:8040/api/v1",
        help="API 서버 기본 URL (기본값: http://localhost:8040/api/v1)",
    )

    args = parser.parse_args()

    # 1) 템플릿 파일 생성
    create_all_templates()

    # 2) --register 옵션이 있으면 API 서버에 등록
    if args.register:
        register_templates_via_api(base_url=args.base_url)
