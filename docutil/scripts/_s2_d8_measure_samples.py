"""Phase 4 S2 D8-b — PPTX A/B 비교 자동 메트릭 수집 스크립트.

`_s2_d8_generate_ab_samples.py` 가 생성한 5 샘플을 python-pptx 로 재오픈해
다음 메트릭을 측정하고 `docs/s2_d8_samples/metrics.json` 에 저장한다.

측정 항목:
    1. 슬라이드 수
    2. 슬라이드별 실제 shape 수 (컴포넌트 스키마 수 대비 "누락" 감지)
    3. 파일 크기(KB)
    4. 빌드 재현 시간(ms) — 동일 schema 로 다시 빌드 후 perf_counter 로 측정.
       (파일 사이즈/생성 시간은 metadata.json 에도 있으나 "측정 독립성" 을
        위해 여기서 다시 수행.)
    5. IDINO 팔레트 일치율
       — 슬라이드 내 모든 shape/font 에서 수집된 RGB 의 분포를
         `KNOWN_IDINO_COLORS` 집합(constants.py 색 + 차트 팔레트 합집합) 과
         비교해 "팔레트 내 비율" 을 계산.
    6. 사용된 레이아웃 이름 (slide.slide_layout.name)

출력:
    docs/s2_d8_samples/metrics.json

주의:
    - 본 스크립트는 **읽기 전용** — 샘플 파일을 수정하지 않는다.
    - Registry 경유로 빌더를 획득 (P1) 하여 build 재현 시간 측정.
    - 파이썬 파일명 / 디렉토리 경로 등의 "상수" 는 한 곳에 모아둔다.
"""

from __future__ import annotations

import asyncio
import json
import sys
import time
from collections import Counter
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Iterable
from uuid import UUID

# ---------------------------------------------------------------------------
# 경로 — backend/ 를 sys.path 에 먼저 추가.
# ---------------------------------------------------------------------------
PROJECT_ROOT = Path(r"D:\workspace\document_utilization")
BACKEND_ROOT = PROJECT_ROOT / "backend"
SAMPLES_DIR = PROJECT_ROOT / "docs" / "s2_d8_samples"
METADATA_PATH = SAMPLES_DIR / "metadata.json"
METRICS_OUT_PATH = SAMPLES_DIR / "metrics.json"

sys.path.insert(0, str(BACKEND_ROOT))

# ruff: noqa: E402 — 의도된 지연 임포트.
from pptx import Presentation  # noqa: E402
from pptx.util import Emu  # noqa: E402

from app.integrations.document_builders import BuilderRegistry  # noqa: E402
from app.integrations.document_builders.pptx import register_pptx_builder  # noqa: E402
from app.integrations.document_builders.pptx.constants import (  # noqa: E402
    CHART_SERIES_PALETTE,
    IDINO_ACCENT,
    IDINO_HEADER_NAVY,
    IDINO_PALETTE,
    IDINO_PRIMARY,
    IDINO_TEXT,
    IDINO_TEXT_MUTED,
    IDINO_WHITE,
    SLIDE_HEIGHT_IN,
    SLIDE_WIDTH_IN,
)


# ---------------------------------------------------------------------------
# 알려진 IDINO 색상 집합 — constants.IDINO_PALETTE ∪ CHART_SERIES_PALETTE.
# "일치율" 의 기준. 대소문자 차이 흡수를 위해 upper() 로 정규화.
# ---------------------------------------------------------------------------


def _normalize_hex(hex_str: str) -> str:
    """'#0a4fc2' / '#0A4FC2' / '0A4FC2' 모두 '0A4FC2' 형식으로 정규화."""
    s = hex_str.strip()
    if s.startswith("#"):
        s = s[1:]
    return s.upper()


KNOWN_IDINO_COLORS: frozenset[str] = frozenset(
    {
        _normalize_hex(IDINO_PRIMARY),
        _normalize_hex(IDINO_ACCENT),
        _normalize_hex(IDINO_TEXT),
        _normalize_hex(IDINO_TEXT_MUTED),
        _normalize_hex(IDINO_WHITE),
        _normalize_hex(IDINO_HEADER_NAVY),
    }
    | {_normalize_hex(v) for v in IDINO_PALETTE.values()}
    | {_normalize_hex(c) for c in CHART_SERIES_PALETTE}
)


# ---------------------------------------------------------------------------
# Helpers — shape/text 트리에서 색상 수집.
# ---------------------------------------------------------------------------


def _safe_hex_from_rgb(rgb_obj: Any) -> str | None:
    """python-pptx `RGBColor` → 6 자리 hex 대문자로 변환. 실패 시 None.

    `RGBColor` 는 str() 시 "0A4FC2" 같은 6 자리 hex 를 반환하지만, type=MSO_THEME_COLOR
    같은 스킴 기반 색은 rgb 속성 접근이 예외를 던진다.
    """
    try:
        return str(rgb_obj).upper()
    except Exception:
        return None


def _collect_font_color(font: Any, sink: list[str]) -> None:
    """font.color.rgb 를 안전하게 추출해 sink 에 추가.

    font 는 python-pptx 내부 타입. color 접근 및 rgb 접근 모두 MSO_THEME_COLOR
    또는 상속 미해결 시 예외를 던질 수 있어 광범위한 방어가 필요하다.
    """
    try:
        rgb_obj = font.color.rgb
    except Exception:
        return
    if rgb_obj is None:
        return
    hex_val = _safe_hex_from_rgb(rgb_obj)
    if hex_val:
        sink.append(hex_val)


def _collect_text_frame_colors(text_frame: Any, sink: list[str]) -> None:
    """text_frame 의 모든 paragraph/run 에서 폰트 색을 수집."""
    try:
        paragraphs = list(text_frame.paragraphs)
    except Exception:
        return
    for paragraph in paragraphs:
        # paragraph.font 는 run 이 없을 때의 기본 폰트.
        _collect_font_color(paragraph.font, sink)
        try:
            runs = list(paragraph.runs)
        except Exception:
            runs = []
        for run in runs:
            _collect_font_color(run.font, sink)


def _collect_shape_fill_color(shape: Any, sink: list[str]) -> None:
    """shape.fill.fore_color.rgb 를 안전하게 추출.

    shape 가 fill 을 가지지 않는 타입이거나 fill.type 이 solid 가 아닐 때는
    조용히 무시.
    """
    try:
        fill = shape.fill
    except Exception:
        return
    try:
        # fill.type == MSO_FILL.SOLID 일 때만 fore_color.rgb 접근이 안전.
        if str(fill.type) != "MSO_FILL_TYPE.SOLID (1)":
            # 비-SOLID(배경 투명, 패턴, 그라데이션, 미설정 등)는 색 정보 없음으로 간주.
            # (문자열 비교는 python-pptx 버전에 따라 enum repr 이 달라질 수 있어
            #  실제 색 접근을 try 로 한 번 더 감싼다.)
            pass
    except Exception:
        return
    try:
        rgb_obj = fill.fore_color.rgb
    except Exception:
        return
    hex_val = _safe_hex_from_rgb(rgb_obj)
    if hex_val:
        sink.append(hex_val)


def _iter_shapes_deep(shapes: Iterable[Any]) -> Iterable[Any]:
    """shape tree 를 평탄화해 순회 (GroupShape 내부도 재귀적으로 파고든다).

    python-pptx `GroupShape` 는 `.shapes` 속성으로 자식 shape 를 노출한다. 본
    샘플들은 group 을 쓰지 않지만, 향후 차트/표 내부의 sub-shape 도 안전하게
    순회할 수 있게 방어적으로 재귀.
    """
    for shape in shapes:
        yield shape
        sub = getattr(shape, "shapes", None)
        if sub is not None and sub is not shape:
            try:
                yield from _iter_shapes_deep(sub)
            except Exception:
                continue


def _collect_colors_from_slide(slide: Any) -> list[str]:
    """한 슬라이드 내 모든 shape/text/table 의 색상을 평탄한 리스트로 수집."""
    sink: list[str] = []
    for shape in _iter_shapes_deep(slide.shapes):
        # 1) shape 자체의 solid fill.
        _collect_shape_fill_color(shape, sink)

        # 2) text frame (shape 가 텍스트를 담고 있을 때).
        tf = getattr(shape, "text_frame", None)
        if tf is not None:
            _collect_text_frame_colors(tf, sink)

        # 3) table 셀 내부 — table 은 shape.has_table 로 판별.
        has_table = getattr(shape, "has_table", False)
        if has_table:
            try:
                table = shape.table
            except Exception:
                table = None
            if table is not None:
                for row in table.rows:
                    for cell in row.cells:
                        # 셀 배경.
                        try:
                            rgb_obj = cell.fill.fore_color.rgb
                        except Exception:
                            rgb_obj = None
                        if rgb_obj is not None:
                            hex_val = _safe_hex_from_rgb(rgb_obj)
                            if hex_val:
                                sink.append(hex_val)
                        # 셀 텍스트.
                        cell_tf = getattr(cell, "text_frame", None)
                        if cell_tf is not None:
                            _collect_text_frame_colors(cell_tf, sink)

        # 4) chart — has_chart 로 판별. 시리즈 색을 추출하기는 내부 XML 이 복잡해
        #    본 단계에서는 chart 자체는 건너뛰고 테이블·텍스트 수집으로 충분히
        #    팔레트 일치 샘플 확보. (chart 시리즈 색은 PPT 내부에 인덱스 기반으로
        #    저장될 수 있어 RGB 직접 추출이 항상 가능하진 않다.)
    return sink


def _shape_count(slide: Any) -> int:
    """슬라이드의 top-level shape 개수."""
    return len(list(slide.shapes))


def _compute_palette_match_ratio(colors: list[str]) -> tuple[float, dict[str, int]]:
    """수집된 색상 리스트에서 IDINO 팔레트 내 색상 비율 + 색상 빈도 반환.

    Returns:
        (ratio ∈ [0, 1], histogram). histogram 은 각 hex → 카운트.
    """
    if not colors:
        return 0.0, {}
    hist = Counter(colors)
    inside = sum(cnt for hex_val, cnt in hist.items() if hex_val in KNOWN_IDINO_COLORS)
    total = sum(hist.values())
    return (inside / total if total else 0.0), dict(hist)


# ---------------------------------------------------------------------------
# 빌드 재현 — metadata.json 의 설정값을 재구성하지 않고, 생성 스크립트를
# 다시 import 해 같은 함수를 호출한다. 이래야 "의도된 빌더 파이프라인" 을
# 완전히 동일하게 재현할 수 있다.
# ---------------------------------------------------------------------------


def _import_build_functions() -> dict[str, Any]:
    """생성 스크립트에서 build_sample_* 함수들을 import 해 반환.

    _s2_d8_generate_ab_samples.py 가 같은 scripts/ 디렉토리에 있어 경로 주입으로
    가져온다. 파일명이 '_' 로 시작하는 private 스크립트라 import 관습에선
    벗어나지만, 메트릭 스크립트 전용의 내부 재사용이라 허용.
    """
    sys.path.insert(0, str(PROJECT_ROOT / "scripts"))
    # ruff: noqa: I001
    from _s2_d8_generate_ab_samples import (
        SAMPLE_A_ID,
        SAMPLE_B_ID,
        SAMPLE_C_ID,
        SAMPLE_D_ID,
        SAMPLE_E_ID,
        build_sample_a,
        build_sample_b,
        build_sample_c,
        build_sample_d,
        build_sample_e,
    )

    return {
        SAMPLE_A_ID: build_sample_a,
        SAMPLE_B_ID: build_sample_b,
        SAMPLE_C_ID: build_sample_c,
        SAMPLE_D_ID: build_sample_d,
        SAMPLE_E_ID: build_sample_e,
    }


async def _rebuild_time_ms(builder: Any, schema_factory: Any) -> float:
    """동일 schema 를 즉시 재빌드해 시간 측정 (ms)."""
    schema = schema_factory()
    start = time.perf_counter()
    _ = await builder.build(schema)
    return round((time.perf_counter() - start) * 1000.0, 2)


# ---------------------------------------------------------------------------
# 메인 측정 로직
# ---------------------------------------------------------------------------


def _measure_sample(sample_path: Path, expected_component_counts: dict[str, int]) -> dict[str, Any]:
    """단일 샘플 파일을 열어 메트릭을 측정.

    Args:
        sample_path: .pptx 파일 경로.
        expected_component_counts: metadata.json 에 기록된 스키마 컴포넌트 수.

    Returns:
        단일 샘플 메트릭 dict.
    """
    prs = Presentation(str(sample_path))
    slide_count = len(prs.slides)

    # 슬라이드 크기(인치 확인) — 16:9 여부 검증.
    slide_width_in = float(Emu(prs.slide_width).inches) if prs.slide_width else 0.0
    slide_height_in = float(Emu(prs.slide_height).inches) if prs.slide_height else 0.0

    per_slide: list[dict[str, Any]] = []
    all_colors: list[str] = []
    for idx, slide in enumerate(prs.slides):
        shape_cnt = _shape_count(slide)
        colors = _collect_colors_from_slide(slide)
        all_colors.extend(colors)
        per_slide.append(
            {
                "slide_index": idx,
                "layout_name": slide.slide_layout.name,
                "shape_count": shape_cnt,
                "color_sample_count": len(colors),
            }
        )

    ratio, histogram = _compute_palette_match_ratio(all_colors)

    # 컴포넌트 누락 감지 — 매우 거친 추정:
    #   스키마 컴포넌트 수 합 vs 슬라이드별 top-level shape 수 합.
    #   KPI 1 개는 3~4 shape (label/value/delta [+bg]) 을 만들고, DataTable 은 1 shape,
    #   Image placeholder 는 2 shape (bg + label) 등으로 1:1 대응이 아니다.
    #   따라서 "shape 합 >= 컴포넌트 합" 이면 누락 없음으로 간주한다 (가장 보수적).
    schema_component_total = sum(expected_component_counts.values())
    total_shapes = sum(s["shape_count"] for s in per_slide)
    components_missing = total_shapes < schema_component_total

    file_size_kb = round(sample_path.stat().st_size / 1024.0, 2)

    return {
        "file_name": sample_path.name,
        "file_size_kb": file_size_kb,
        "slide_count": slide_count,
        "slide_width_in": round(slide_width_in, 3),
        "slide_height_in": round(slide_height_in, 3),
        "is_16_9": (
            abs(slide_width_in - SLIDE_WIDTH_IN) < 0.01
            and abs(slide_height_in - SLIDE_HEIGHT_IN) < 0.01
        ),
        "per_slide": per_slide,
        "schema_component_counts": expected_component_counts,
        "schema_component_total": schema_component_total,
        "total_top_level_shapes": total_shapes,
        "components_likely_missing": components_missing,
        "idino_palette_match_ratio": round(ratio, 4),
        "color_samples_total": len(all_colors),
        "color_histogram_top10": dict(
            sorted(histogram.items(), key=lambda kv: kv[1], reverse=True)[:10]
        ),
    }


async def run() -> int:
    if not METADATA_PATH.exists():
        print(
            f"[ERROR] {METADATA_PATH} 가 없습니다. 먼저 _s2_d8_generate_ab_samples.py 를 실행하세요.",
        )
        return 1

    metadata_payload = json.loads(METADATA_PATH.read_text(encoding="utf-8"))
    register_pptx_builder()
    builder = BuilderRegistry.get("pptx")
    factories = _import_build_functions()

    metrics_per_sample: list[dict[str, Any]] = []
    for sample_meta in metadata_payload["samples"]:
        sample_id = sample_meta["sample_id"]
        sample_path = SAMPLES_DIR / sample_meta["file_name"]
        if not sample_path.exists():
            print(f"[WARN] 샘플 파일 없음 — skip: {sample_path}")
            continue
        entry = _measure_sample(sample_path, sample_meta["component_counts"])
        entry["sample_id"] = sample_id
        # 재현 빌드 시간 측정.
        factory = factories.get(sample_id)
        if factory is not None:
            entry["rebuild_time_ms"] = await _rebuild_time_ms(builder, factory)
        else:
            entry["rebuild_time_ms"] = None
        metrics_per_sample.append(entry)
        print(
            f"[OK] {sample_id:20s} | slides={entry['slide_count']} | "
            f"shapes={entry['total_top_level_shapes']} | "
            f"palette={entry['idino_palette_match_ratio']:.2%} | "
            f"size={entry['file_size_kb']} KB"
        )

    payload: dict[str, Any] = {
        "measured_at_utc": datetime.now(tz=timezone.utc).isoformat(),
        "known_idino_colors": sorted(KNOWN_IDINO_COLORS),
        "expected_slide_size_in": [SLIDE_WIDTH_IN, SLIDE_HEIGHT_IN],
        "samples": metrics_per_sample,
        "summary": {
            "sample_count": len(metrics_per_sample),
            "all_16_9": all(e["is_16_9"] for e in metrics_per_sample),
            "any_component_missing": any(e["components_likely_missing"] for e in metrics_per_sample),
            "avg_palette_match_ratio": round(
                sum(e["idino_palette_match_ratio"] for e in metrics_per_sample)
                / max(1, len(metrics_per_sample)),
                4,
            ),
            "avg_rebuild_time_ms": round(
                sum(e.get("rebuild_time_ms") or 0.0 for e in metrics_per_sample)
                / max(1, len(metrics_per_sample)),
                2,
            ),
        },
    }
    METRICS_OUT_PATH.write_text(
        json.dumps(payload, ensure_ascii=False, indent=2, default=_json_default),
        encoding="utf-8",
    )
    print(f"[DONE] metrics.json 저장 → {METRICS_OUT_PATH}")
    return 0


def _json_default(obj: Any) -> Any:
    """json.dumps 에서 UUID 등 기본 직렬화되지 않는 객체를 처리."""
    if isinstance(obj, UUID):
        return str(obj)
    if isinstance(obj, datetime):
        return obj.isoformat()
    raise TypeError(f"지원하지 않는 JSON 직렬화 타입: {type(obj).__name__}")


if __name__ == "__main__":
    raise SystemExit(asyncio.run(run()))
