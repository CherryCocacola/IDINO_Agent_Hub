"""
Generate professional DOCX and PPTX template files with {{variable}} placeholders.

These templates are designed for the DocUtil report generator, which fills
{{variable}} patterns using RAG-extracted content from source documents.

Usage:
    python scripts/create_sample_templates.py
"""

from __future__ import annotations

import os
from pathlib import Path

# ---------------------------------------------------------------------------
# DOCX imports
# ---------------------------------------------------------------------------
from docx import Document
from docx.shared import Inches, Pt, RGBColor, Cm, Emu
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.section import WD_ORIENT
from docx.oxml.ns import qn, nsdecls
from docx.oxml import parse_xml

# ---------------------------------------------------------------------------
# PPTX imports
# ---------------------------------------------------------------------------
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt, Emu as PptxEmu
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUTPUT_DIR = Path(__file__).resolve().parent.parent / "sample_templates"


# ===================================================================
# Shared helpers
# ===================================================================

def _set_cell_shading(cell, color_hex: str):
    """Set background shading on a table cell."""
    shading = parse_xml(
        f'<w:shd {nsdecls("w")} w:fill="{color_hex}" w:val="clear"/>'
    )
    cell._tc.get_or_add_tcPr().append(shading)


def _set_cell_border(cell, **kwargs):
    """Set borders on a cell. kwargs like top=("single","4","000000")."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    tcBorders = parse_xml(f'<w:tcBorders {nsdecls("w")}/>')
    for edge, (style, sz, color) in kwargs.items():
        element = parse_xml(
            f'<w:{edge} {nsdecls("w")} w:val="{style}" w:sz="{sz}" '
            f'w:space="0" w:color="{color}"/>'
        )
        tcBorders.append(element)
    tcPr.append(tcBorders)


def _add_styled_paragraph(doc, text, style_name=None, font_name="맑은 고딕",
                          font_size=None, bold=False, color=None,
                          alignment=None, space_before=None, space_after=None):
    """Add a paragraph with fine-grained formatting."""
    p = doc.add_paragraph()
    if style_name:
        p.style = doc.styles[style_name]
    if alignment is not None:
        p.alignment = alignment
    pf = p.paragraph_format
    if space_before is not None:
        pf.space_before = Pt(space_before)
    if space_after is not None:
        pf.space_after = Pt(space_after)

    run = p.add_run(text)
    run.font.name = font_name
    run._element.rPr.rFonts.set(qn("w:eastAsia"), font_name)
    if font_size:
        run.font.size = Pt(font_size)
    if bold:
        run.font.bold = True
    if color:
        run.font.color.rgb = RGBColor(*color)
    return p


def _configure_docx_styles(doc, primary_color, heading_font="맑은 고딕",
                            body_font="맑은 고딕"):
    """Configure heading and body styles with consistent branding."""
    style = doc.styles["Normal"]
    style.font.name = body_font
    style.font.size = Pt(10)
    style.element.rPr.rFonts.set(qn("w:eastAsia"), body_font)
    style.paragraph_format.space_after = Pt(6)
    style.paragraph_format.line_spacing = 1.15

    for level, (size, is_bold) in enumerate(
        [(18, True), (14, True), (12, True)], start=1
    ):
        heading_style = doc.styles[f"Heading {level}"]
        heading_style.font.name = heading_font
        heading_style.font.size = Pt(size)
        heading_style.font.bold = is_bold
        heading_style.font.color.rgb = RGBColor(*primary_color)
        heading_style.element.rPr.rFonts.set(qn("w:eastAsia"), heading_font)
        pf = heading_style.paragraph_format
        pf.space_before = Pt(18 if level == 1 else 12)
        pf.space_after = Pt(6)


def _set_page_margins(section, top=2.54, bottom=2.54, left=2.54, right=2.54):
    """Set page margins in cm."""
    section.top_margin = Cm(top)
    section.bottom_margin = Cm(bottom)
    section.left_margin = Cm(left)
    section.right_margin = Cm(right)


def _add_header_footer(section, header_text: str, primary_color):
    """Add header with text and footer with page number."""
    # Header
    header = section.header
    header.is_linked_to_previous = False
    hp = header.paragraphs[0]
    hp.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    run = hp.add_run(header_text)
    run.font.size = Pt(8)
    run.font.color.rgb = RGBColor(*primary_color)
    run.font.name = "맑은 고딕"

    # Add bottom border to header paragraph
    pPr = hp._p.get_or_add_pPr()
    pBdr = parse_xml(
        f'<w:pBdr {nsdecls("w")}>'
        f'  <w:bottom w:val="single" w:sz="4" w:space="1" '
        f'   w:color="{primary_color[0]:02X}{primary_color[1]:02X}{primary_color[2]:02X}"/>'
        f'</w:pBdr>'
    )
    pPr.append(pBdr)

    # Footer with page number
    footer = section.footer
    footer.is_linked_to_previous = False
    fp = footer.paragraphs[0]
    fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = fp.add_run()
    run.font.size = Pt(8)
    run.font.color.rgb = RGBColor(128, 128, 128)
    # Add page number field
    fldChar1 = parse_xml(f'<w:fldChar {nsdecls("w")} w:fldCharType="begin"/>')
    run._r.append(fldChar1)
    instrText = parse_xml(f'<w:instrText {nsdecls("w")} xml:space="preserve"> PAGE </w:instrText>')
    run._r.append(instrText)
    fldChar2 = parse_xml(f'<w:fldChar {nsdecls("w")} w:fldCharType="end"/>')
    run._r.append(fldChar2)


# ===================================================================
# PPTX helpers
# ===================================================================

def _add_pptx_shape(slide, left, top, width, height, fill_color=None,
                    line_color=None, line_width=None):
    """Add a rectangle shape to a slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.line.fill.background()  # no line by default
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = PptxRGBColor(*fill_color)
    if line_color:
        shape.line.color.rgb = PptxRGBColor(*line_color)
        shape.line.width = PptxPt(line_width or 1)
    return shape


def _add_pptx_textbox(slide, left, top, width, height, text, font_size=14,
                      font_color=(255, 255, 255), bold=False,
                      alignment=PP_ALIGN.LEFT, font_name="맑은 고딕",
                      anchor=MSO_ANCHOR.TOP):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = PptxPt(font_size)
    run.font.color.rgb = PptxRGBColor(*font_color)
    run.font.bold = bold
    run.font.name = font_name
    return txBox


def _set_slide_bg(slide, color):
    """Set solid background color on a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PptxRGBColor(*color)


def _add_slide_number(slide, color=(180, 180, 180)):
    """Add slide number to bottom-right."""
    txBox = slide.shapes.add_textbox(
        PptxInches(8.5), PptxInches(6.8), PptxInches(1.0), PptxInches(0.4)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    # Add slide number field via XML
    fld = parse_xml(
        '<a:fld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" type="slidenum">'
        '<a:rPr lang="ko-KR" dirty="0"/>'
        '<a:t>&lt;#&gt;</a:t>'
        '</a:fld>'
    )
    # Set a unique id
    import uuid
    fld.set("id", f"{{{uuid.uuid4()}}}")
    p._p.append(fld)


def _add_accent_bar(slide, left, top, width, height, color):
    """Add a thin colored accent bar."""
    shape = _add_pptx_shape(slide, left, top, width, height, fill_color=color)
    return shape


def _create_pptx_content_slide(prs, slide_number, title_text, body_placeholder,
                                bg_color, title_color, body_color,
                                accent_color, slide_width, slide_height):
    """Create a standard content slide for PPTX templates."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    _set_slide_bg(slide, bg_color)

    # Top accent bar
    _add_accent_bar(
        slide, PptxInches(0), PptxInches(0),
        slide_width, PptxPt(4), accent_color
    )

    # Slide title
    _add_pptx_textbox(
        slide, PptxInches(0.8), PptxInches(0.5),
        PptxInches(8.5), PptxInches(0.7),
        title_text, font_size=28, font_color=title_color,
        bold=True, alignment=PP_ALIGN.LEFT
    )

    # Divider line under title
    _add_accent_bar(
        slide, PptxInches(0.8), PptxInches(1.2),
        PptxInches(1.5), PptxPt(3), accent_color
    )

    # Body placeholder
    _add_pptx_textbox(
        slide, PptxInches(0.8), PptxInches(1.6),
        PptxInches(8.4), PptxInches(5.0),
        body_placeholder, font_size=16, font_color=body_color,
        bold=False, alignment=PP_ALIGN.LEFT
    )

    _add_slide_number(slide)
    return slide


# ===================================================================
# Template 1: 격식체 요약 보고서 (Formal Summary Report)
# ===================================================================

def create_report_summary_formal() -> Path:
    PRIMARY = (31, 78, 121)      # #1F4E79
    SECONDARY = (68, 114, 148)   # #447294
    LIGHT_BG = (234, 241, 248)   # #EAF1F8

    doc = Document()
    _configure_docx_styles(doc, PRIMARY)

    section = doc.sections[0]
    _set_page_margins(section, top=2.5, bottom=2.5, left=2.8, right=2.8)

    # -- Cover page --
    # Add vertical spacing before title
    for _ in range(6):
        doc.add_paragraph()

    # Top accent line
    p_line = doc.add_paragraph()
    p_line.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p_line.add_run("━" * 50)
    run.font.color.rgb = RGBColor(*PRIMARY)
    run.font.size = Pt(8)

    # Company logo placeholder
    p_logo = doc.add_paragraph()
    p_logo.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p_logo.add_run("[Company Logo]")
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(160, 160, 160)
    p_logo.paragraph_format.space_after = Pt(30)

    # Title
    p_title = doc.add_paragraph()
    p_title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p_title.add_run("{{제목}}")
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*PRIMARY)
    run.font.name = "맑은 고딕"
    run._element.rPr.rFonts.set(qn("w:eastAsia"), "맑은 고딕")
    p_title.paragraph_format.space_after = Pt(40)

    # Bottom accent line
    p_line2 = doc.add_paragraph()
    p_line2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p_line2.add_run("━" * 50)
    run.font.color.rgb = RGBColor(*PRIMARY)
    run.font.size = Pt(8)

    doc.add_paragraph()  # spacing

    # Info table on cover page
    info_table = doc.add_table(rows=2, cols=2)
    info_table.alignment = WD_TABLE_ALIGNMENT.CENTER
    info_data = [
        ("작 성 자", "{{작성자}}"),
        ("작 성 일", "{{작성일}}"),
    ]
    for i, (label, value) in enumerate(info_data):
        row = info_table.rows[i]
        cell_label = row.cells[0]
        cell_value = row.cells[1]
        cell_label.width = Cm(4)
        cell_value.width = Cm(8)

        _set_cell_shading(cell_label, "1F4E79")
        p = cell_label.paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(label)
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.name = "맑은 고딕"

        p2 = cell_value.paragraphs[0]
        p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run2 = p2.add_run(value)
        run2.font.size = Pt(10)
        run2.font.name = "맑은 고딕"

    # Page break after cover
    doc.add_page_break()

    # -- Header / Footer (applied to all sections) --
    _add_header_footer(section, "{{제목}}", PRIMARY)

    # -- Table of Contents placeholder --
    _add_styled_paragraph(
        doc, "목 차", font_size=16, bold=True,
        color=PRIMARY, alignment=WD_ALIGN_PARAGRAPH.CENTER,
        space_after=20
    )

    toc_items = [
        "1. 개요",
        "2. 주요 내용",
        "3. 분석 결과",
        "4. 결론 및 제언",
    ]
    for item in toc_items:
        p = doc.add_paragraph()
        run = p.add_run(item)
        run.font.size = Pt(11)
        run.font.name = "맑은 고딕"
        run._element.rPr.rFonts.set(qn("w:eastAsia"), "맑은 고딕")
        # Add dot leader and page placeholder
        run2 = p.add_run("  " + "·" * 60 + "  ")
        run2.font.size = Pt(8)
        run2.font.color.rgb = RGBColor(180, 180, 180)

    doc.add_page_break()

    # -- Section 1: 개요 --
    doc.add_heading("1. 개요", level=1)
    _add_styled_paragraph(doc, "{{개요}}", font_size=10.5)

    # -- Section 2: 주요 내용 --
    doc.add_heading("2. 주요 내용", level=1)
    _add_styled_paragraph(doc, "{{주요내용}}", font_size=10.5)

    # -- Section 3: 분석 결과 --
    doc.add_heading("3. 분석 결과", level=1)
    _add_styled_paragraph(doc, "{{분석결과}}", font_size=10.5)

    # -- Section 4: 결론 및 제언 --
    doc.add_heading("4. 결론 및 제언", level=1)
    _add_styled_paragraph(doc, "{{결론}}", font_size=10.5)

    path = OUTPUT_DIR / "report_summary_formal.docx"
    doc.save(str(path))
    return path


# ===================================================================
# Template 2: 비즈니스 분석 보고서 (Business Analysis Report)
# ===================================================================

def create_report_analysis_business() -> Path:
    PRIMARY = (46, 64, 87)       # #2E4057
    ACCENT = (4, 138, 129)       # #048A81
    LIGHT_BG = (240, 248, 247)   # #F0F8F7

    doc = Document()
    _configure_docx_styles(doc, PRIMARY)

    section = doc.sections[0]
    _set_page_margins(section, top=2.5, bottom=2.5, left=2.8, right=2.8)

    # -- Cover page --
    for _ in range(4):
        doc.add_paragraph()

    # Accent bar (simulated with colored text line)
    p_bar = doc.add_paragraph()
    p_bar.alignment = WD_ALIGN_PARAGRAPH.LEFT
    run = p_bar.add_run("█" * 12)
    run.font.color.rgb = RGBColor(*ACCENT)
    run.font.size = Pt(14)

    # Document type label
    _add_styled_paragraph(
        doc, "비즈니스 분석 보고서", font_size=12,
        color=ACCENT, bold=True, space_before=12, space_after=8
    )

    # Title
    p_title = doc.add_paragraph()
    run = p_title.add_run("{{프로젝트명}}")
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*PRIMARY)
    run.font.name = "맑은 고딕"
    run._element.rPr.rFonts.set(qn("w:eastAsia"), "맑은 고딕")
    p_title.paragraph_format.space_after = Pt(30)

    # Separator
    p_sep = doc.add_paragraph()
    run = p_sep.add_run("─" * 60)
    run.font.color.rgb = RGBColor(200, 200, 200)
    run.font.size = Pt(6)

    # Info block
    info_items = [
        ("부서", "{{부서명}}"),
        ("작성자", "{{작성자}}"),
        ("작성일", "{{작성일}}"),
    ]

    info_table = doc.add_table(rows=3, cols=2)
    info_table.alignment = WD_TABLE_ALIGNMENT.LEFT
    for i, (label, value) in enumerate(info_items):
        row = info_table.rows[i]
        cell_l = row.cells[0]
        cell_v = row.cells[1]
        cell_l.width = Cm(3)
        cell_v.width = Cm(10)

        p = cell_l.paragraphs[0]
        run = p.add_run(label)
        run.font.bold = True
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(*ACCENT)
        run.font.name = "맑은 고딕"

        p2 = cell_v.paragraphs[0]
        run2 = p2.add_run(value)
        run2.font.size = Pt(10)
        run2.font.name = "맑은 고딕"

    doc.add_page_break()

    # -- Header / Footer --
    _add_header_footer(section, "비즈니스 분석 보고서", PRIMARY)

    # -- Sections --
    sections_data = [
        ("1. 현황 분석", "{{현황분석}}"),
        ("2. 문제점 도출", "{{문제점}}"),
        ("3. 개선 방안", "{{개선방안}}"),
        ("4. 기대 효과", "{{기대효과}}"),
        ("5. 실행 계획", "{{실행계획}}"),
    ]

    for heading_text, placeholder in sections_data:
        doc.add_heading(heading_text, level=1)

        # Section accent underline
        p_accent = doc.add_paragraph()
        run = p_accent.add_run("█" * 6)
        run.font.color.rgb = RGBColor(*ACCENT)
        run.font.size = Pt(6)
        p_accent.paragraph_format.space_after = Pt(8)

        _add_styled_paragraph(doc, placeholder, font_size=10.5)

    path = OUTPUT_DIR / "report_analysis_business.docx"
    doc.save(str(path))
    return path


# ===================================================================
# Template 3: 사업 제안서 PPT (Business Proposal)
# ===================================================================

def create_proposal_business() -> Path:
    BG_DARK = (15, 32, 65)       # #0F2041 - dark navy
    BG_MAIN = (20, 40, 80)       # #142850
    BG_LIGHT = (240, 244, 250)   # #F0F4FA
    ACCENT = (0, 120, 215)       # #0078D7
    WHITE = (255, 255, 255)
    DARK_TEXT = (30, 40, 60)     # #1E283C
    LIGHT_TEXT = (200, 210, 230)

    prs = Presentation()
    prs.slide_width = PptxInches(10)
    prs.slide_height = PptxInches(7.5)
    SW = prs.slide_width
    SH = prs.slide_height

    # -- Slide 1: Title --
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG_DARK)

    # Left accent bar
    _add_accent_bar(slide, PptxInches(0), PptxInches(0),
                    PptxPt(6), SH, ACCENT)

    # Title
    _add_pptx_textbox(
        slide, PptxInches(1.5), PptxInches(2.0),
        PptxInches(7.0), PptxInches(1.5),
        "{{제안서제목}}", font_size=36, font_color=WHITE,
        bold=True, alignment=PP_ALIGN.LEFT
    )

    # Divider
    _add_accent_bar(slide, PptxInches(1.5), PptxInches(3.6),
                    PptxInches(2.0), PptxPt(3), ACCENT)

    # Company and date
    _add_pptx_textbox(
        slide, PptxInches(1.5), PptxInches(4.0),
        PptxInches(7.0), PptxInches(0.5),
        "{{회사명}}", font_size=18, font_color=LIGHT_TEXT,
        bold=False, alignment=PP_ALIGN.LEFT
    )
    _add_pptx_textbox(
        slide, PptxInches(1.5), PptxInches(4.6),
        PptxInches(7.0), PptxInches(0.4),
        "{{작성일}}", font_size=14, font_color=LIGHT_TEXT,
        bold=False, alignment=PP_ALIGN.LEFT
    )

    # -- Slide 2: Table of Contents --
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_accent_bar(slide, PptxInches(0), PptxInches(0), SW, PptxPt(4), ACCENT)

    _add_pptx_textbox(
        slide, PptxInches(0.8), PptxInches(0.5),
        PptxInches(8.5), PptxInches(0.7),
        "목 차", font_size=28, font_color=BG_DARK,
        bold=True, alignment=PP_ALIGN.LEFT
    )
    _add_accent_bar(slide, PptxInches(0.8), PptxInches(1.2),
                    PptxInches(1.2), PptxPt(3), ACCENT)

    toc_items = [
        "01  사업 개요",
        "02  추진 배경",
        "03  제안 내용",
        "04  기술 구성",
        "05  추진 일정",
        "06  기대 효과",
        "07  예산",
    ]
    toc_text = "\n\n".join(toc_items)
    _add_pptx_textbox(
        slide, PptxInches(1.2), PptxInches(1.8),
        PptxInches(7.0), PptxInches(5.0),
        toc_text, font_size=16, font_color=DARK_TEXT,
        bold=False, alignment=PP_ALIGN.LEFT
    )
    _add_slide_number(slide, color=DARK_TEXT)

    # -- Content slides --
    content_slides = [
        ("사업 개요", "{{사업개요}}"),
        ("추진 배경", "{{추진배경}}"),
        ("제안 내용", "{{제안내용}}"),
        ("기술 구성", "{{기술구성}}"),
        ("추진 일정", "{{추진일정}}"),
        ("기대 효과", "{{기대효과}}"),
        ("예산", "{{예산}}"),
    ]

    for title, placeholder in content_slides:
        _create_pptx_content_slide(
            prs, 0, title, placeholder,
            bg_color=WHITE, title_color=BG_DARK, body_color=DARK_TEXT,
            accent_color=ACCENT, slide_width=SW, slide_height=SH
        )

    # -- Slide 10: Thank You --
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG_DARK)
    _add_accent_bar(slide, PptxInches(0), PptxInches(0),
                    PptxPt(6), SH, ACCENT)

    _add_pptx_textbox(
        slide, PptxInches(1.0), PptxInches(2.5),
        PptxInches(8.0), PptxInches(1.2),
        "Q & A", font_size=44, font_color=WHITE,
        bold=True, alignment=PP_ALIGN.CENTER
    )
    _add_pptx_textbox(
        slide, PptxInches(1.0), PptxInches(4.0),
        PptxInches(8.0), PptxInches(0.6),
        "감사합니다", font_size=20, font_color=LIGHT_TEXT,
        bold=False, alignment=PP_ALIGN.CENTER
    )
    _add_accent_bar(slide, PptxInches(4.0), PptxInches(3.8),
                    PptxInches(2.0), PptxPt(3), ACCENT)

    path = OUTPUT_DIR / "proposal_business.pptx"
    prs.save(str(path))
    return path


# ===================================================================
# Template 4: 기술 제안서 PPT (Technical Proposal)
# ===================================================================

def create_proposal_tech() -> Path:
    BG_DARK = (27, 40, 56)      # #1B2838
    BG_CARD = (35, 52, 72)      # #233448
    ACCENT = (0, 180, 216)      # #00B4D8
    WHITE = (255, 255, 255)
    LIGHT_TEXT = (180, 200, 220)
    GRID_LINE = (50, 70, 90)

    prs = Presentation()
    prs.slide_width = PptxInches(10)
    prs.slide_height = PptxInches(7.5)
    SW = prs.slide_width
    SH = prs.slide_height

    # -- Slide 1: Title --
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG_DARK)

    # Top accent bar
    _add_accent_bar(slide, PptxInches(0), PptxInches(0),
                    SW, PptxPt(4), ACCENT)

    # Decorative element - thin vertical accent
    _add_accent_bar(slide, PptxInches(1.0), PptxInches(1.5),
                    PptxPt(3), PptxInches(4.5), ACCENT)

    # Title
    _add_pptx_textbox(
        slide, PptxInches(1.5), PptxInches(2.0),
        PptxInches(7.5), PptxInches(1.5),
        "{{제목}}", font_size=34, font_color=WHITE,
        bold=True, alignment=PP_ALIGN.LEFT
    )

    # Subtitle / presenter info
    _add_pptx_textbox(
        slide, PptxInches(1.5), PptxInches(3.8),
        PptxInches(7.0), PptxInches(0.5),
        "{{제안사}}", font_size=18, font_color=ACCENT,
        bold=False, alignment=PP_ALIGN.LEFT
    )
    _add_pptx_textbox(
        slide, PptxInches(1.5), PptxInches(4.4),
        PptxInches(7.0), PptxInches(0.4),
        "{{제안일}}", font_size=13, font_color=LIGHT_TEXT,
        bold=False, alignment=PP_ALIGN.LEFT
    )

    # Bottom accent bar
    _add_accent_bar(slide, PptxInches(0), PptxInches(7.2),
                    SW, PptxPt(4), ACCENT)

    # -- Content slides --
    tech_slides = [
        ("기술 개요", "{{기술개요}}"),
        ("시스템 아키텍처", "{{아키텍처}}"),
        ("핵심 기술", "{{핵심기술}}"),
        ("구현 방안", "{{구현방안}}"),
        ("품질 관리", "{{품질관리}}"),
        ("프로젝트 관리", "{{프로젝트관리}}"),
        ("지원 체계", "{{지원체계}}"),
    ]

    for title, placeholder in tech_slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_slide_bg(slide, BG_DARK)

        # Top accent line
        _add_accent_bar(slide, PptxInches(0), PptxInches(0),
                        SW, PptxPt(3), ACCENT)

        # Section number area (left accent block)
        _add_pptx_shape(
            slide, PptxInches(0.5), PptxInches(0.5),
            PptxInches(0.15), PptxInches(0.6),
            fill_color=ACCENT
        )

        # Title
        _add_pptx_textbox(
            slide, PptxInches(1.0), PptxInches(0.4),
            PptxInches(8.5), PptxInches(0.7),
            title, font_size=26, font_color=WHITE,
            bold=True, alignment=PP_ALIGN.LEFT
        )

        # Thin divider
        _add_accent_bar(slide, PptxInches(1.0), PptxInches(1.15),
                        PptxInches(1.2), PptxPt(2), ACCENT)

        # Content card background
        _add_pptx_shape(
            slide, PptxInches(0.5), PptxInches(1.5),
            PptxInches(9.0), PptxInches(5.5),
            fill_color=BG_CARD
        )

        # Content text
        _add_pptx_textbox(
            slide, PptxInches(0.9), PptxInches(1.8),
            PptxInches(8.2), PptxInches(4.8),
            placeholder, font_size=15, font_color=LIGHT_TEXT,
            bold=False, alignment=PP_ALIGN.LEFT
        )

        # Bottom bar
        _add_accent_bar(slide, PptxInches(0), PptxInches(7.2),
                        SW, PptxPt(2), ACCENT)

        _add_slide_number(slide)

    path = OUTPUT_DIR / "proposal_tech.pptx"
    prs.save(str(path))
    return path


# ===================================================================
# Template 5: 회의록 (Meeting Minutes)
# ===================================================================

def create_report_meeting_minutes() -> Path:
    PRIMARY = (55, 65, 81)       # #374151
    ACCENT = (59, 130, 246)      # #3B82F6
    LIGHT_BG = (243, 244, 246)   # #F3F4F6

    doc = Document()
    _configure_docx_styles(doc, PRIMARY)

    section = doc.sections[0]
    _set_page_margins(section, top=2.0, bottom=2.0, left=2.5, right=2.5)

    # -- Header / Footer --
    _add_header_footer(section, "회 의 록", PRIMARY)

    # -- Title --
    p_title = doc.add_paragraph()
    p_title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p_title.add_run("{{회의제목}}")
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*PRIMARY)
    run.font.name = "맑은 고딕"
    run._element.rPr.rFonts.set(qn("w:eastAsia"), "맑은 고딕")
    p_title.paragraph_format.space_before = Pt(12)
    p_title.paragraph_format.space_after = Pt(6)

    # Subtitle accent
    p_accent = doc.add_paragraph()
    p_accent.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p_accent.add_run("━" * 30)
    run.font.color.rgb = RGBColor(*ACCENT)
    run.font.size = Pt(6)
    p_accent.paragraph_format.space_after = Pt(16)

    # -- Info table --
    info_table = doc.add_table(rows=4, cols=2)
    info_table.alignment = WD_TABLE_ALIGNMENT.CENTER

    info_data = [
        ("일  시", "{{회의일시}}"),
        ("장  소", "{{회의장소}}"),
        ("참석자", "{{참석자}}"),
        ("작성자", "{{작성자}}"),
    ]

    for i, (label, value) in enumerate(info_data):
        row = info_table.rows[i]
        cell_label = row.cells[0]
        cell_value = row.cells[1]
        cell_label.width = Cm(3.5)
        cell_value.width = Cm(12)

        # Style label cell
        _set_cell_shading(cell_label, "F3F4F6")
        p = cell_label.paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(label)
        run.font.bold = True
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(*PRIMARY)
        run.font.name = "맑은 고딕"

        # Style value cell
        p2 = cell_value.paragraphs[0]
        run2 = p2.add_run(value)
        run2.font.size = Pt(10)
        run2.font.name = "맑은 고딕"

        # Add borders to all cells
        for cell in [cell_label, cell_value]:
            _set_cell_border(
                cell,
                top=("single", "4", "D1D5DB"),
                bottom=("single", "4", "D1D5DB"),
                left=("single", "4", "D1D5DB"),
                right=("single", "4", "D1D5DB"),
            )

    doc.add_paragraph()  # spacing

    # -- Sections --
    sections_data = [
        ("회의 안건", "{{안건}}"),
        ("논의 내용", "{{논의내용}}"),
        ("결정 사항", "{{결정사항}}"),
        ("향후 일정", "{{향후일정}}"),
    ]

    for heading_text, placeholder in sections_data:
        # Section heading with accent
        p_heading = doc.add_paragraph()
        p_heading.paragraph_format.space_before = Pt(16)
        p_heading.paragraph_format.space_after = Pt(8)

        # Blue accent marker
        run_marker = p_heading.add_run("■ ")
        run_marker.font.color.rgb = RGBColor(*ACCENT)
        run_marker.font.size = Pt(12)

        run_text = p_heading.add_run(heading_text)
        run_text.font.size = Pt(14)
        run_text.font.bold = True
        run_text.font.color.rgb = RGBColor(*PRIMARY)
        run_text.font.name = "맑은 고딕"
        run_text._element.rPr.rFonts.set(qn("w:eastAsia"), "맑은 고딕")

        # Light separator
        p_sep = doc.add_paragraph()
        run = p_sep.add_run("─" * 60)
        run.font.color.rgb = RGBColor(220, 220, 220)
        run.font.size = Pt(5)
        p_sep.paragraph_format.space_after = Pt(4)

        # Content placeholder
        _add_styled_paragraph(doc, placeholder, font_size=10.5, space_after=6)

    path = OUTPUT_DIR / "report_meeting_minutes.docx"
    doc.save(str(path))
    return path


# ===================================================================
# Main
# ===================================================================

def main():
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    print(f"Output directory: {OUTPUT_DIR}\n")

    generators = [
        ("Template 1 - 격식체 요약 보고서 (DOCX)", create_report_summary_formal),
        ("Template 2 - 비즈니스 분석 보고서 (DOCX)", create_report_analysis_business),
        ("Template 3 - 사업 제안서 (PPTX)", create_proposal_business),
        ("Template 4 - 기술 제안서 (PPTX)", create_proposal_tech),
        ("Template 5 - 회의록 (DOCX)", create_report_meeting_minutes),
    ]

    for label, func in generators:
        path = func()
        print(f"[OK] {label}")
        print(f"     -> {path}\n")

    print("All templates created successfully.")


if __name__ == "__main__":
    main()
