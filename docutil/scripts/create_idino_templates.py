#!/usr/bin/env python3
"""
IDINO 기업용 문서 템플릿 생성기
================================
PPTX 8종 + DOCX 8종 = 총 16개의 전문적 디자인 템플릿을 생성합니다.

사용법:
    python scripts/create_idino_templates.py

의존성:
    pip install python-pptx python-docx

출력:
    sample_templates/ 디렉토리에 16개 파일 생성
"""

import os
from pathlib import Path

# ─── python-docx: Word 문서 생성 라이브러리 ─────────────────────────────────
from docx import Document
from docx.shared import Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn, nsdecls
from docx.oxml import parse_xml

# ─── python-pptx: PowerPoint 프레젠테이션 생성 라이브러리 ───────────────────
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt, Emu as PptxEmu
from pptx.dml.color import RGBColor as PptxRGB
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ═══════════════════════════════════════════════════════════════════════════════
#  IDINO 디자인 시스템 — 색상, 폰트, 크기 상수
# ═══════════════════════════════════════════════════════════════════════════════

HEADER_BAR    = "34495E"   # 진한 남색 — 헤더/푸터 배경, 타이틀 슬라이드 배경
ACCENT_BLUE   = "2980B9"   # 밝은 파랑 — 강조선, 라벨, 링크
HIGHLIGHT_BG  = "EBF5FB"   # 밝은 파랑 배경 — 콘텐츠 박스 배경
DARK_ACCENT   = "1A5276"   # 진한 파랑 — 부제, 강조 텍스트
BODY_TEXT      = "2C3E50"   # 진한 회색 — 본문 텍스트
SUBTITLE_GRAY = "D5D8DC"   # 밝은 회색 — 부제, 날짜
WHITE          = "FFFFFF"   # 흰색 — 헤더바 위 텍스트

FONT_NAME = "Malgun Gothic"  # 맑은 고딕 — 전체 문서에 사용

# 출력 디렉토리 (프로젝트 루트 기준)
OUTPUT_DIR = Path(__file__).resolve().parent.parent / "sample_templates"


# ═══════════════════════════════════════════════════════════════════════════════
#  PPTX 헬퍼 함수들
# ═══════════════════════════════════════════════════════════════════════════════

# 슬라이드 크기: 16:9 비율 (13.333" x 7.5")
SLIDE_W = PptxEmu(12192000)  # 13.333인치를 EMU(English Metric Units)로 변환
SLIDE_H = PptxEmu(6858000)   # 7.5인치를 EMU로 변환

# 헤더바와 푸터바 높이
HEADER_BAR_H = PptxInches(0.42)   # 상단 헤더바 높이
FOOTER_BAR_H = PptxInches(0.18)   # 하단 푸터바 높이
FOOTER_TOP   = SLIDE_H - FOOTER_BAR_H  # 푸터바 Y 좌표 (슬라이드 높이 - 푸터 높이)


def _new_pptx():
    """
    새 프레젠테이션 객체를 생성합니다.
    슬라이드 크기를 16:9로 설정합니다.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _add_shape_rect(slide, left, top, width, height, fill_hex):
    """
    슬라이드에 색이 채워진 사각형을 추가합니다.

    매개변수:
        slide: 슬라이드 객체
        left, top: 좌측 상단 좌표
        width, height: 사각형 크기
        fill_hex: 채우기 색상 (예: "34495E")
    """
    # MSO_SHAPE.RECTANGLE = 1 (직사각형 도형 타입)
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()                                  # 단색 채우기 설정
    shape.fill.fore_color.rgb = PptxRGB.from_string(fill_hex)  # 색상 적용
    shape.line.fill.background()                        # 테두리 없음 (투명)
    return shape


def _add_text_box(slide, left, top, width, height, text,
                  size=12, bold=False, color=BODY_TEXT,
                  alignment=PP_ALIGN.LEFT, font=FONT_NAME,
                  anchor=MSO_ANCHOR.TOP):
    """
    슬라이드에 텍스트 상자를 추가합니다.

    매개변수:
        slide: 슬라이드 객체
        left, top, width, height: 위치와 크기
        text: 표시할 텍스트
        size: 글자 크기 (pt)
        bold: 굵게 여부
        color: 글자 색상
        alignment: 텍스트 정렬 (왼쪽, 가운데, 오른쪽)
        font: 폰트 이름
        anchor: 세로 정렬 (상단, 중앙, 하단)
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True       # 텍스트 자동 줄바꿈 활성화
    tf.auto_size = None       # 자동 크기 조절 비활성화 (고정 크기)

    p = tf.paragraphs[0]      # 기본 문단에 텍스트 설정
    p.text = text
    p.font.size = PptxPt(size)
    p.font.bold = bold
    p.font.color.rgb = PptxRGB.from_string(color)
    p.font.name = font
    p.alignment = alignment
    return txBox


def _pptx_add_header_bar(slide, section_title):
    """
    콘텐츠 슬라이드 상단에 헤더바를 추가합니다.
    - 전체 너비 x 0.42" 높이의 진한 남색 배경
    - 왼쪽 0.5" 여백에 섹션 제목 (흰색, 13pt, 굵게)

    매개변수:
        slide: 슬라이드 객체
        section_title: 헤더바에 표시할 섹션 제목
    """
    # 배경 사각형 추가
    _add_shape_rect(slide, 0, 0, SLIDE_W, HEADER_BAR_H, HEADER_BAR)
    # 섹션 제목 텍스트
    _add_text_box(
        slide, PptxInches(0.5), 0,
        PptxInches(10), HEADER_BAR_H,
        section_title,
        size=13, bold=True, color=WHITE,
        alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE
    )


def _pptx_add_footer_bar(slide):
    """
    콘텐츠 슬라이드 하단에 푸터바를 추가합니다.
    - 전체 너비 x 0.18" 높이의 진한 남색 배경
    - 우측 1.2" 위치에 "IDINO" 텍스트 (7pt, 굵게, 흰색)
    """
    # 배경 사각형 추가
    _add_shape_rect(slide, 0, FOOTER_TOP, SLIDE_W, FOOTER_BAR_H, HEADER_BAR)
    # IDINO 로고 텍스트 (우측 정렬)
    _add_text_box(
        slide, SLIDE_W - PptxInches(1.2), FOOTER_TOP,
        PptxInches(1.0), FOOTER_BAR_H,
        "IDINO",
        size=7, bold=True, color=WHITE,
        alignment=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE
    )


def _pptx_add_title_slide(prs, title, subtitle, date_text="{{작성일}}"):
    """
    타이틀 슬라이드를 추가합니다 (전체 배경 #34495E).

    구성:
        - 전체 배경: 진한 남색
        - 상단 강조선: 파란색, 전체 너비, 0.03" 높이
        - 제목: 32pt, Bold, White, 중앙 정렬, Y=2.5"
        - 부제/{{변수}}: 14pt, 밝은 회색, 중앙, Y=3.5"
        - 날짜: 11pt, 밝은 회색, 중앙, Y=4.5"
        - 하단 IDINO 로고: 11pt, 파란색, 중앙, Y=6.5"

    매개변수:
        prs: Presentation 객체
        title: 메인 제목 텍스트
        subtitle: 부제 텍스트
        date_text: 날짜 텍스트 (기본값: "{{작성일}}")
    """
    slide_layout = prs.slide_layouts[6]       # 빈 레이아웃 사용 (수동 배치)
    slide = prs.slides.add_slide(slide_layout)

    # 1) 전체 배경: 진한 남색 사각형
    _add_shape_rect(slide, 0, 0, SLIDE_W, SLIDE_H, HEADER_BAR)

    # 2) 상단 강조선: 파란색 가로 줄 (전체 너비, 0.03" 높이)
    _add_shape_rect(slide, 0, 0, SLIDE_W, PptxInches(0.03), ACCENT_BLUE)

    # 3) 제목: 32pt, 굵게, 흰색, 가운데 정렬, Y=2.5"
    _add_text_box(
        slide, PptxInches(1), PptxInches(2.5),
        PptxInches(11.333), PptxInches(1.0),
        title,
        size=32, bold=True, color=WHITE,
        alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP
    )

    # 4) 부제: 14pt, 밝은 회색, 가운데 정렬, Y=3.5"
    _add_text_box(
        slide, PptxInches(1), PptxInches(3.5),
        PptxInches(11.333), PptxInches(0.5),
        subtitle,
        size=14, color=SUBTITLE_GRAY,
        alignment=PP_ALIGN.CENTER
    )

    # 5) 날짜: 11pt, 밝은 회색, 가운데 정렬, Y=4.5"
    _add_text_box(
        slide, PptxInches(1), PptxInches(4.5),
        PptxInches(11.333), PptxInches(0.4),
        date_text,
        size=11, color=SUBTITLE_GRAY,
        alignment=PP_ALIGN.CENTER
    )

    # 6) 하단 IDINO 로고 텍스트: 11pt, 파란색, 가운데, Y=6.5"
    _add_text_box(
        slide, PptxInches(5), PptxInches(6.5),
        PptxInches(3.333), PptxInches(0.4),
        "IDINO",
        size=11, bold=True, color=ACCENT_BLUE,
        alignment=PP_ALIGN.CENTER
    )

    return slide


def _pptx_add_content_slide(prs, section_title, items):
    """
    콘텐츠 슬라이드를 추가합니다.

    구성:
        - 상단 헤더바 (진한 남색 + 섹션 제목)
        - 콘텐츠 영역 (Y=0.7" ~ Y=6.8"):
            각 항목마다:
            - "■ 항목명" 라벨 (12pt, 파란색, 굵게)
            - 밝은 파란 배경 박스 + 좌측 파란 테두리
            - {{변수}} 텍스트 (11pt, 본문색)
        - 하단 푸터바 (IDINO 로고)

    매개변수:
        prs: Presentation 객체
        section_title: 헤더바에 표시될 섹션 제목
        items: [(라벨, 변수)] 리스트. 예: [("개요", "{{개요}}")]
    """
    slide_layout = prs.slide_layouts[6]       # 빈 레이아웃
    slide = prs.slides.add_slide(slide_layout)

    # 상단 헤더바
    _pptx_add_header_bar(slide, section_title)
    # 하단 푸터바
    _pptx_add_footer_bar(slide)

    # ── 콘텐츠 영역 레이아웃 계산 ──
    # 사용 가능한 Y 범위: 0.7" ~ 6.6" (약 5.9" 높이)
    # 각 항목은 라벨(0.3") + 박스 + 간격으로 구성됨
    content_left = PptxInches(0.5)   # 콘텐츠 좌측 여백
    content_width = PptxInches(12.3) # 콘텐츠 너비 (슬라이드 너비 - 양쪽 여백)
    y_start = PptxInches(0.7)        # 콘텐츠 시작 Y 좌표
    y_end = PptxInches(6.6)          # 콘텐츠 끝 Y 좌표

    n = len(items)
    if n == 0:
        return slide

    # 각 항목에 할당되는 총 높이 계산 (라벨 + 박스 + 간격)
    total_height_emu = y_end - y_start
    item_height_emu = total_height_emu // n          # 항목 하나의 전체 높이
    label_h = PptxInches(0.3)                        # 라벨 높이 (고정)
    gap = PptxInches(0.1)                            # 항목 간 간격 (고정)
    box_h = item_height_emu - label_h - gap          # 나머지가 박스 높이

    # 박스 최소/최대 높이 제한 (너무 작거나 크지 않도록)
    min_box = PptxInches(0.4)
    max_box = PptxInches(1.5)
    if box_h < min_box:
        box_h = min_box
    elif box_h > max_box:
        box_h = max_box

    y = y_start
    for label, variable in items:
        # ── "■ 항목명" 라벨 ──
        _add_text_box(
            slide, content_left, y,
            content_width, label_h,
            f"\u25a0 {label}",       # ■ + 항목명
            size=12, bold=True, color=ACCENT_BLUE
        )
        y += label_h

        # ── 콘텐츠 배경 박스 (밝은 파랑 배경) ──
        _add_shape_rect(
            slide, content_left, y,
            content_width, box_h, HIGHLIGHT_BG
        )

        # ── 좌측 파란 테두리 (0.15" 폭의 가느다란 사각형) ──
        _add_shape_rect(
            slide, content_left, y,
            PptxInches(0.08), box_h, ACCENT_BLUE
        )

        # ── {{변수}} 텍스트 (배경 박스 안, 좌측 0.15" 여백) ──
        _add_text_box(
            slide, content_left + PptxInches(0.2), y + PptxInches(0.05),
            content_width - PptxInches(0.35), box_h - PptxInches(0.1),
            variable,
            size=11, color=BODY_TEXT
        )

        y += box_h + gap

    return slide


def _pptx_add_closing_slide(prs):
    """
    마침(감사합니다) 슬라이드를 추가합니다.

    구성:
        - 전체 배경: 진한 남색
        - "감사합니다": 36pt, Bold, White, 중앙
        - 중앙 강조선: #2980B9, 2" 너비, 0.03" 높이
        - 연락처 안내문: 12pt, 밝은 회색, 중앙
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 전체 배경
    _add_shape_rect(slide, 0, 0, SLIDE_W, SLIDE_H, HEADER_BAR)

    # "감사합니다" 텍스트 (중앙)
    _add_text_box(
        slide, PptxInches(1), PptxInches(2.5),
        PptxInches(11.333), PptxInches(1.2),
        "감사합니다",
        size=36, bold=True, color=WHITE,
        alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE
    )

    # 중앙 강조선 (2" 너비, 가운데 배치)
    line_w = PptxInches(2)
    line_left = (SLIDE_W - line_w) // 2  # 수평 중앙 계산
    _add_shape_rect(
        slide, line_left, PptxInches(3.8),
        line_w, PptxInches(0.03), ACCENT_BLUE
    )

    # 연락처 안내문
    _add_text_box(
        slide, PptxInches(2), PptxInches(4.2),
        PptxInches(9.333), PptxInches(0.5),
        "IDINO | 문의사항은 담당자에게 연락 바랍니다",
        size=12, color=SUBTITLE_GRAY,
        alignment=PP_ALIGN.CENTER
    )

    return slide


def _build_pptx(filename, title, subtitle, slides_data):
    """
    PPTX 템플릿 하나를 생성하는 통합 함수입니다.

    매개변수:
        filename: 저장할 파일명 (예: "report_summary_pptx.pptx")
        title: 타이틀 슬라이드의 제목
        subtitle: 타이틀 슬라이드의 부제
        slides_data: [(섹션제목, [(라벨, 변수), ...])] 콘텐츠 슬라이드 목록

    생성 순서:
        1. 타이틀 슬라이드
        2. 콘텐츠 슬라이드들 (slides_data 순서대로)
        3. 마침 슬라이드 (감사합니다)
    """
    prs = _new_pptx()

    # 1) 타이틀 슬라이드
    _pptx_add_title_slide(prs, title, subtitle)

    # 2) 각 콘텐츠 슬라이드 추가
    for section_title, items in slides_data:
        _pptx_add_content_slide(prs, section_title, items)

    # 3) 마침 슬라이드
    _pptx_add_closing_slide(prs)

    # 저장
    path = OUTPUT_DIR / filename
    prs.save(str(path))
    print(f"  [PPTX] {filename} ({len(prs.slides)}슬라이드)")


# ═══════════════════════════════════════════════════════════════════════════════
#  DOCX 헬퍼 함수들
# ═══════════════════════════════════════════════════════════════════════════════

def _set_cell_shading(cell, color_hex):
    """
    표 셀의 배경색을 설정합니다.

    매개변수:
        cell: python-docx 테이블 셀 객체
        color_hex: 배경 색상 (예: "34495E")
    """
    # XML로 배경 음영(shading) 속성 생성 후 셀에 추가
    shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{color_hex}"/>')
    cell._tc.get_or_add_tcPr().append(shading)


def _set_cell_border(cell, **kwargs):
    """
    표 셀의 테두리를 설정합니다.

    매개변수:
        cell: python-docx 테이블 셀 객체
        kwargs: 테두리 위치별 설정. 예: left=("12", "2980B9")
                키: top, bottom, left, right
                값: (두께, 색상) 튜플
    """
    tc_pr = cell._tc.get_or_add_tcPr()
    # 테두리 컨테이너 XML 생성
    borders = parse_xml(f'<w:tcBorders {nsdecls("w")}/>')
    for edge, (sz, color) in kwargs.items():
        # 각 변(edge)에 대한 테두리 속성 추가
        el = parse_xml(
            f'<w:{edge} {nsdecls("w")} w:val="single" w:sz="{sz}" '
            f'w:space="0" w:color="{color}"/>'
        )
        borders.append(el)
    tc_pr.append(borders)


def _docx_run(paragraph, text, size=11, bold=False, color=BODY_TEXT, font=FONT_NAME):
    """
    문단에 서식이 적용된 텍스트 조각(run)을 추가합니다.

    매개변수:
        paragraph: 문단 객체
        text: 추가할 텍스트
        size: 글자 크기 (pt)
        bold: 굵게 여부
        color: 글자 색상
        font: 폰트 이름
    """
    r = paragraph.add_run(text)
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = RGBColor.from_string(color)
    r.font.name = font
    # 동아시아 폰트도 같은 폰트로 설정 (한글 표시에 필요)
    r._element.rPr.rFonts.set(qn("w:eastAsia"), font)
    return r


def _docx_add_thin_line(doc, color=ACCENT_BLUE):
    """
    가는 수평 구분선을 추가합니다 (문단 하단 테두리 방식).

    매개변수:
        doc: Document 객체
        color: 구분선 색상
    """
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(2)
    p.paragraph_format.space_after = Pt(6)
    # XML로 문단 하단 테두리 추가
    pPr = p._element.get_or_add_pPr()
    pBdr = parse_xml(
        f'<w:pBdr {nsdecls("w")}>'
        f'  <w:bottom w:val="single" w:sz="6" w:space="1" w:color="{color}"/>'
        f'</w:pBdr>'
    )
    pPr.append(pBdr)
    return p


def _docx_setup():
    """
    새 DOCX 문서를 생성하고 기본 스타일을 설정합니다.

    설정 내용:
        - A4 용지 (21cm x 29.7cm)
        - 여백 2.5cm (상하좌우)
        - 기본 폰트: 맑은 고딕 11pt
        - 행간: 1.15
        - Heading 1 스타일: 진한 남색 (#34495E) + 파란 언더라인
    """
    doc = Document()

    # ── 페이지 크기 및 여백 설정 ──
    section = doc.sections[0]
    section.page_width = Cm(21)       # A4 가로
    section.page_height = Cm(29.7)    # A4 세로
    section.top_margin = Cm(2.5)      # 상단 여백
    section.bottom_margin = Cm(2.5)   # 하단 여백
    section.left_margin = Cm(2.5)     # 좌측 여백
    section.right_margin = Cm(2.5)    # 우측 여백

    # ── 기본(Normal) 스타일: 맑은 고딕 11pt, 행간 1.15 ──
    style = doc.styles["Normal"]
    style.font.name = FONT_NAME
    style.font.size = Pt(11)
    style.font.color.rgb = RGBColor.from_string(BODY_TEXT)
    style.paragraph_format.line_spacing = 1.15
    # 동아시아 폰트 설정 (한글에 필요)
    rpr = style._element.rPr
    rpr.rFonts.set(qn("w:eastAsia"), FONT_NAME)

    # ── Heading 스타일 커스터마이징 ──
    # Heading 1: 16pt, 굵게, 진한 남색 (#34495E)
    # Heading 2: 14pt, 굵게, 파란색 (#2980B9)
    # Heading 3: 12pt, 굵게, 파란색 (#2980B9)
    for level, sz, clr in [(1, 16, HEADER_BAR), (2, 14, ACCENT_BLUE), (3, 12, ACCENT_BLUE)]:
        hs = doc.styles[f"Heading {level}"]
        hs.font.name = FONT_NAME
        hs.font.size = Pt(sz)
        hs.font.bold = True
        hs.font.color.rgb = RGBColor.from_string(clr)
        hs._element.rPr.rFonts.set(qn("w:eastAsia"), FONT_NAME)

    return doc


def _docx_add_cover_page(doc, title, subtitle_lines):
    """
    커버 페이지를 추가합니다 (표 기반 어두운 배너).

    구성:
        - 어두운 배너 (표 1행 1열, #34495E 배경)
        - 제목: 28pt, 굵게, 흰색, 가운데 정렬
        - 부제 줄들: 13pt, 밝은 회색, 가운데 정렬
        - 파란색 구분선 (━ 문자 반복)
        - 페이지 넘김

    매개변수:
        doc: Document 객체
        title: 문서 제목
        subtitle_lines: 부제 줄 리스트 (예: ["작성자: {{작성자}}", ...])
    """
    # ── 어두운 배너를 표(1x1)로 구현 ──
    # python-docx에서는 전체 페이지 배경을 직접 지정하기 어려우므로
    # 표의 셀 배경색을 사용하여 배너 효과를 냅니다.
    tbl = doc.add_table(rows=1, cols=1)
    tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
    cell = tbl.rows[0].cells[0]
    _set_cell_shading(cell, HEADER_BAR)  # 진한 남색 배경
    cell.width = Cm(16)

    # 셀 높이 설정 (약 1.8" = 3600 twips)
    tr = tbl.rows[0]._tr
    trPr = tr.get_or_add_trPr()
    trHeight = parse_xml(
        f'<w:trHeight {nsdecls("w")} w:val="3600" w:hRule="atLeast"/>'
    )
    trPr.append(trHeight)

    # ── 제목 텍스트 ──
    p = cell.paragraphs[0]
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(60)  # 상단 여백으로 수직 중앙 효과
    _docx_run(p, title, size=28, bold=True, color=WHITE)

    # ── 부제 줄들 ──
    for line in subtitle_lines:
        p2 = cell.add_paragraph()
        p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p2.paragraph_format.space_before = Pt(6)
        _docx_run(p2, line, size=13, color=SUBTITLE_GRAY)

    # ── 파란색 구분선 (━ 문자 30개) ──
    p3 = cell.add_paragraph()
    p3.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p3.paragraph_format.space_before = Pt(20)
    _docx_run(p3, "\u2501" * 30, size=11, color=ACCENT_BLUE)

    # 배너 뒤에 빈 줄 + 페이지 넘김
    doc.add_paragraph()
    doc.add_page_break()


def _docx_add_header(doc, title_text):
    """
    문서 상단에 헤더를 추가합니다.
    - 좌측: 문서 제목 (9pt, 굵게, 진한 남색)
    - 하단: 구분선

    매개변수:
        doc: Document 객체
        title_text: 헤더에 표시할 문서 제목
    """
    section = doc.sections[0]
    header = section.header
    header.is_linked_to_previous = False  # 이전 섹션과 연결 해제

    p = header.paragraphs[0] if header.paragraphs else header.add_paragraph()
    _docx_run(p, title_text, size=9, bold=True, color=HEADER_BAR)

    # 헤더 하단에 가느다란 구분선 추가
    pPr = p._element.get_or_add_pPr()
    pBdr = parse_xml(
        f'<w:pBdr {nsdecls("w")}>'
        f'  <w:bottom w:val="single" w:sz="4" w:space="4" w:color="{HEADER_BAR}"/>'
        f'</w:pBdr>'
    )
    pPr.append(pBdr)


def _docx_add_footer(doc):
    """
    문서 하단에 푸터(페이지 번호)를 추가합니다.
    - 중앙: "— 1 —" 형식의 페이지 번호
    - 상단: 가느다란 구분선
    """
    section = doc.sections[0]
    footer = section.footer
    footer.is_linked_to_previous = False

    p = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # "— " 텍스트 추가
    _docx_run(p, "\u2014 ", size=9, color=HEADER_BAR)

    # 페이지 번호 필드 (Word의 PAGE 필드 코드 사용)
    fld_xml = (
        f'<w:fldSimple {nsdecls("w")} w:instr=" PAGE \\* MERGEFORMAT ">'
        f'  <w:r><w:rPr><w:sz w:val="18"/><w:color w:val="{HEADER_BAR}"/></w:rPr>'
        f'  <w:t>1</w:t></w:r></w:fldSimple>'
    )
    p._element.append(parse_xml(fld_xml))

    # " —" 텍스트 추가
    _docx_run(p, " \u2014", size=9, color=HEADER_BAR)

    # 푸터 상단 구분선
    pPr = p._element.get_or_add_pPr()
    pBdr = parse_xml(
        f'<w:pBdr {nsdecls("w")}>'
        f'  <w:top w:val="single" w:sz="4" w:space="4" w:color="{HEADER_BAR}"/>'
        f'</w:pBdr>'
    )
    pPr.append(pBdr)


def _docx_add_section_heading(doc, text):
    """
    섹션 제목(Heading 1)과 파란 구분선을 추가합니다.

    매개변수:
        doc: Document 객체
        text: 섹션 제목 (예: "1. 개요")
    """
    doc.add_heading(text, level=1)
    _docx_add_thin_line(doc, ACCENT_BLUE)  # 제목 아래 파란 구분선


def _docx_add_content_block(doc, variable):
    """
    콘텐츠 블록을 추가합니다 (밝은 파란 배경 + 좌측 파란 테두리).

    표(1x1)를 사용하여 배경색과 좌측 테두리를 구현합니다.
    {{변수}} 텍스트가 안에 들어갑니다.

    매개변수:
        doc: Document 객체
        variable: 변수 텍스트 (예: "{{개요}}")
    """
    tbl = doc.add_table(rows=1, cols=1)
    cell = tbl.rows[0].cells[0]
    _set_cell_shading(cell, HIGHLIGHT_BG)  # 밝은 파란 배경

    p = cell.paragraphs[0]
    p.paragraph_format.space_before = Pt(8)
    p.paragraph_format.space_after = Pt(8)
    _docx_run(p, variable, size=11, color=BODY_TEXT)

    # 좌측 파란 테두리 (두께 12, 색상 파란색)
    _set_cell_border(cell, left=("12", ACCENT_BLUE))

    # 블록 뒤에 빈 줄 추가 (간격 확보)
    doc.add_paragraph()


def _docx_add_info_table(doc, rows_data):
    """
    정보 테이블을 추가합니다 (라벨-값 쌍).
    라벨 셀은 진한 남색 배경 + 흰색 글씨,
    값 셀은 기본 배경 + 본문색 글씨.

    매개변수:
        doc: Document 객체
        rows_data: [(라벨, 값)] 리스트. 예: [("회의 제목", "{{회의제목}}")]
    """
    tbl = doc.add_table(rows=len(rows_data), cols=2)
    tbl.alignment = WD_TABLE_ALIGNMENT.CENTER

    for i, (label, value) in enumerate(rows_data):
        # 라벨 셀 (왼쪽): 진한 남색 배경, 흰색 텍스트
        lc = tbl.rows[i].cells[0]
        lc.width = Cm(3.5)
        _set_cell_shading(lc, HEADER_BAR)
        p = lc.paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        _docx_run(p, label, size=10, bold=True, color=WHITE)

        # 값 셀 (오른쪽): 기본 배경, 본문색 텍스트
        vc = tbl.rows[i].cells[1]
        vc.width = Cm(12.5)
        p2 = vc.paragraphs[0]
        _docx_run(p2, value, size=10, color=BODY_TEXT)
        p2.paragraph_format.left_indent = Cm(0.3)

    doc.add_paragraph()  # 테이블 뒤 간격


def _build_docx(filename, doc_title, cover_title, cover_subtitles, sections):
    """
    DOCX 템플릿 하나를 생성하는 통합 함수입니다.

    매개변수:
        filename: 저장할 파일명 (예: "report_summary_docx.docx")
        doc_title: 헤더에 표시될 문서 유형명 (예: "요약 보고서")
        cover_title: 커버 페이지 제목 (예: "{{제목}}")
        cover_subtitles: 커버 페이지 부제 줄 리스트
        sections: [(섹션번호+제목, 변수)] 리스트
                  예: [("1. 개요", "{{개요}}"), ...]

    생성 순서:
        1. 문서 초기화 (A4, 스타일 설정)
        2. 커버 페이지
        3. 헤더 + 푸터
        4. 각 섹션 (제목 + 콘텐츠 블록)
    """
    doc = _docx_setup()

    # 1) 커버 페이지
    _docx_add_cover_page(doc, cover_title, cover_subtitles)

    # 2) 헤더와 푸터
    _docx_add_header(doc, doc_title)
    _docx_add_footer(doc)

    # 3) 각 섹션 추가
    for heading, variable in sections:
        _docx_add_section_heading(doc, heading)
        _docx_add_content_block(doc, variable)

    # 저장
    path = OUTPUT_DIR / filename
    doc.save(str(path))
    print(f"  [DOCX] {filename} ({len(sections)}섹션)")


def _build_docx_meeting(filename):
    """
    회의록 템플릿을 생성합니다.
    커버 페이지 대신 회의 정보 테이블이 상단에 위치하는 특수한 구조입니다.

    구성:
        1. "회 의 록" 제목 + 구분선
        2. 회의 정보 테이블 (제목, 일시, 장소, 참석자)
        3. 안건, 논의내용, 결정사항, 향후계획, 비고 섹션
    """
    doc = _docx_setup()
    _docx_add_header(doc, "회의록")
    _docx_add_footer(doc)

    # ── "회 의 록" 대제목 ──
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    _docx_run(p, "회 의 록", size=22, bold=True, color=HEADER_BAR)
    _docx_add_thin_line(doc, HEADER_BAR)

    # ── 회의 정보 테이블 ──
    _docx_add_info_table(doc, [
        ("회의 제목", "{{회의명}}"),
        ("회의 일시", "{{회의일시}}"),
        ("회의 장소", "{{회의장소}}"),
        ("참 석 자", "{{참석자}}"),
    ])

    # ── 내용 섹션들 ──
    sections = [
        ("1. 안건", "{{안건}}"),
        ("2. 논의 내용", "{{논의내용}}"),
        ("3. 결정 사항", "{{결정사항}}"),
        ("4. 향후 계획", "{{향후계획}}"),
        ("5. 비고", "{{비고}}"),
    ]
    for heading, variable in sections:
        _docx_add_section_heading(doc, heading)
        _docx_add_content_block(doc, variable)

    path = OUTPUT_DIR / filename
    doc.save(str(path))
    print(f"  [DOCX] {filename} (회의록, 5섹션)")


# ═══════════════════════════════════════════════════════════════════════════════
#  PPTX 템플릿 8종 정의
# ═══════════════════════════════════════════════════════════════════════════════

def create_all_pptx():
    """
    PPTX 템플릿 8종을 모두 생성합니다.
    각 템플릿은 타이틀 + 콘텐츠 슬라이드들 + 마침 슬라이드로 구성됩니다.
    """
    print("\n[PPTX 템플릿 생성 중...]")

    # ── 1. 요약 보고서 (8슬라이드) ──
    _build_pptx(
        "report_summary_pptx.pptx",
        title="{{제목}}",
        subtitle="요약 보고서  |  {{부서}}",
        slides_data=[
            ("1. 개요", [
                ("보고 목적", "{{보고목적}}"),
                ("보고 범위", "{{보고범위}}"),
                ("작성 배경", "{{작성배경}}"),
            ]),
            ("2. 주요 내용", [
                ("핵심 요약", "{{핵심요약}}"),
                ("주요 사항", "{{주요사항}}"),
            ]),
            ("3. 분석 결과", [
                ("분석 방법", "{{분석방법}}"),
                ("분석 결과", "{{분석결과}}"),
            ]),
            ("4. 세부 내용", [
                ("세부 항목 1", "{{세부항목1}}"),
                ("세부 항목 2", "{{세부항목2}}"),
            ]),
            ("5. 시사점", [
                ("주요 시사점", "{{주요시사점}}"),
                ("개선 방향", "{{개선방향}}"),
            ]),
            ("6. 결론", [
                ("종합 결론", "{{종합결론}}"),
                ("향후 계획", "{{향후계획}}"),
            ]),
        ]
    )

    # ── 2. 비교 분석 (8슬라이드) ──
    _build_pptx(
        "report_analysis_pptx.pptx",
        title="{{제목}}",
        subtitle="비교 분석 보고서  |  {{부서}}",
        slides_data=[
            ("1. 분석 대상", [
                ("분석 대상", "{{분석대상}}"),
                ("분석 목적", "{{분석목적}}"),
            ]),
            ("2. 비교 항목", [
                ("비교 기준", "{{비교기준}}"),
                ("비교 항목", "{{비교항목}}"),
            ]),
            ("3. 분석 결과", [
                ("정량 분석", "{{정량분석}}"),
                ("정성 분석", "{{정성분석}}"),
            ]),
            ("4. 강점", [
                ("주요 강점", "{{주요강점}}"),
                ("차별화 요소", "{{차별화요소}}"),
            ]),
            ("5. 약점", [
                ("주요 약점", "{{주요약점}}"),
                ("개선 필요사항", "{{개선필요사항}}"),
            ]),
            ("6. 결론", [
                ("종합 평가", "{{종합평가}}"),
                ("권고 사항", "{{권고사항}}"),
            ]),
        ]
    )

    # ── 3. 성과 보고서 (8슬라이드) ──
    _build_pptx(
        "report_performance_pptx.pptx",
        title="{{제목}}",
        subtitle="성과 보고서  |  {{부서}}",
        slides_data=[
            ("1. 목표", [
                ("사업 목표", "{{사업목표}}"),
                ("핵심 지표(KPI)", "{{핵심지표}}"),
            ]),
            ("2. 달성 현황", [
                ("목표 대비 달성률", "{{달성률}}"),
                ("주요 실적", "{{주요실적}}"),
            ]),
            ("3. 핵심 성과", [
                ("정량 성과", "{{정량성과}}"),
                ("정성 성과", "{{정성성과}}"),
            ]),
            ("4. 성과 상세", [
                ("분야별 성과", "{{분야별성과}}"),
                ("우수 사례", "{{우수사례}}"),
            ]),
            ("5. 개선 사항", [
                ("미달 항목", "{{미달항목}}"),
                ("개선 방안", "{{개선방안}}"),
            ]),
            ("6. 향후 계획", [
                ("차기 목표", "{{차기목표}}"),
                ("실행 계획", "{{실행계획}}"),
            ]),
        ]
    )

    # ── 4. 주간 보고 (6슬라이드) ──
    _build_pptx(
        "report_weekly_pptx.pptx",
        title="{{제목}}",
        subtitle="주간 업무 보고  |  {{부서}}",
        slides_data=[
            ("1. 주요 업무", [
                ("금주 업무 요약", "{{금주업무요약}}"),
                ("완료 업무", "{{완료업무}}"),
                ("진행중 업무", "{{진행중업무}}"),
            ]),
            ("2. 이슈 사항", [
                ("주요 이슈", "{{주요이슈}}"),
                ("대응 현황", "{{대응현황}}"),
            ]),
            ("3. 다음주 계획", [
                ("예정 업무", "{{예정업무}}"),
                ("목표", "{{다음주목표}}"),
            ]),
            ("4. 요청 사항", [
                ("지원 요청", "{{지원요청}}"),
                ("협조 사항", "{{협조사항}}"),
            ]),
        ]
    )

    # ── 5. 사업 제안서 (10슬라이드) ──
    _build_pptx(
        "proposal_business_pptx.pptx",
        title="{{제목}}",
        subtitle="사업 제안서  |  {{회사명}}",
        slides_data=[
            ("1. 사업 개요", [
                ("사업명", "{{사업명}}"),
                ("사업 목적", "{{사업목적}}"),
            ]),
            ("2. 추진 배경", [
                ("현황 및 문제점", "{{현황및문제점}}"),
                ("추진 필요성", "{{추진필요성}}"),
            ]),
            ("3. 제안 내용", [
                ("핵심 제안", "{{핵심제안}}"),
                ("제안 범위", "{{제안범위}}"),
            ]),
            ("4. 기술 구성", [
                ("기술 아키텍처", "{{기술아키텍처}}"),
                ("핵심 기술", "{{핵심기술}}"),
            ]),
            ("5. 추진 일정", [
                ("단계별 일정", "{{단계별일정}}"),
                ("마일스톤", "{{마일스톤}}"),
            ]),
            ("6. 기대 효과", [
                ("정량적 효과", "{{정량적효과}}"),
                ("정성적 효과", "{{정성적효과}}"),
            ]),
            ("7. 예산", [
                ("총 예산", "{{총예산}}"),
                ("항목별 내역", "{{항목별내역}}"),
            ]),
            ("8. 요약", [
                ("핵심 요약", "{{핵심요약}}"),
                ("기대 성과", "{{기대성과}}"),
            ]),
        ]
    )

    # ── 6. 기술 제안서 (10슬라이드) ──
    _build_pptx(
        "proposal_tech_pptx.pptx",
        title="{{제목}}",
        subtitle="기술 제안서  |  {{회사명}}",
        slides_data=[
            ("1. 기술 개요", [
                ("기술명", "{{기술명}}"),
                ("기술 개요", "{{기술개요}}"),
            ]),
            ("2. 현황 분석", [
                ("현재 시스템 현황", "{{현재시스템현황}}"),
                ("문제점 및 한계", "{{문제점및한계}}"),
            ]),
            ("3. 제안 기술", [
                ("핵심 기술", "{{핵심기술}}"),
                ("기술 차별점", "{{기술차별점}}"),
            ]),
            ("4. 아키텍처", [
                ("시스템 아키텍처", "{{시스템아키텍처}}"),
                ("데이터 흐름", "{{데이터흐름}}"),
            ]),
            ("5. 구현 방안", [
                ("개발 방법론", "{{개발방법론}}"),
                ("구현 전략", "{{구현전략}}"),
            ]),
            ("6. 추진 일정", [
                ("단계별 일정", "{{단계별일정}}"),
                ("주요 마일스톤", "{{주요마일스톤}}"),
            ]),
            ("7. 기대 효과", [
                ("기술적 효과", "{{기술적효과}}"),
                ("업무 효과", "{{업무효과}}"),
            ]),
            ("8. 요약", [
                ("제안 요약", "{{제안요약}}"),
                ("차별화 포인트", "{{차별화포인트}}"),
            ]),
        ]
    )

    # ── 7. 서비스 제안서 (8슬라이드) ──
    _build_pptx(
        "proposal_service_pptx.pptx",
        title="{{제목}}",
        subtitle="서비스 제안서  |  {{회사명}}",
        slides_data=[
            ("1. 서비스 개요", [
                ("서비스명", "{{서비스명}}"),
                ("서비스 개요", "{{서비스개요}}"),
                ("대상 고객", "{{대상고객}}"),
            ]),
            ("2. 특장점", [
                ("핵심 특장점", "{{핵심특장점}}"),
                ("경쟁 우위", "{{경쟁우위}}"),
            ]),
            ("3. 서비스 구성도", [
                ("서비스 구성", "{{서비스구성}}"),
                ("운영 체계", "{{운영체계}}"),
            ]),
            ("4. 도입 효과", [
                ("기대 효과", "{{기대효과}}"),
                ("도입 사례", "{{도입사례}}"),
            ]),
            ("5. 운영 방안", [
                ("운영 계획", "{{운영계획}}"),
                ("SLA 기준", "{{SLA기준}}"),
            ]),
            ("6. 요약", [
                ("제안 요약", "{{제안요약}}"),
                ("Next Step", "{{NextStep}}"),
            ]),
        ]
    )

    # ── 8. 교육 제안서 (8슬라이드) ──
    _build_pptx(
        "proposal_edu_pptx.pptx",
        title="{{제목}}",
        subtitle="교육 제안서  |  {{기관명}}",
        slides_data=[
            ("1. 교육 개요", [
                ("교육명", "{{교육명}}"),
                ("교육 목적", "{{교육목적}}"),
                ("교육 대상", "{{교육대상}}"),
            ]),
            ("2. 교육 과정", [
                ("커리큘럼", "{{커리큘럼}}"),
                ("교육 내용", "{{교육내용}}"),
            ]),
            ("3. 교육 방법", [
                ("교육 방식", "{{교육방식}}"),
                ("교육 도구", "{{교육도구}}"),
            ]),
            ("4. 교육 일정", [
                ("전체 일정", "{{전체일정}}"),
                ("차시별 계획", "{{차시별계획}}"),
            ]),
            ("5. 기대 효과", [
                ("학습 성과", "{{학습성과}}"),
                ("업무 적용 효과", "{{업무적용효과}}"),
            ]),
            ("6. 요약", [
                ("교육 요약", "{{교육요약}}"),
                ("수료 기준", "{{수료기준}}"),
            ]),
        ]
    )


# ═══════════════════════════════════════════════════════════════════════════════
#  DOCX 템플릿 8종 정의
# ═══════════════════════════════════════════════════════════════════════════════

def create_all_docx():
    """
    DOCX 템플릿 8종을 모두 생성합니다.
    각 템플릿은 커버 페이지 + 섹션들로 구성됩니다 (회의록만 예외).
    """
    print("\n[DOCX 템플릿 생성 중...]")

    # ── 1. 요약 보고서 (6섹션) ──
    _build_docx(
        "report_summary_docx.docx",
        doc_title="요약 보고서",
        cover_title="{{제목}}",
        cover_subtitles=[
            "작성자: {{작성자}}  |  부서: {{부서}}",
            "작성일: {{작성일}}",
        ],
        sections=[
            ("1. 개요", "{{개요}}"),
            ("2. 주요 내용", "{{주요내용}}"),
            ("3. 분석 결과", "{{분석결과}}"),
            ("4. 세부 내용", "{{세부내용}}"),
            ("5. 시사점", "{{시사점}}"),
            ("6. 결론", "{{결론}}"),
        ]
    )

    # ── 2. 비교 분석 (6섹션) ──
    _build_docx(
        "report_analysis_docx.docx",
        doc_title="비교 분석 보고서",
        cover_title="{{제목}}",
        cover_subtitles=[
            "비교 분석 보고서",
            "작성자: {{작성자}}  |  작성일: {{작성일}}",
        ],
        sections=[
            ("1. 분석 대상", "{{분석대상}}"),
            ("2. 비교 항목", "{{비교항목}}"),
            ("3. 분석 결과", "{{분석결과}}"),
            ("4. 강점", "{{강점}}"),
            ("5. 약점", "{{약점}}"),
            ("6. 결론", "{{결론}}"),
        ]
    )

    # ── 3. 회의록 (특수 구조: 정보 테이블 + 5섹션) ──
    _build_docx_meeting("report_meeting_docx.docx")

    # ── 4. 주간 보고 (4섹션) ──
    _build_docx(
        "report_weekly_docx.docx",
        doc_title="주간 업무 보고",
        cover_title="{{제목}}",
        cover_subtitles=[
            "주간 업무 보고",
            "작성자: {{작성자}}  |  부서: {{부서}}",
            "보고 기간: {{보고기간}}",
        ],
        sections=[
            ("1. 주요 업무", "{{주요업무}}"),
            ("2. 이슈 사항", "{{이슈사항}}"),
            ("3. 다음주 계획", "{{다음주계획}}"),
            ("4. 요청 사항", "{{요청사항}}"),
        ]
    )

    # ── 5. 사업 제안서 (8섹션) ──
    _build_docx(
        "proposal_business_docx.docx",
        doc_title="사업 제안서",
        cover_title="{{제목}}",
        cover_subtitles=[
            "사업 제안서",
            "{{회사명}}  |  {{작성일}}",
        ],
        sections=[
            ("1. 사업 개요", "{{사업개요}}"),
            ("2. 추진 배경", "{{추진배경}}"),
            ("3. 제안 내용", "{{제안내용}}"),
            ("4. 기술 구성", "{{기술구성}}"),
            ("5. 추진 일정", "{{추진일정}}"),
            ("6. 기대 효과", "{{기대효과}}"),
            ("7. 예산", "{{예산}}"),
            ("8. 요약", "{{요약}}"),
        ]
    )

    # ── 6. 기술 제안서 (8섹션) ──
    _build_docx(
        "proposal_tech_docx.docx",
        doc_title="기술 제안서",
        cover_title="{{제목}}",
        cover_subtitles=[
            "기술 제안서",
            "{{회사명}}  |  {{작성일}}",
        ],
        sections=[
            ("1. 기술 개요", "{{기술개요}}"),
            ("2. 현황 분석", "{{현황분석}}"),
            ("3. 제안 기술", "{{제안기술}}"),
            ("4. 아키텍처", "{{아키텍처}}"),
            ("5. 구현 방안", "{{구현방안}}"),
            ("6. 추진 일정", "{{일정}}"),
            ("7. 기대 효과", "{{기대효과}}"),
            ("8. 요약", "{{요약}}"),
        ]
    )

    # ── 7. R&D 제안서 (6섹션) ──
    _build_docx(
        "proposal_rd_docx.docx",
        doc_title="R&D 제안서",
        cover_title="{{제목}}",
        cover_subtitles=[
            "R&D 제안서",
            "{{기관명}}  |  {{작성일}}",
        ],
        sections=[
            ("1. 연구 목적", "{{연구목적}}"),
            ("2. 연구 내용", "{{연구내용}}"),
            ("3. 연구 방법", "{{연구방법}}"),
            ("4. 기대 효과", "{{기대효과}}"),
            ("5. 예산", "{{예산}}"),
            ("6. 요약", "{{요약}}"),
        ]
    )

    # ── 8. 교육 제안서 (6섹션) ──
    _build_docx(
        "proposal_edu_docx.docx",
        doc_title="교육 제안서",
        cover_title="{{제목}}",
        cover_subtitles=[
            "교육 제안서",
            "{{기관명}}  |  {{작성일}}",
        ],
        sections=[
            ("1. 교육 개요", "{{교육개요}}"),
            ("2. 교육 과정", "{{교육과정}}"),
            ("3. 교육 방법", "{{교육방법}}"),
            ("4. 교육 일정", "{{교육일정}}"),
            ("5. 기대 효과", "{{기대효과}}"),
            ("6. 요약", "{{요약}}"),
        ]
    )


# ═══════════════════════════════════════════════════════════════════════════════
#  메인 실행
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    """
    메인 함수: 출력 디렉토리를 준비하고 모든 템플릿을 생성합니다.

    실행 순서:
        1. sample_templates/ 디렉토리 생성 (없으면)
        2. 기존 파일 삭제
        3. PPTX 8종 생성
        4. DOCX 8종 생성
        5. 결과 요약 출력
    """
    print("=" * 60)
    print("  IDINO 기업용 문서 템플릿 생성기")
    print("=" * 60)

    # 출력 디렉토리 준비
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    # 기존 파일 삭제 (깨끗한 상태에서 시작)
    existing = list(OUTPUT_DIR.glob("*.pptx")) + list(OUTPUT_DIR.glob("*.docx"))
    if existing:
        print(f"\n기존 파일 {len(existing)}개 삭제 중...")
        for f in existing:
            f.unlink()
            print(f"  삭제: {f.name}")

    # PPTX 템플릿 8종 생성
    create_all_pptx()

    # DOCX 템플릿 8종 생성
    create_all_docx()

    # 결과 요약
    created = list(OUTPUT_DIR.glob("*.pptx")) + list(OUTPUT_DIR.glob("*.docx"))
    print(f"\n{'=' * 60}")
    print(f"  완료! 총 {len(created)}개 파일 생성됨")
    print(f"  위치: {OUTPUT_DIR}")
    print(f"{'=' * 60}")

    # 생성된 파일 목록 출력
    for f in sorted(OUTPUT_DIR.iterdir()):
        size_kb = f.stat().st_size / 1024
        print(f"  {f.name:45s} {size_kb:6.1f} KB")


if __name__ == "__main__":
    main()
