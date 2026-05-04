"""
Jinja2 템플릿 엔진 모듈.

DOCX/PPTX 파일 내 Jinja2 변수를 추출하고, 템플릿을 렌더링하며,
AI 분석 결과를 기반으로 일반 문서를 Jinja2 템플릿으로 변환하는 기능을 제공한다.

주요 기능:
  1. extract_docx_variables  - DOCX 파일에서 Jinja2 변수 추출
  2. extract_pptx_variables  - PPTX 파일에서 Jinja2 변수 추출
  3. render_docx_jinja2      - DOCX 템플릿에 데이터를 채워 렌더링
  4. render_pptx_jinja2      - PPTX 템플릿에 데이터를 채워 렌더링
  5. convert_example_to_template - 일반 문서를 Jinja2 템플릿으로 변환
  6. analyze_blank_form      - 빈 양식 문서의 구조를 분석
  7. auto_generate_jinja2_from_structure - 빈 양식에 Jinja2 변수 자동 삽입
"""

from __future__ import annotations

import io
import logging
import re
from typing import Any

import jinja2
from docx.shared import Mm
from docxtpl import DocxTemplate, InlineImage
from pptx import Presentation

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# 정규식 패턴 정의
# ---------------------------------------------------------------------------

# {{ variable }} 형태의 변수 패턴 (점 표기법 포함: item.name 등)
VAR_PATTERN = re.compile(r"\{\{\s*(\w+(?:\.\w+)*)\s*\}\}")

# {% for item in items %} 형태의 반복문 패턴
FOR_PATTERN = re.compile(r"\{%[-\s]*for\s+\w+\s+in\s+(\w+(?:\.\w+)*)\s*[-\s]*%\}")

# {% if condition %} 형태의 조건문 패턴
IF_PATTERN = re.compile(r"\{%[-\s]*if\s+(\w+(?:\.\w+)*)\s*[-\s]*%\}")


# ---------------------------------------------------------------------------
# 변수 카테고리 분류 상수 및 함수
# ---------------------------------------------------------------------------
# 변수를 3가지 카테고리로 분류한다:
#   - "session_auto": 로그인한 사용자 세션 정보에서 자동으로 채울 수 있는 변수
#     (예: 소속 부서명, 작성자 이름 등)
#   - "user_input": 사용자가 보고서 생성 시 직접 입력해야 하는 최소 항목
#     (예: 보고서명, 기간, 참석자 등)
#   - "ai_generated": AI가 소스 문서 내용을 분석하여 자동 생성하는 항목
#     (예: 보고내용, 문제점, 기대효과 등)
#
# 분류는 변수 이름에 포함된 한글 키워드를 기반으로 자동 판별한다.

# session_auto: 로그인 세션(사용자/부서 정보)에서 자동 채울 수 있는 키워드
SESSION_AUTO_KEYWORDS = ["소속", "작성자", "성명", "부서", "이름"]

# user_input: 사용자가 직접 입력해야 하는 최소 항목 키워드
USER_INPUT_KEYWORDS = [
    "보고서명",
    "사업명",
    "제목",
    "회의제목",
    "주제",
    "참석자",
    "일시",
    "회의일자",
    "장소",
    "날짜",
    "기간",
    "보고기간",
]

# date: 날짜 입력이 필요한 변수 (캘린더 picker 사용)
DATE_KEYWORDS = ["일시", "일자", "날짜", "기간", "보고기간", "검토일", "제안일", "보고일"]


def classify_variable_category(var_name: str) -> str:
    """변수 이름을 분석하여 카테고리를 자동 분류한다.

    변수 이름에 포함된 한글 키워드를 검사하여 3가지 카테고리 중 하나를 반환한다.
    키워드 매칭은 '포함(in)' 방식으로, 변수 이름 어디에나 키워드가 있으면 매칭된다.
    예: "보고기간" → "보고기간" 키워드 매칭 → "user_input"
        "소속부서" → "소속" 키워드 매칭 → "session_auto"

    우선순위: session_auto > user_input > ai_generated (기본값)

    Args:
        var_name: Jinja2 변수 이름 (예: "소속", "보고서명", "보고내용과의견")

    Returns:
        "session_auto", "user_input", 또는 "ai_generated" 중 하나
    """
    # 점 표기법(item.속성) 변수는 반복 항목의 속성이므로 AI 생성으로 분류
    if "." in var_name:
        return "ai_generated"

    # 1순위: session_auto — 세션 정보로 자동 채울 수 있는 변수인지 확인
    for keyword in SESSION_AUTO_KEYWORDS:
        if keyword in var_name:
            return "session_auto"

    # 2순위: user_input — 사용자가 직접 입력해야 하는 변수인지 확인
    for keyword in USER_INPUT_KEYWORDS:
        if keyword in var_name:
            return "user_input"

    # 3순위: 위 키워드에 해당하지 않으면 AI가 자동 생성할 내용으로 분류
    return "ai_generated"


def classify_variable_type(var_name: str, default_type: str = "string") -> str:
    """변수 이름으로 입력 타입을 자동 분류한다.

    날짜 관련 키워드가 포함되면 "date" 타입으로 반환하여
    프론트에서 캘린더 picker를 사용하도록 한다.

    Args:
        var_name: 변수 이름
        default_type: 기본 타입 (키워드 미매칭 시)

    Returns:
        "date" 또는 default_type
    """
    for keyword in DATE_KEYWORDS:
        if keyword in var_name:
            return "date"
    return default_type


# ---------------------------------------------------------------------------
# 내부 헬퍼 함수들
# ---------------------------------------------------------------------------


def _classify_variable_type(name: str, source: str) -> str:
    """변수가 어떤 맥락에서 발견되었는지에 따라 타입을 추정한다.

    Args:
        name: 변수 이름 (예: "items", "title")
        source: 변수가 발견된 맥락 ("for" = 반복문, "if" = 조건문, "var" = 일반 변수)

    Returns:
        추정된 타입 문자열 ("array", "boolean", "string")
    """
    if source == "for":
        # for 반복문에 사용되면 배열(리스트) 타입
        return "array"
    elif source == "if":
        # if 조건문에 사용되면 불린 타입 (단, 단순 존재 체크일 수도 있음)
        return "boolean"
    else:
        # 그 외 일반 변수는 문자열 타입으로 추정
        return "string"


def _extract_variables_from_text(text: str) -> list[dict[str, str]]:
    """텍스트 문자열에서 모든 Jinja2 변수를 추출하여 중복 없이 반환한다.

    {{ var }}, {% for x in list %}, {% if cond %} 세 가지 패턴을 모두 검색한다.
    동일 변수가 여러 번 등장하면 첫 번째 발견 시의 타입을 사용한다.

    Args:
        text: Jinja2 패턴이 포함된 전체 텍스트

    Returns:
        [{"name": "변수명", "type": "string|array|boolean"}, ...] 형태의 리스트
    """
    # 발견된 변수명 → 타입 매핑 (중복 방지용, 삽입 순서 보존)
    seen: dict[str, str] = {}

    # 1) {% for item in items %} 에서 "items"를 array 타입으로 추출
    for match in FOR_PATTERN.finditer(text):
        var_name = match.group(1)
        if var_name not in seen:
            seen[var_name] = _classify_variable_type(var_name, "for")

    # 2) {% if condition %} 에서 "condition"을 boolean 타입으로 추출
    for match in IF_PATTERN.finditer(text):
        var_name = match.group(1)
        if var_name not in seen:
            seen[var_name] = _classify_variable_type(var_name, "if")

    # 3) {{ variable }} 에서 일반 변수를 string 타입으로 추출
    for match in VAR_PATTERN.finditer(text):
        var_name = match.group(1)
        # 점 표기법(item.name)의 경우 루트 변수(item)만 등록하지 않고
        # 전체 경로를 그대로 등록한다 (UI에서 중첩 구조를 보여줄 수 있도록)
        if var_name not in seen:
            seen[var_name] = _classify_variable_type(var_name, "var")

    # 각 변수에 카테고리(category) 정보를 추가하여 반환한다
    # 카테고리는 변수 이름의 키워드를 분석하여 자동 분류된다
    return [
        {
            "name": name,
            "type": vtype,
            "category": classify_variable_category(name),
        }
        for name, vtype in seen.items()
    ]


def _collect_docx_text(template_bytes: bytes) -> str:
    """DOCX 파일에서 모든 텍스트를 추출하여 하나의 문자열로 합친다.

    paragraphs(본문 문단)와 tables(표) 양쪽 모두에서 텍스트를 수집한다.
    docxtpl이 아닌 python-docx를 사용하여 원본 XML 텍스트를 그대로 읽는다.

    Args:
        template_bytes: DOCX 파일의 바이트 데이터

    Returns:
        모든 텍스트가 줄바꿈으로 연결된 문자열
    """
    from docx import Document

    doc = Document(io.BytesIO(template_bytes))
    parts: list[str] = []

    # 본문 문단에서 텍스트 수집
    for paragraph in doc.paragraphs:
        parts.append(paragraph.text)

    # 표(table)의 각 셀에서 텍스트 수집
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    parts.append(paragraph.text)

    return "\n".join(parts)


def _collect_pptx_text(prs: Presentation) -> str:
    """PPTX Presentation 객체에서 모든 텍스트를 추출하여 하나의 문자열로 합친다.

    각 슬라이드의 모든 shape에서 text_frame과 table을 검색한다.
    XML run 분리 문제를 해결하기 위해 paragraph 단위로 전체 텍스트를 조합한다.

    Args:
        prs: python-pptx Presentation 객체

    Returns:
        모든 텍스트가 줄바꿈으로 연결된 문자열
    """
    parts: list[str] = []

    for slide in prs.slides:
        for shape in slide.shapes:
            # text_frame이 있는 shape에서 텍스트 추출
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    # paragraph.text는 모든 run의 텍스트를 합쳐서 반환
                    # 이렇게 하면 {{ var }}가 여러 run으로 분리된 경우에도 올바르게 인식
                    parts.append(paragraph.text)

            # 표(table)가 있는 shape에서 텍스트 추출
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            parts.append(paragraph.text)

    return "\n".join(parts)


# ---------------------------------------------------------------------------
# 공개 API 함수들
# ---------------------------------------------------------------------------


def extract_docx_variables(file_bytes: bytes) -> dict:
    """DOCX 파일에서 Jinja2 템플릿 변수를 추출한다.

    docxtpl(DocxTemplate)로 파일을 로드한 뒤, 문서 내 모든 텍스트
    (paragraphs + tables)에서 Jinja2 변수 패턴을 찾아 반환한다.

    지원하는 패턴:
      - {{ variable }}         → type: "string"
      - {% for item in list %} → type: "array"
      - {% if condition %}     → type: "boolean"

    Args:
        file_bytes: DOCX 파일의 바이트 데이터

    Returns:
        {"variables": [{"name": "title", "type": "string"}, ...]}

    Raises:
        ValueError: DOCX 파일을 읽을 수 없는 경우
    """
    try:
        logger.info("DOCX 파일에서 Jinja2 변수 추출 시작")

        # DOCX 파일의 모든 텍스트를 수집
        full_text = _collect_docx_text(file_bytes)

        # 텍스트에서 변수 패턴 추출
        variables = _extract_variables_from_text(full_text)

        logger.info("DOCX 변수 추출 완료: %d개 변수 발견", len(variables))
        return {"variables": variables}

    except Exception as e:
        logger.error("DOCX 변수 추출 실패: %s", str(e))
        raise ValueError(f"DOCX 파일에서 변수를 추출할 수 없습니다: {e}") from e


def extract_pptx_variables(file_bytes: bytes) -> dict:
    """PPTX 파일에서 Jinja2 템플릿 변수를 추출한다.

    python-pptx로 Presentation을 로드한 뒤, 모든 슬라이드의
    텍스트프레임과 표 셀에서 Jinja2 변수 패턴을 찾아 반환한다.

    paragraph 단위로 전체 텍스트를 조합한 뒤 파싱하여,
    XML run이 분리되어 {{ 와 variable }}이 다른 run에 있는 경우에도
    올바르게 변수를 인식한다.

    Args:
        file_bytes: PPTX 파일의 바이트 데이터

    Returns:
        {"variables": [{"name": "title", "type": "string"}, ...]}

    Raises:
        ValueError: PPTX 파일을 읽을 수 없는 경우
    """
    try:
        logger.info("PPTX 파일에서 Jinja2 변수 추출 시작")

        # PPTX 파일 로드
        prs = Presentation(io.BytesIO(file_bytes))

        # 모든 슬라이드 텍스트 수집
        full_text = _collect_pptx_text(prs)

        # 텍스트에서 변수 패턴 추출
        variables = _extract_variables_from_text(full_text)

        logger.info("PPTX 변수 추출 완료: %d개 변수 발견", len(variables))
        return {"variables": variables}

    except Exception as e:
        logger.error("PPTX 변수 추출 실패: %s", str(e))
        raise ValueError(f"PPTX 파일에서 변수를 추출할 수 없습니다: {e}") from e


def render_docx_jinja2(template_bytes: bytes, context: dict) -> bytes:
    """DOCX Jinja2 템플릿에 데이터를 채워 렌더링한다.

    docxtpl 라이브러리를 사용하여 템플릿을 로드하고,
    context 딕셔너리의 값으로 변수를 치환한다.

    context에 이미지 데이터가 포함된 경우 (bytes 타입 값),
    자동으로 docxtpl.InlineImage로 변환하여 문서에 삽입한다.

    Args:
        template_bytes: DOCX 템플릿 파일의 바이트 데이터
        context: 변수명 → 값 매핑 딕셔너리
                 예: {"title": "보고서 제목", "items": [{"name": "항목1"}, ...]}

    Returns:
        렌더링된 DOCX 파일의 바이트 데이터

    Raises:
        ValueError: 템플릿 로드 또는 렌더링에 실패한 경우
    """
    try:
        logger.info("DOCX Jinja2 렌더링 시작")

        # 1단계: docxtpl로 템플릿 로드
        tpl = DocxTemplate(io.BytesIO(template_bytes))

        # 2단계: context에서 이미지 데이터(bytes)를 InlineImage로 변환
        processed_context = _process_context_images(tpl, context)

        # 3단계: 템플릿 렌더링 (변수 치환)
        tpl.render(processed_context)

        # 4단계: 렌더링 결과를 바이트로 저장
        output = io.BytesIO()
        tpl.save(output)
        output.seek(0)

        logger.info("DOCX Jinja2 렌더링 완료")
        return output.read()

    except Exception as e:
        logger.error("DOCX 렌더링 실패: %s", str(e))
        raise ValueError(f"DOCX 템플릿 렌더링에 실패했습니다: {e}") from e


def _process_context_images(tpl: DocxTemplate, context: dict) -> dict:
    """context 딕셔너리에서 이미지 데이터를 찾아 InlineImage 객체로 변환한다.

    이미지 데이터는 다음 형태로 전달될 수 있다:
      - bytes 타입: 이미지 바이너리 데이터
      - dict 타입 {"image_data": bytes, "width_mm": int}: 크기 지정 포함

    Args:
        tpl: docxtpl DocxTemplate 객체 (InlineImage 생성에 필요)
        context: 원본 context 딕셔너리

    Returns:
        이미지가 InlineImage로 변환된 새 context 딕셔너리
    """
    processed = {}

    for key, value in context.items():
        if isinstance(value, bytes):
            # bytes 데이터를 이미지로 간주하여 InlineImage로 변환
            # 기본 너비 100mm
            try:
                processed[key] = InlineImage(tpl, io.BytesIO(value), width=Mm(100))
                logger.debug("이미지 변수 변환: %s (기본 크기)", key)
            except Exception as img_err:
                logger.warning("이미지 변환 실패 (%s), 원본 값 유지: %s", key, img_err)
                processed[key] = value

        elif isinstance(value, dict) and "image_data" in value:
            # {"image_data": bytes, "width_mm": int} 형태의 이미지 데이터
            try:
                width = Mm(value.get("width_mm", 100))
                processed[key] = InlineImage(tpl, io.BytesIO(value["image_data"]), width=width)
                logger.debug("이미지 변수 변환: %s (너비 %dmm)", key, value.get("width_mm", 100))
            except Exception as img_err:
                logger.warning("이미지 변환 실패 (%s), 원본 값 유지: %s", key, img_err)
                processed[key] = value

        elif isinstance(value, list):
            # 리스트 내부에도 이미지가 있을 수 있으므로 재귀 처리
            processed[key] = [_process_context_images(tpl, item) if isinstance(item, dict) else item for item in value]

        elif isinstance(value, dict):
            # 중첩 딕셔너리도 재귀 처리
            processed[key] = _process_context_images(tpl, value)

        else:
            # 일반 값은 그대로 전달
            processed[key] = value

    return processed


def render_pptx_jinja2(template_bytes: bytes, context: dict) -> bytes:
    """PPTX 템플릿에 Jinja2 변수를 치환하여 렌더링한다.

    python-pptx로 Presentation을 로드한 뒤, 모든 슬라이드를 순회하면서
    각 shape의 text_frame에서 Jinja2 패턴을 찾아 치환한다.

    서식 보존 전략:
      - paragraph의 전체 텍스트를 조합하여 Jinja2 렌더링 수행
      - 렌더링 결과를 첫 번째 run에 할당하고, 나머지 run은 빈 문자열로 설정
      - 이렇게 하면 첫 번째 run의 서식(폰트, 크기, 색상 등)이 유지됨

    Args:
        template_bytes: PPTX 템플릿 파일의 바이트 데이터
        context: 변수명 → 값 매핑 딕셔너리

    Returns:
        렌더링된 PPTX 파일의 바이트 데이터

    Raises:
        ValueError: 템플릿 로드 또는 렌더링에 실패한 경우
    """
    try:
        logger.info("PPTX Jinja2 렌더링 시작")

        # 1단계: PPTX 파일 로드
        prs = Presentation(io.BytesIO(template_bytes))

        # 2단계: 모든 슬라이드 순회하며 텍스트 치환
        for _slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                # text_frame이 있는 shape 처리
                if shape.has_text_frame:
                    _render_text_frame(shape.text_frame, context)

                # 표(table)가 있는 shape 처리
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            _render_text_frame(cell.text_frame, context)

        # 3단계: 렌더링 결과를 바이트로 저장
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)

        logger.info("PPTX Jinja2 렌더링 완료")
        return output.read()

    except Exception as e:
        logger.error("PPTX 렌더링 실패: %s", str(e))
        raise ValueError(f"PPTX 템플릿 렌더링에 실패했습니다: {e}") from e


def _render_text_frame(text_frame: Any, context: dict) -> None:
    """하나의 text_frame 내 모든 paragraph에 대해 Jinja2 렌더링을 수행한다.

    각 paragraph의 전체 텍스트를 조합한 뒤 Jinja2 패턴이 포함되어 있는지 확인하고,
    패턴이 있으면 렌더링하여 결과를 첫 번째 run에 할당한다.

    Args:
        text_frame: python-pptx TextFrame 객체
        context: 변수명 → 값 매핑 딕셔너리
    """
    for paragraph in text_frame.paragraphs:
        # paragraph 내 모든 run의 텍스트를 합쳐서 전체 텍스트 생성
        full_text = "".join(run.text for run in paragraph.runs)

        # Jinja2 패턴이 포함되어 있는지 빠르게 확인 (성능 최적화)
        if "{{" not in full_text and "{%" not in full_text:
            continue

        try:
            # Jinja2 렌더링 수행
            # undefined=jinja2.Undefined: 정의되지 않은 변수는 빈 문자열로 처리
            env = jinja2.Environment(undefined=jinja2.Undefined)
            template = env.from_string(full_text)
            rendered = template.render(context)

            # 렌더링 결과를 run에 반영 (서식 보존)
            if paragraph.runs:
                # 첫 번째 run에 렌더링 결과 전체를 할당
                paragraph.runs[0].text = rendered
                # 나머지 run은 빈 문자열로 설정 (중복 방지)
                for run in paragraph.runs[1:]:
                    run.text = ""

        except jinja2.TemplateError as te:
            # Jinja2 문법 오류 시 원본 텍스트 유지
            logger.warning(
                "PPTX 텍스트 렌더링 건너뜀 (Jinja2 오류): %s / 원본: %s",
                str(te),
                full_text[:100],
            )
        except Exception as e:
            logger.warning(
                "PPTX 텍스트 렌더링 건너뜀 (일반 오류): %s / 원본: %s",
                str(e),
                full_text[:100],
            )


def convert_example_to_template(
    file_bytes: bytes,
    format: str,
    ai_analysis: dict,
) -> bytes:
    """일반 문서를 Jinja2 템플릿으로 변환한다.

    AI 분석 결과(ai_analysis)에 포함된 텍스트 → 변수 매핑 정보를 사용하여,
    문서 내 특정 텍스트를 {{ variable }} 패턴으로 치환한다.

    ai_analysis 형식 예시:
        {
            "replacements": [
                {"original": "2026년 1분기", "variable": "period"},
                {"original": "홍길동", "variable": "author_name"},
                {"original": "프로젝트 현황 보고", "variable": "report_title"},
            ]
        }

    Args:
        file_bytes: 원본 DOCX 또는 PPTX 파일의 바이트 데이터
        format: 파일 형식 ("docx" 또는 "pptx")
        ai_analysis: AI가 분석한 텍스트 → 변수 매핑 정보

    Returns:
        Jinja2 변수로 치환된 템플릿 파일의 바이트 데이터

    Raises:
        ValueError: 지원하지 않는 파일 형식이거나 변환에 실패한 경우
    """
    # ai_analysis에서 치환 매핑 목록을 가져온다
    replacements = ai_analysis.get("replacements", [])
    if not replacements:
        logger.warning("AI 분석 결과에 치환 매핑이 없습니다. 원본 파일을 그대로 반환합니다.")
        return file_bytes

    logger.info(
        "문서 → 템플릿 변환 시작 (형식: %s, 치환 항목: %d개)",
        format,
        len(replacements),
    )

    # 파일 형식에 따라 적절한 변환 함수 호출
    normalized_format = format.lower().strip().lstrip(".")

    if normalized_format == "docx":
        return _convert_docx_to_template(file_bytes, replacements)
    elif normalized_format == "pptx":
        return _convert_pptx_to_template(file_bytes, replacements)
    else:
        raise ValueError(f"지원하지 않는 파일 형식입니다: {format} (docx 또는 pptx만 지원)")


def _convert_docx_to_template(
    file_bytes: bytes,
    replacements: list[dict[str, str]],
) -> bytes:
    """DOCX 파일 내 지정된 텍스트를 {{ variable }} 패턴으로 치환한다.

    python-docx를 사용하여 문서의 paragraphs와 tables를 순회하면서
    원본 텍스트를 Jinja2 변수 패턴으로 교체한다.

    Args:
        file_bytes: 원본 DOCX 파일 바이트 데이터
        replacements: [{"original": "원본 텍스트", "variable": "변수명"}, ...]

    Returns:
        변환된 DOCX 템플릿 파일의 바이트 데이터
    """
    from docx import Document

    try:
        doc = Document(io.BytesIO(file_bytes))

        # 본문 문단에서 텍스트 치환
        for paragraph in doc.paragraphs:
            _replace_in_paragraph(paragraph, replacements)

        # 표(table)의 각 셀에서 텍스트 치환
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        _replace_in_paragraph(paragraph, replacements)

        # 결과를 바이트로 저장
        output = io.BytesIO()
        doc.save(output)
        output.seek(0)

        logger.info("DOCX → 템플릿 변환 완료")
        return output.read()

    except Exception as e:
        logger.error("DOCX → 템플릿 변환 실패: %s", str(e))
        raise ValueError(f"DOCX 템플릿 변환에 실패했습니다: {e}") from e


def _convert_pptx_to_template(
    file_bytes: bytes,
    replacements: list[dict[str, str]],
) -> bytes:
    """PPTX 파일 내 지정된 텍스트를 {{ variable }} 패턴으로 치환한다.

    python-pptx를 사용하여 모든 슬라이드의 shape를 순회하면서
    원본 텍스트를 Jinja2 변수 패턴으로 교체한다.

    Args:
        file_bytes: 원본 PPTX 파일 바이트 데이터
        replacements: [{"original": "원본 텍스트", "variable": "변수명"}, ...]

    Returns:
        변환된 PPTX 템플릿 파일의 바이트 데이터
    """
    try:
        prs = Presentation(io.BytesIO(file_bytes))

        for slide in prs.slides:
            for shape in slide.shapes:
                # text_frame이 있는 shape 처리
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        _replace_in_paragraph(paragraph, replacements)

                # 표(table)가 있는 shape 처리
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                _replace_in_paragraph(paragraph, replacements)

        # 결과를 바이트로 저장
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)

        logger.info("PPTX → 템플릿 변환 완료")
        return output.read()

    except Exception as e:
        logger.error("PPTX → 템플릿 변환 실패: %s", str(e))
        raise ValueError(f"PPTX 템플릿 변환에 실패했습니다: {e}") from e


def analyze_blank_form(file_bytes: bytes, format: str) -> dict:
    """빈 양식 문서의 구조를 분석하여 섹션 목록을 반환한다.

    DOCX/PPTX 파일에서 제목(heading), 빈 영역, 표 구조를 파악하여
    어떤 내용이 채워져야 하는지 구조 정보를 JSON으로 반환한다.

    사용자가 항목/제목만 있는 빈 양식을 업로드하면 AI가 구조를 분석하여
    자동으로 Jinja2 변수를 삽입할 수 있도록 구조 정보를 제공한다.

    Args:
        file_bytes: 문서 바이트 데이터
        format: 'docx' 또는 'pptx'

    Returns:
        {
            "sections": [
                {"heading": "1. 사업 개요", "level": 1, "has_content": false,
                 "has_table": false},
                {"heading": "2. 추진 전략", "level": 1, "has_content": false,
                 "has_table": true, "table_headers": ["항목", "내용", "비고"]},
                ...
            ],
            "title": "사업명",
            "metadata": {"total_sections": 4, "empty_sections": 4}
        }

    Raises:
        ValueError: 지원하지 않는 파일 형식이거나 분석에 실패한 경우
    """
    # 파일 형식을 정규화 (앞뒤 공백, 점 제거, 소문자화)
    normalized = format.lower().strip().lstrip(".")

    if normalized == "docx":
        return _analyze_blank_docx(file_bytes)
    elif normalized == "pptx":
        return _analyze_blank_pptx(file_bytes)
    else:
        raise ValueError(f"지원하지 않는 파일 형식입니다: {format} (docx 또는 pptx만 지원)")


# ---------------------------------------------------------------------------
# AI 기반 빈 양식 분석 함수들
# ---------------------------------------------------------------------------
# 휴리스틱 분석(analyze_blank_form)이 부족할 때 GPT-4o를 사용하여
# 더 정확한 양식 구조 분석을 수행한다.
# ---------------------------------------------------------------------------


def _extract_docx_structure_text(file_bytes: bytes) -> str:
    """DOCX 문서를 열어서 GPT-4o에 보낼 구조 텍스트를 생성한다.

    문서의 문단(paragraph)과 표(table)를 텍스트로 변환하여
    AI가 양식 구조를 이해할 수 있도록 명확하게 표현한다.

    문단은 스타일과 텍스트를, 표는 행/열 수와 각 셀의
    병합 정보(gridSpan, vMerge)까지 포함하여 출력한다.

    출력 예시:
        [문서 구조]
        [paragraph 0] style=Title text="회의록"
        [table 0] 8행 x 6열
          [row 0]
            [cell 0,0] gridSpan=4 vMerge=restart text=""
            [cell 0,4] gridSpan=1 text="PM"

    Args:
        file_bytes: DOCX 파일의 바이트 데이터

    Returns:
        AI에 전달할 문서 구조 텍스트 문자열
    """
    from docx import Document

    doc = Document(io.BytesIO(file_bytes))

    # XML 네임스페이스 — python-docx 내부 XML에 접근할 때 필요하다
    _ns = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"

    lines: list[str] = ["[문서 구조]"]

    # --- 문단(paragraph) 구조 추출 ---
    # 각 문단의 스타일명과 텍스트를 한 줄로 표현한다
    for idx, para in enumerate(doc.paragraphs):
        style_name = para.style.name if para.style else "Normal"
        text = para.text.strip()
        lines.append(f'[paragraph {idx}] style={style_name} text="{text}"')

    # --- 표(table) 구조 추출 ---
    # 각 표의 행/열 수와 셀 정보를 상세히 출력한다
    for t_idx, table in enumerate(doc.tables):
        row_count = len(table.rows)
        # 열 수는 첫 번째 행의 셀 수로 추정한다
        col_count = len(table.rows[0].cells) if table.rows else 0
        lines.append(f"\n[table {t_idx}] {row_count}행 x {col_count}열")

        for r_idx, row in enumerate(table.rows):
            lines.append(f"  [row {r_idx}]")

            # 같은 셀 객체가 여러 번 반복되는 경우(병합 셀) 중복 제거
            seen_ids: set[int] = set()

            for c_idx, cell in enumerate(row.cells):
                # 파이썬 id()로 중복 셀 판별 — 병합된 셀은 같은 객체를 가리킨다
                if id(cell) in seen_ids:
                    continue
                seen_ids.add(id(cell))

                # gridSpan: 가로 병합 셀 수 (기본값 1)
                tc = cell._tc
                grid_span_el = tc.find(f".//{_ns}gridSpan")
                grid_span = 1
                if grid_span_el is not None:
                    grid_span = int(grid_span_el.get(f"{_ns}val", "1"))

                # vMerge: 세로 병합 상태
                #   - vMerge="restart" → 병합의 시작 셀
                #   - vMerge="" (속성만 있고 값 없음) → 병합 연속 셀
                #   - vMerge 없음 → 병합 아님
                v_merge_el = tc.find(f".//{_ns}vMerge")
                v_merge = None
                if v_merge_el is not None:
                    v_merge_val = v_merge_el.get(f"{_ns}val", "")
                    v_merge = v_merge_val if v_merge_val else "continue"

                text = cell.text.strip()

                # 셀 정보를 한 줄로 포맷팅한다
                parts = [f"    [cell {r_idx},{c_idx}]"]
                parts.append(f"gridSpan={grid_span}")
                if v_merge is not None:
                    parts.append(f"vMerge={v_merge}")
                parts.append(f'text="{text}"')
                lines.append(" ".join(parts))

    return "\n".join(lines)


def analyze_blank_form_with_ai(
    file_bytes: bytes,
    format: str,
    api_key: str,
    model: str = "gpt-4o",
) -> dict:
    """GPT-4o를 사용하여 빈 양식 문서의 구조를 분석한다.

    휴리스틱 분석(analyze_blank_form)이 양식 구조를 제대로 감지하지 못할 때
    AI가 문서 구조를 직접 분석하여 라벨/값 필드를 식별하고,
    각 필드의 카테고리(session_auto, user_input, ai_generated)까지 분류한다.

    처리 흐름:
      1. _extract_docx_structure_text()로 문서 구조를 텍스트로 추출
      2. OpenAI Structured Outputs API로 양식 분석 요청
      3. AI 응답을 기존 structure 형식으로 변환

    Args:
        file_bytes: 문서 바이트 데이터
        format: 파일 형식 ('docx' 또는 'pptx')
        api_key: OpenAI API 키
        model: 사용할 GPT 모델명 (기본값: 'gpt-4o')

    Returns:
        기존 analyze_blank_form()과 동일한 형식의 구조 딕셔너리:
        {
            "sections": [...],
            "title": "...",
            "metadata": {"total_sections": N, "empty_sections": N}
        }

    Raises:
        ValueError: AI 분석 호출 실패 시
    """
    # 현재 DOCX만 지원한다 (PPTX는 향후 확장 가능)
    normalized = format.lower().strip().lstrip(".")
    if normalized != "docx":
        raise ValueError(f"AI 양식 분석은 현재 DOCX만 지원합니다. (입력: {format})")

    try:
        # 1단계: DOCX 문서에서 구조 텍스트를 추출한다
        structure_text = _extract_docx_structure_text(file_bytes)
        logger.info("AI 양식 분석 — 구조 텍스트 추출 완료 (%d자)", len(structure_text))

        # 2단계: LLMClient 추상화를 통해 Structured Outputs로 양식 분석을 요청한다
        from app.integrations.llm.factory import create_llm_client, get_provider_for_task

        # 템플릿 분석용 프로바이더를 설정에서 가져온다 (기본: openai)
        provider = get_provider_for_task("template")
        llm_client = create_llm_client(provider, api_key=api_key)

        # 시스템 프롬프트: AI에게 양식 분석 전문가 역할을 부여한다
        system_prompt = (
            "당신은 한국어 업무 양식(보고서, 회의록, 제안서 등) 분석 전문가입니다.\n"
            "주어진 DOCX 문서 구조를 분석하여 각 필드의 라벨과 변수 정보를 추출하세요.\n\n"
            "## 분석 규칙\n"
            "1. **라벨 vs 값 구분**: 표에서 짧은 텍스트(2~8자)가 있는 셀은 '라벨',\n"
            "   비어있거나 긴 텍스트가 있는 셀은 '값 입력란'입니다.\n"
            "2. **variable_name**: 라벨에서 공백/괄호/번호를 제거한 한글 변수명을 생성하세요.\n"
            "   예: '장 소' → '장소', '회 의 일 자' → '회의일자'\n"
            "3. **field_type 판단**:\n"
            "   - 'short': 한 줄짜리 입력란 (이름, 날짜, 장소 등)\n"
            "   - 'long': 여러 줄 입력란 (내용, 의견, 사항, 비고 등)\n"
            "4. **category 분류**:\n"
            "   - 'session_auto': 소속, 작성자, 성명, 부서, 이름 — 로그인 세션에서 자동 채움\n"
            "   - 'user_input': 보고서명, 제목, 주제, 참석자, 일시, 장소, 날짜, 기간 등\n"
            "     — 사용자가 직접 입력하는 최소 항목\n"
            "   - 'ai_generated': 보고내용, 문제점, 기대효과, 정보출처, 의견 등\n"
            "     — AI가 소스 문서 기반으로 자동 생성\n"
            "5. **var_type 분류**:\n"
            "   - 'string': 일반 텍스트 (기본값)\n"
            "   - 'date': 일시, 일자, 날짜, 기간 관련 필드\n"
            "   - 'array': 목록이나 반복 항목\n"
            "6. **title**: 문서 제목(Title 스타일이나 상단 큰 텍스트)이 있으면 추출,\n"
            "   없으면 null\n"
            "7. vMerge=continue 셀, 날짜 형식 텍스트('20 년 월 일'), "
            "보안등급 텍스트('기밀 보통')는 건너뛰세요.\n"
            "8. description에는 해당 필드가 어떤 정보를 담는지 한국어로 간단히 설명하세요."
        )

        # 사용자 프롬프트: 추출한 문서 구조를 전달한다
        user_prompt = f"아래 DOCX 문서 구조를 분석하여 양식 필드를 추출하세요.\n\n{structure_text}"

        # OpenAI Structured Outputs용 JSON 스키마 정의
        # strict=True로 AI가 반드시 이 형식에 맞춰 응답하도록 강제한다
        json_schema = {
            "name": "form_analysis",
            "strict": True,
            "schema": {
                "type": "object",
                "properties": {
                    "title": {
                        "type": ["string", "null"],
                    },
                    "fields": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "label": {"type": "string"},
                                "variable_name": {"type": "string"},
                                "field_type": {
                                    "type": "string",
                                    "enum": ["short", "long"],
                                },
                                "category": {
                                    "type": "string",
                                    "enum": [
                                        "session_auto",
                                        "user_input",
                                        "ai_generated",
                                    ],
                                },
                                "var_type": {
                                    "type": "string",
                                    "enum": ["string", "date", "array"],
                                },
                                "description": {"type": "string"},
                            },
                            "required": [
                                "label",
                                "variable_name",
                                "field_type",
                                "category",
                                "var_type",
                                "description",
                            ],
                            "additionalProperties": False,
                        },
                    },
                },
                "required": ["title", "fields"],
                "additionalProperties": False,
            },
        }

        # LLMClient 동기 호출 — Celery worker 환경이므로 동기 메서드를 사용한다
        # generate_structured_sync()는 JSON 스키마에 맞는 dict를 직접 반환한다
        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ]
        ai_result = llm_client.generate_structured_sync(messages, json_schema, temperature=0.2, max_tokens=4096)
        logger.info(
            "AI 양식 분석 완료: title=%s, 필드 %d개",
            ai_result.get("title"),
            len(ai_result.get("fields", [])),
        )

        # 3단계: AI 결과를 기존 structure 형식으로 변환한다
        return _convert_ai_result_to_structure(ai_result)

    except Exception as e:
        logger.error("AI 양식 분석 실패: %s", str(e))
        raise ValueError(f"AI 양식 분석에 실패했습니다: {e}") from e


def _convert_ai_result_to_structure(ai_result: dict) -> dict:
    """AI 분석 결과를 기존 analyze_blank_form() 형식으로 변환한다.

    AI가 반환한 fields 배열을 sections 배열로 변환하여,
    기존 auto_generate_jinja2_from_structure()가 그대로 사용할 수 있도록 한다.

    AI 결과의 각 field는 다음과 같이 매핑된다:
      - field.label → section.heading (원본 라벨 텍스트)
      - field.variable_name → section.heading (변수명으로 사용)
      - field.field_type → section.field_type ("short" 또는 "long")
      - field.category → section.category (기존 structure에는 없던 필드, AI가 분류)

    Args:
        ai_result: AI가 반환한 분석 결과
            {"title": "...", "fields": [{"label": "...", ...}, ...]}

    Returns:
        기존 structure 형식의 딕셔너리:
        {
            "sections": [
                {"heading": "장소", "level": 1, "has_content": False,
                 "field_type": "short", "category": "user_input"},
                ...
            ],
            "title": "...",
            "metadata": {"total_sections": N, "empty_sections": N}
        }
    """
    title = ai_result.get("title")
    fields = ai_result.get("fields", [])

    sections: list[dict] = []
    for field in fields:
        # heading에는 변수명을 사용한다
        # (auto_generate_jinja2_from_structure에서 heading을 변수명으로 변환하므로)
        heading = field.get("variable_name", field.get("label", ""))

        section = {
            "heading": heading,
            "level": 1,
            "has_content": False,  # 빈 양식이므로 내용 없음
            "field_type": field.get("field_type", "short"),
            # AI가 분류한 카테고리를 포함한다
            # 이 값은 _auto_fill_docx()에서 변수 메타데이터 생성 시 참조된다
            "category": field.get("category", "ai_generated"),
        }
        sections.append(section)

    empty_count = sum(1 for s in sections if not s["has_content"])

    logger.info(
        "AI 결과 → structure 변환 완료: 제목=%s, 섹션 %d개 (빈 섹션 %d개)",
        title,
        len(sections),
        empty_count,
    )

    return {
        "sections": sections,
        "title": title,
        "metadata": {
            "total_sections": len(sections),
            "empty_sections": empty_count,
        },
    }


# ---------------------------------------------------------------------------
# 에디터용 구조 추출 함수
# ---------------------------------------------------------------------------


def extract_docx_structure_for_editor(
    file_bytes: bytes,
    existing_variables: dict | None = None,
) -> dict:
    """원본 DOCX를 파싱하여 에디터에서 렌더링할 수 있는 JSON 구조를 반환한다.

    프론트엔드의 변수 매핑 에디터에서 문서의 표 구조를 시각적으로 표시하고,
    사용자가 클릭으로 셀과 변수를 매핑할 수 있도록 상세 구조 정보를 제공한다.

    기능:
      - 문단: 인덱스, 텍스트, 스타일
      - 표: 각 셀의 행/열 좌표, 텍스트, 가로/세로 병합 정보, 라벨 여부
      - 기존 변수 매핑이 있으면 해당 셀에 variable_name을 미리 설정

    Args:
        file_bytes: DOCX 파일의 바이트 데이터
        existing_variables: 기존에 저장된 변수 정보 (jinja2_variables 형식)
            {"variables": [{"name": "장소", "label": "장 소", ...}, ...]}

    Returns:
        에디터용 구조 딕셔너리:
        {
            "paragraphs": [...],
            "tables": [...],
            "existing_variables": [...]
        }
    """
    from docx import Document

    doc = Document(io.BytesIO(file_bytes))

    # XML 네임스페이스
    _ns = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"

    # --- 기존 변수 매핑 정보를 준비한다 ---
    # 라벨 → 변수 정보 매핑 딕셔너리를 만들어서
    # 표 셀의 텍스트와 매칭할 때 사용한다
    var_by_label: dict[str, dict] = {}
    var_list: list[dict] = []
    if existing_variables and "variables" in existing_variables:
        for var in existing_variables["variables"]:
            var_name = var.get("name", "")
            label = var.get("label", "")
            if label:
                var_by_label[label] = var
            # _heading_to_variable_name() 결과로도 매핑한다
            # (라벨에 공백이 포함된 경우를 대비)
            cleaned_label = _heading_to_variable_name(label) if label else ""
            if cleaned_label:
                var_by_label[cleaned_label] = var

            var_list.append(
                {
                    "name": var_name,
                    "var_type": var.get("type", var.get("var_type", "string")),
                    "label": label,
                    "category": var.get("category", "ai_generated"),
                    "field_type": var.get("field_type", "short"),
                }
            )

    # --- 라벨 판별 함수 ---
    # 짧고 한글/영문으로 된 텍스트를 라벨로 인식한다
    def _is_label_cell(text: str) -> bool:
        """셀 텍스트가 양식의 라벨(항목명)인지 판별한다.

        라벨의 특징:
          - 2~10자 정도의 짧은 텍스트
          - 콜론(:)으로 끝나는 경우
          - 숫자로 시작하지 않는 경우
        """
        cleaned = re.sub(r"\s+", "", text)
        if not cleaned:
            return False
        # 콜론으로 끝나면 라벨
        if text.rstrip().endswith(":") or text.rstrip().endswith("："):
            return True
        # 너무 짧거나 너무 긴 텍스트는 라벨 아님
        if len(cleaned) < 2 or len(cleaned) > 10:
            return False
        # 숫자로 시작하면 값 (날짜, 수량 등)
        return not cleaned[0].isdigit()

    # --- 문단(paragraph) 구조 추출 ---
    paragraphs_out: list[dict] = []
    for idx, para in enumerate(doc.paragraphs):
        style_name = para.style.name if para.style else "Normal"
        text = para.text.strip()

        # 기존 변수와 매칭되는지 확인
        var_name_match = None
        if text in var_by_label:
            var_name_match = var_by_label[text].get("name")

        paragraphs_out.append(
            {
                "index": idx,
                "text": text,
                "style": style_name,
                "variable_name": var_name_match,
            }
        )

    # --- 표(table) 구조 추출 ---
    tables_out: list[dict] = []
    for t_idx, table in enumerate(doc.tables):
        total_rows = len(table.rows)
        total_cols = len(table.rows[0].cells) if table.rows else 0

        rows_out: list[dict] = []
        for r_idx, row in enumerate(table.rows):
            cells_out: list[dict] = []

            # 병합된 셀은 같은 객체를 여러 번 참조하므로 id()로 중복 제거
            seen_ids: set[int] = set()

            for c_idx, cell in enumerate(row.cells):
                if id(cell) in seen_ids:
                    continue
                seen_ids.add(id(cell))

                # gridSpan: 가로 병합 셀 수
                tc = cell._tc
                grid_span_el = tc.find(f".//{_ns}gridSpan")
                grid_span = 1
                if grid_span_el is not None:
                    grid_span = int(grid_span_el.get(f"{_ns}val", "1"))

                # vMerge: 세로 병합 상태
                v_merge_el = tc.find(f".//{_ns}vMerge")
                row_span = 1
                if v_merge_el is not None:
                    v_merge_val = v_merge_el.get(f"{_ns}val", "")
                    if v_merge_val == "restart":
                        # 병합 시작 셀 — 아래 행을 순회하여 실제 rowSpan을 계산한다
                        row_span = 1
                        for next_r_idx in range(r_idx + 1, total_rows):
                            next_row = table.rows[next_r_idx]
                            # 같은 열 위치의 셀에서 vMerge를 확인한다
                            if c_idx < len(next_row.cells):
                                next_tc = next_row.cells[c_idx]._tc
                                next_vm = next_tc.find(f".//{_ns}vMerge")
                                if next_vm is not None:
                                    next_val = next_vm.get(f"{_ns}val", "")
                                    if next_val != "restart":
                                        # continue 상태 — 병합이 이어지고 있다
                                        row_span += 1
                                    else:
                                        break
                                else:
                                    break
                            else:
                                break
                    elif v_merge_val == "" or v_merge_val is None:
                        # continue 상태 — 이 셀은 위의 셀에 병합되어 있으므로
                        # 건너뛴다 (시작 셀에서 이미 rowSpan으로 표현됨)
                        continue

                text = cell.text.strip()

                # 라벨 여부 판별
                is_label = _is_label_cell(text)

                # 기존 변수와 매칭
                # 주의: 라벨 셀에는 변수를 매핑하지 않는다.
                # "장 소" 라벨과 "장소" 변수가 매칭되어 라벨 셀에 변수가 들어가는 것을 방지.
                # 변수는 라벨 옆의 빈 값 셀에만 매핑된다.
                var_name_match = None
                if not is_label:
                    cleaned_text = _heading_to_variable_name(text) if text else ""
                    if text in var_by_label:
                        var_name_match = var_by_label[text].get("name")
                    elif cleaned_text in var_by_label:
                        var_name_match = var_by_label[cleaned_text].get("name")

                cells_out.append(
                    {
                        "row": r_idx,
                        "col": c_idx,
                        "text": text,
                        "gridSpan": grid_span,
                        "rowSpan": row_span,
                        "is_label": is_label,
                        "variable_name": var_name_match,
                    }
                )

            # --- 라벨 셀의 변수를 옆 빈 값 셀에 자동 배치 ---
            # 기존 변수가 라벨과 이름이 매칭되지만 라벨 셀에는 매핑되지 않았으므로,
            # 라벨 바로 옆의 빈 셀에 해당 변수를 배치한다.
            for ci, cell_info in enumerate(cells_out):
                if cell_info["is_label"] and cell_info["variable_name"] is None:
                    cleaned = _heading_to_variable_name(cell_info["text"]) if cell_info["text"] else ""
                    matched_var = var_by_label.get(cell_info["text"]) or var_by_label.get(cleaned)
                    if matched_var and ci + 1 < len(cells_out):
                        next_cell = cells_out[ci + 1]
                        if not next_cell["is_label"] and next_cell["variable_name"] is None:
                            next_cell["variable_name"] = matched_var.get("name")

            rows_out.append(
                {
                    "index": r_idx,
                    "cells": cells_out,
                }
            )

        tables_out.append(
            {
                "index": t_idx,
                "total_rows": total_rows,
                "total_cols": total_cols,
                "rows": rows_out,
            }
        )

    logger.info(
        "에디터용 DOCX 구조 추출 완료: 문단 %d개, 표 %d개, 기존 변수 %d개",
        len(paragraphs_out),
        len(tables_out),
        len(var_list),
    )

    return {
        "paragraphs": paragraphs_out,
        "tables": tables_out,
        "existing_variables": var_list,
    }


def _analyze_blank_docx(file_bytes: bytes) -> dict:
    """DOCX 빈 양식의 구조를 분석한다.

    Heading 스타일을 가진 paragraph를 섹션 시작으로 인식하고,
    그 아래 paragraph가 비어 있으면 해당 섹션을 빈(empty) 섹션으로 표시한다.
    표(table)가 있으면 헤더 행을 추출하고, 데이터 행이 비어있는지 확인한다.

    Args:
        file_bytes: DOCX 파일의 바이트 데이터

    Returns:
        구조 분석 결과 딕셔너리
    """
    from docx import Document

    try:
        doc = Document(io.BytesIO(file_bytes))

        # 문서 제목을 감지한다 (Title 스타일이나 첫 번째 텍스트)
        title: str | None = None
        sections: list[dict] = []
        # 현재 섹션 정보를 추적하는 변수
        current_section: dict | None = None

        # --- 모든 paragraph를 순회하며 구조를 파악한다 ---
        for para in doc.paragraphs:
            style_name = (para.style.name or "").lower() if para.style else ""
            text = para.text.strip()

            # Title 스타일인 경우 문서 제목으로 인식
            if "title" in style_name and text:
                title = text
                continue

            # Heading 스타일인 경우 새 섹션 시작
            if "heading" in style_name and text:
                # 이전 섹션이 있으면 sections 리스트에 추가
                if current_section is not None:
                    sections.append(current_section)

                # heading 레벨을 추출한다 (예: "Heading 2" → 2)
                level = 1  # 기본 레벨
                for char in style_name:
                    if char.isdigit():
                        level = int(char)
                        break

                # 새 섹션 정보를 초기화한다
                current_section = {
                    "heading": text,
                    "level": level,
                    "has_content": False,
                    "has_table": False,
                }
                continue

            # 일반 paragraph: 현재 섹션에 내용이 있는지 확인
            if current_section is not None and text:
                current_section["has_content"] = True

            # Title 스타일이 없는 경우, 첫 번째 비어있지 않은 텍스트를 제목으로
            # 사용한다 (대괄호로 감싸인 텍스트도 제목 후보)
            if title is None and text:
                # [사업명] 형태의 텍스트를 제목으로 인식
                bracket_match = re.match(r"^\[(.+)\]$", text)
                if bracket_match:
                    title = bracket_match.group(1)
                elif current_section is None and len(text) <= 50:
                    # 첫 번째 heading 이전의 짧은 텍스트(50자 이하)를 제목으로 간주
                    # 긴 텍스트는 AI 생성 내용일 수 있으므로 제목으로 사용하지 않음
                    title = text

        # 마지막 섹션을 추가한다
        if current_section is not None:
            sections.append(current_section)

        # --- 표(table) 분석 ---
        # 각 표를 순서대로 분석하여 가장 가까운 섹션에 연결한다
        for table in doc.tables:
            # 표의 헤더 행 (첫 번째 행)에서 열 이름을 추출한다
            headers: list[str] = []
            if table.rows:
                for cell in table.rows[0].cells:
                    header_text = cell.text.strip()
                    if header_text:
                        headers.append(header_text)

            # 데이터 행(두 번째 행부터)이 비어 있는지 확인한다
            table_is_empty = True
            for row in table.rows[1:]:
                for cell in row.cells:
                    if cell.text.strip():
                        table_is_empty = False
                        break
                if not table_is_empty:
                    break

            # 마지막 섹션에 표 정보를 연결한다
            if sections:
                last_section = sections[-1]
                last_section["has_table"] = True
                if headers:
                    last_section["table_headers"] = headers
                # 표 데이터가 비어 있으면 내용이 없는 것으로 간주
                if table_is_empty:
                    last_section["has_content"] = False

        # Heading 스타일이 없는 문서도 처리한다
        # (번호 매기기 패턴으로 섹션을 감지: "1. 제목", "2) 제목" 등)
        if not sections:
            sections = _detect_sections_by_numbering(doc)

        # ---------------------------------------------------------------
        # 번호 패턴으로도 섹션이 없고 표가 있으면 → 표 기반 양식 분석
        # 업무 양식은 대부분 표로 구성되어 있기 때문에,
        # heading이나 번호 없이 표만 있는 문서도 올바르게 분석한다.
        # ---------------------------------------------------------------
        if not sections and doc.tables:
            for table in doc.tables:
                table_fields = _analyze_table_form_fields(table)
                sections.extend(table_fields)

            # 표의 첫 행에서 문서 제목을 감지한다
            # 예: "보 고 서" 같은 패턴에서 공백을 제거하면 "보고서"
            if title is None and doc.tables:
                try:
                    first_row_text = doc.tables[0].rows[0].cells[0].text.strip()
                    # 공백 제거 후 날짜 부분("20년월일" 등)을 제거하여 제목만 추출
                    cleaned = re.sub(r"\s+", "", first_row_text)
                    cleaned = re.sub(r"\d+년.*", "", cleaned)  # 날짜 부분 제거
                    if cleaned and 2 <= len(cleaned) <= 10:
                        title = cleaned
                except (IndexError, AttributeError):
                    # 표 구조가 비정상적인 경우 제목 감지를 건너뛴다
                    pass

        # 빈 섹션 수 계산
        empty_count = sum(1 for s in sections if not s["has_content"])

        logger.info(
            "DOCX 빈 양식 분석 완료: 제목=%s, 섹션 %d개 (빈 섹션 %d개)",
            title,
            len(sections),
            empty_count,
        )

        return {
            "sections": sections,
            "title": title,
            "metadata": {
                "total_sections": len(sections),
                "empty_sections": empty_count,
            },
        }

    except Exception as e:
        logger.error("DOCX 빈 양식 분석 실패: %s", str(e))
        raise ValueError(f"DOCX 빈 양식을 분석할 수 없습니다: {e}") from e


def _detect_sections_by_numbering(doc: Any) -> list[dict]:
    """Heading 스타일이 없는 DOCX에서 번호 매기기 패턴으로 섹션을 감지한다.

    "1. 사업 개요", "2) 추진 전략", "제1장 개요" 등의 번호 패턴을 찾아
    섹션으로 인식한다. 실무에서 자주 쓰이는 한글 문서 양식을 지원한다.

    Args:
        doc: python-docx Document 객체

    Returns:
        감지된 섹션 리스트
    """
    # 번호 매기기 패턴: "1. ", "1) ", "제1장 ", "가. ", "I. " 등
    numbering_pattern = re.compile(r"^(\d+[\.\)]\s+|제\d+[장절항]\s+|[가-힣][\.\)]\s+|[IVX]+[\.\)]\s+)")

    sections: list[dict] = []
    current_section: dict | None = None

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        # 번호 매기기 패턴과 일치하면 새 섹션 시작
        if numbering_pattern.match(text):
            if current_section is not None:
                sections.append(current_section)

            current_section = {
                "heading": text,
                "level": 1,
                "has_content": False,
                "has_table": False,
            }
        elif current_section is not None and text:
            # 섹션 아래에 내용이 있으면 표시
            current_section["has_content"] = True

    # 마지막 섹션 추가
    if current_section is not None:
        sections.append(current_section)

    return sections


def _analyze_table_form_fields(table: Any) -> list[dict]:
    """표에서 '라벨:값' 패턴의 폼 필드를 감지한다.

    업무 양식에서 흔히 쓰이는 표 기반 입력 폼을 분석한다.
    표의 각 행을 순회하면서, 라벨 셀과 빈 셀의 조합을 감지하여
    어떤 항목에 값을 채워야 하는지 구조 정보를 반환한다.

    감지 전략:
      1. 행의 고유 셀 텍스트를 추출 (병합 셀 중복 제거)
      2. 패턴 A: "라벨 셀 + 빈 셀" 인접 패턴 → 짧은 입력란 (short)
      3. 패턴 B: 행 전체를 차지하는 단독 텍스트 → 큰 입력 영역 라벨 (long)
      4. 패턴 C: "라벨:" 형태 (콜론으로 끝나는 텍스트) → 짧은 입력란 (short)

    Args:
        table: python-docx Table 객체

    Returns:
        감지된 폼 필드 리스트. 각 항목은 다음 형태:
        [{"heading": "소속", "level": 1, "has_content": False, "field_type": "short"},
         {"heading": "보고내용과의견", "level": 1, "has_content": False, "field_type": "long"},
         ...]
    """
    fields: list[dict] = []

    # -----------------------------------------------------------------------
    # 건너뛸 텍스트 패턴 정의
    # -----------------------------------------------------------------------
    # 날짜 패턴: "20 년 월 일", "2026년", "년 월 일" 등
    # 선택 항목 패턴: "기밀 보통" 같이 여러 등급이 나열된 경우
    date_pattern = re.compile(
        r"(20\s*\d{0,2}\s*년|년\s*월\s*일|\d{4}년)"  # 날짜 관련 텍스트
    )
    choice_pattern = re.compile(
        r"(기밀|보통|대외비|극비|비밀)"  # 보안 등급 선택 항목
    )

    # -----------------------------------------------------------------------
    # 큰 텍스트 영역(long)으로 판별할 키워드 패턴
    # 한국어 양식에서 "내용", "사항" 등이 붙은 필드는 여러 줄 입력 영역이다
    # -----------------------------------------------------------------------
    long_field_keywords = re.compile(r"(내용|사항|설명|비고|의견|결과|요약|결론|소견|경과|조치|대책)")

    # -----------------------------------------------------------------------
    # 행(row) 단위로 분석
    # -----------------------------------------------------------------------
    rows = list(table.rows)
    for row_idx, row in enumerate(rows):
        # -----------------------------------------------------------------
        # 수직 병합(vMerge) 셀이 포함된 행은 건너뛴다.
        # 양식 상단의 로고/문서번호/PM 헤더 영역은 vMerge로 병합되어 있는데,
        # 이 영역은 입력 필드가 아닌 고정 헤더이므로 변수로 감지하지 않는다.
        # -----------------------------------------------------------------
        _ns = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
        has_vmerge = False
        for cell in row.cells:
            tc = cell._tc
            vm_el = tc.find(f".//{_ns}vMerge")
            if vm_el is not None:
                has_vmerge = True
                break
        if has_vmerge:
            continue

        # 병합된 셀은 같은 객체를 여러 번 참조하므로, id()로 중복을 제거한다
        seen_cell_ids: set[int] = set()
        unique_cells: list[Any] = []
        for cell in row.cells:
            cell_id = id(cell)
            if cell_id not in seen_cell_ids:
                seen_cell_ids.add(cell_id)
                unique_cells.append(cell)

        # 각 고유 셀의 텍스트를 추출한다 (앞뒤 공백 제거)
        cell_texts = [cell.text.strip() for cell in unique_cells]

        # 모든 셀이 비어있으면 건너뛴다 (빈 행)
        if not any(cell_texts):
            continue

        # -----------------------------------------------------------------------
        # 패턴 B: 행 전체가 하나의 라벨인 경우 (넓은 영역 입력란)
        # 예: "보고내용과 의견", "정보(자료)출처", "문제점"
        # 조건: 고유 셀이 1개이고 텍스트가 있으며, 다음 행이 비어있거나 없음
        # -----------------------------------------------------------------------
        non_empty_texts = [t for t in cell_texts if t]
        if len(non_empty_texts) == 1 and len(unique_cells) <= 3:
            label_text = non_empty_texts[0]

            # 날짜/선택 항목은 건너뛴다
            cleaned_label = re.sub(r"\s+", "", label_text)
            if date_pattern.search(label_text):
                continue
            if choice_pattern.search(label_text):
                continue
            # 너무 짧은 텍스트(1자)는 라벨로 보지 않는다
            if len(cleaned_label) <= 1:
                continue

            # 라벨+값이 같은 셀에 섞여있는 경우 (줄바꿈으로 구분)
            # 예1: "미결\n사항\n현재까지 논의된..." → 라벨="미결사항"
            # 예2: "회의\n내용\n회의에서는..." → 라벨="회의내용"
            # 짧은 줄(5자 이하)이 연속되면 합쳐서 라벨로, 긴 줄부터는 값으로 판단
            lines = label_text.split("\n")
            if len(lines) > 1:
                label_parts = []
                for line in lines:
                    line_clean = re.sub(r"\s+", "", line.strip())
                    if len(line_clean) <= 5 and line_clean:
                        label_parts.append(line_clean)
                    else:
                        break  # 긴 줄부터는 값
                if label_parts:
                    label_text = "".join(label_parts)
                    cleaned_label = label_text

            # "long" 영역인지 판단한다
            # 3가지 조건 중 하나라도 충족하면 큰 입력 영역(long)으로 판단:
            #   1) 라벨에 줄바꿈(\n)이 있음 — 셀 안에서 명시적으로 줄을 나눈 라벨
            #      (예: "회의\n내용", "특기\n사항")
            #   2) 라벨에 "내용/사항/설명/비고/의견/결과" 등 long 키워드 포함
            #   3) 다음 행이 비어있거나 마지막 행 (기존 조건)
            is_long = False
            original_label = non_empty_texts[0]
            if "\n" in original_label or long_field_keywords.search(cleaned_label):
                is_long = True
            elif row_idx + 1 < len(rows):
                next_row_cells = rows[row_idx + 1].cells
                next_texts = [c.text.strip() for c in next_row_cells]
                # 다음 행이 전부 비어있으면 큰 입력 영역이다
                if not any(next_texts):
                    is_long = True
            else:
                # 마지막 행이면 큰 입력 영역으로 간주한다
                is_long = True

            # 라벨에서 변수명을 생성한다
            var_label = _heading_to_variable_name(label_text)

            fields.append(
                {
                    "heading": var_label,
                    "level": 1,
                    "has_content": False,
                    "field_type": "long" if is_long else "short",
                }
            )
            continue

        # -----------------------------------------------------------------------
        # 패턴 A: 가로 "라벨:값" 교차 배치 감지
        # 예: ["장소", "회의실", "일시", "2026-03-25"]
        #     → 짝수 인덱스(0,2) = 라벨, 홀수 인덱스(1,3) = 값
        # 예: ["소속:", "홍길동"]
        #     → 콜론으로 끝나는 셀 = 라벨
        # -----------------------------------------------------------------------

        # 라벨 판별 함수: 짧고(2~8자) 공백 제거한 텍스트가 명사형인지
        def _is_label(txt: str) -> bool:
            """셀 텍스트가 폼 라벨인지 판별한다."""
            cleaned = re.sub(r"\s+", "", txt)
            # 콜론으로 끝나면 라벨
            if txt.rstrip().endswith(":") or txt.rstrip().endswith("："):
                return True
            # 너무 짧거나 너무 길면 라벨 아님
            if len(cleaned) < 2 or len(cleaned) > 10:
                return False
            # 날짜 패턴이면 값임
            if date_pattern.search(txt):
                return False
            if choice_pattern.search(txt):
                return False
            # 숫자로 시작하면 값임 (예: "2026-03-25", "3명")
            return not cleaned[0].isdigit()

        def _is_value(txt: str) -> bool:
            """셀 텍스트가 값(내용)인지 판별한다."""
            cleaned = re.sub(r"\s+", "", txt)
            # 긴 텍스트(15자+)는 값
            if len(cleaned) > 15:
                return True
            # 날짜/숫자 포함이면 값
            if date_pattern.search(txt) or (cleaned and cleaned[0].isdigit()):
                return True
            # "외N명" 패턴은 값 (예: "조예주 외3명")
            return bool(re.search(r"외\d+명", cleaned))

        # 가로 교차 패턴 감지: [라벨, 값, 라벨, 값, ...]
        # 짝수 인덱스가 라벨이고 홀수 인덱스가 값/빈셀이면 교차 패턴
        if len(non_empty_texts) >= 2:
            is_alternating = True
            for ci in range(0, len(cell_texts) - 1, 2):
                left = cell_texts[ci]
                cell_texts[ci + 1] if ci + 1 < len(cell_texts) else ""
                if left and not _is_label(left):
                    is_alternating = False
                    break

            if is_alternating:
                # 짝수 인덱스 셀만 라벨로 추출
                for ci in range(0, len(cell_texts), 2):
                    text = cell_texts[ci]
                    if not text:
                        continue
                    if not _is_label(text):
                        continue
                    var_label = _heading_to_variable_name(text)
                    if any(f["heading"] == var_label for f in fields):
                        continue
                    fields.append(
                        {
                            "heading": var_label,
                            "level": 1,
                            "has_content": False,
                            "field_type": "short",
                        }
                    )
                continue  # 이 행은 처리 완료

        # 교차 패턴이 아닌 경우: 개별 셀 분석
        for cell_idx, text in enumerate(cell_texts):
            if not text:
                continue
            # 값으로 보이는 셀은 건너뛴다
            if _is_value(text):
                continue
            if not _is_label(text):
                continue

            # 바로 옆 셀이 빈 셀이거나 값인 경우에만 라벨로 인정
            next_is_empty_or_value = (
                cell_idx + 1 >= len(cell_texts)  # 마지막 셀
                or not cell_texts[cell_idx + 1]  # 옆이 빈 셀
                or _is_value(cell_texts[cell_idx + 1])  # 옆이 값
            )
            if not next_is_empty_or_value:
                continue

            var_label = _heading_to_variable_name(text)
            if any(f["heading"] == var_label for f in fields):
                continue

            fields.append(
                {
                    "heading": var_label,
                    "level": 1,
                    "has_content": False,
                    "field_type": "short",
                }
            )

    logger.debug(
        "표 폼 필드 분석 완료: %d개 필드 감지 (%s)",
        len(fields),
        ", ".join(f["heading"] for f in fields),
    )
    return fields


def _analyze_blank_pptx(file_bytes: bytes) -> dict:
    """PPTX 빈 양식의 구조를 분석한다.

    각 슬라이드에서 title shape를 섹션 제목으로 인식하고,
    본문 영역(body/subtitle/content placeholder)이 비어 있으면
    해당 슬라이드를 빈 섹션으로 표시한다.

    Args:
        file_bytes: PPTX 파일의 바이트 데이터

    Returns:
        구조 분석 결과 딕셔너리
    """
    from pptx.util import Pt  # noqa: F401 — 사용되지 않지만 pptx 패키지 확인용

    try:
        prs = Presentation(io.BytesIO(file_bytes))

        title: str | None = None
        sections: list[dict] = []

        for slide_idx, slide in enumerate(prs.slides):
            slide_title: str | None = None
            has_content = False
            has_table = False
            table_headers: list[str] = []

            for shape in slide.shapes:
                # 제목(title) shape 감지
                if shape.has_text_frame:
                    # placeholder 타입으로 title을 구분한다
                    is_title = False
                    if hasattr(shape, "placeholder_format") and shape.placeholder_format:
                        # placeholder idx 0 = title, 1 = subtitle/body
                        ph_idx = shape.placeholder_format.idx
                        if ph_idx == 0:
                            is_title = True

                    if is_title:
                        text = shape.text_frame.text.strip()
                        if text:
                            slide_title = text
                            # 첫 번째 슬라이드의 제목을 문서 전체 제목으로 사용
                            if title is None and slide_idx == 0:
                                title = text
                    else:
                        # title이 아닌 shape에서 내용이 있는지 확인
                        text = shape.text_frame.text.strip()
                        if text:
                            has_content = True

                # 표(table) 감지
                if shape.has_table:
                    has_table = True
                    tbl = shape.table
                    # 첫 번째 행에서 헤더를 추출
                    if tbl.rows:
                        for cell in tbl.rows[0].cells:
                            h = cell.text_frame.text.strip()
                            if h:
                                table_headers.append(h)

                    # 데이터 행이 비어 있는지 확인
                    table_is_empty = True
                    for row in tbl.rows[1:]:
                        for cell in row.cells:
                            if cell.text_frame.text.strip():
                                table_is_empty = False
                                break
                        if not table_is_empty:
                            break

                    if not table_is_empty:
                        has_content = True

            # 슬라이드에 제목이 있으면 섹션으로 등록
            if slide_title:
                section_info: dict = {
                    "heading": slide_title,
                    "level": 1,
                    "has_content": has_content,
                    "has_table": has_table,
                }
                if table_headers:
                    section_info["table_headers"] = table_headers
                sections.append(section_info)

        # 빈 섹션 수 계산
        empty_count = sum(1 for s in sections if not s["has_content"])

        logger.info(
            "PPTX 빈 양식 분석 완료: 제목=%s, 섹션 %d개 (빈 섹션 %d개)",
            title,
            len(sections),
            empty_count,
        )

        return {
            "sections": sections,
            "title": title,
            "metadata": {
                "total_sections": len(sections),
                "empty_sections": empty_count,
            },
        }

    except Exception as e:
        logger.error("PPTX 빈 양식 분석 실패: %s", str(e))
        raise ValueError(f"PPTX 빈 양식을 분석할 수 없습니다: {e}") from e


def _heading_to_variable_name(heading: str) -> str:
    """섹션 heading 텍스트를 Jinja2 변수명으로 변환한다.

    앞의 번호/기호/공백을 제거하고, 남은 텍스트를 밑줄로 연결한다.
    예시:
      "1. 사업 개요"   → "사업개요"
      "2) 추진 전략"   → "추진전략"
      "제1장 개요"     → "개요"
      "[사업명]"       → "사업명"

    Args:
        heading: 원본 heading 텍스트

    Returns:
        Jinja2에서 사용 가능한 변수명 문자열
    """
    # 앞의 번호 패턴을 제거한다 (예: "1. ", "2) ", "제1장 ", "가. ", "I. ")
    cleaned = re.sub(
        r"^(\d+[\.\)]\s*|제\d+[장절항]\s*|[가-힣][\.\)]\s*|[IVX]+[\.\)]\s*)",
        "",
        heading.strip(),
    )
    # 대괄호를 제거한다 (예: "[사업명]" → "사업명")
    cleaned = re.sub(r"[\[\]]", "", cleaned)
    # 괄호와 괄호 안의 내용을 제거한다 (예: "보고처(부서)" → "보고처")
    cleaned = re.sub(r"\([^)]*\)", "", cleaned)
    # 콜론을 제거한다 (예: "소 속 :" → "소 속 ")
    cleaned = re.sub(r"[:：]", "", cleaned)
    # 공백을 제거하여 변수명으로 만든다 (예: "사업 개요" → "사업개요", "회 의 일 자" → "회의일자")
    cleaned = re.sub(r"\s+", "", cleaned)
    # 변수명이 비어있으면 원본에서 숫자/기호를 제거한 결과를 사용
    if not cleaned:
        cleaned = re.sub(r"[^가-힣a-zA-Z0-9_]", "", heading)
    # 그래도 비어있으면 기본값
    if not cleaned:
        cleaned = "항목"
    return cleaned


def auto_generate_jinja2_from_structure(
    file_bytes: bytes,
    format: str,
    structure: dict,
) -> tuple[bytes, dict]:
    """빈 양식 구조를 기반으로 Jinja2 변수를 자동 삽입한다.

    analyze_blank_form()의 결과를 받아, 각 빈 섹션에 적절한 Jinja2 변수를
    자동으로 삽입한 새 DOCX/PPTX 파일을 생성한다.

    변수 생성 규칙:
      - 빈 텍스트 섹션: {{ 변수명 }} 삽입 (변수명은 heading에서 자동 생성)
      - 빈 표 섹션: {% for row in 변수명_항목 %} + {{ row.열이름 }} 삽입
      - 문서 제목: {{ 제목 }} 삽입

    Args:
        file_bytes: 원본 문서 바이트
        format: 'docx' 또는 'pptx'
        structure: analyze_blank_form()의 반환값

    Returns:
        (변환된 파일 bytes, 추출된 변수 메타데이터 dict)

    Raises:
        ValueError: 지원하지 않는 파일 형식이거나 변환에 실패한 경우
    """
    normalized = format.lower().strip().lstrip(".")

    if normalized == "docx":
        return _auto_fill_docx(file_bytes, structure)
    elif normalized == "pptx":
        return _auto_fill_pptx(file_bytes, structure)
    else:
        raise ValueError(f"지원하지 않는 파일 형식입니다: {format} (docx 또는 pptx만 지원)")


def _auto_fill_docx(file_bytes: bytes, structure: dict) -> tuple[bytes, dict]:
    """DOCX 빈 양식에 Jinja2 변수를 자동 삽입한다.

    각 빈 섹션의 heading 바로 다음 paragraph에 {{ 변수명 }} 을 삽입하고,
    표가 있는 섹션에는 for 루프 변수를 삽입한다.

    Args:
        file_bytes: DOCX 파일의 바이트 데이터
        structure: analyze_blank_form()의 반환값

    Returns:
        (변환된 DOCX 바이트, 변수 메타데이터 dict)
    """
    from docx import Document

    try:
        doc = Document(io.BytesIO(file_bytes))
        sections = structure.get("sections", [])
        doc_title = structure.get("title")

        # 생성된 변수 메타데이터를 수집하는 리스트
        variables_meta: list[dict] = []

        # --- 제목 변수 처리 ---
        # 문서 제목이 감지된 경우, 해당 텍스트를 {{ 문서제목 }}로 치환
        # 버그 수정: "회의록" 같은 모호한 변수명이면 AI가 긴 내용을 생성하므로
        #           변수명을 "문서제목"으로 고정하고, category를 user_input으로 강제 지정
        if doc_title:
            title_var = "문서제목"
            variables_meta.append(
                {
                    "name": title_var,
                    "type": "string",
                    "label": doc_title,
                    "description": f"문서 제목 — 짧은 제목만 입력 ({doc_title})",
                    "required": True,
                    # 제목은 항상 user_input — AI가 생성하면 안 됨
                    "category": "user_input",
                }
            )
            # 문서 내에서 제목 텍스트를 찾아 치환
            for para in doc.paragraphs:
                if para.text.strip() == doc_title or (
                    re.match(r"^\[.+\]$", para.text.strip()) and doc_title in para.text
                ):
                    if para.runs:
                        para.runs[0].text = "{{ " + title_var + " }}"
                        for run in para.runs[1:]:
                            run.text = ""
                    break

        # --- 섹션 기반 heading 매칭 및 변수 삽입 ---
        # 각 섹션의 heading을 문서에서 찾아, 빈 섹션이면 변수를 삽입한다
        for section in sections:
            heading_text = section["heading"]
            var_name = _heading_to_variable_name(heading_text)
            has_content = section.get("has_content", False)
            has_table = section.get("has_table", False)
            field_type = section.get("field_type")  # 표 기반 폼 필드인 경우

            # ---------------------------------------------------------------
            # 표 기반 폼 필드 처리 (field_type이 있는 경우)
            # _analyze_table_form_fields()에서 감지된 라벨에 대해
            # 표의 빈 셀에 {{ 변수명 }} 을 삽입한다.
            # ---------------------------------------------------------------
            if field_type:
                # 카테고리 결정: AI 분석 결과 우선, 없으면 키워드 기반 자동 분류
                category = section.get("category") or classify_variable_category(var_name)
                variables_meta.append(
                    {
                        "name": var_name,
                        "type": classify_variable_type(var_name, "string"),
                        "label": heading_text,
                        "description": (
                            f"표 입력란: {heading_text}" + (" (여러 줄)" if field_type == "long" else " (한 줄)")
                        ),
                        # user_input만 필수 — session_auto/ai_generated는 자동 채워지므로 선택
                        "required": category == "user_input",
                        "field_type": field_type,
                        "category": category,
                    }
                )

                # 실제 표 셀에서 라벨 텍스트를 찾아 옆/아래 빈 셀에 변수를 삽입한다
                _insert_table_form_variable(doc, heading_text, var_name, field_type)
                continue

            # 이미 내용이 있는 섹션은 건너뛴다
            if has_content and not has_table:
                continue

            # heading을 찾아서 바로 다음 빈 paragraph에 변수를 삽입한다
            for i, para in enumerate(doc.paragraphs):
                # heading 텍스트와 일치하는 paragraph를 찾는다
                if para.text.strip() == heading_text:
                    # 표가 있는 섹션은 별도 처리 (아래 표 처리 로직에서)
                    if has_table:
                        break

                    # 다음 paragraph를 찾아서 변수를 삽입한다
                    content_var = var_name + "_내용"
                    _insert_variable_after_heading(doc, i, content_var)

                    # 변수 메타데이터 등록
                    variables_meta.append(
                        {
                            "name": content_var,
                            "type": "string",
                            "label": heading_text,
                            "description": f"{heading_text} 섹션의 내용",
                            # heading 기반 섹션은 AI가 생성하는 내용이므로 필수 아님
                            "required": False,
                            # AI 분석 결과의 category를 우선 사용하고,
                            # 없으면 변수 이름 키워드 기반으로 자동 분류한다
                            "category": section.get("category") or classify_variable_category(content_var),
                        }
                    )
                    break

            # --- 표 변수 처리 ---
            if has_table:
                table_headers = section.get("table_headers", [])
                if table_headers:
                    # 표 반복 변수명 생성
                    loop_var = var_name + "_항목"
                    variables_meta.append(
                        {
                            "name": loop_var,
                            "type": "array",
                            "label": f"{heading_text} 표 데이터",
                            "description": (f"열 구조: {', '.join(table_headers)}"),
                            # 표 데이터는 AI가 자동 생성하므로 필수 아님
                            "required": False,
                            # 표 데이터 항목은 AI가 자동 생성
                            "category": classify_variable_category(loop_var),
                        }
                    )

        # --- 표(table)에 for 루프 + row 변수 삽입 ---
        # 실제 DOCX 표 안에 Jinja2 for 구문을 넣는다
        _insert_table_jinja2_variables(doc, sections)

        # 변환 결과를 바이트로 저장
        output = io.BytesIO()
        doc.save(output)
        output.seek(0)

        logger.info(
            "DOCX 빈 양식 → Jinja2 자동 변환 완료: 변수 %d개 생성",
            len(variables_meta),
        )

        return output.read(), {"variables": variables_meta}

    except Exception as e:
        logger.error("DOCX 빈 양식 자동 변환 실패: %s", str(e))
        raise ValueError(f"DOCX 빈 양식 자동 변환에 실패했습니다: {e}") from e


def _insert_variable_after_heading(
    doc: Any,
    heading_index: int,
    variable_name: str,
) -> None:
    """heading 바로 다음의 빈 paragraph에 {{ 변수명 }}을 삽입한다.

    heading_index 다음 paragraph가 비어 있으면 거기에 변수를 넣고,
    비어 있지 않으면 새 paragraph를 추가한다.

    Args:
        doc: python-docx Document 객체
        heading_index: heading paragraph의 인덱스
        variable_name: 삽입할 Jinja2 변수명
    """
    jinja2_text = "{{ " + variable_name + " }}"

    # heading 바로 다음 paragraph를 확인한다
    next_idx = heading_index + 1
    if next_idx < len(doc.paragraphs):
        next_para = doc.paragraphs[next_idx]
        if not next_para.text.strip():
            # 빈 paragraph에 변수 텍스트를 삽입한다
            if next_para.runs:
                next_para.runs[0].text = jinja2_text
            else:
                next_para.add_run(jinja2_text)
            return

    # 빈 paragraph가 없으면, heading paragraph 뒤에 변수 텍스트를 가진
    # 새 run을 추가한다 (실제로는 다음 paragraph에 넣는 것이 더 자연스럽지만
    # python-docx에서는 paragraph 삽입이 복잡하므로 기존 빈 paragraph에 넣는
    # 방식을 우선한다)
    heading_para = doc.paragraphs[heading_index]
    # heading 다음에 내용이 있는 경우, heading 자체 뒤에 줄바꿈+변수를 추가
    # 하지 않고 건너뛴다 (이미 내용이 있으므로)
    logger.debug(
        "heading '%s' 다음에 빈 paragraph 없음 — 변수 삽입 건너뜀",
        heading_para.text[:50],
    )


def _insert_table_form_variable(
    doc: Any,
    label_text: str,
    variable_name: str,
    field_type: str,
) -> None:
    """표 기반 폼 필드의 빈 셀에 {{ 변수명 }}을 삽입한다.

    _analyze_table_form_fields()에서 감지된 라벨에 대응하는 빈 셀을 찾아
    Jinja2 변수를 삽입한다.

    field_type별 처리:
      - "short": 라벨 셀 바로 오른쪽의 빈 셀에 변수를 삽입한다.
      - "long": 라벨 행 바로 아래의 빈 행에 변수를 삽입한다.

    Args:
        doc: python-docx Document 객체
        label_text: 라벨 원본 텍스트 (_heading_to_variable_name 적용 전의 텍스트)
        variable_name: 삽입할 Jinja2 변수명 (_heading_to_variable_name 적용 후)
        field_type: "short" 또는 "long"
    """
    jinja2_text = "{{ " + variable_name + " }}"

    for table in doc.tables:
        rows = list(table.rows)
        for row_idx, row in enumerate(rows):
            # 병합된 셀 중복 제거
            seen_cell_ids: set[int] = set()
            unique_cells: list[Any] = []
            for cell in row.cells:
                cell_id = id(cell)
                if cell_id not in seen_cell_ids:
                    seen_cell_ids.add(cell_id)
                    unique_cells.append(cell)

            for cell_idx, cell in enumerate(unique_cells):
                cell_text = cell.text.strip()
                if not cell_text:
                    continue

                # 셀 텍스트에서 변수명을 생성하여 일치하는지 확인한다
                cell_var_name = _heading_to_variable_name(cell_text)
                if cell_var_name != variable_name:
                    continue

                # --- 공통: 오른쪽 빈 셀에 변수 삽입 시도 (short/long 모두) ---
                # 회의록 양식처럼 라벨이 첫 열, 값이 오른쪽 병합 셀인 경우
                # short/long 모두 오른쪽 셀에 먼저 삽입을 시도한다.
                if cell_idx + 1 < len(unique_cells):
                    target_cell = unique_cells[cell_idx + 1]
                    if not target_cell.text.strip():
                        if target_cell.paragraphs:
                            para = target_cell.paragraphs[0]
                            if para.runs:
                                para.runs[0].text = jinja2_text
                            else:
                                para.add_run(jinja2_text)
                        logger.debug(
                            "표 %s 필드 변수 삽입: '%s' → %s (오른쪽 셀)",
                            field_type,
                            cell_text,
                            jinja2_text,
                        )
                        return

                # --- long 타입: 오른쪽 셀이 없으면 아래 행 시도 ---
                if field_type == "long":
                    inserted = False
                    if row_idx + 1 < len(rows):
                        next_row = rows[row_idx + 1]
                        next_cells = list(next_row.cells)
                        if next_cells:
                            target_cell = next_cells[0]
                            if not target_cell.text.strip():
                                if target_cell.paragraphs:
                                    para = target_cell.paragraphs[0]
                                    if para.runs:
                                        para.runs[0].text = jinja2_text
                                    else:
                                        para.add_run(jinja2_text)
                                logger.debug(
                                    "표 long 필드 변수 삽입: '%s' → %s (다음 행)",
                                    cell_text,
                                    jinja2_text,
                                )
                                inserted = True

                    if not inserted:
                        cell.paragraphs[-1].add_run("\n" + jinja2_text)
                        logger.debug(
                            "표 long 필드 변수 삽입: '%s' → %s (같은 셀에 추가)",
                            cell_text,
                            jinja2_text,
                        )
                    return

                # --- short 타입인데 오른쪽 빈 셀이 없는 경우 ---
                if field_type == "short":
                    cell.paragraphs[-1].add_run("\n" + jinja2_text)
                    logger.debug(
                        "표 short 필드 변수 삽입: '%s' → %s (같은 셀, 오른쪽 없음)",
                        cell_text,
                        jinja2_text,
                    )
                    return

    # 매칭 실패 시 경고 로그만 출력하고, 에러는 발생시키지 않는다
    logger.debug(
        "표에서 라벨 '%s'에 대응하는 빈 셀을 찾지 못함 — 변수 삽입 건너뜀",
        label_text,
    )


def _insert_table_jinja2_variables(doc: Any, sections: list[dict]) -> None:
    """DOCX 문서의 표에 Jinja2 for 루프 변수를 삽입한다.

    표가 있는 섹션에 대해, 표의 데이터 행(두 번째 행)에
    {% for row in 변수명_항목 %} / {{ row.열이름 }} / {% endfor %}를 삽입한다.

    Args:
        doc: python-docx Document 객체
        sections: 구조 분석에서 얻은 섹션 리스트
    """
    # 표가 있는 섹션 목록을 만든다
    table_sections = [s for s in sections if s.get("has_table")]

    if not table_sections or not doc.tables:
        return

    # 단순 매핑: 섹션 순서와 표 순서를 1:1 매칭한다
    # (하나의 섹션에 하나의 표가 대응된다고 가정)
    for _idx, (section, table) in enumerate(zip(table_sections, doc.tables, strict=False)):
        headers = section.get("table_headers", [])
        if not headers:
            continue

        var_name = _heading_to_variable_name(section["heading"])
        loop_var = var_name + "_항목"

        # 표에 데이터 행이 1개 이상 있어야 변수를 삽입할 수 있다
        if len(table.rows) < 2:
            continue

        # 두 번째 행(데이터 행)에 row 변수를 삽입한다
        data_row = table.rows[1]
        for cell_idx, cell in enumerate(data_row.cells):
            if cell_idx < len(headers):
                # 열 이름에서 공백을 제거하여 변수명으로 사용
                col_var = re.sub(r"\s+", "", headers[cell_idx])
                cell_text = "{{ row." + col_var + " }}"
            else:
                cell_text = ""

            # 셀 내용을 변수로 교체한다
            for para in cell.paragraphs:
                if para.runs:
                    para.runs[0].text = cell_text
                    for run in para.runs[1:]:
                        run.text = ""
                elif not para.text.strip():
                    para.add_run(cell_text)

        # 첫 번째 셀에 for 루프 시작 태그를 prefix로 추가한다
        first_cell = data_row.cells[0]
        if first_cell.paragraphs:
            first_para = first_cell.paragraphs[0]
            current_text = first_para.text
            for_tag = "{%- for row in " + loop_var + " %}"
            if first_para.runs:
                first_para.runs[0].text = for_tag + current_text
            else:
                first_para.add_run(for_tag + current_text)

        # 마지막 셀에 endfor 태그를 suffix로 추가한다
        last_cell = data_row.cells[-1]
        if last_cell.paragraphs:
            last_para = last_cell.paragraphs[-1]
            current_text = last_para.text
            endfor_tag = "{%- endfor %}"
            if last_para.runs:
                last_para.runs[-1].text = last_para.runs[-1].text + endfor_tag
            else:
                last_para.add_run(current_text + endfor_tag)

        logger.debug(
            "표 Jinja2 변수 삽입: 섹션='%s', 열=%s",
            section["heading"],
            headers,
        )


def _auto_fill_pptx(file_bytes: bytes, structure: dict) -> tuple[bytes, dict]:
    """PPTX 빈 양식에 Jinja2 변수를 자동 삽입한다.

    각 슬라이드의 빈 본문 영역에 {{ 변수명 }} 을 삽입하고,
    표가 있는 슬라이드에는 for 루프 변수를 삽입한다.

    Args:
        file_bytes: PPTX 파일의 바이트 데이터
        structure: analyze_blank_form()의 반환값

    Returns:
        (변환된 PPTX 바이트, 변수 메타데이터 dict)
    """
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        sections = structure.get("sections", [])
        doc_title = structure.get("title")

        # 생성된 변수 메타데이터를 수집하는 리스트
        variables_meta: list[dict] = []

        # 섹션 인덱스 → 슬라이드 매핑을 위해 title이 있는 슬라이드를 순서대로 추적

        for slide in prs.slides:
            # 이 슬라이드의 title을 찾는다
            slide_title: str | None = None
            title_shape = None
            body_shapes: list = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    is_title = False
                    if (
                        hasattr(shape, "placeholder_format")
                        and shape.placeholder_format
                        and shape.placeholder_format.idx == 0
                    ):
                        is_title = True

                    if is_title:
                        slide_title = shape.text_frame.text.strip()
                        title_shape = shape
                    else:
                        body_shapes.append(shape)

            if slide_title is None:
                continue

            # 해당 슬라이드에 대응하는 section을 찾는다
            matched_section = None
            for s in sections:
                if s["heading"] == slide_title:
                    matched_section = s
                    break

            if matched_section is None:
                continue

            has_content = matched_section.get("has_content", False)
            has_table = matched_section.get("has_table", False)

            # --- 제목 슬라이드(첫 슬라이드)의 제목을 변수화 ---
            if doc_title and slide_title == doc_title and title_shape:
                title_var = _heading_to_variable_name(doc_title)
                variables_meta.append(
                    {
                        "name": title_var,
                        "type": "string",
                        "label": doc_title,
                        "description": f"프레젠테이션 제목 ({doc_title})",
                        "required": True,
                        # 제목은 사용자가 직접 입력하는 항목으로 분류
                        "category": classify_variable_category(title_var),
                    }
                )
                # title shape의 텍스트를 변수로 치환
                if title_shape.has_text_frame:
                    for para in title_shape.text_frame.paragraphs:
                        if para.runs:
                            para.runs[0].text = "{{ " + title_var + " }}"
                            for run in para.runs[1:]:
                                run.text = ""
                            break

            # --- 빈 본문 영역에 변수 삽입 ---
            if not has_content and not has_table:
                var_name = _heading_to_variable_name(slide_title)
                content_var = var_name + "_내용"

                # body shape 중 빈 것을 찾아 변수를 삽입한다
                inserted = False
                for shape in body_shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if not text:
                            # 빈 body shape에 변수 삽입
                            if shape.text_frame.paragraphs:
                                para = shape.text_frame.paragraphs[0]
                                if para.runs:
                                    para.runs[0].text = "{{ " + content_var + " }}"
                                else:
                                    para.add_run("{{ " + content_var + " }}")
                            inserted = True
                            break

                if not inserted and body_shapes:
                    # 빈 shape가 없으면 첫 번째 body shape에 삽입
                    shape = body_shapes[0]
                    if shape.has_text_frame and shape.text_frame.paragraphs:
                        para = shape.text_frame.paragraphs[0]
                        if para.runs:
                            para.runs[0].text = "{{ " + content_var + " }}"
                            for run in para.runs[1:]:
                                run.text = ""
                        else:
                            para.add_run("{{ " + content_var + " }}")
                    inserted = True

                variables_meta.append(
                    {
                        "name": content_var,
                        "type": "string",
                        "label": slide_title,
                        "description": f"{slide_title} 슬라이드의 내용",
                        "required": True,
                        # 슬라이드 내용은 대부분 AI가 생성할 항목
                        "category": classify_variable_category(content_var),
                    }
                )

            # --- 표 변수 처리 ---
            if has_table:
                table_headers = matched_section.get("table_headers", [])
                if table_headers:
                    var_name = _heading_to_variable_name(slide_title)
                    loop_var = var_name + "_항목"

                    variables_meta.append(
                        {
                            "name": loop_var,
                            "type": "array",
                            "label": f"{slide_title} 표 데이터",
                            "description": f"열 구조: {', '.join(table_headers)}",
                            "required": True,
                            # 표 데이터 항목은 AI가 자동 생성
                            "category": classify_variable_category(loop_var),
                        }
                    )

                    # 슬라이드의 표에 변수를 삽입한다
                    for shape in slide.shapes:
                        if shape.has_table:
                            tbl = shape.table
                            if len(tbl.rows) >= 2:
                                _fill_pptx_table_row(
                                    tbl,
                                    table_headers,
                                    loop_var,
                                )
                            break

        # 변환 결과를 바이트로 저장
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)

        logger.info(
            "PPTX 빈 양식 → Jinja2 자동 변환 완료: 변수 %d개 생성",
            len(variables_meta),
        )

        return output.read(), {"variables": variables_meta}

    except Exception as e:
        logger.error("PPTX 빈 양식 자동 변환 실패: %s", str(e))
        raise ValueError(f"PPTX 빈 양식 자동 변환에 실패했습니다: {e}") from e


def _fill_pptx_table_row(
    tbl: Any,
    headers: list[str],
    loop_var: str,
) -> None:
    """PPTX 표의 데이터 행에 Jinja2 for 루프 변수를 삽입한다.

    Args:
        tbl: python-pptx Table 객체
        headers: 표 헤더 열 이름 리스트
        loop_var: for 루프의 반복 대상 변수명 (예: "사업개요_항목")
    """
    # 두 번째 행(데이터 행)에 변수를 삽입한다
    data_row = tbl.rows[1]

    for cell_idx, cell in enumerate(data_row.cells):
        if cell_idx < len(headers):
            col_var = re.sub(r"\s+", "", headers[cell_idx])
            cell_text = "{{ row." + col_var + " }}"
        else:
            cell_text = ""

        # 셀 텍스트를 교체한다
        if cell.text_frame.paragraphs:
            para = cell.text_frame.paragraphs[0]
            if para.runs:
                para.runs[0].text = cell_text
                for run in para.runs[1:]:
                    run.text = ""
            else:
                para.add_run(cell_text)

    # 첫 번째 셀에 for 루프 시작 태그를 추가한다
    first_cell = data_row.cells[0]
    if first_cell.text_frame.paragraphs:
        para = first_cell.text_frame.paragraphs[0]
        for_tag = "{%- for row in " + loop_var + " %}"
        if para.runs:
            para.runs[0].text = for_tag + para.runs[0].text
        else:
            para.add_run(for_tag)

    # 마지막 셀에 endfor 태그를 추가한다
    last_cell = data_row.cells[-1]
    if last_cell.text_frame.paragraphs:
        para = last_cell.text_frame.paragraphs[-1]
        endfor_tag = "{%- endfor %}"
        if para.runs:
            para.runs[-1].text = para.runs[-1].text + endfor_tag
        else:
            para.add_run(endfor_tag)


def _replace_in_paragraph(
    paragraph: Any,
    replacements: list[dict[str, str]],
) -> None:
    """하나의 paragraph 내에서 지정된 텍스트를 Jinja2 변수로 치환한다.

    XML run 분리 문제를 해결하기 위해:
      1. 모든 run의 텍스트를 합쳐서 전체 텍스트를 만든다.
      2. 전체 텍스트에서 치환을 수행한다.
      3. 치환이 발생했으면, 결과를 첫 번째 run에 할당하고 나머지 run은 비운다.

    Args:
        paragraph: python-docx Paragraph 또는 python-pptx TextFrame의 paragraph
        replacements: [{"original": "원본 텍스트", "variable": "변수명"}, ...]
    """
    # paragraph의 runs에서 전체 텍스트 조합
    runs = paragraph.runs
    if not runs:
        return

    full_text = "".join(run.text for run in runs)
    if not full_text.strip():
        return

    # 각 치환 항목을 적용
    modified = False
    new_text = full_text

    for item in replacements:
        original = item.get("original", "")
        variable = item.get("variable", "")

        if not original or not variable:
            continue

        # 원본 텍스트가 포함되어 있으면 {{ variable }}로 치환
        jinja2_var = "{{ " + variable + " }}"
        if original in new_text:
            new_text = new_text.replace(original, jinja2_var)
            modified = True

    # 치환이 발생한 경우에만 run 텍스트 업데이트
    if modified:
        # 첫 번째 run에 전체 치환 결과를 할당
        runs[0].text = new_text
        # 나머지 run은 빈 문자열로 설정 (서식 보존, 중복 방지)
        for run in runs[1:]:
            run.text = ""
