"""
Report generation Celery worker.

Generates reports by:
  1. Loading the report record and template from the database / MinIO.
  2. Parsing template variables ({{var}} patterns).
  3. Using RAG to extract content for each variable from source documents.
  4. Filling the template using python-docx, python-pptx, or WeasyPrint.
  5. Uploading the finished report to MinIO.
  6. Updating the report status.
"""

from __future__ import annotations

import contextlib
import io
import json
import logging
import re
import tempfile
from pathlib import Path
from typing import Any

from app.core.config import get_settings
from app.integrations.image_generation.service import ImageGenerationService
from app.workers.celery_app import celery_app
from app.workers.jinja2_engine import (
    classify_variable_category,
    render_docx_jinja2,
    render_pptx_jinja2,
)

logger = logging.getLogger(__name__)
settings = get_settings()


def _get_worker_session():
    """Celery Worker 전용 async session factory."""
    from sqlalchemy.ext.asyncio import AsyncSession, async_sessionmaker, create_async_engine

    engine = create_async_engine(settings.database_url, echo=False, pool_pre_ping=True)
    return async_sessionmaker(bind=engine, class_=AsyncSession, expire_on_commit=False)


def _get_agent_sync(agent_id: str) -> dict | None:
    """에이전트 설정을 DB에서 동기적으로 로드한다.

    Celery worker에서 Agent의 system_prompt, LLM 모델, temperature,
    max_tokens 등을 가져와서 보고서 생성 시 적용한다.

    Args:
        agent_id: 에이전트 UUID 문자열

    Returns:
        에이전트 정보 딕셔너리 또는 None (에이전트를 찾을 수 없는 경우)
    """

    async def _load():
        from sqlalchemy import select

        from app.modules.agents.models import Agent

        async with _get_worker_session()() as session:
            result = await session.execute(
                select(Agent).where(Agent.id == agent_id, Agent.is_active == True)  # noqa: E712
            )
            agent = result.scalars().first()
            if agent is None:
                return None
            return {
                "id": str(agent.id),
                "name": agent.name,
                "system_prompt": agent.system_prompt,
                "llm_provider": getattr(agent, "llm_provider", None),
                "llm_model": agent.llm_model,
                "temperature": agent.temperature,
                "max_tokens": agent.max_tokens,
            }

    return _run_async(_load())


def _get_document_template_sync(template_id: str) -> dict | None:
    """document_templates 테이블에서 템플릿 정보를 동기적으로 로드한다.

    Celery worker에서 Jinja2 렌더링 분기에 필요한 템플릿 정보
    (jinja2_variables 스키마, template_storage_path, image_generation_config 등)를
    DB에서 조회한다.  _get_agent_sync()와 동일한 패턴(_run_async)을 사용한다.

    Args:
        template_id: 문서 템플릿의 UUID 문자열

    Returns:
        템플릿 정보 딕셔너리 또는 None (템플릿을 찾을 수 없는 경우)
    """

    async def _load():
        from sqlalchemy import select

        from app.modules.templates.models import DocumentTemplate

        async with _get_worker_session()() as session:
            result = await session.execute(
                select(DocumentTemplate).where(
                    DocumentTemplate.id == template_id,
                    DocumentTemplate.is_active == True,  # noqa: E712
                )
            )
            tmpl = result.scalars().first()
            if tmpl is None:
                return None
            return {
                "id": str(tmpl.id),
                "name": tmpl.name,
                "description": tmpl.description,
                "template_type": tmpl.template_type,
                "output_format": tmpl.output_format,
                "template_storage_path": tmpl.template_storage_path,
                "jinja2_variables": tmpl.jinja2_variables,
                "rendering_mode": tmpl.rendering_mode,
                "image_generation_config": tmpl.image_generation_config,
                "schema": tmpl.schema_,
                "sample_prompt": tmpl.sample_prompt,
            }

    return _run_async(_load())


# Pattern to match template variables: {{variable_name}}
TEMPLATE_VAR_PATTERN = re.compile(r"\{\{(\w+(?:\.\w+)*)\}\}")


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _run_async(coro):
    """Run an async coroutine from synchronous Celery task context."""
    import asyncio

    try:
        loop = asyncio.get_running_loop()
    except RuntimeError:
        loop = None

    if loop and loop.is_running():
        import concurrent.futures

        with concurrent.futures.ThreadPoolExecutor() as pool:
            return pool.submit(asyncio.run, coro).result()
    else:
        return asyncio.run(coro)


async def _get_report(report_id: str) -> dict[str, Any] | None:
    """리포트 레코드를 DB에서 조회한다.

    ORM 모델로 올바른 테이블(tb_generated_reports)을 참조한다.
    """
    from sqlalchemy import select

    from app.modules.reports.models import GeneratedReport

    async with _get_worker_session()() as session:
        result = await session.execute(select(GeneratedReport).where(GeneratedReport.id == report_id))
        report = result.scalars().first()
        if report is None:
            return None
        return {
            "id": str(report.id),
            "template_id": str(report.template_id) if report.template_id else None,
            "organization_id": str(report.organization_id),
            "title": report.title,
            "status": report.status,
            "output_format": report.output_format,
            "output_storage_path": report.output_storage_path,
            "source_document_ids": [str(sid) for sid in report.source_document_ids]
            if report.source_document_ids
            else [],
            "source_chat_session_id": str(report.source_chat_session_id) if report.source_chat_session_id else None,
            "generation_params": report.generation_params,
            "generated_by": str(report.generated_by),
        }


async def _update_report_status(
    report_id: str,
    status: str,
    output_path: str | None = None,
    error_message: str | None = None,
) -> None:
    """리포트 상태를 업데이트한다.

    ORM update를 사용하여 올바른 테이블(tb_generated_reports)을 참조한다.
    """
    from sqlalchemy import update

    from app.modules.reports.models import GeneratedReport

    async with _get_worker_session()() as session:
        values: dict[str, Any] = {"status": status}
        if output_path is not None:
            values["output_storage_path"] = output_path
        if error_message is not None:
            values["error_message"] = error_message[:1000]

        await session.execute(update(GeneratedReport).where(GeneratedReport.id == report_id).values(**values))
        await session.commit()


async def _get_report_template(template_id: str) -> dict[str, Any] | None:
    """리포트 템플릿 레코드를 DB에서 조회한다.

    ORM 모델로 올바른 테이블(tb_report_templates)을 참조한다.
    """
    from sqlalchemy import select

    from app.modules.reports.models import ReportTemplate

    async with _get_worker_session()() as session:
        result = await session.execute(select(ReportTemplate).where(ReportTemplate.id == template_id))
        template = result.scalars().first()
        if template is None:
            return None
        return {
            "id": str(template.id),
            "organization_id": str(template.organization_id),
            "name": template.name,
            "description": template.description,
            "format": template.format,
            "file_path": template.template_storage_path,
        }


def _extract_template_variables(content: str | bytes) -> list[str]:
    """Extract all {{variable}} patterns from template content."""
    if isinstance(content, bytes):
        # Try to decode for text-based searching
        try:
            text = content.decode("utf-8")
        except UnicodeDecodeError:
            text = content.decode("latin-1")
    else:
        text = content

    variables = TEMPLATE_VAR_PATTERN.findall(text)
    return list(dict.fromkeys(variables))  # deduplicate, preserving order


def _extract_docx_variables(file_data: bytes) -> list[str]:
    """Extract template variables from a DOCX file."""
    from docx import Document

    doc = Document(io.BytesIO(file_data))
    all_text: list[str] = []

    for paragraph in doc.paragraphs:
        all_text.append(paragraph.text)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                all_text.append(cell.text)

    full_text = "\n".join(all_text)
    return _extract_template_variables(full_text)


def _extract_pptx_variables(file_data: bytes) -> list[str]:
    """Extract template variables from a PPTX file."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_data))
    all_text: list[str] = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                all_text.append(shape.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        all_text.append(cell.text)

    full_text = "\n".join(all_text)
    return _extract_template_variables(full_text)


def _rag_extract_content(
    variable_name: str,
    source_document_ids: list[str],
    org_id: str,
) -> str:
    """Use RAG to extract content for a template variable from source documents.

    Performs a hybrid search in Qdrant for the variable name, then uses the
    LLM to synthesise the answer from retrieved chunks.
    """
    import httpx

    from app.integrations.llm.prompts import REPORT_SECTION_PROMPT
    from app.integrations.vector_store.qdrant_client import QdrantService

    # Generate a query embedding for the variable name
    query_text = variable_name.replace("_", " ").replace(".", " ")

    # Get dense embedding for the query
    embed_url = f"{settings.vllm_url}/embeddings"
    headers = {}
    if settings.openai_api_key:
        headers["Authorization"] = f"Bearer {settings.openai_api_key}"

    try:
        embed_response = httpx.post(
            embed_url,
            json={"model": settings.embedding_model, "input": [query_text]},
            headers=headers,
            timeout=60.0,
        )
        embed_response.raise_for_status()
        dense_vector = embed_response.json()["data"][0]["embedding"]
    except Exception as exc:
        logger.warning("Embedding generation for query failed: %s", exc)
        return f"[Could not generate content for: {variable_name}]"

    # Search Qdrant for relevant chunks
    qdrant = QdrantService(
        url=f"http://{settings.qdrant_host}:{settings.qdrant_port}",
        api_key=settings.qdrant_api_key,
    )

    # Build document filter
    doc_filter = None
    if source_document_ids:
        from qdrant_client.models import FieldCondition, Filter, MatchAny

        doc_filter = Filter(
            must=[
                FieldCondition(
                    key="document_id",
                    match=MatchAny(any=source_document_ids),
                ),
            ]
        )

    results = qdrant.hybrid_search(
        org_id=org_id,
        dense_vector=dense_vector,
        sparse_vector=None,
        filters=doc_filter,
        limit=5,
    )

    if not results:
        return f"[No relevant content found for: {variable_name}]"

    # Build context from search results
    context_parts: list[str] = []
    for i, result in enumerate(results, 1):
        payload = result.payload if hasattr(result, "payload") else result
        content = payload.get("content", "")
        section = payload.get("section_title", "")
        prefix = f"[Source {i}]"
        if section:
            prefix += f" ({section})"
        context_parts.append(f"{prefix}\n{content}")

    context = "\n\n---\n\n".join(context_parts)

    # Use LLM to generate the section content
    prompt = REPORT_SECTION_PROMPT.format(
        variable_name=variable_name,
        context=context,
    )

    try:
        llm_response = httpx.post(
            f"{settings.vllm_url}/chat/completions",
            json={
                "model": settings.llm_model,
                "messages": [
                    {"role": "system", "content": "You are a report generation assistant."},
                    {"role": "user", "content": prompt},
                ],
                "temperature": 0.2,
                "max_tokens": settings.llm_max_tokens,
            },
            headers=headers,
            timeout=120.0,
        )
        llm_response.raise_for_status()
        return llm_response.json()["choices"][0]["message"]["content"]
    except Exception as exc:
        logger.warning("LLM generation failed for variable '%s': %s", variable_name, exc)
        return f"[Content generation failed for: {variable_name}]"


def _fill_docx_template(
    file_data: bytes,
    variables: dict[str, str],
) -> bytes:
    """Fill a DOCX template by replacing {{var}} patterns with content."""
    from docx import Document

    doc = Document(io.BytesIO(file_data))

    def _replace_in_paragraph(paragraph):
        full_text = paragraph.text
        for var_name, value in variables.items():
            pattern = "{{" + var_name + "}}"
            if pattern in full_text:
                full_text = full_text.replace(pattern, value)

        if full_text != paragraph.text:
            # Clear existing runs and set new text
            for run in paragraph.runs:
                run.text = ""
            if paragraph.runs:
                paragraph.runs[0].text = full_text
            else:
                paragraph.add_run(full_text)

    for paragraph in doc.paragraphs:
        _replace_in_paragraph(paragraph)

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    _replace_in_paragraph(paragraph)

    output = io.BytesIO()
    doc.save(output)
    return output.getvalue()


def _fill_pptx_template(
    file_data: bytes,
    variables: dict[str, str],
) -> bytes:
    """Fill a PPTX template by replacing {{var}} patterns with content."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_data))

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    full_text = paragraph.text
                    replaced = False
                    for var_name, value in variables.items():
                        pattern = "{{" + var_name + "}}"
                        if pattern in full_text:
                            full_text = full_text.replace(pattern, value)
                            replaced = True

                    if replaced:
                        for run in paragraph.runs:
                            run.text = ""
                        if paragraph.runs:
                            paragraph.runs[0].text = full_text
                        else:
                            paragraph.add_run().text = full_text

            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            full_text = paragraph.text
                            replaced = False
                            for var_name, value in variables.items():
                                pattern = "{{" + var_name + "}}"
                                if pattern in full_text:
                                    full_text = full_text.replace(pattern, value)
                                    replaced = True
                            if replaced and paragraph.runs:
                                for run in paragraph.runs:
                                    run.text = ""
                                paragraph.runs[0].text = full_text

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


def _fill_html_template(
    template_content: str,
    variables: dict[str, str],
) -> bytes:
    """Fill an HTML template and convert to PDF via WeasyPrint."""
    for var_name, value in variables.items():
        pattern = "{{" + var_name + "}}"
        # Escape HTML in the value to prevent injection
        import html

        safe_value = html.escape(value).replace("\n", "<br>")
        template_content = template_content.replace(pattern, safe_value)

    try:
        from weasyprint import HTML

        pdf_bytes = HTML(string=template_content).write_pdf()
        return pdf_bytes
    except ImportError:
        logger.error("WeasyPrint is not installed. Install with: pip install weasyprint")
        raise


async def _get_document_template_name(template_id: str) -> str | None:
    """document_templates 테이블에서 이름을 조회한다."""
    from sqlalchemy import select

    from app.modules.templates.models import DocumentTemplate

    async with _get_worker_session()() as session:
        result = await session.execute(select(DocumentTemplate.name).where(DocumentTemplate.id == template_id))
        row = result.scalar_one_or_none()
        return row


async def _find_report_template_by_name(name: str, org_id: str) -> dict | None:
    """report_templates에서 이름으로 검색하여 파일 경로를 반환한다."""
    from sqlalchemy import select

    from app.modules.reports.models import ReportTemplate

    async with _get_worker_session()() as session:
        # 정확 매칭 먼저, 없으면 ILIKE 부분 매칭
        result = await session.execute(
            select(ReportTemplate)
            .where(
                ReportTemplate.name == name,
                ReportTemplate.organization_id == org_id,
            )
            .limit(1)
        )
        tmpl = result.scalar_one_or_none()
        if tmpl is None:
            # 부분 매칭 시도 (limit 1로 첫 번째만)
            result2 = await session.execute(
                select(ReportTemplate)
                .where(
                    ReportTemplate.name.ilike(f"%{name}%"),
                    ReportTemplate.organization_id == org_id,
                )
                .limit(1)
            )
            tmpl = result2.scalar_one_or_none()
        if tmpl is None:
            return None
        return {
            "id": str(tmpl.id),
            "name": tmpl.name,
            "file_path": tmpl.template_storage_path,
            "format": tmpl.format,
        }


async def _load_source_chunks(
    document_ids: list[str],
    total_limit: int = 400,
    per_doc_min: int = 20,
) -> str:
    """소스 문서의 청크 텍스트를 문서별 균등 쿼터로 로드하여 연결한다.

    H5 핫픽스 배경
    --------------
    구버전은 ``ORDER BY (document_id UUID, chunk_index) LIMIT 100`` 이었다.
    UUID 정렬 특성상 앞쪽 document_id가 100개 청크를 전부 잠식해 다중 문서를
    지정해도 뒤쪽 문서가 통째로 누락되는 문제가 있었다.

    수정 방향
    --------
    1) 문서 ID가 여러 개이면 각 문서마다 ``ceil(total_limit / N)``개 또는
       ``per_doc_min`` 중 큰 값만큼 청크를 로드한다.
    2) 각 문서 내부에서는 ``chunk_index`` 오름차순을 유지해 논리적 순서를
       보존한다.
    3) 반환 문자열은 문서 블록을 순서대로 이어 붙이되, 각 블록 사이에
       빈 줄 두 개를 넣는다(기존 포맷과 호환).

    Args:
        document_ids: 로드할 문서 UUID 문자열 목록.
        total_limit: 전체 청크 상한. 기본 400.
        per_doc_min: 문서당 최소 보장 청크 수. 기본 20.
    """
    from sqlalchemy import select

    from app.modules.documents.models import DocumentChunk

    if not document_ids:
        return ""

    # 중복 제거 (호출 측이 실수로 중복 id를 넘겨도 쿼터 배분이 왜곡되지 않게).
    unique_ids = list(dict.fromkeys(document_ids))
    n_docs = len(unique_ids)

    # 단일 문서: 구버전처럼 한 번에 조회(단순·빠름).
    if n_docs == 1:
        async with _get_worker_session()() as session:
            stmt = (
                select(DocumentChunk.content)
                .where(DocumentChunk.document_id == unique_ids[0])
                .order_by(DocumentChunk.chunk_index)
                .limit(total_limit)
            )
            result = await session.execute(stmt)
            chunks = [row[0] for row in result.fetchall()]
        logger.info("_load_source_chunks: 단일 문서 %d개 청크 로드", len(chunks))
        return "\n\n".join(chunks)

    # 다중 문서: 문서별 쿼터 계산.
    # (ceil 연산을 정수 산술로 구현: (a + b - 1) // b)
    per_doc_quota = max((total_limit + n_docs - 1) // n_docs, per_doc_min)

    blocks: list[str] = []
    per_doc_counts: list[tuple[str, int]] = []

    async with _get_worker_session()() as session:
        for doc_id in unique_ids:
            stmt = (
                select(DocumentChunk.content)
                .where(DocumentChunk.document_id == doc_id)
                .order_by(DocumentChunk.chunk_index)
                .limit(per_doc_quota)
            )
            result = await session.execute(stmt)
            doc_chunks = [row[0] for row in result.fetchall()]
            per_doc_counts.append((str(doc_id), len(doc_chunks)))
            if doc_chunks:
                blocks.append("\n\n".join(doc_chunks))

    logger.info(
        "_load_source_chunks: %d개 문서에서 문서별 청크 수 = %s (quota=%d, total_limit=%d)",
        n_docs,
        per_doc_counts,
        per_doc_quota,
        total_limit,
    )

    return "\n\n".join(blocks)


def _generate_docx(title: str, content: str) -> bytes:
    """AI 생성 텍스트를 DOCX로 변환한다."""
    from docx import Document
    from docx.shared import Pt

    doc = Document()
    doc.add_heading(title, level=1)

    for paragraph_text in content.split("\n"):
        text = paragraph_text.strip()
        if not text:
            continue
        if text.startswith("# "):
            doc.add_heading(text[2:], level=1)
        elif text.startswith("## "):
            doc.add_heading(text[3:], level=2)
        elif text.startswith("### "):
            doc.add_heading(text[4:], level=3)
        else:
            p = doc.add_paragraph(text)
            for run in p.runs:
                run.font.size = Pt(11)

    output = io.BytesIO()
    doc.save(output)
    return output.getvalue()


def _get_layout_by_name(prs, name: str):
    """레이아웃 이름으로 찾기."""
    for layout in prs.slide_master.slide_layouts:
        if layout.name == name:
            return layout
    return None


def _fill_placeholder(slide, ph_idx: int, text: str):
    """placeholder idx로 텍스트 주입."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == ph_idx:
            ph.text = text
            return
    # idx가 없으면 무시


def _add_slide_with_design(prs, layout):
    """레이아웃의 디자인 도형+이미지를 포함하여 슬라이드를 추가한다.

    python-pptx의 add_slide()는 placeholder만 복사하고 레이아웃의
    배경 도형(선, 이미지, 도형 등)은 복사하지 않는다.
    이 함수는:
    1. 레이아웃 XML의 non-placeholder 요소를 슬라이드 spTree에 deepcopy
    2. 이미지 relationship(rId→image파일)도 슬라이드 part에 복사
    """
    from copy import deepcopy

    from pptx.oxml.ns import qn

    slide = prs.slides.add_slide(layout)

    # 레이아웃의 spTree에서 non-placeholder 요소를 XML 레벨로 복사
    layout_cSld = layout.part._element.find(qn("p:cSld"))
    layout_spTree = layout_cSld.find(qn("p:spTree"))
    slide_cSld = slide.part._element.find(qn("p:cSld"))
    slide_spTree = slide_cSld.find(qn("p:spTree"))

    # 레이아웃의 이미지 relationship을 슬라이드에 복사 (rId 매핑)
    rId_map = {}
    for rId, rel in layout.part.rels.items():
        if "image" in str(rel.reltype):
            # 슬라이드에 같은 이미지 파일에 대한 relationship 추가
            new_rId = slide.part.relate_to(rel.target_part, rel.reltype)
            rId_map[rId] = new_rId

    for child in layout_spTree:
        tag = child.tag.split("}")[-1]
        if tag in ("nvGrpSpPr", "grpSpPr"):
            continue

        # placeholder 여부 확인
        nvSpPr = child.find(qn("p:nvSpPr"))
        is_ph = False
        if nvSpPr is not None:
            nvPr = nvSpPr.find(qn("p:nvPr"))
            if nvPr is not None and nvPr.find(qn("p:ph")) is not None:
                is_ph = True

        if not is_ph:
            elem_copy = deepcopy(child)

            # 이미지 참조(r:embed) rId 재매핑
            for blip in elem_copy.iter(qn("a:blip")):
                old_embed = blip.get(qn("r:embed"))
                if old_embed and old_embed in rId_map:
                    blip.set(qn("r:embed"), rId_map[old_embed])

            slide_spTree.insert(1, elem_copy)  # grpSpPr 뒤에 삽입

    return slide


def _build_from_slide_master(template_data: bytes, title: str, content: str, params: dict) -> bytes:
    """슬라이드마스터 기반 PPTX 생성.

    Tech Spec 핵심 원칙: 파일을 처음부터 만들지 않는다.
    1. 원본 pptx를 그대로 복사
    2. python-pptx로 파일 열기
    3. 원하는 레이아웃으로 슬라이드 추가 (add_slide)
    4. placeholder에 텍스트만 주입
    5. 저장
    """
    import tempfile

    from pptx import Presentation
    from pptx.oxml.ns import qn

    # 1. 원본을 임시 파일로 복사 (shutil.copy 원칙)
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp.write(template_data)
        tmp_path = tmp.name

    # 2. 복사본 열기 (마스터/테마 100% 보존)
    prs = Presentation(tmp_path)

    # 레이아웃 이름 목록
    layout_names = [lo.name for lo in prs.slide_master.slide_layouts]
    logger.info("Slide master layouts: %s", layout_names)

    # 기존 슬라이드 개수 기억 (나중에 제거)
    original_slide_count = len(prs.slides)

    # AI 텍스트를 섹션으로 파싱
    sections: list[tuple[str, str]] = []
    current_heading = ""
    current_body: list[str] = []
    for line in content.split("\n"):
        stripped = line.strip()
        if stripped.startswith("# ") or stripped.startswith("## ") or stripped.startswith("### "):
            if current_heading or current_body:
                sections.append((current_heading, "\n".join(current_body)))
            current_heading = stripped.lstrip("#").strip()
            current_body = []
        elif stripped:
            # 마크다운 서식 제거
            clean = stripped.replace("**", "").replace("*", "").replace("`", "")
            current_body.append(clean)
    if current_heading or current_body:
        sections.append((current_heading, "\n".join(current_body)))

    date_str = params.get("date", "")
    subtitle = params.get("subtitle", "")
    company = params.get("company", "")

    # 가로형 감지: 슬라이드 너비 > 높이
    is_landscape = prs.slide_width > prs.slide_height

    if is_landscape:
        # ── 가로형 ──
        layout = _get_layout_by_name(prs, "1_표지")
        if layout:
            slide = _add_slide_with_design(prs, layout)
            _fill_placeholder(slide, 0, title)
            _fill_placeholder(slide, 1, subtitle or company)
            _fill_placeholder(slide, 10, date_str)

        layout = _get_layout_by_name(prs, "1_인덱스")
        if layout:
            slide = _add_slide_with_design(prs, layout)
            _fill_placeholder(slide, 11, "\n".join(str(i + 1).zfill(2) for i in range(len(sections))))
            _fill_placeholder(slide, 1, "\n".join(h for h, _ in sections))

        section_layouts = ["Ⅰ. 본문", "Ⅱ. 본문", "Ⅲ. 본문", "Ⅳ. 본문", "Ⅴ. 본문", "Ⅵ. 본문", "Ⅶ. 본문", "Ⅷ. 본문"]
        divider_layouts = [
            "2_소단원 표지",
            "3_소단원 표지",
            "4_소단원 표지",
            "5_소단원 표지",
            "6_소단원 표지",
            "7_소단원 표지",
        ]

        for i, (heading, body) in enumerate(sections[:8]):
            div = _get_layout_by_name(prs, divider_layouts[i % len(divider_layouts)])
            if div:
                ds = _add_slide_with_design(prs, div)
                _fill_placeholder(ds, 12, heading)
                _fill_placeholder(ds, 11, f"{i + 1:02d}")

            bl = _get_layout_by_name(prs, section_layouts[i % len(section_layouts)])
            if bl:
                bs = _add_slide_with_design(prs, bl)
                _fill_placeholder(bs, 0, heading)
                _fill_placeholder(bs, 1, body[:2000])

        layout = _get_layout_by_name(prs, "9_마침화면")
        if layout:
            slide = _add_slide_with_design(prs, layout)
            _fill_placeholder(slide, 10, "감사합니다")

    else:
        # ── 세로형 ──
        layout = _get_layout_by_name(prs, "표지-앞면")
        if layout:
            slide = _add_slide_with_design(prs, layout)
            _fill_placeholder(slide, 0, title)
            _fill_placeholder(slide, 1, subtitle or company)
            _fill_placeholder(slide, 10, date_str)

        layout = _get_layout_by_name(prs, "목차 및 조견표")
        if layout:
            slide = _add_slide_with_design(prs, layout)
            _fill_placeholder(slide, 0, "목  차")
            toc = "\n".join(f"{i + 1}. {h}" for i, (h, _) in enumerate(sections))
            _fill_placeholder(slide, 16, toc)

        body_layouts = ["본문 - Ⅰ", "본문 - Ⅱ", "본문 - Ⅲ", "본문 - Ⅳ", "본문 - Ⅴ", "본문 - Ⅵ", "본문 - Ⅶ", "본문 - Ⅷ"]
        for i, (heading, body) in enumerate(sections[:8]):
            bl = _get_layout_by_name(prs, body_layouts[i % len(body_layouts)])
            if bl:
                bs = _add_slide_with_design(prs, bl)
                _fill_placeholder(bs, 0, heading)
                _fill_placeholder(bs, 16, body[:2000])

        layout = _get_layout_by_name(prs, "맺음말")
        if layout:
            slide = _add_slide_with_design(prs, layout)
            _fill_placeholder(slide, 0, "감사합니다")
            _fill_placeholder(slide, 1, company or "")

    # 5. 원본에 있던 기존 슬라이드 제거 (앞에서부터 original_slide_count개)
    #    add_slide는 항상 끝에 추가되므로, 앞쪽 슬라이드가 원본
    if original_slide_count > 0:
        sldIdLst = prs.part._element.find(qn("p:sldIdLst"))
        for _ in range(original_slide_count):
            if len(sldIdLst) > 0:
                first = sldIdLst[0]
                rId = first.get(qn("r:id"))
                if rId:
                    with contextlib.suppress(Exception):
                        prs.part.drop_rel(rId)
                sldIdLst.remove(first)

    output = io.BytesIO()
    prs.save(output)

    # 임시 파일 정리
    with contextlib.suppress(Exception):
        Path(tmp_path).unlink()

    return output.getvalue()


# ---------------------------------------------------------------------------
# Structured Output 관련 함수들
# ---------------------------------------------------------------------------


def _call_llm_structured(
    system_msg: str,
    user_msg: str,
    json_schema: dict,
    api_key: str,
    model: str = "gpt-4o",
    temperature: float = 0.3,
    max_tokens: int = 16384,
    provider_override: str | None = None,
) -> dict:
    """LLMClient 추상화를 통해 Structured Outputs JSON 응답을 받는다.

    팩토리에서 report 태스크용 프로바이더를 자동 해석하고,
    동기 메서드(generate_structured_sync)를 사용하여 Celery worker에서 안전하게 호출한다.

    Args:
        system_msg: 시스템 프롬프트 (역할 지시)
        user_msg: 사용자 메시지 (생성 요청 내용)
        json_schema: OpenAI Structured Outputs 형식의 JSON 스키마
        api_key: LLM API 인증 키
        model: 사용할 LLM 모델명 (기본: gpt-4o)
        temperature: 생성 온도 (0에 가까울수록 결정적)
        max_tokens: 최대 생성 토큰 수
        provider_override: 에이전트가 지정한 프로바이더 (None이면 시스템 기본값)

    Returns:
        파싱된 JSON 딕셔너리

    Raises:
        Exception: LLM 호출 실패 또는 JSON 파싱 실패 시
    """
    from app.integrations.llm.factory import create_llm_client, get_provider_for_task

    # 에이전트 프로바이더 우선, 없으면 시스템 기본값 (기능별 오버라이드 적용)
    provider = provider_override or get_provider_for_task("report")
    client = create_llm_client(provider, api_key=api_key, model=model)

    # 메시지 구성
    messages = [
        {"role": "system", "content": system_msg},
        {"role": "user", "content": user_msg},
    ]

    # LLMClient의 동기 Structured Outputs 메서드를 호출한다
    return client.generate_structured_sync(
        messages=messages,
        json_schema=json_schema,
        temperature=temperature,
        max_tokens=max_tokens,
    )


# ---------------------------------------------------------------------------
# 표(Table) 삽입 함수들
# ---------------------------------------------------------------------------


def _set_cell_bg(cell, color_hex: str):
    """PPTX 테이블 셀의 배경색을 설정한다.

    python-pptx에서 셀 배경색은 XML을 직접 조작해야 한다.
    tcPr(테이블 셀 속성) 요소에 solidFill을 추가하는 방식이다.

    Args:
        cell: python-pptx의 테이블 셀 객체
        color_hex: 6자리 HEX 색상 코드 (예: "34495E")
    """
    from lxml import etree
    from pptx.oxml.ns import qn

    # 셀의 tcPr(테이블 셀 속성) 요소를 가져오거나 생성한다
    tc = cell._tc
    tcPr = tc.find(qn("a:tcPr"))
    if tcPr is None:
        tcPr = etree.SubElement(tc, qn("a:tcPr"))

    # 기존 solidFill이 있으면 제거하고 새로 추가한다
    for old_fill in tcPr.findall(qn("a:solidFill")):
        tcPr.remove(old_fill)

    solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", color_hex)


def _add_table_to_slide(slide, table_data: dict):
    """PPTX 슬라이드에 표를 추가한다.

    IDINO 브랜드 스타일(#34495E 짙은 남색 헤더, 흰색 텍스트)을 적용하여
    전문적인 표를 생성한다.

    Args:
        slide: python-pptx 슬라이드 객체
        table_data: {"headers": [...], "rows": [[...]]} 형태의 표 데이터
    """
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    headers = table_data.get("headers", [])
    rows = table_data.get("rows", [])
    if not headers:
        return

    # 표의 크기와 위치를 설정한다 (슬라이드 중앙 하단)
    col_count = len(headers)
    row_count = len(rows) + 1  # 헤더 행 포함
    left = Inches(0.5)
    top = Inches(3.0)
    width = Inches(9.0)
    height = Inches(0.4) * row_count

    # 표 도형을 슬라이드에 추가한다
    table_shape = slide.shapes.add_table(
        row_count,
        col_count,
        left,
        top,
        width,
        height,
    )
    table = table_shape.table

    # 헤더 행 스타일 적용: IDINO 짙은 남색 배경 + 흰색 볼드 텍스트
    for col_idx, header_text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header_text
        _set_cell_bg(cell, "34495E")
        # 헤더 텍스트를 흰색 볼드로 설정한다
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                run.font.bold = True
                run.font.size = Pt(10)

    # 데이터 행을 채운다
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            if col_idx < col_count:
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(cell_text)
                # 데이터 셀 폰트 크기 설정
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(9)


def _add_table_to_docx(doc, table_data: dict):
    """DOCX 문서에 표를 추가한다.

    IDINO 스타일 헤더(짙은 남색 배경 + 흰색 텍스트)를 적용하여
    전문적인 표를 생성한다.

    Args:
        doc: python-docx Document 객체
        table_data: {"headers": [...], "rows": [[...]]} 형태의 표 데이터
    """
    from docx.oxml.ns import qn as docx_qn
    from docx.shared import Pt, RGBColor
    from lxml import etree

    headers = table_data.get("headers", [])
    rows = table_data.get("rows", [])
    if not headers:
        return

    # 표를 생성한다 (헤더 + 데이터 행)
    col_count = len(headers)
    table = doc.add_table(rows=1 + len(rows), cols=col_count)
    table.style = "Table Grid"

    # 헤더 행을 채운다
    header_row = table.rows[0]
    for col_idx, header_text in enumerate(headers):
        cell = header_row.cells[col_idx]
        cell.text = header_text

        # 셀 배경색을 IDINO 짙은 남색(#34495E)으로 설정한다
        tc = cell._tc
        tcPr = tc.find(docx_qn("w:tcPr"))
        if tcPr is None:
            tcPr = etree.SubElement(tc, docx_qn("w:tcPr"))
        shading = etree.SubElement(tcPr, docx_qn("w:shd"))
        shading.set(docx_qn("w:fill"), "34495E")
        shading.set(docx_qn("w:val"), "clear")

        # 헤더 텍스트를 흰색 볼드로 설정한다
        for paragraph in cell.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                run.font.bold = True
                run.font.size = Pt(10)

    # 데이터 행을 채운다
    for row_idx, row_data in enumerate(rows):
        row = table.rows[row_idx + 1]
        for col_idx, cell_text in enumerate(row_data):
            if col_idx < col_count:
                cell = row.cells[col_idx]
                cell.text = str(cell_text)
                for paragraph in cell.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(10)

    # 표 뒤에 빈 문단을 추가하여 간격을 확보한다
    doc.add_paragraph("")


# ---------------------------------------------------------------------------
# 차트(Chart) 삽입 함수들
# ---------------------------------------------------------------------------


def _add_chart_to_slide(slide, chart_data: dict):
    """PPTX 슬라이드에 차트를 추가한다 (bar/pie/line).

    python-pptx의 내장 차트 기능을 사용하여 카테고리 차트를 생성한다.
    chart_type에 따라 막대/원형/꺾은선 차트가 생성된다.

    Args:
        slide: python-pptx 슬라이드 객체
        chart_data: {chart_type, title, categories, series} 형태의 차트 데이터
    """
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    chart_type_str = chart_data.get("chart_type", "bar")
    title = chart_data.get("title", "")
    categories = chart_data.get("categories", [])
    series_list = chart_data.get("series", [])

    if not categories or not series_list:
        return

    # 차트 유형을 python-pptx 열거형으로 매핑한다
    chart_type_map = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,  # 묶은 세로 막대형
        "pie": XL_CHART_TYPE.PIE,  # 원형
        "line": XL_CHART_TYPE.LINE_MARKERS,  # 꺾은선(마커 포함)
    }
    xl_chart_type = chart_type_map.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

    # 카테고리 차트 데이터를 구성한다
    chart_data_obj = CategoryChartData()
    chart_data_obj.categories = categories

    for s in series_list:
        chart_data_obj.add_series(s.get("name", ""), s.get("values", []))

    # 차트를 슬라이드 중앙 하단에 추가한다
    chart_shape = slide.shapes.add_chart(
        xl_chart_type,
        Inches(1.0),  # 왼쪽 여백
        Inches(3.0),  # 상단 여백
        Inches(8.0),  # 너비
        Inches(4.0),  # 높이
        chart_data_obj,
    )

    # 차트 제목을 설정한다
    chart = chart_shape.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = title


def _add_chart_to_docx(doc, chart_data: dict):
    """DOCX 문서에 차트를 matplotlib 이미지로 삽입한다.

    python-docx는 네이티브 차트를 지원하지 않으므로,
    matplotlib로 차트를 렌더링한 뒤 PNG 이미지로 삽입한다.
    Malgun Gothic 폰트를 사용하여 한글이 깨지지 않도록 한다.

    Args:
        doc: python-docx Document 객체
        chart_data: {chart_type, title, categories, series} 형태의 차트 데이터
    """
    import matplotlib

    matplotlib.use("Agg")  # GUI 없는 백엔드 사용 (서버 환경)
    import matplotlib.pyplot as plt
    from docx.shared import Inches

    chart_type_str = chart_data.get("chart_type", "bar")
    title = chart_data.get("title", "")
    categories = chart_data.get("categories", [])
    series_list = chart_data.get("series", [])

    if not categories or not series_list:
        return

    # 한글 폰트 설정 (Malgun Gothic = Windows 기본 한글 폰트)
    plt.rcParams["font.family"] = "Malgun Gothic"
    plt.rcParams["axes.unicode_minus"] = False  # 마이너스 기호 깨짐 방지

    fig, ax = plt.subplots(figsize=(8, 5))

    # 차트 유형에 따라 적절한 matplotlib 차트를 그린다
    if chart_type_str == "pie":
        # 원형 차트: 첫 번째 시리즈만 사용한다
        values = series_list[0].get("values", [])
        ax.pie(values, labels=categories, autopct="%1.1f%%", startangle=90)
    elif chart_type_str == "line":
        # 꺾은선 차트: 모든 시리즈를 선으로 표시한다
        for s in series_list:
            ax.plot(categories, s.get("values", []), marker="o", label=s.get("name", ""))
        ax.legend()
        ax.grid(True, alpha=0.3)
    else:
        # 막대 차트 (기본): 시리즈별 막대를 나란히 표시한다
        import numpy as np

        x = np.arange(len(categories))
        bar_width = 0.8 / max(len(series_list), 1)
        for idx, s in enumerate(series_list):
            offset = (idx - len(series_list) / 2 + 0.5) * bar_width
            ax.bar(x + offset, s.get("values", []), bar_width, label=s.get("name", ""))
        ax.set_xticks(x)
        ax.set_xticklabels(categories)
        ax.legend()

    ax.set_title(title)
    plt.tight_layout()

    # 차트를 바이트 버퍼에 PNG로 저장한다
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)

    # DOCX에 이미지로 삽입한다
    doc.add_picture(buf, width=Inches(6.0))
    doc.add_paragraph("")  # 이미지 뒤 간격


# ---------------------------------------------------------------------------
# 이미지 검색/삽입 함수들
# ---------------------------------------------------------------------------


def _fetch_stock_image(query: str) -> bytes | None:
    """Unsplash API에서 무료 스톡 이미지를 검색한다.

    UNSPLASH_ACCESS_KEY 환경변수가 설정되어 있어야 동작한다.
    API 키가 없거나 검색 결과가 없으면 None을 반환한다.

    Args:
        query: 영어 이미지 검색 키워드

    Returns:
        이미지 바이트 데이터 또는 None
    """
    import httpx

    # Unsplash API 키가 없으면 이미지 검색을 건너뛴다
    access_key = settings.unsplash_access_key
    if not access_key:
        logger.debug("Unsplash API key not configured, skipping image fetch")
        return None

    try:
        # Unsplash 검색 API로 이미지를 검색한다
        search_resp = httpx.get(
            "https://api.unsplash.com/search/photos",
            params={"query": query, "per_page": 1, "orientation": "landscape"},
            headers={"Authorization": f"Client-ID {access_key}"},
            timeout=15.0,
        )
        search_resp.raise_for_status()
        results = search_resp.json().get("results", [])
        if not results:
            return None

        # 검색된 첫 번째 이미지의 일반 크기 URL로 다운로드한다
        image_url = results[0]["urls"]["regular"]
        img_resp = httpx.get(image_url, timeout=15.0)
        img_resp.raise_for_status()
        return img_resp.content

    except Exception as exc:
        logger.warning("Failed to fetch stock image for '%s': %s", query, exc)
        return None


def _add_image_to_slide(slide, image_query: str):
    """슬라이드에 Unsplash 스톡 이미지를 추가한다.

    image_query 키워드로 이미지를 검색하여 슬라이드에 삽입한다.
    이미지를 찾지 못하면 아무 동작도 하지 않는다.

    Args:
        slide: python-pptx 슬라이드 객체
        image_query: 영어 이미지 검색 키워드
    """
    from pptx.util import Inches

    image_bytes = _fetch_stock_image(image_query)
    if image_bytes is None:
        return

    # 이미지를 슬라이드 중앙 하단에 배치한다
    img_stream = io.BytesIO(image_bytes)
    slide.shapes.add_picture(
        img_stream,
        Inches(1.5),  # 왼쪽 여백
        Inches(3.0),  # 상단 여백
        Inches(7.0),  # 너비
        Inches(4.0),  # 높이
    )


# ---------------------------------------------------------------------------
# 레이아웃 카탈로그 빌드
# ---------------------------------------------------------------------------


def _build_layout_catalog(prs) -> dict:
    """IDINO 슬라이드마스터의 레이아웃을 layout_type으로 매핑한다.

    IDINO 슬라이드마스터에는 가로형과 세로형 두 가지가 있다.
    각 레이아웃 이름을 Structured Output의 layout_type 열거형 값에 매핑하여
    슬라이드 생성 시 올바른 레이아웃을 선택할 수 있게 한다.

    가로형 매핑:
        1_표지 → title, 1_인덱스 → index,
        2~7_소단원 표지 → section_divider,
        Ⅰ~Ⅷ. 본문 → body_text/body_with_table/body_with_chart/body_with_image/two_column,
        9_마침화면 → closing

    세로형 매핑:
        표지-앞면 → title, 목차 및 조견표 → index,
        본문-Ⅰ~Ⅷ → body_text 계열,
        맺음말 → closing

    Args:
        prs: python-pptx Presentation 객체

    Returns:
        {layout_type: layout_object} 매핑 딕셔너리.
        body_text/body_with_table/body_with_chart/body_with_image/two_column은
        모두 같은 본문 레이아웃을 공유한다.
    """
    catalog: dict = {}

    # 슬라이드 너비가 높이보다 크면 가로형, 아니면 세로형
    is_landscape = prs.slide_width > prs.slide_height

    # 레이아웃 이름 → layout_type 매핑 규칙을 정의한다
    if is_landscape:
        # ── 가로형 IDINO 슬라이드마스터 ──
        name_to_type = {
            "1_표지": "title",
            "1_인덱스": "index",
            "9_마침화면": "closing",
        }
        # 소단원 표지(2~7)는 모두 section_divider에 매핑
        for i in range(2, 8):
            name_to_type[f"{i}_소단원 표지"] = "section_divider"

        # 본문 레이아웃(Ⅰ~Ⅷ)은 body 계열 레이아웃에 매핑
        # 첫 번째 본문 레이아웃을 기본으로 사용한다
        body_layout_names = [
            "Ⅰ. 본문",
            "Ⅱ. 본문",
            "Ⅲ. 본문",
            "Ⅳ. 본문",
            "Ⅴ. 본문",
            "Ⅵ. 본문",
            "Ⅶ. 본문",
            "Ⅷ. 본문",
        ]
    else:
        # ── 세로형 IDINO 슬라이드마스터 ──
        name_to_type = {
            "표지-앞면": "title",
            "목차 및 조견표": "index",
            "맺음말": "closing",
        }
        body_layout_names = [
            "본문 - Ⅰ",
            "본문 - Ⅱ",
            "본문 - Ⅲ",
            "본문 - Ⅳ",
            "본문 - Ⅴ",
            "본문 - Ⅵ",
            "본문 - Ⅶ",
            "본문 - Ⅷ",
        ]

    # 슬라이드마스터의 모든 레이아웃을 순회하며 매핑한다
    for layout in prs.slide_master.slide_layouts:
        layout_name = layout.name

        # 1. 직접 매핑 (title, index, closing, section_divider)
        if layout_name in name_to_type:
            catalog[name_to_type[layout_name]] = layout

        # 2. 본문 레이아웃은 body 계열 모든 타입에 공유
        if layout_name in body_layout_names and "body_text" not in catalog:
            # 본문 계열 레이아웃은 모두 같은 레이아웃을 사용한다
            for body_type in [
                "body_text",
                "body_with_table",
                "body_with_chart",
                "body_with_image",
                "two_column",
            ]:
                catalog[body_type] = layout

    logger.info("Layout catalog built: %s", list(catalog.keys()))
    return catalog


# ---------------------------------------------------------------------------
# Structured Output 기반 문서 빌드 함수들
# ---------------------------------------------------------------------------


def _build_pptx_from_structured(
    data: dict,
    slide_master_data: bytes | None,
    generation_params: dict,
) -> bytes:
    """Structured Output JSON으로 PPTX를 생성한다.

    GPT-4o가 생성한 구조화된 프레젠테이션 데이터를 기반으로
    실제 PPTX 파일을 만든다. 슬라이드마스터가 있으면 IDINO 디자인을
    적용하고, 없으면 기본 python-pptx 템플릿을 사용한다.

    Args:
        data: {"presentation": {"title", "subtitle", "slides": [...]}} 형태
        slide_master_data: IDINO 슬라이드마스터 PPTX 바이트 (없으면 None)
        generation_params: 생성 옵션 (include_charts, include_tables 등)

    Returns:
        완성된 PPTX 파일의 바이트 데이터
    """
    from pptx import Presentation
    from pptx.oxml.ns import qn
    from pptx.util import Inches, Pt

    presentation = data.get("presentation", {})
    slides_data = presentation.get("slides", [])

    # 생성 옵션: 표/차트/이미지 포함 여부 (기본: True)
    include_tables = generation_params.get("include_tables", True)
    include_charts = generation_params.get("include_charts", True)
    include_images = generation_params.get("include_images", True)

    if slide_master_data:
        # ── 슬라이드마스터 기반 생성 ──
        # 원본 파일을 임시 파일로 복사하여 디자인을 보존한다
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp.write(slide_master_data)
            tmp_path = tmp.name

        prs = Presentation(tmp_path)
        original_slide_count = len(prs.slides)

        # 레이아웃 카탈로그를 빌드하여 layout_type → 실제 레이아웃 매핑
        catalog = _build_layout_catalog(prs)

        for slide_info in slides_data:
            layout_type = slide_info.get("layout_type", "body_text")
            layout = catalog.get(layout_type)

            if layout is None:
                # 매핑되지 않는 layout_type은 body_text로 대체한다
                layout = catalog.get("body_text")
            if layout is None:
                # 그래도 없으면 첫 번째 레이아웃을 사용한다
                layout = prs.slide_master.slide_layouts[0]

            # 디자인 요소를 포함하여 슬라이드를 추가한다
            slide = _add_slide_with_design(prs, layout)

            # placeholder에 텍스트를 주입한다
            heading = slide_info.get("heading", "")
            body = slide_info.get("body", "")
            bullet_points = slide_info.get("bullet_points", [])

            # 제목 placeholder (idx=0)
            _fill_placeholder(slide, 0, heading)

            # 본문 placeholder (idx=1 또는 idx=16, 세로형에서는 16 사용)
            body_text = body
            if bullet_points:
                body_text += "\n" + "\n".join(f"• {bp}" for bp in bullet_points)
            _fill_placeholder(slide, 1, body_text[:2000])
            _fill_placeholder(slide, 16, body_text[:2000])

            # two_column 레이아웃 처리
            if layout_type == "two_column":
                col_left = slide_info.get("column_left", "")
                col_right = slide_info.get("column_right", "")
                if col_left:
                    _fill_placeholder(slide, 1, col_left[:1000])
                if col_right:
                    _fill_placeholder(slide, 2, col_right[:1000])

            # 표 삽입 (include_tables가 True이고 데이터가 있을 때만)
            table_data = slide_info.get("table")
            if include_tables and table_data:
                try:
                    _add_table_to_slide(slide, table_data)
                except Exception:
                    logger.warning("Failed to add table to slide", exc_info=True)

            # 차트 삽입 (include_charts가 True이고 데이터가 있을 때만)
            chart_data = slide_info.get("chart")
            if include_charts and chart_data:
                try:
                    _add_chart_to_slide(slide, chart_data)
                except Exception:
                    logger.warning("Failed to add chart to slide", exc_info=True)

            # 이미지 삽입 (include_images가 True이고 쿼리가 있을 때만)
            image_query = slide_info.get("image_query")
            if include_images and image_query:
                try:
                    _add_image_to_slide(slide, image_query)
                except Exception:
                    logger.warning("Failed to add image to slide", exc_info=True)

            # 발표자 노트 추가
            speaker_notes = slide_info.get("speaker_notes", "")
            if speaker_notes:
                try:
                    slide.notes_slide.notes_text_frame.text = speaker_notes
                except Exception:
                    logger.debug("Failed to add speaker notes", exc_info=True)

        # 원본 슬라이드 제거 (새로 추가한 슬라이드만 남긴다)
        if original_slide_count > 0:
            sldIdLst = prs.part._element.find(qn("p:sldIdLst"))
            for _ in range(original_slide_count):
                if len(sldIdLst) > 0:
                    first = sldIdLst[0]
                    rId = first.get(qn("r:id"))
                    if rId:
                        with contextlib.suppress(Exception):
                            prs.part.drop_rel(rId)
                    sldIdLst.remove(first)

        output = io.BytesIO()
        prs.save(output)

        # 임시 파일 정리
        with contextlib.suppress(Exception):
            Path(tmp_path).unlink()

        return output.getvalue()

    else:
        # ── IDINO 디자인 기본 생성 (슬라이드마스터 없음) ──
        # 슬라이드마스터가 없어도 IDINO 브랜드 디자인을 적용한다.
        # 빈 레이아웃(slide_layouts[6]) 위에 도형+텍스트박스를 수동 배치한다.
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN

        # ── IDINO 디자인 시스템 색상 정의 ──
        CLR_HEADER = RGBColor(0x34, 0x49, 0x5E)  # 진한 남색 - 헤더/푸터 배경
        CLR_ACCENT = RGBColor(0x29, 0x80, 0xB9)  # 밝은 파랑 - 강조
        RGBColor(0xEB, 0xF5, 0xFB)  # 밝은 파랑 배경 - 콘텐츠 박스
        RGBColor(0x1A, 0x52, 0x76)  # 진한 파랑 - 라벨
        CLR_BODY = RGBColor(0x2C, 0x3E, 0x50)  # 진한 회색 - 본문
        CLR_SUBTITLE = RGBColor(0xD5, 0xD8, 0xDC)  # 밝은 회색 - 부제
        CLR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)  # 흰색
        FONT = "Malgun Gothic"  # 맑은 고딕 폰트

        # 16:9 비율의 프레젠테이션을 생성한다
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # 빈 레이아웃 참조 (도형을 수동으로 배치하기 위해 빈 레이아웃 사용)
        blank_layout = prs.slide_layouts[6]

        # -----------------------------------------------------------------
        # IDINO 헬퍼: 도형에 배경색을 채우는 유틸리티
        # -----------------------------------------------------------------
        def _set_shape_fill(shape, color: RGBColor):
            """도형의 배경색을 단색으로 설정한다."""
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            # 테두리를 제거하여 깔끔한 디자인을 만든다
            shape.line.fill.background()

        def _set_tf_text(tf, text: str, font_size: int, color: RGBColor, bold: bool = False, alignment=PP_ALIGN.LEFT):
            """텍스트 프레임에 서식이 적용된 텍스트를 설정한다.

            Args:
                tf: 텍스트 프레임 객체
                text: 표시할 텍스트
                font_size: 폰트 크기 (pt 단위)
                color: 글자 색상
                bold: 굵게 표시 여부
                alignment: 텍스트 정렬 방식
            """
            tf.clear()
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.bold = bold
            p.font.name = FONT
            p.alignment = alignment

        def _add_idino_header_bar(slide):
            """콘텐츠 슬라이드 상단에 IDINO 헤더바(진한 남색)를 추가한다.

            Returns:
                헤더바 도형 (제목 텍스트를 추가할 때 사용)
            """
            # 상단 헤더바: 전체 너비, 높이 0.42인치, 진한 남색 배경
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0),
                Inches(0),  # 좌측 상단에 배치
                Inches(13.333),
                Inches(0.42),  # 전체 너비, 높이 0.42인치
            )
            _set_shape_fill(bar, CLR_HEADER)
            return bar

        def _add_idino_header_title(slide, heading: str):
            """헤더바 위에 제목 텍스트를 추가한다.

            Args:
                slide: 슬라이드 객체
                heading: 제목 텍스트
            """
            # 제목 텍스트박스: 좌측 0.4인치 여백, 헤더바 안에 배치
            txbox = slide.shapes.add_textbox(
                Inches(0.4),
                Inches(0.04),  # 좌측 여백 0.4인치
                Inches(12),
                Inches(0.36),  # 넉넉한 너비
            )
            tf = txbox.text_frame
            tf.word_wrap = True
            _set_tf_text(tf, heading, 13, CLR_WHITE, bold=True)

        def _add_idino_footer_bar(slide):
            """콘텐츠 슬라이드 하단에 IDINO 푸터바를 추가한다.

            하단에 얇은 남색 바를 배치하고, 우측에 "IDINO" 로고 텍스트를 넣는다.
            """
            # 하단 푸터바: 전체 너비, 높이 0.18인치
            footer_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0),
                Inches(7.32),  # 슬라이드 하단에 배치
                Inches(13.333),
                Inches(0.18),
            )
            _set_shape_fill(footer_bar, CLR_HEADER)
            # 우측에 "IDINO" 로고 텍스트를 추가한다
            logo_box = slide.shapes.add_textbox(
                Inches(12.2),
                Inches(7.32),  # 우측에 배치
                Inches(1),
                Inches(0.18),
            )
            tf = logo_box.text_frame
            tf.word_wrap = False
            _set_tf_text(tf, "IDINO", 7, CLR_WHITE, bold=True, alignment=PP_ALIGN.RIGHT)

        # -----------------------------------------------------------------
        # IDINO 슬라이드 생성 헬퍼 함수들
        # -----------------------------------------------------------------

        def _add_idino_title_slide(prs, heading: str, subtitle: str = ""):
            """IDINO 스타일 타이틀 슬라이드를 추가한다.

            전체 배경이 진한 남색이고, 상단에 파란 강조선이 있으며,
            가운데에 제목과 부제를 표시한다.

            Args:
                prs: 프레젠테이션 객체
                heading: 메인 제목 텍스트
                subtitle: 부제 텍스트 (선택)

            Returns:
                생성된 슬라이드 객체
            """
            slide = prs.slides.add_slide(blank_layout)

            # 전체 배경을 진한 남색(#34495E)으로 채운다
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0),
                Inches(0),
                Inches(13.333),
                Inches(7.5),
            )
            _set_shape_fill(bg, CLR_HEADER)

            # 상단 강조선: 밝은 파랑색 얇은 줄 (높이 0.03인치)
            accent_line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0),
                Inches(0),
                Inches(13.333),
                Inches(0.03),
            )
            _set_shape_fill(accent_line, CLR_ACCENT)

            # 제목 텍스트: 32pt, Bold, 흰색, 중앙 정렬
            title_box = slide.shapes.add_textbox(
                Inches(1.5),
                Inches(2.5),
                Inches(10.333),
                Inches(1.5),
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            _set_tf_text(tf, heading, 32, CLR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

            # 부제 텍스트: 14pt, 밝은 회색, 중앙 정렬
            if subtitle:
                sub_box = slide.shapes.add_textbox(
                    Inches(2),
                    Inches(4.2),
                    Inches(9.333),
                    Inches(0.8),
                )
                tf_sub = sub_box.text_frame
                tf_sub.word_wrap = True
                _set_tf_text(tf_sub, subtitle, 14, CLR_SUBTITLE, alignment=PP_ALIGN.CENTER)

            # 하단 "IDINO" 로고 텍스트: 11pt, 밝은 파랑, 중앙 정렬
            logo_box = slide.shapes.add_textbox(
                Inches(5.5),
                Inches(6.5),
                Inches(2.333),
                Inches(0.5),
            )
            tf_logo = logo_box.text_frame
            _set_tf_text(tf_logo, "IDINO", 11, CLR_ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

            return slide

        def _add_idino_content_slide(prs, heading: str, body: str = "", bullet_points: list | None = None):
            """IDINO 스타일 콘텐츠 슬라이드를 추가한다.

            상단 헤더바, 본문 영역, 하단 푸터바로 구성된다.
            bullet_points가 있으면 "•" 접두사를 붙여 표시한다.

            Args:
                prs: 프레젠테이션 객체
                heading: 슬라이드 제목
                body: 본문 텍스트
                bullet_points: 글머리 기호 항목 목록 (선택)

            Returns:
                생성된 슬라이드 객체
            """
            slide = prs.slides.add_slide(blank_layout)

            # 상단 헤더바 + 제목 텍스트
            _add_idino_header_bar(slide)
            _add_idino_header_title(slide, heading)

            # 본문 텍스트를 조합한다 (body + bullet_points)
            content_parts = []
            if body:
                content_parts.append(body)
            if bullet_points:
                for bp in bullet_points:
                    content_parts.append(f"• {bp}")
            full_text = "\n".join(content_parts)[:2000]

            # 본문 텍스트박스: 콘텐츠 영역 (0.7" ~ 6.5" 세로 범위)
            if full_text:
                content_box = slide.shapes.add_textbox(
                    Inches(0.6),
                    Inches(0.7),
                    Inches(12.133),
                    Inches(5.8),
                )
                tf = content_box.text_frame
                tf.word_wrap = True
                # 본문에 여러 줄이 있으면 각 줄을 별도 문단으로 추가한다
                lines = full_text.split("\n")
                for i, line in enumerate(lines):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = line
                    p.font.size = Pt(11)
                    p.font.color.rgb = CLR_BODY
                    p.font.name = FONT
                    p.space_after = Pt(4)

            # 하단 푸터바
            _add_idino_footer_bar(slide)

            return slide

        def _add_idino_index_slide(prs, heading: str, items: list):
            """IDINO 스타일 인덱스(목차) 슬라이드를 추가한다.

            상단 헤더바에 "목차" 제목을 표시하고,
            번호가 매겨진 항목 목록을 본문에 나열한다.

            Args:
                prs: 프레젠테이션 객체
                heading: 목차 제목 (예: "목차")
                items: 목차 항목 텍스트 목록

            Returns:
                생성된 슬라이드 객체
            """
            slide = prs.slides.add_slide(blank_layout)

            # 상단 헤더바 + 제목
            _add_idino_header_bar(slide)
            _add_idino_header_title(slide, heading or "목차")

            # 번호 매기기된 항목 목록을 본문에 추가한다
            if items:
                content_box = slide.shapes.add_textbox(
                    Inches(1.0),
                    Inches(1.0),
                    Inches(11.333),
                    Inches(5.5),
                )
                tf = content_box.text_frame
                tf.word_wrap = True
                for i, item in enumerate(items):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    # 1. 2. 3. 형식으로 번호를 붙인다
                    p.text = f"{i + 1}. {item}"
                    p.font.size = Pt(14)
                    p.font.color.rgb = CLR_BODY
                    p.font.name = FONT
                    p.font.bold = False
                    p.space_after = Pt(8)

            # 하단 푸터바
            _add_idino_footer_bar(slide)

            return slide

        def _add_idino_section_slide(prs, heading: str):
            """IDINO 스타일 섹션 구분 슬라이드를 추가한다.

            전체 배경이 진한 남색이고, 가운데에 섹션 제목이 크게 표시된다.
            하단에 파란색 강조선이 있다.

            Args:
                prs: 프레젠테이션 객체
                heading: 섹션 제목

            Returns:
                생성된 슬라이드 객체
            """
            slide = prs.slides.add_slide(blank_layout)

            # 전체 배경을 진한 남색으로 채운다
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0),
                Inches(0),
                Inches(13.333),
                Inches(7.5),
            )
            _set_shape_fill(bg, CLR_HEADER)

            # 섹션 제목: 28pt, Bold, 흰색, 중앙 정렬
            title_box = slide.shapes.add_textbox(
                Inches(2),
                Inches(2.8),
                Inches(9.333),
                Inches(1.5),
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            _set_tf_text(tf, heading, 28, CLR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

            # 하단 강조선: 밝은 파랑색 (가운데 배치, 너비 4인치)
            accent_line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(4.667),
                Inches(4.5),
                Inches(4),
                Inches(0.03),
            )
            _set_shape_fill(accent_line, CLR_ACCENT)

            return slide

        def _add_idino_closing_slide(prs, heading: str = "감사합니다", contact: str = ""):
            """IDINO 스타일 마침 슬라이드를 추가한다.

            전체 배경이 진한 남색이고, 가운데에 감사 메시지와
            파란색 강조선, 연락처 정보를 표시한다.

            Args:
                prs: 프레젠테이션 객체
                heading: 마침 메시지 (기본: "감사합니다")
                contact: 연락처 정보 (선택)

            Returns:
                생성된 슬라이드 객체
            """
            slide = prs.slides.add_slide(blank_layout)

            # 전체 배경을 진한 남색으로 채운다
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0),
                Inches(0),
                Inches(13.333),
                Inches(7.5),
            )
            _set_shape_fill(bg, CLR_HEADER)

            # 감사 메시지: 36pt, Bold, 흰색, 중앙 정렬
            title_box = slide.shapes.add_textbox(
                Inches(2),
                Inches(2.5),
                Inches(9.333),
                Inches(1.5),
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            _set_tf_text(tf, heading or "감사합니다", 36, CLR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

            # 중앙 강조선: 밝은 파랑색, 너비 2인치
            accent_line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(5.667),
                Inches(4.2),
                Inches(2),
                Inches(0.03),
            )
            _set_shape_fill(accent_line, CLR_ACCENT)

            # 연락처 정보: 12pt, 밝은 회색, 중앙 정렬
            if contact:
                contact_box = slide.shapes.add_textbox(
                    Inches(3),
                    Inches(4.8),
                    Inches(7.333),
                    Inches(0.6),
                )
                tf_c = contact_box.text_frame
                tf_c.word_wrap = True
                _set_tf_text(tf_c, contact, 12, CLR_SUBTITLE, alignment=PP_ALIGN.CENTER)

            return slide

        def _add_idino_two_column_slide(prs, heading: str, left_text: str, right_text: str):
            """IDINO 스타일 2단 레이아웃 슬라이드를 추가한다.

            상단 헤더바 아래에 좌우 두 개의 텍스트 영역을 배치한다.
            각 영역은 약 6인치 너비로, 좌우 균등하게 분할된다.

            Args:
                prs: 프레젠테이션 객체
                heading: 슬라이드 제목
                left_text: 왼쪽 영역 텍스트
                right_text: 오른쪽 영역 텍스트

            Returns:
                생성된 슬라이드 객체
            """
            slide = prs.slides.add_slide(blank_layout)

            # 상단 헤더바 + 제목
            _add_idino_header_bar(slide)
            _add_idino_header_title(slide, heading)

            # 왼쪽 텍스트 영역: 좌측 절반 (약 6인치 너비)
            if left_text:
                left_box = slide.shapes.add_textbox(
                    Inches(0.5),
                    Inches(0.7),
                    Inches(6.0),
                    Inches(5.8),
                )
                tf_l = left_box.text_frame
                tf_l.word_wrap = True
                for i, line in enumerate(left_text[:1000].split("\n")):
                    p = tf_l.paragraphs[0] if i == 0 else tf_l.add_paragraph()
                    p.text = line
                    p.font.size = Pt(11)
                    p.font.color.rgb = CLR_BODY
                    p.font.name = FONT

            # 오른쪽 텍스트 영역: 우측 절반 (약 6인치 너비)
            if right_text:
                right_box = slide.shapes.add_textbox(
                    Inches(6.833),
                    Inches(0.7),
                    Inches(6.0),
                    Inches(5.8),
                )
                tf_r = right_box.text_frame
                tf_r.word_wrap = True
                for i, line in enumerate(right_text[:1000].split("\n")):
                    p = tf_r.paragraphs[0] if i == 0 else tf_r.add_paragraph()
                    p.text = line
                    p.font.size = Pt(11)
                    p.font.color.rgb = CLR_BODY
                    p.font.name = FONT

            # 하단 푸터바
            _add_idino_footer_bar(slide)

            return slide

        # -----------------------------------------------------------------
        # 슬라이드 데이터를 순회하며 적절한 IDINO 헬퍼를 호출한다
        # -----------------------------------------------------------------
        for slide_info in slides_data:
            layout_type = slide_info.get("layout_type", "body_text")
            heading = slide_info.get("heading", "")
            body = slide_info.get("body", "")
            bullet_points = slide_info.get("bullet_points", [])

            # layout_type에 따라 IDINO 디자인 슬라이드를 생성한다
            if layout_type == "title":
                # 타이틀 슬라이드: 전체 남색 배경 + 제목 + 부제
                slide = _add_idino_title_slide(prs, heading, subtitle=body)

            elif layout_type == "index":
                # 인덱스(목차) 슬라이드: 헤더바 + 번호 목록
                items = bullet_points if bullet_points else []
                slide = _add_idino_index_slide(prs, heading, items)

            elif layout_type == "section_divider":
                # 섹션 구분 슬라이드: 전체 남색 배경 + 큰 제목
                slide = _add_idino_section_slide(prs, heading)

            elif layout_type == "closing":
                # 마침 슬라이드: 전체 남색 배경 + 감사 메시지
                slide = _add_idino_closing_slide(prs, heading=heading or "감사합니다", contact=body)

            elif layout_type == "two_column":
                # 2단 레이아웃: 좌우 분할 텍스트
                col_left = slide_info.get("column_left", body)
                col_right = slide_info.get("column_right", "")
                slide = _add_idino_two_column_slide(prs, heading, col_left, col_right)

            else:
                # body_text, body_with_table, body_with_chart, body_with_image 등
                # 모든 콘텐츠 타입은 동일한 IDINO 콘텐츠 슬라이드 사용
                slide = _add_idino_content_slide(prs, heading, body, bullet_points)

            # 표 삽입 (include_tables가 True이고 데이터가 있을 때만)
            table_data = slide_info.get("table")
            if include_tables and table_data:
                try:
                    _add_table_to_slide(slide, table_data)
                except Exception:
                    logger.warning("Failed to add table to slide", exc_info=True)

            # 차트 삽입 (include_charts가 True이고 데이터가 있을 때만)
            chart_data = slide_info.get("chart")
            if include_charts and chart_data:
                try:
                    _add_chart_to_slide(slide, chart_data)
                except Exception:
                    logger.warning("Failed to add chart to slide", exc_info=True)

            # 이미지 삽입 (include_images가 True이고 쿼리가 있을 때만)
            image_query = slide_info.get("image_query")
            if include_images and image_query:
                try:
                    _add_image_to_slide(slide, image_query)
                except Exception:
                    logger.warning("Failed to add image to slide", exc_info=True)

            # 발표자 노트 추가
            speaker_notes = slide_info.get("speaker_notes", "")
            if speaker_notes:
                with contextlib.suppress(Exception):
                    slide.notes_slide.notes_text_frame.text = speaker_notes

        output = io.BytesIO()
        prs.save(output)
        return output.getvalue()


def _build_docx_from_structured(data: dict) -> bytes:
    """Structured Output JSON으로 IDINO 스타일 DOCX를 생성한다.

    GPT-4o가 생성한 구조화된 문서 데이터를 기반으로 실제 DOCX 파일을 만든다.
    IDINO 브랜드 디자인(커버 페이지, 헤더/푸터, 색상, 폰트)을 적용하여
    전문적인 보고서를 생성한다.

    Args:
        data: {"document": {"title", "sections": [...]}} 형태

    Returns:
        완성된 DOCX 파일의 바이트 데이터
    """
    from docx import Document
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml import parse_xml
    from docx.oxml.ns import nsdecls
    from docx.shared import Pt, RGBColor

    # IDINO 디자인 시스템 색상 (DOCX용 RGBColor)
    CLR_HEADER_HEX = "34495E"  # 진한 남색 - 헤더/배너 배경
    CLR_ACCENT_HEX = "2980B9"  # 밝은 파랑 - 강조
    FONT = "Malgun Gothic"  # 맑은 고딕 폰트

    doc = Document()
    document = data.get("document", {})
    title = document.get("title", "")
    sections_data = document.get("sections", [])

    # -----------------------------------------------------------------
    # 기본 스타일 설정: Normal 스타일에 IDINO 폰트와 색상 적용
    # -----------------------------------------------------------------
    style = doc.styles["Normal"]
    style.font.name = FONT
    style.font.size = Pt(11)
    style.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    # 행간 1.15 설정 (가독성 향상)
    style.paragraph_format.line_spacing = 1.15

    # Heading 스타일에도 IDINO 폰트 적용
    for level in range(1, 5):
        style_name = f"Heading {level}"
        if style_name in doc.styles:
            h_style = doc.styles[style_name]
            h_style.font.name = FONT
            h_style.font.color.rgb = RGBColor(0x1A, 0x52, 0x76)  # 진한 파랑

    # -----------------------------------------------------------------
    # 커버 페이지: 다크 배너(표 이용) + 제목 + 부제
    # -----------------------------------------------------------------
    if title:
        # 상단 여백을 줄여 커버 페이지를 넓게 사용한다
        doc.add_paragraph()  # 빈 줄 (상단 여백)

        # 다크 배너: 1행 1열 표를 사용하여 배경색이 있는 배너를 만든다
        banner_table = doc.add_table(rows=1, cols=1)
        banner_table.alignment = WD_TABLE_ALIGNMENT.CENTER
        # 표 너비를 페이지 전체로 설정한다
        banner_cell = banner_table.cell(0, 0)

        # 셀 배경색을 진한 남색(#34495E)으로 설정한다
        shading_elm = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{CLR_HEADER_HEX}"/>')
        banner_cell._tc.get_or_add_tcPr().append(shading_elm)

        # 셀 높이를 설정하여 배너처럼 보이게 한다 (상하 패딩)
        # 배너 안에 제목 텍스트를 추가한다
        banner_p = banner_cell.paragraphs[0]
        banner_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        # 위쪽 여백 (빈 줄)
        banner_p.add_run("\n")
        # 제목 텍스트
        title_run = banner_p.add_run(title)
        title_run.font.size = Pt(26)
        title_run.font.bold = True
        title_run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        title_run.font.name = FONT
        # 아래쪽 여백 (빈 줄)
        banner_p.add_run("\n")

        # 배너 아래에 파란색 강조선을 추가한다 (얇은 표)
        accent_table = doc.add_table(rows=1, cols=1)
        accent_table.alignment = WD_TABLE_ALIGNMENT.CENTER
        accent_cell = accent_table.cell(0, 0)
        accent_shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{CLR_ACCENT_HEX}"/>')
        accent_cell._tc.get_or_add_tcPr().append(accent_shading)
        # 강조선 셀은 최소 높이로 설정한다
        accent_cell.text = ""

        # 커버 페이지 아래 빈 줄 추가
        doc.add_paragraph()

        # IDINO 로고 텍스트
        logo_p = doc.add_paragraph()
        logo_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        logo_run = logo_p.add_run("IDINO")
        logo_run.font.size = Pt(11)
        logo_run.font.bold = True
        logo_run.font.color.rgb = RGBColor(0x29, 0x80, 0xB9)
        logo_run.font.name = FONT

        # 페이지 나누기: 커버 페이지 다음에 본문이 시작되도록 한다
        doc.add_page_break()

    # -----------------------------------------------------------------
    # 헤더/푸터 설정: 문서 제목 + 페이지 번호
    # -----------------------------------------------------------------
    # 기본 섹션의 헤더에 문서 제목을 추가한다
    doc_section = doc.sections[0]
    header = doc_section.header
    header.is_linked_to_previous = False
    header_p = header.paragraphs[0]
    header_p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    header_run = header_p.add_run(title)
    header_run.font.size = Pt(8)
    header_run.font.color.rgb = RGBColor(0x29, 0x80, 0xB9)
    header_run.font.name = FONT

    # 헤더 하단에 파란색 하단 라인을 추가한다
    # (paragraph의 하단 border를 XML로 직접 설정)
    pPr = header_p._p.get_or_add_pPr()
    pBdr = parse_xml(
        f'<w:pBdr {nsdecls("w")}>  <w:bottom w:val="single" w:sz="6" w:space="1" w:color="{CLR_ACCENT_HEX}"/></w:pBdr>'
    )
    pPr.append(pBdr)

    # 푸터에 페이지 번호를 추가한다
    footer = doc_section.footer
    footer.is_linked_to_previous = False
    footer_p = footer.paragraphs[0]
    footer_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    # 페이지 번호 필드를 XML로 삽입한다
    page_run = footer_p.add_run()
    page_run.font.size = Pt(8)
    page_run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    page_run.font.name = FONT
    fldChar_begin = parse_xml(f'<w:fldChar {nsdecls("w")} w:fldCharType="begin"/>')
    page_run._r.append(fldChar_begin)
    instrText = parse_xml(f'<w:instrText {nsdecls("w")} xml:space="preserve"> PAGE </w:instrText>')
    page_run._r.append(instrText)
    fldChar_end = parse_xml(f'<w:fldChar {nsdecls("w")} w:fldCharType="end"/>')
    page_run._r.append(fldChar_end)

    # -----------------------------------------------------------------
    # 각 섹션을 순서대로 문서에 추가한다
    # -----------------------------------------------------------------
    for section in sections_data:
        section_heading = section.get("heading", "")
        heading_level = section.get("heading_level", 1)
        paragraphs = section.get("paragraphs", [])
        bullet_points = section.get("bullet_points", [])
        table_data = section.get("table")
        chart_data = section.get("chart")

        # 섹션 제목 추가: IDINO 스타일 (파란색 언더라인 포함)
        if section_heading:
            h = doc.add_heading(section_heading, level=heading_level)
            for run in h.runs:
                run.font.name = FONT
                run.font.color.rgb = RGBColor(0x1A, 0x52, 0x76)

            # H1 레벨에만 파란색 하단 라인(언더라인)을 추가한다
            if heading_level == 1:
                h_pPr = h._p.get_or_add_pPr()
                h_pBdr = parse_xml(
                    f"<w:pBdr {nsdecls('w')}>"
                    f'  <w:bottom w:val="single" w:sz="8" w:space="2" w:color="{CLR_ACCENT_HEX}"/>'
                    f"</w:pBdr>"
                )
                h_pPr.append(h_pBdr)

        # 본문 문단 추가: Malgun Gothic 11pt, 1.15 행간
        for para_text in paragraphs:
            p = doc.add_paragraph(para_text)
            for run in p.runs:
                run.font.name = FONT
                run.font.size = Pt(11)
                run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

        # 글머리 기호 항목 추가
        for bullet in bullet_points:
            p = doc.add_paragraph(bullet, style="List Bullet")
            for run in p.runs:
                run.font.name = FONT
                run.font.size = Pt(11)
                run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

        # 표 삽입: IDINO 스타일 헤더 적용
        if table_data:
            try:
                _add_table_to_docx(doc, table_data)
                # 마지막으로 추가된 표에 IDINO 스타일 헤더 색상을 적용한다
                if doc.tables:
                    last_table = doc.tables[-1]
                    # 첫 번째 행(헤더)의 각 셀에 진한 남색 배경을 적용한다
                    if len(last_table.rows) > 0:
                        for cell in last_table.rows[0].cells:
                            cell_shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{CLR_HEADER_HEX}"/>')
                            cell._tc.get_or_add_tcPr().append(cell_shading)
                            # 헤더 셀 텍스트를 흰색 Bold로 변경한다
                            for paragraph in cell.paragraphs:
                                for run in paragraph.runs:
                                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                                    run.font.bold = True
                                    run.font.name = FONT
            except Exception:
                logger.warning("Failed to add table to DOCX", exc_info=True)

        # 차트 삽입 (matplotlib 이미지로)
        if chart_data:
            try:
                _add_chart_to_docx(doc, chart_data)
            except Exception:
                logger.warning("Failed to add chart to DOCX", exc_info=True)

    output = io.BytesIO()
    doc.save(output)
    return output.getvalue()


def _generate_pptx(title: str, content: str) -> bytes:
    """AI 생성 텍스트를 PPTX로 변환한다."""
    from pptx import Presentation

    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title

    # Split content into sections by headings
    sections: list[tuple[str, str]] = []
    current_heading = ""
    current_body: list[str] = []

    for line in content.split("\n"):
        stripped = line.strip()
        if stripped.startswith("# ") or stripped.startswith("## "):
            if current_heading or current_body:
                sections.append((current_heading, "\n".join(current_body)))
            current_heading = stripped.lstrip("#").strip()
            current_body = []
        elif stripped:
            current_body.append(stripped)

    if current_heading or current_body:
        sections.append((current_heading, "\n".join(current_body)))

    # Create a slide per section
    for heading, body in sections[:20]:  # max 20 slides
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = heading or title
        if slide.placeholders[1]:
            slide.placeholders[1].text = body[:1000]

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


# ---------------------------------------------------------------------------
# Jinja2 렌더링 관련 함수들
# ---------------------------------------------------------------------------


def _generate_jinja2_context(
    variables_schema: dict,
    source_chunks: list[str],
    ai_prompt: str | None,
    api_key: str,
    agent_config: dict | None = None,
) -> dict:
    """Jinja2 변수 스키마를 기반으로 AI가 context dict를 생성한다.

    템플릿의 jinja2_variables 스키마(예: {"title": "string", "items": "list"})를
    분석하여 OpenAI Structured Outputs JSON Schema를 동적으로 만들고,
    소스 문서 청크를 참고하여 AI가 각 변수에 맞는 값을 JSON으로 생성한다.

    동작 순서:
      1. jinja2_variables 스키마의 각 변수 타입을 JSON Schema 속성으로 변환
      2. 소스 문서 텍스트를 참고 자료로 시스템 프롬프트에 포함
      3. _call_llm_structured()로 Structured Outputs API 호출
      4. 에이전트 설정이 있으면 system_prompt를 병합하여 톤/스타일 반영

    Args:
        variables_schema: 템플릿의 jinja2_variables 딕셔너리
            예: {"title": "string", "author": "string", "sections": "list",
                 "cover_image": "image"}
        source_chunks: 소스 문서에서 추출된 텍스트 청크 목록
        ai_prompt: 사용자가 추가로 입력한 생성 지시사항 (선택)
        api_key: LLM API 인증 키
        agent_config: 에이전트 설정 딕셔너리 (선택)
            {"system_prompt": "...", "llm_model": "...", "temperature": ...}

    Returns:
        Jinja2 렌더링에 사용할 context 딕셔너리
        예: {"title": "AI 동향 보고서", "author": "홍길동", "sections": [...]}
    """
    # ── Step 1: jinja2_variables 스키마를 JSON Schema properties로 변환한다 ──
    # variables_schema 구조: {"variables": [{"name": "...", "type": "..."}, ...]}
    # 각 변수의 타입을 OpenAI Structured Outputs JSON Schema 형식으로 매핑한다
    properties: dict[str, Any] = {}
    required_keys: list[str] = []

    # variables_schema에서 실제 변수 목록을 추출한다
    var_list = variables_schema.get("variables", [])
    if not var_list and isinstance(variables_schema, dict):
        # 이전 형식 호환: {"변수명": "타입"} 형태일 수도 있음
        var_list = [{"name": k, "type": v} for k, v in variables_schema.items() if k != "variables"]

    # ── 점(.) 표기 변수를 객체 배열로 변환 ──
    # 예: "task.업무명", "task.담당자" → 부모 "완료업무"를 객체 배열로 만든다
    # 1단계: 점 표기 변수에서 속성 이름을 수집한다
    dot_properties: dict[str, list[str]] = {}  # {"task": ["업무명", "담당자", "진행률"]}
    # 어떤 배열 변수가 어떤 점 변수의 부모인지 매핑
    array_vars = set()
    for var_info in var_list:
        if isinstance(var_info, dict):
            name = var_info.get("name", "")
            vtype = var_info.get("type", "string")
            if "." in name:
                # "task.업무명" → prefix="task", prop="업무명"
                prefix, prop = name.split(".", 1)
                if prefix not in dot_properties:
                    dot_properties[prefix] = []
                dot_properties[prefix].append(prop)
            if vtype == "array":
                array_vars.add(name)

    for var_info in var_list:
        # var_info가 dict이면 name/type 추출, 아니면 건너뜀
        if isinstance(var_info, dict):
            var_name = var_info.get("name", "")
            var_type = var_info.get("type", "string")
        else:
            continue
        if not var_name:
            continue
        # 점(.) 표기 변수는 부모 배열에 포함되므로 개별 처리하지 않는다
        if "." in var_name:
            continue
        # "image" 타입은 이미지 생성용이므로 AI에게는 설명 텍스트를 요청한다
        if var_type == "image":
            properties[var_name] = {
                "type": "string",
                "description": (
                    f"'{var_name}'에 들어갈 이미지를 설명하는 짧은 영어 프롬프트 "
                    "(예: 'modern office meeting with diverse team')"
                ),
            }
        elif var_type in ("list", "array") and var_name in array_vars:
            # 배열 변수: 점 표기 속성이 있으면 객체 배열로 변환
            # 예: "완료업무" + dot_properties["task"] → [{업무명, 담당자, 진행률}]
            # for 루프에서 사용하는 iterator 이름과 속성을 매칭
            item_props = {}
            for _prefix, props in dot_properties.items():
                for prop in props:
                    item_props[prop] = {"type": "string"}

            if item_props:
                properties[var_name] = {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": item_props,
                        "required": list(item_props.keys()),
                        "additionalProperties": False,
                    },
                    "description": f"'{var_name}'에 들어갈 항목 목록 (각 항목은 {list(item_props.keys())} 속성을 가짐)",
                }
            else:
                properties[var_name] = {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": f"'{var_name}'에 들어갈 항목 목록",
                }
        elif var_type in ("list", "array"):
            # 점 표기 없는 단순 배열
            properties[var_name] = {
                "type": "array",
                "items": {"type": "string"},
                "description": f"'{var_name}'에 들어갈 항목 목록",
            }
        elif var_type == "number":
            properties[var_name] = {
                "type": "number",
                "description": f"'{var_name}'에 들어갈 숫자 값",
            }
        elif var_type == "boolean":
            properties[var_name] = {
                "type": "boolean",
                "description": f"'{var_name}' 참/거짓 값",
            }
        elif isinstance(var_type, dict):
            # 중첩 객체 타입: 그대로 JSON Schema에 포함한다
            properties[var_name] = var_type
        else:
            # 기본값은 string 타입으로 처리한다
            properties[var_name] = {
                "type": "string",
                "description": f"'{var_name}'에 들어갈 텍스트 내용 (상세하게 작성)",
            }
        required_keys.append(var_name)

    # ── Step 2: Structured Outputs용 JSON Schema를 구성한다 ──
    json_schema = {
        "name": "jinja2_context",
        "strict": True,
        "schema": {
            "type": "object",
            "properties": properties,
            "required": required_keys,
            "additionalProperties": False,
        },
    }

    # ── Step 3: 시스템 프롬프트를 구성한다 ──
    # 소스 문서 청크를 참고 자료로 포함하여 AI가 문맥에 맞는 값을 생성하도록 한다
    source_text = "\n\n---\n\n".join(source_chunks[:50]) if source_chunks else "(참고 문서 없음)"

    system_msg = (
        "당신은 문서 템플릿의 변수에 적절한 값을 채워넣는 AI 어시스턴트입니다.\n"
        "아래 제공되는 참고 문서 내용을 기반으로 각 변수에 맞는 값을 JSON으로 생성하세요.\n\n"
        "중요 규칙:\n"
        "- 각 변수에 충분히 상세한 내용을 작성하세요 (최소 2~3문장 이상)\n"
        "- 문서에서 관련 정보를 찾아 구체적으로 기술하세요\n"
        "- '보고내용', '의견', '내용' 등의 변수는 문단 수준의 상세한 텍스트를 작성하세요\n"
        "- 배열(array) 타입은 구체적인 항목을 3~5개 이상 생성하세요\n"
        "- 한국어로 작성하되, image 타입의 변수는 영어 프롬프트로 작성하세요\n\n"
        f"참고 문서 내용:\n{source_text[:25000]}"
    )

    # 에이전트 설정이 있으면 에이전트의 system_prompt를 기존 메시지에 병합한다
    if agent_config and agent_config.get("system_prompt"):
        system_msg = f"{agent_config['system_prompt']}\n\n--- 추가 지시 ---\n{system_msg}"

    # ── Step 4: 사용자 메시지를 구성한다 ──
    user_msg = "위 참고 문서를 바탕으로 각 변수에 적합한 값을 JSON으로 생성해주세요."
    if ai_prompt:
        user_msg += f"\n\n추가 지시사항:\n{ai_prompt}"

    # ── Step 5: LLM 모델/온도/프로바이더 설정 (에이전트 설정 우선) ──
    llm_model = settings.llm_model
    llm_temperature = 0.3
    # 에이전트의 llm_provider가 있으면 우선 사용, 없으면 시스템 기본값
    agent_provider = None
    if agent_config:
        llm_model = agent_config.get("llm_model", llm_model)
        llm_temperature = agent_config.get("temperature", llm_temperature)
        agent_provider = agent_config.get("llm_provider")

    # ── Step 6: _call_llm_structured()로 Structured Outputs API를 호출한다 ──
    logger.info(
        "Jinja2 context 생성 호출: 변수 %d개, 모델=%s, 프로바이더=%s",
        len(properties),
        llm_model,
        agent_provider or "(시스템 기본값)",
    )
    context = _call_llm_structured(
        system_msg=system_msg,
        user_msg=user_msg,
        json_schema=json_schema,
        api_key=api_key,
        model=llm_model,
        temperature=llm_temperature,
        provider_override=agent_provider,
    )
    logger.info("Jinja2 context 생성 완료: 키=%s", list(context.keys()))
    return context


def _process_image_variables(
    context: dict,
    variables_schema: dict,
    image_config: dict | None,
) -> dict:
    """context에서 type이 'image'인 변수를 찾아 이미지를 생성/다운로드한다.

    _generate_jinja2_context()에서 image 타입 변수에는 이미지 설명 프롬프트가
    문자열로 들어 있다. 이 함수는 그 프롬프트를 사용하여 실제 이미지를
    DALL-E 3 또는 Unsplash에서 생성/검색한 뒤, context의 해당 값을
    이미지 bytes로 교체한다.

    Celery worker는 동기 환경이므로 _run_async()를 사용하여
    비동기 ImageGenerationService를 호출한다.

    Args:
        context: Jinja2 context 딕셔너리 (image 변수에는 프롬프트 문자열이 들어있음)
        variables_schema: 템플릿의 jinja2_variables 스키마
            예: {"cover_image": "image", "title": "string"}
        image_config: 이미지 생성 설정 (DocumentTemplate.image_generation_config)
            예: {"provider": "dalle3", "size": "1024x1024", "style": "natural"}
            None이면 기본 설정 사용

    Returns:
        이미지 bytes가 삽입된 context 딕셔너리
        (image 타입 변수의 값이 str → bytes로 교체됨)
    """
    # image 타입 변수가 없으면 바로 반환한다
    image_vars = [k for k, v in variables_schema.items() if v == "image"]
    if not image_vars:
        return context

    # 이미지 생성 서비스 인스턴스를 만든다
    img_service = ImageGenerationService()

    # 이미지 설정에서 provider, size, style을 추출한다
    config = image_config or {}
    provider = config.get("provider")  # "dalle3" 또는 "unsplash"
    size = config.get("size", "1024x1024")
    style = config.get("style", "natural")

    for var_name in image_vars:
        prompt = context.get(var_name, "")
        if not prompt or not isinstance(prompt, str):
            logger.warning("이미지 변수 '%s'에 프롬프트가 없음, 건너뜀", var_name)
            continue

        logger.info("이미지 생성 시작: 변수='%s', 프롬프트='%s'", var_name, prompt[:80])
        try:
            # 비동기 ImageGenerationService.generate()를 동기적으로 호출한다
            image_bytes = _run_async(
                img_service.generate(
                    prompt=prompt,
                    provider=provider,
                    size=size,
                    style=style,
                )
            )
            if image_bytes:
                # 프롬프트 문자열을 이미지 bytes로 교체한다
                context[var_name] = image_bytes
                logger.info(
                    "이미지 생성 완료: 변수='%s', 크기=%d bytes",
                    var_name,
                    len(image_bytes),
                )
            else:
                logger.warning("이미지 생성 결과가 None: 변수='%s'", var_name)
                # 이미지 실패 시 빈 문자열로 대체 (렌더링 오류 방지)
                context[var_name] = ""
        except Exception:
            logger.warning(
                "이미지 생성 실패: 변수='%s'",
                var_name,
                exc_info=True,
            )
            # 실패해도 보고서 생성은 계속 진행한다 — 빈 문자열로 대체
            context[var_name] = ""

    return context


async def _load_chat_transcript(chat_session_id: str, max_chars: int = 8000) -> str:
    """Chat 세션의 최근 메시지를 "역할: 내용" 문자열로 직렬화한다.

    H7 핫픽스 배경
    --------------
    회의록·요약 보고서를 Chat 기반으로 생성하려면 실제 대화 내용이 LLM
    컨텍스트로 전달되어야 한다. 이전 코드는 ``source_chat_session_id`` 를
    읽기만 하고 ``tb_chat_messages`` 를 조회하지 않아 컨텍스트가 비어 있었다.

    구현 방침
    --------
    1. 최신 메시지부터 ``chunk_index`` 역순으로 로드한다.
    2. 누적 글자수가 ``max_chars`` 를 넘지 않는 선에서 최근 대화를 포함한다.
    3. 반환 시점에는 시간 오름차순(=대화 순서)으로 다시 정렬한다.
    4. LLM 토큰 폭주를 막기 위해 기본 상한을 8000자(대략 8K 토큰 이내)로 둔다.

    Args:
        chat_session_id: Chat 세션 UUID 문자열.
        max_chars: 반환 문자열 상한. 기본 8000.

    Returns:
        "user: ...\nassistant: ..." 형식의 다중 줄 문자열. 메시지가 없으면 빈 문자열.
    """
    from sqlalchemy import select

    from app.modules.chat.models import ChatMessage

    async with _get_worker_session()() as session:
        # 최신 메시지부터 역순 로드하여 앞쪽 대화가 잘려도 최근 대화는 살린다.
        stmt = (
            select(ChatMessage)
            .where(ChatMessage.session_id == chat_session_id)
            .order_by(ChatMessage.ins_dt.desc())
            .limit(200)
        )
        result = await session.execute(stmt)
        messages = list(result.scalars().all())

    if not messages:
        return ""

    # 오래된 것부터 읽히도록 뒤집는다.
    messages.reverse()

    lines: list[str] = []
    total = 0
    # 토큰·글자수 초과를 방지하기 위해 뒤에서부터 포함 여부를 판단한다.
    # (최근 메시지 우선 보존)
    kept_rev: list[str] = []
    for msg in reversed(messages):
        role_label = {
            "user": "사용자",
            "assistant": "AI",
            "system": "시스템",
        }.get(str(msg.role), str(msg.role))
        line = f"{role_label}: {msg.content or ''}"
        if total + len(line) + 1 > max_chars and kept_rev:
            # 상한 초과: 그만 추가한다.
            break
        kept_rev.append(line)
        total += len(line) + 1

    # kept_rev 는 최신 → 과거 순이므로 다시 뒤집어 자연스러운 순서로 반환.
    lines = list(reversed(kept_rev))
    return "\n".join(lines)


def _minutes_to_docx_document(data: dict, fallback_title: str) -> dict:
    """회의록 JSON 을 DOCX 구조 스키마로 변환한다 (H6 핫픽스).

    ``MINUTES_STRUCTURED_SCHEMA`` 로 받은 JSON 을 ``DOCX_DOCUMENT_SCHEMA``
    의 {"document": {"title", "sections": [...]}} 형태로 매핑하여
    기존 ``_build_docx_from_structured`` 빌더를 재사용한다.

    섹션 구성:
      1. 회의 개요 (일시·장소·참석자)
      2. 안건
      3. 논의 내용
      4. 결정사항
      5. 액션아이템 (표)
      6. 차기 회의 (있는 경우)

    Args:
        data: LLM 이 반환한 {"minutes": {...}} 형태의 dict.
        fallback_title: ``meeting_title`` 이 비어있을 때 사용할 문서 제목.

    Returns:
        {"document": {"title": ..., "sections": [...]}} 형태의 dict.
    """
    minutes = data.get("minutes", {}) or {}
    title = minutes.get("meeting_title") or fallback_title or "회의록"
    meeting_date = minutes.get("meeting_date") or "일자 미상"
    location = minutes.get("location")
    attendees: list[str] = minutes.get("attendees") or []
    agenda: list[str] = minutes.get("agenda") or []
    discussion: list[str] = minutes.get("discussion_points") or []
    decisions: list[str] = minutes.get("decisions") or []
    action_items: list[dict] = minutes.get("action_items") or []
    next_meeting = minutes.get("next_meeting")

    sections: list[dict] = []

    # 1. 회의 개요
    overview_paragraphs = [
        f"일시: {meeting_date}",
        f"장소: {location}" if location else "장소: (미상)",
    ]
    sections.append(
        {
            "heading": "회의 개요",
            "heading_level": 1,
            "paragraphs": overview_paragraphs,
            "table": None,
            "chart": None,
            "bullet_points": [f"참석자: {a}" for a in attendees] if attendees else [],
        }
    )

    # 2. 안건
    sections.append(
        {
            "heading": "안건",
            "heading_level": 1,
            "paragraphs": [],
            "table": None,
            "chart": None,
            "bullet_points": agenda,
        }
    )

    # 3. 논의 내용
    sections.append(
        {
            "heading": "논의 내용",
            "heading_level": 1,
            "paragraphs": discussion,
            "table": None,
            "chart": None,
            "bullet_points": [],
        }
    )

    # 4. 결정사항
    sections.append(
        {
            "heading": "결정사항",
            "heading_level": 1,
            "paragraphs": [],
            "table": None,
            "chart": None,
            "bullet_points": decisions,
        }
    )

    # 5. 액션아이템 — 표로 정리
    if action_items:
        rows: list[list[str]] = []
        for item in action_items:
            if not isinstance(item, dict):
                continue
            assignee = str(item.get("assignee") or "미정")
            task = str(item.get("task") or "")
            due = str(item.get("due_date") or "") or "-"
            rows.append([assignee, task, due])

        action_table = {
            "headers": ["담당자", "업무", "마감일"],
            "rows": rows,
        }
        sections.append(
            {
                "heading": "액션아이템",
                "heading_level": 1,
                "paragraphs": [],
                "table": action_table,
                "chart": None,
                "bullet_points": [],
            }
        )
    else:
        sections.append(
            {
                "heading": "액션아이템",
                "heading_level": 1,
                "paragraphs": ["별도의 액션아이템은 도출되지 않았습니다."],
                "table": None,
                "chart": None,
                "bullet_points": [],
            }
        )

    # 6. 차기 회의
    if next_meeting:
        sections.append(
            {
                "heading": "차기 회의",
                "heading_level": 1,
                "paragraphs": [next_meeting],
                "table": None,
                "chart": None,
                "bullet_points": [],
            }
        )

    return {
        "document": {
            "title": title,
            "sections": sections,
        }
    }


async def _save_jinja2_context(report_id: str, context: dict, rendering_mode: str) -> None:
    """생성된 Jinja2 context를 DB에 저장한다 (디버깅 및 재생성 용도).

    GeneratedReport의 jinja2_context, rendering_mode 컬럼에
    실제 사용된 변수 값과 렌더링 방식을 기록한다.

    Args:
        report_id: 리포트 UUID 문자열
        context: Jinja2 context 딕셔너리 (이미지 bytes는 제외됨)
        rendering_mode: 렌더링 방식 ("jinja2", "structured", "regex")
    """
    from sqlalchemy import update

    from app.modules.reports.models import GeneratedReport

    # context에서 bytes 값(이미지)은 JSON 직렬화 불가하므로 제거한다
    serializable_ctx = {}
    for k, v in context.items():
        if isinstance(v, bytes):
            serializable_ctx[k] = f"[image: {len(v)} bytes]"
        else:
            serializable_ctx[k] = v

    async with _get_worker_session()() as session:
        await session.execute(
            update(GeneratedReport)
            .where(GeneratedReport.id == report_id)
            .values(
                jinja2_context=serializable_ctx,
                rendering_mode=rendering_mode,
            )
        )
        await session.commit()


# ---------------------------------------------------------------------------
# Celery task
# ---------------------------------------------------------------------------


@celery_app.task(
    bind=True,
    name="workers.report_generator.generate_report",
    max_retries=3,
    default_retry_delay=60,
    acks_late=True,
    track_started=True,
)
def generate_report(self, report_id: str) -> dict[str, Any]:
    """Generate a report from a template using RAG-extracted content.

    Steps:
      1. Load the report record and associated template.
      2. Download the template file from MinIO.
      3. Parse template variables ({{var}} patterns).
      4. For each variable, use RAG to extract content from source documents.
      5. Fill the template with generated content.
      6. Upload the completed report to MinIO.
      7. Update the report status.
    """
    logger.info("Starting report generation: %s", report_id)

    try:
        # Step 1: Load report record
        _run_async(_update_report_status(report_id, "processing"))

        report = _run_async(_get_report(report_id))
        if report is None:
            raise ValueError(f"Report not found: {report_id}")

        template_id = report.get("template_id")  # None if no template
        org_id = str(report.get("organization_id", ""))
        source_document_ids = report.get("source_document_ids", [])
        generation_params = report.get("generation_params") or {}
        if isinstance(source_document_ids, str):
            source_document_ids = json.loads(source_document_ids)

        # Determine output format from report record
        desired_format = report.get("output_format", "docx")
        output_ext = f".{desired_format}"

        # rendering_mode를 generation_params에서 가져온다 (기본값: "structured")
        # - "jinja2": document_template의 Jinja2 템플릿으로 렌더링
        # - "regex": 기존 {{var}} 패턴 치환 (report_template 기반)
        # - "structured": AI Structured Outputs → 자유형 문서 생성
        rendering_mode = generation_params.get("rendering_mode", "structured")

        # document_template_id가 있으면 document_templates 테이블에서 로드한다
        doc_template_id_param = generation_params.get("document_template_id")
        document_template = None
        if doc_template_id_param:
            try:
                document_template = _get_document_template_sync(str(doc_template_id_param))
                if document_template:
                    logger.info(
                        "Document template loaded: %s (mode=%s)",
                        document_template["name"],
                        document_template["rendering_mode"],
                    )
                    # document_template 자체에 rendering_mode가 지정되어 있으면 우선 사용한다
                    if document_template.get("rendering_mode"):
                        rendering_mode = document_template["rendering_mode"]
            except Exception:
                logger.warning(
                    "Failed to load document template %s",
                    doc_template_id_param,
                    exc_info=True,
                )

        # ── Jinja2 렌더링 분기 ──
        # rendering_mode가 "jinja2"이고 document_template이 있으면
        # Jinja2 엔진으로 DOCX/PPTX 템플릿을 렌더링한다
        if rendering_mode == "jinja2" and document_template:
            logger.info("Jinja2 렌더링 경로 진입: template=%s", document_template["name"])

            from app.core.llm_keys import resolve_api_key as _resolve_key
            from app.integrations.object_storage.minio_client import MinIOService

            minio = MinIOService(
                endpoint=settings.minio_endpoint,
                access_key=settings.minio_access_key,
                secret_key=settings.minio_secret_key,
                secure=settings.minio_secure,
            )

            # 1. document_template의 template_storage_path에서 템플릿 파일을 다운로드한다
            tmpl_storage_path = document_template.get("template_storage_path", "")
            if not tmpl_storage_path:
                raise ValueError(
                    f"Document template '{document_template['name']}'에 template_storage_path가 설정되지 않았습니다."
                )
            logger.info("Jinja2 템플릿 파일 다운로드: %s", tmpl_storage_path)
            jinja2_template_data = minio.download_file(
                bucket=settings.minio_bucket,
                object_name=tmpl_storage_path,
            )
            jinja2_template_format = Path(tmpl_storage_path).suffix.lower()

            # 2. jinja2_variables 스키마를 로드한다
            jinja2_variables_schema = document_template.get("jinja2_variables") or {}
            if not jinja2_variables_schema:
                logger.warning("jinja2_variables 스키마가 비어있음, 빈 context로 렌더링 시도")

            # 3. 소스 문서 청크를 로드한다 (AI context 생성에 사용)
            source_chunks_text = ""
            source_chunks_list: list[str] = []
            if source_document_ids:
                source_chunks_text = _run_async(_load_source_chunks(source_document_ids))
                # 각 청크를 개별 항목으로 분리한다
                source_chunks_list = [chunk.strip() for chunk in source_chunks_text.split("\n\n") if chunk.strip()]
            logger.info(
                "소스 청크 로드 완료: %d개 청크, %d자",
                len(source_chunks_list),
                len(source_chunks_text),
            )

            # 5. 에이전트 설정을 로드한다 (있으면) — API 키 해석보다 먼저 로드
            agent_id = generation_params.get("agent_id")
            agent_config = None
            if agent_id:
                try:
                    agent_config = _get_agent_sync(str(agent_id))
                    if agent_config:
                        logger.info(
                            "Jinja2 에이전트 로드: %s (provider=%s)",
                            agent_config["name"],
                            agent_config.get("llm_provider") or "(시스템 기본값)",
                        )
                except Exception:
                    logger.warning("에이전트 로드 실패, 기본 설정 사용", exc_info=True)

            # 4. LLM API 키를 조회한다 (DB 저장 키 > 환경변수 fallback)
            # 에이전트의 llm_provider가 있으면 해당 프로바이더의 키를 사용한다
            _jinja2_provider = (agent_config.get("llm_provider") if agent_config else None) or settings.llm_provider
            api_key = _run_async(_resolve_key(None, None, _jinja2_provider))
            if org_id:
                try:

                    async def _get_key_jinja2():
                        async with _get_worker_session()() as _sess:
                            return await _resolve_key(
                                _sess,
                                __import__("uuid").UUID(org_id),
                                _jinja2_provider,
                            )

                    api_key = _run_async(_get_key_jinja2())
                except Exception:
                    pass  # 실패 시 환경변수 키를 사용한다

            # ── 6a. session_auto 변수: 세션 사용자 정보로 자동 채움 ──
            # router에서 주입한 _session_user 정보를 사용하여
            # "소속", "작성자", "성명" 등의 변수를 자동으로 채운다.
            # 이 변수들은 AI에게 생성을 요청하지 않고 사용자 세션에서 바로 채운다.
            session_user = generation_params.get("_session_user", {})

            # session_auto 키워드 → 세션 정보 매핑 테이블
            # 변수 이름에 해당 키워드가 포함되면 매핑된 값을 자동으로 사용한다
            # "소속"/"부서"는 조직명 + 부서명을 결합 (예: "(주)아이디노 플랫폼사업팀")
            _org_name = session_user.get("organization_name", "")
            # H3: 라우터에서 주입한 새 키(``department_name``)와 기존 키(``department``)를
            # 모두 지원한다. full_name 도 성명 후보로 허용해 별도 프로필 설정이 있어도
            # 자연스럽게 매핑되도록 한다.
            _dept_name = session_user.get("department_name") or session_user.get("department", "")
            _affiliation = f"{_org_name} {_dept_name}".strip() if _org_name else _dept_name
            _user_display = session_user.get("full_name") or session_user.get("username") or ""

            SESSION_MAPPING = {
                "소속": _affiliation,
                "부서": _affiliation,
                "작성자": _user_display,
                "성명": _user_display,
                "이름": _user_display,
            }

            # session_auto 변수에 값을 자동 채운 context 딕셔너리
            session_auto_ctx: dict = {}
            # AI 생성 대상에서 제외할 변수 이름 집합
            # (session_auto + user_input 변수는 AI 스키마에서 제외)
            exclude_from_ai: set[str] = set()

            # jinja2_variables_schema에서 각 변수의 카테고리를 확인한다
            var_list_for_category = jinja2_variables_schema.get("variables", []) if jinja2_variables_schema else []
            for var_info in var_list_for_category:
                if not isinstance(var_info, dict):
                    continue
                vname = var_info.get("name", "")
                if not vname:
                    continue

                # 변수에 저장된 category를 사용하되, 없으면 이름 기반으로 자동 분류
                # (하위호환: category 필드가 없는 기존 데이터도 정상 동작)
                category = var_info.get("category") or classify_variable_category(vname)

                if category == "session_auto":
                    # session_auto 변수: 키워드 매핑으로 값을 채운다
                    for keyword, value in SESSION_MAPPING.items():
                        if keyword in vname and value:
                            session_auto_ctx[vname] = value
                            break
                    # 매핑되지 않아도 AI 대상에서는 제외한다
                    exclude_from_ai.add(vname)
                elif category == "user_input":
                    # user_input 변수: 사용자가 custom_context로 입력한 값을 사용
                    # AI 생성 대상에서 제외한다
                    exclude_from_ai.add(vname)

            if session_auto_ctx:
                logger.info(
                    "session_auto 변수 자동 채움 (%d개): %s",
                    len(session_auto_ctx),
                    list(session_auto_ctx.keys()),
                )
            if exclude_from_ai:
                logger.info(
                    "AI 생성 대상 제외 변수 (%d개): %s",
                    len(exclude_from_ai),
                    list(exclude_from_ai),
                )

            # ── 6b. AI 생성 대상 변수만 포함하는 스키마를 만든다 ──
            # session_auto와 user_input 변수는 AI에게 요청하지 않는다.
            # 이렇게 하면 AI 토큰을 절약하고, AI가 불필요한 내용을 생성하지 않는다.
            ai_only_schema: dict | None = None
            if jinja2_variables_schema and exclude_from_ai:
                # 원본 스키마를 복사하고 제외 변수를 필터링한다
                filtered_vars = [
                    v for v in var_list_for_category if isinstance(v, dict) and v.get("name", "") not in exclude_from_ai
                ]
                ai_only_schema = {"variables": filtered_vars} if filtered_vars else None
            else:
                # 제외할 변수가 없으면 원본 스키마를 그대로 사용한다
                ai_only_schema = jinja2_variables_schema

            # ── 6c. _generate_jinja2_context()로 AI가 ai_generated 변수 값을 생성한다 ──
            # generation_params에서 컨텐츠 옵션 읽기 (표/차트/이미지 포함 여부)
            include_charts = generation_params.get("include_charts", False)
            include_tables = generation_params.get("include_tables", True)
            include_images = generation_params.get("include_images", False)

            ai_prompt = generation_params.get("ai_prompt", "")

            # AI 프롬프트에 컨텐츠 옵션 지시사항 추가
            content_instructions = []
            if include_tables:
                content_instructions.append("표(table) 데이터가 적절한 곳에 표 형태로 정리해주세요")
            if include_charts:
                content_instructions.append("통계/수치 데이터가 있으면 차트 설명을 포함해주세요")
            if include_images:
                content_instructions.append("이미지가 필요한 곳에 이미지 설명 프롬프트를 포함해주세요")

            if content_instructions:
                ai_prompt = (ai_prompt or "") + "\n\n" + "\n".join(content_instructions)
            jinja2_ctx: dict = {}
            if ai_only_schema:
                jinja2_ctx = _generate_jinja2_context(
                    variables_schema=ai_only_schema,
                    source_chunks=source_chunks_list,
                    ai_prompt=ai_prompt,
                    api_key=api_key,
                    agent_config=agent_config,
                )

            # ── 7. 최종 context 조립: session_auto + ai_generated + custom_context(user_input) ──
            # 우선순위: custom_context > session_auto > ai_generated
            # 즉, 사용자가 직접 입력한 값이 가장 높은 우선순위를 갖는다.

            # 먼저 session_auto 값을 AI 생성 context에 병합한다
            if session_auto_ctx:
                jinja2_ctx.update(session_auto_ctx)

            # custom_context(user_input 등 사용자 입력값)를 최종 병합한다
            custom_context = generation_params.get("custom_context")
            if custom_context and isinstance(custom_context, dict):
                # 빈 값이 아닌 필드만 필터링하여 병합
                non_empty = {
                    k: v for k, v in custom_context.items() if v is not None and v != "" and v != [] and v != {}
                }
                if non_empty:
                    logger.info("custom_context 병합 (%d개): %s", len(non_empty), list(non_empty.keys()))
                    jinja2_ctx.update(non_empty)

            # 8. _process_image_variables()로 image 타입 변수를 실제 이미지로 교체한다
            # include_images가 False이면 이미지 처리를 건너뛴다
            if include_images:
                image_config = document_template.get("image_generation_config")
                jinja2_ctx = _process_image_variables(
                    context=jinja2_ctx,
                    variables_schema=jinja2_variables_schema,
                    image_config=image_config,
                )

            # 9. Jinja2 context를 DB에 저장한다 (디버깅/재생성 용도)
            try:
                _run_async(_save_jinja2_context(report_id, jinja2_ctx, "jinja2"))
            except Exception:
                logger.warning("Jinja2 context DB 저장 실패 (보고서 생성은 계속)", exc_info=True)

            # 10. 템플릿 형식에 따라 적절한 Jinja2 렌더러를 호출한다
            logger.info(
                "Jinja2 렌더링 시작: 형식=%s, context 키=%s",
                jinja2_template_format,
                list(jinja2_ctx.keys()),
            )
            if jinja2_template_format == ".docx":
                filled_data = render_docx_jinja2(jinja2_template_data, jinja2_ctx)
                output_ext = ".docx"
            elif jinja2_template_format == ".pptx":
                filled_data = render_pptx_jinja2(jinja2_template_data, jinja2_ctx)
                output_ext = ".pptx"
            else:
                raise ValueError(f"Jinja2 렌더링은 DOCX/PPTX만 지원합니다 (현재: {jinja2_template_format})")
            logger.info("Jinja2 렌더링 완료: %d bytes", len(filled_data))

        elif template_id:
            # Load template record from DB
            template = _run_async(_get_report_template(template_id))
            if template is None:
                raise ValueError(f"Report template not found: {template_id}")

            template_file_path = template.get("file_path", "")
            template_format = Path(template_file_path).suffix.lower()

            # Step 2: Download template from MinIO
            logger.info("Downloading template: %s", template_file_path)
            from app.integrations.object_storage.minio_client import MinIOService

            minio = MinIOService(
                endpoint=settings.minio_endpoint,
                access_key=settings.minio_access_key,
                secret_key=settings.minio_secret_key,
                secure=settings.minio_secure,
            )
            template_data = minio.download_file(
                bucket=settings.minio_bucket,
                object_name=template_file_path,
            )

            # Step 3: Parse template variables
            logger.info("Extracting template variables")
            if template_format == ".docx":
                variable_names = _extract_docx_variables(template_data)
            elif template_format == ".pptx":
                variable_names = _extract_pptx_variables(template_data)
            else:
                variable_names = _extract_template_variables(template_data)

            logger.info("Found %d template variables: %s", len(variable_names), variable_names)

            # Step 4: Use RAG to extract content for each variable
            variables: dict[str, str] = {}
            for i, var_name in enumerate(variable_names):
                logger.info(
                    "Generating content for variable '%s' (%d/%d)",
                    var_name,
                    i + 1,
                    len(variable_names),
                )
                content = _rag_extract_content(
                    variable_name=var_name,
                    source_document_ids=source_document_ids,
                    org_id=org_id,
                )
                variables[var_name] = content

            # Step 5: Fill the template
            logger.info("Filling template with generated content")
            if template_format == ".docx":
                filled_data = _fill_docx_template(template_data, variables)
            elif template_format == ".pptx":
                filled_data = _fill_pptx_template(template_data, variables)
            elif template_format in (".html", ".htm"):
                filled_data = _fill_html_template(template_data.decode("utf-8"), variables)
                output_ext = ".pdf"
            else:
                text_content = template_data.decode("utf-8")
                for var_name, value in variables.items():
                    text_content = text_content.replace("{{" + var_name + "}}", value)
                filled_data = text_content.encode("utf-8")

        else:
            # --- Free-form AI generation (no DB template) ---
            # 템플릿 없이 AI가 직접 콘텐츠를 생성하는 경로이다.
            # Structured Outputs를 사용하면 표/차트/이미지가 포함된 구조화된
            # 문서를 생성할 수 있고, 실패 시 기존 plain text 경로로 fallback한다.
            logger.info("Free-form AI generation (no template)")
            from app.integrations.object_storage.minio_client import MinIOService

            minio = MinIOService(
                endpoint=settings.minio_endpoint,
                access_key=settings.minio_access_key,
                secret_key=settings.minio_secret_key,
                secure=settings.minio_secure,
            )

            # 소스 문서의 청크 텍스트를 로드한다 (RAG 컨텍스트)
            doc_context = ""
            if source_document_ids:
                doc_context = _run_async(_load_source_chunks(source_document_ids))

            # H7: 지정된 Chat 세션이 있으면 대화 기록을 로드한다.
            # (회의록·요약 생성 시 LLM 컨텍스트와 plain text fallback 양쪽에서 공용으로 사용)
            chat_transcript = ""
            source_chat_session_id_local = report.get("source_chat_session_id")
            if source_chat_session_id_local:
                try:
                    chat_transcript = _run_async(
                        _load_chat_transcript(str(source_chat_session_id_local), max_chars=8000)
                    )
                    logger.info(
                        "Chat 메시지 로드 완료: session=%s, %d자",
                        source_chat_session_id_local,
                        len(chat_transcript),
                    )
                except Exception:
                    logger.warning(
                        "Chat transcript 로드 실패 (session_id=%s)",
                        source_chat_session_id_local,
                        exc_info=True,
                    )

            # 기본 프롬프트 파라미터를 추출한다
            ai_prompt = generation_params.get("ai_prompt", "")
            template_name = generation_params.get("template_name", "보고서")
            document_type = generation_params.get("document_type", "report")
            title = report.get("title", "문서")

            type_label = "보고서" if document_type == "report" else "제안서"
            system_msg = (
                f"당신은 전문 {type_label} 작성 AI입니다. "
                f"주어진 참고 문서 내용을 바탕으로 '{title}' 제목의 {type_label}를 작성하세요. "
                f"템플릿 유형: {template_name}. "
                "한국어로 작성하고, 구조적이며 전문적인 문서 형식을 따르세요. "
                "마크다운 형식(# 제목, ## 소제목, - 목록)을 사용하여 구조화하세요."
            )
            if ai_prompt:
                system_msg += (
                    f"\n\n**반드시 준수해야 할 지시사항:**\n{ai_prompt}\n"
                    "위 지시사항의 페이지 수, 분량, 형식 등 구체적인 요구사항을 최우선으로 따르세요."
                )

            # ── Agent 로드 ──
            # generation_params에 agent_id가 있으면 해당 에이전트의 설정을
            # 시스템 프롬프트, LLM 모델, temperature, max_tokens에 반영한다.
            agent_id = generation_params.get("agent_id")
            agent_config = None
            llm_model = settings.llm_model  # 기본 LLM 모델
            llm_temperature = 0.3  # 기본 온도
            llm_max_tokens_override = None  # Agent가 지정한 max_tokens

            if agent_id:
                try:
                    agent_config = _get_agent_sync(str(agent_id))
                    if agent_config:
                        logger.info(
                            "Agent loaded: %s (model=%s, temp=%s)",
                            agent_config["name"],
                            agent_config["llm_model"],
                            agent_config["temperature"],
                        )
                        # 에이전트의 system_prompt를 기존 시스템 메시지에 병합한다
                        agent_system_prompt = agent_config.get("system_prompt", "")
                        if agent_system_prompt:
                            system_msg = f"{agent_system_prompt}\n\n--- 추가 지시 ---\n{system_msg}"
                        # 에이전트의 LLM 설정을 적용한다
                        llm_model = agent_config.get("llm_model", llm_model)
                        llm_temperature = agent_config.get("temperature", llm_temperature)
                        llm_max_tokens_override = agent_config.get("max_tokens")
                except Exception:
                    logger.warning(
                        "Failed to load agent %s, using defaults",
                        agent_id,
                        exc_info=True,
                    )

            # 페이지 수 요청 감지 → 동적 max_tokens 조절
            import re as _re

            page_match = _re.search(r"(\d+)\s*(?:페이지|page|쪽)", ai_prompt or "", _re.IGNORECASE)
            if page_match:
                requested_pages = int(page_match.group(1))
                # 한국어 기준 1페이지 ≈ 1500자 ≈ 800토큰
                report_max_tokens = min(max(requested_pages * 800, 4096), 16384)
            elif llm_max_tokens_override:
                # 에이전트가 지정한 max_tokens를 사용한다
                report_max_tokens = llm_max_tokens_override
            else:
                report_max_tokens = max(settings.llm_max_tokens, 8192)

            # H7: chat_transcript 가 있으면 plain text fallback 경로에도 포함한다.
            user_msg_parts = [
                "참고 문서 내용:",
                "",
                doc_context[:30000] if doc_context else "(참고 문서 없음)",
            ]
            if chat_transcript:
                user_msg_parts.extend(["", "다음은 관련된 대화 기록입니다:", chat_transcript])
            user_msg_parts.extend(["", f"위 내용을 바탕으로 {type_label}를 작성해주세요."])
            user_msg = "\n".join(user_msg_parts)

            # LLM API 키 조회 (DB 저장 키 > 환경변수 fallback)
            # 에이전트의 llm_provider가 있으면 해당 프로바이더의 키를 사용한다

            from app.core.llm_keys import resolve_api_key as _resolve_key

            _report_provider = (agent_config.get("llm_provider") if agent_config else None) or settings.llm_provider
            api_key = _run_async(_resolve_key(None, None, _report_provider))
            # Celery worker에서는 DB 세션이 없으므로 org_id로 조회 시도
            if org_id:
                try:

                    async def _get_key():
                        async with _get_worker_session()() as _sess:
                            return await _resolve_key(_sess, __import__("uuid").UUID(org_id), _report_provider)

                    api_key = _run_async(_get_key())
                except Exception:
                    pass  # 실패 시 환경변수 키를 사용한다

            # ── Structured Outputs 경로 (pptx/docx) ──
            # output_format이 pptx 또는 docx이면 Structured Outputs로 구조화된
            # JSON을 먼저 생성하고, 해당 JSON으로 문서를 빌드한다.
            # 실패하면 기존 plain text 경로로 fallback한다.
            structured_success = False

            if desired_format in ("pptx", "docx"):
                try:
                    from app.integrations.llm.prompts import (
                        STRUCTURED_DOCX_SYSTEM_PROMPT,
                        STRUCTURED_MINUTES_SYSTEM_PROMPT,
                        STRUCTURED_PPTX_SYSTEM_PROMPT,
                    )
                    from app.workers.structured_schemas import (
                        DOCX_DOCUMENT_SCHEMA,
                        MINUTES_STRUCTURED_SCHEMA,
                        PPTX_SLIDE_SCHEMA,
                    )

                    # H6: 회의록 분기
                    # document_type == "minutes" 이면 회의록 전용 스키마/프롬프트를 사용한다.
                    # 회의록은 보통 DOCX 로 출력되므로 PPTX 요청이 들어와도 여기서는
                    # DOCX 빌더로 전환한다(사용자가 명시적으로 pptx를 원하면 회의록이 아니라
                    # 발표 자료로 간주하므로 기존 경로를 유지).
                    is_minutes = document_type == "minutes" and desired_format == "docx"

                    if is_minutes:
                        json_schema = MINUTES_STRUCTURED_SCHEMA
                        structured_system = STRUCTURED_MINUTES_SYSTEM_PROMPT
                    elif desired_format == "pptx":
                        json_schema = PPTX_SLIDE_SCHEMA
                        structured_system = STRUCTURED_PPTX_SYSTEM_PROMPT
                    else:
                        json_schema = DOCX_DOCUMENT_SCHEMA
                        structured_system = STRUCTURED_DOCX_SYSTEM_PROMPT

                    # 에이전트/사용자 지시사항을 Structured 프롬프트에 병합한다
                    full_structured_system = structured_system
                    if agent_config and agent_config.get("system_prompt"):
                        full_structured_system = (
                            f"{agent_config['system_prompt']}\n\n--- 문서 구조 생성 규칙 ---\n{structured_system}"
                        )
                    if ai_prompt:
                        full_structured_system += f"\n\n--- 사용자 추가 지시사항 ---\n{ai_prompt}"

                    # Structured Outputs API를 호출하여 JSON 구조를 생성한다.
                    # H7: chat_transcript 가 있으면 별도 섹션으로 붙여 준다.
                    structured_user_parts: list[str] = [
                        f"제목: {title}",
                        f"유형: {type_label}",
                        "",
                        "참고 문서 내용:",
                        doc_context[:25000] if doc_context else "(참고 문서 없음)",
                    ]
                    if chat_transcript:
                        structured_user_parts.extend(
                            [
                                "",
                                "다음은 관련된 대화 기록입니다 (회의록·요약 시 우선 참조):",
                                chat_transcript,
                            ]
                        )
                    structured_user_parts.extend(
                        [
                            "",
                            (
                                "위 내용을 바탕으로 회의록 JSON을 생성해주세요."
                                if is_minutes
                                else f"위 내용을 바탕으로 {type_label} 구조를 JSON으로 생성해주세요."
                            ),
                        ]
                    )
                    structured_user_msg = "\n".join(structured_user_parts)

                    # 에이전트 프로바이더 우선, 없으면 시스템 기본값
                    _structured_provider = (agent_config.get("llm_provider") if agent_config else None) or None
                    logger.info(
                        "Calling Structured Outputs (format=%s, model=%s, provider=%s)",
                        desired_format,
                        llm_model,
                        _structured_provider or "(시스템 기본값)",
                    )
                    structured_data = _call_llm_structured(
                        system_msg=full_structured_system,
                        user_msg=structured_user_msg,
                        json_schema=json_schema,
                        api_key=api_key,
                        model=llm_model,
                        temperature=llm_temperature,
                        max_tokens=report_max_tokens,
                        provider_override=_structured_provider,
                    )
                    logger.info("Structured Outputs response received successfully")

                    # 슬라이드마스터 로드 (PPTX인 경우)
                    # 1순위: document_template(tb_document_templates)에 업로드된 PPTX 파일
                    # 2순위: report_template(tb_report_templates)에서 이름으로 검색
                    slide_master_data = None
                    if desired_format == "pptx":
                        # 1순위: document_template에서 직접 로드
                        if document_template and document_template.get("template_storage_path"):
                            try:
                                slide_master_data = minio.download_file(
                                    bucket=settings.minio_bucket,
                                    object_name=document_template["template_storage_path"],
                                )
                                logger.info(
                                    "document_template에서 슬라이드마스터 로드: %s",
                                    document_template["template_storage_path"],
                                )
                            except Exception:
                                logger.warning(
                                    "document_template 슬라이드마스터 로드 실패",
                                    exc_info=True,
                                )

                        # 2순위: report_template에서 이름으로 검색 (fallback)
                        if slide_master_data is None:
                            doc_template_id = generation_params.get("document_template_id")
                            template_name_hint = generation_params.get("template_name", "")
                            try:
                                rt_name = None
                                if doc_template_id:
                                    rt_name = _run_async(_get_document_template_name(doc_template_id))
                                search_name = rt_name or template_name_hint
                                if search_name:
                                    rt = _run_async(_find_report_template_by_name(search_name, org_id))
                                    if rt and rt.get("file_path"):
                                        slide_master_data = minio.download_file(
                                            bucket=settings.minio_bucket,
                                            object_name=rt["file_path"],
                                        )
                                        logger.info(
                                            "report_template에서 슬라이드마스터 로드: %s",
                                            rt["file_path"],
                                        )
                            except Exception:
                                logger.warning(
                                    "report_template 슬라이드마스터 로드 실패",
                                    exc_info=True,
                                )

                    # Structured JSON으로 문서를 빌드한다
                    if is_minutes:
                        # 회의록 JSON 을 DOCX 구조 스키마로 변환해 기존 빌더를 재사용한다.
                        docx_payload = _minutes_to_docx_document(structured_data, title)
                        filled_data = _build_docx_from_structured(docx_payload)
                    elif desired_format == "pptx":
                        filled_data = _build_pptx_from_structured(
                            structured_data,
                            slide_master_data,
                            generation_params,
                        )
                    else:
                        filled_data = _build_docx_from_structured(structured_data)

                    structured_success = True
                    logger.info("Structured Output build completed successfully")

                except Exception as structured_exc:
                    # Structured Outputs 실패 시 기존 plain text 경로로 fallback
                    logger.warning(
                        "Structured Outputs failed, falling back to plain text: %s",
                        structured_exc,
                        exc_info=True,
                    )

            # ── Plain text fallback 경로 ──
            # Structured Outputs가 비활성화되었거나 실패한 경우 기존 방식을 사용한다
            if not structured_success:
                try:
                    from app.integrations.llm.factory import create_llm_client, get_provider_for_task

                    # 에이전트 프로바이더 우선, 없으면 시스템 기본값
                    provider = (agent_config.get("llm_provider") if agent_config else None) or get_provider_for_task(
                        "report"
                    )
                    llm_client = create_llm_client(provider, api_key=api_key, model=llm_model)

                    # LLMClient의 동기 텍스트 생성 메서드를 호출한다
                    messages = [
                        {"role": "system", "content": system_msg},
                        {"role": "user", "content": user_msg},
                    ]
                    generated_text = llm_client.generate_sync(
                        messages=messages,
                        temperature=llm_temperature,
                        max_tokens=report_max_tokens,
                    )
                except Exception as exc:
                    logger.exception("LLM generation failed: %s", exc)
                    generated_text = f"[AI 생성 실패: {exc}]"

                # 출력 포맷에 맞게 문서를 변환한다
                # document_template_id가 있으면 해당 파일 기반 슬라이드마스터 사용
                doc_template_id = generation_params.get("document_template_id")
                template_name_hint = generation_params.get("template_name", "")
                slide_master_data = None
                logger.info(
                    "doc_template_id=%s, template_name=%s, format=%s",
                    doc_template_id,
                    template_name_hint,
                    desired_format,
                )

                if desired_format == "pptx":
                    try:
                        # 방법 1: document_template_id로 이름 조회 → report_template 파일 검색
                        rt_name = None
                        if doc_template_id:
                            rt_name = _run_async(_get_document_template_name(doc_template_id))
                            logger.info("Document template name from DB: %s", rt_name)

                        # 방법 2: template_name으로 직접 검색 (fallback)
                        search_name = rt_name or template_name_hint
                        if search_name:
                            rt = _run_async(_find_report_template_by_name(search_name, org_id))
                            logger.info("Report template match: %s", rt)
                            if rt and rt.get("file_path"):
                                slide_master_data = minio.download_file(
                                    bucket=settings.minio_bucket,
                                    object_name=rt["file_path"],
                                )
                                logger.info(
                                    "Loaded slide master: %s (%d bytes)", rt["file_path"], len(slide_master_data)
                                )
                    except Exception:
                        logger.warning("Failed to load slide master, using default", exc_info=True)

                if desired_format == "docx":
                    filled_data = _generate_docx(title, generated_text)
                elif desired_format == "pptx" and slide_master_data:
                    filled_data = _build_from_slide_master(slide_master_data, title, generated_text, generation_params)
                elif desired_format == "pptx":
                    filled_data = _generate_pptx(title, generated_text)
                else:
                    # HTML/PDF/TXT fallback
                    filled_data = generated_text.encode("utf-8")

        # Step 6: Upload completed report to MinIO
        report_filename = f"reports/{org_id}/{report_id}{output_ext}"
        content_type_map = {
            ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ".pdf": "application/pdf",
            ".html": "text/html",
            ".txt": "text/plain",
        }
        content_type = content_type_map.get(output_ext, "application/octet-stream")

        logger.info("Uploading report to: %s", report_filename)
        minio.upload_file(
            bucket=settings.minio_bucket,
            object_name=report_filename,
            file_data=filled_data,
            content_type=content_type,
        )

        # Step 7: Update report status
        _run_async(
            _update_report_status(
                report_id,
                status="completed",
                output_path=report_filename,
            )
        )

        logger.info("Report generation complete: %s", report_id)
        return {
            "report_id": report_id,
            "status": "completed",
            "output_path": report_filename,
        }

    except Exception as exc:
        logger.exception("Report generation failed for %s: %s", report_id, exc)
        _run_async(
            _update_report_status(
                report_id,
                status="error",
                error_message=str(exc)[:500],
            )
        )

        try:
            raise self.retry(exc=exc)
        except self.MaxRetriesExceededError:
            logger.error("Max retries exceeded for report generation: %s", report_id)
            return {
                "report_id": report_id,
                "status": "error",
                "error": str(exc)[:500],
            }
