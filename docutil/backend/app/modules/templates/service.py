"""Business logic for document template management.

기본 CRUD 외에 Jinja2 템플릿 파일 업로드, 변환, 변수 관리, 파일 다운로드 기능을 제공한다.
"""

from __future__ import annotations

import io
import json
import logging
import uuid as uuid_mod
from pathlib import PurePosixPath
from typing import TYPE_CHECKING
from uuid import UUID

from fastapi import HTTPException, status
from sqlalchemy import func, select

from app.core.config import get_settings
from app.integrations.llm.factory import create_llm_client, get_provider_for_task
from app.integrations.object_storage.minio_client import MinIOService
from app.modules.documents.models import DocumentChunk
from app.workers.jinja2_engine import (
    analyze_blank_form,
    analyze_blank_form_with_ai,
    auto_generate_jinja2_from_structure,
    convert_example_to_template,
    extract_docx_structure_for_editor,
    extract_docx_variables,
    extract_pptx_variables,
)

from .models import DocumentTemplate

if TYPE_CHECKING:
    from sqlalchemy.ext.asyncio import AsyncSession

    from .schemas import TemplateCreate, TemplateUpdate

logger = logging.getLogger(__name__)
settings = get_settings()

# 파일 확장자 → MIME 타입 매핑 (업로드 시 사용)
_CONTENT_TYPES: dict[str, str] = {
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "pdf": "application/pdf",
    "html": "text/html",
}


class TemplateService:
    """Stateless service methods for document template CRUD + Jinja2 확장."""

    # ------------------------------------------------------------------
    # Create
    # ------------------------------------------------------------------

    @staticmethod
    async def create_template(
        db: AsyncSession,
        org_id: UUID,
        user_id: UUID,
        data: TemplateCreate,
    ) -> DocumentTemplate:
        """새 문서 템플릿을 생성한다."""
        template = DocumentTemplate(
            organization_id=org_id,
            name=data.name,
            description=data.description,
            template_type=data.template_type,
            tone=data.tone,
            output_format=data.output_format,
            schema_=data.schema_,
            sample_prompt=data.sample_prompt,
            rendering_mode=data.rendering_mode,
            is_active=True,
            created_by=user_id,
        )
        db.add(template)
        await db.flush()
        await db.refresh(template)
        return template

    # ------------------------------------------------------------------
    # List
    # ------------------------------------------------------------------

    @staticmethod
    async def get_templates(
        db: AsyncSession,
        org_id: UUID,
        template_type: str | None = None,
        page: int = 1,
        size: int = 20,
    ) -> tuple[list[DocumentTemplate], int]:
        """조직별 템플릿 목록을 페이지네이션하여 반환한다.

        template_type이 주어지면 해당 유형만 필터링한다.
        """
        base_query = select(DocumentTemplate).where(
            DocumentTemplate.organization_id == org_id,
        )

        if template_type is not None:
            base_query = base_query.where(
                DocumentTemplate.template_type == template_type,
            )

        # 전체 건수 조회
        count_query = select(func.count()).select_from(base_query.subquery())
        total = (await db.execute(count_query)).scalar_one()

        # 페이지네이션 적용
        offset = (page - 1) * size
        items_query = base_query.order_by(DocumentTemplate.ins_dt.desc()).offset(offset).limit(size)
        result = await db.execute(items_query)
        items = list(result.scalars().all())

        return items, total

    # ------------------------------------------------------------------
    # Get single
    # ------------------------------------------------------------------

    @staticmethod
    async def get_template(
        db: AsyncSession,
        template_id: UUID,
        org_id: UUID,
    ) -> DocumentTemplate:
        """ID로 단일 템플릿을 조회한다. 없으면 404 예외를 발생시킨다."""
        stmt = select(DocumentTemplate).where(
            DocumentTemplate.id == template_id,
            DocumentTemplate.organization_id == org_id,
        )
        result = await db.execute(stmt)
        template = result.scalar_one_or_none()

        if template is None:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail=f"Document template '{template_id}' not found.",
            )
        return template

    # ------------------------------------------------------------------
    # Update
    # ------------------------------------------------------------------

    @staticmethod
    async def update_template(
        db: AsyncSession,
        template_id: UUID,
        org_id: UUID,
        data: TemplateUpdate,
    ) -> DocumentTemplate:
        """기존 템플릿의 필드를 부분 업데이트한다."""
        template = await TemplateService.get_template(db, template_id, org_id)

        update_fields = data.model_dump(exclude_unset=True, by_alias=False)
        for field, value in update_fields.items():
            setattr(template, field, value)

        await db.flush()
        await db.refresh(template)
        return template

    # ------------------------------------------------------------------
    # Delete
    # ------------------------------------------------------------------

    @staticmethod
    async def delete_template(
        db: AsyncSession,
        template_id: UUID,
        org_id: UUID,
    ) -> None:
        """ID로 템플릿을 삭제한다. 존재하지 않으면 404 예외."""
        template = await TemplateService.get_template(db, template_id, org_id)
        await db.delete(template)
        await db.flush()

    # ==================================================================
    # Jinja2 템플릿 확장 메서드들
    # ==================================================================

    # ------------------------------------------------------------------
    # upload_template — 템플릿 파일 업로드 + 변수 자동 추출
    # ------------------------------------------------------------------

    @staticmethod
    async def upload_template(
        db: AsyncSession,
        org_id: UUID,
        user_id: UUID,
        file_bytes: bytes,
        filename: str,
        template_type: str,
        tone: str,
        output_format: str,
        name: str | None = None,
        description: str | None = None,
    ) -> DocumentTemplate:
        """Jinja2 템플릿 파일을 MinIO에 업로드하고 변수를 자동 추출한다.

        처리 순서:
          1. 새 template_id(UUID) 생성
          2. MinIO에 파일 업로드 (경로: templates/{org_id}/{template_id}/{filename})
          3. 파일 확장자(docx/pptx)에 따라 Jinja2 변수 추출
          4. DB에 DocumentTemplate 레코드 저장
          5. 생성된 템플릿 반환

        Args:
            db: 비동기 DB 세션
            org_id: 조직 ID
            user_id: 업로드한 사용자 ID
            file_bytes: 업로드된 파일의 바이트 데이터
            filename: 원본 파일 이름 (확장자 포함)
            template_type: 템플릿 유형 (자유 텍스트)
            tone: 문서 어조
            output_format: 출력 형식
            name: 템플릿 이름 (미지정 시 파일명 사용)
            description: 템플릿 설명

        Returns:
            생성된 DocumentTemplate ORM 객체
        """
        # 0. 중복 체크: 같은 이름 + 같은 org_id의 템플릿이 이미 있는지 확인한다
        effective_name = name or filename
        existing_stmt = select(DocumentTemplate).where(
            DocumentTemplate.organization_id == org_id,
            DocumentTemplate.name == effective_name,
        )
        existing_result = await db.execute(existing_stmt)
        existing_template = existing_result.scalar_one_or_none()

        # 1. 기존 템플릿이 있으면 해당 ID 재사용, 없으면 새 UUID 생성
        template_id = existing_template.id if existing_template else uuid_mod.uuid4()

        # 2. MinIO에 파일 업로드
        # 경로 형식: templates/{org_id}/{template_id}/{filename}
        object_name = f"templates/{org_id}/{template_id}/{filename}"
        # 파일 확장자로 MIME 타입 결정
        ext = PurePosixPath(filename).suffix.lower().lstrip(".")
        content_type = _CONTENT_TYPES.get(ext, "application/octet-stream")

        minio = MinIOService()
        minio.upload_file(
            bucket=settings.minio_bucket,
            object_name=object_name,
            file_data=file_bytes,
            content_type=content_type,
        )
        logger.info("템플릿 파일 MinIO 업로드 완료: %s", object_name)

        # 3. 파일 형식에 따라 Jinja2 변수 추출
        jinja2_vars: dict | None = None
        try:
            if ext == "docx":
                jinja2_vars = extract_docx_variables(file_bytes)
            elif ext == "pptx":
                jinja2_vars = extract_pptx_variables(file_bytes)
            else:
                logger.info("변수 추출 미지원 형식: %s — 건너뜀", ext)
        except ValueError as e:
            logger.warning("Jinja2 변수 추출 실패 (계속 진행): %s", e)

        # 변수가 추출되지 않았으면 슬라이드마스터 양식 등으로 판단 → structured 모드
        has_variables = bool(jinja2_vars and jinja2_vars.get("variables"))
        rendering_mode = "jinja2" if has_variables else "structured"
        if not has_variables:
            logger.info("변수 0개 → structured 모드로 자동 전환")

        # 4. 기존 템플릿이 있으면 업데이트, 없으면 새로 생성
        if existing_template:
            existing_template.template_type = template_type
            existing_template.tone = tone
            existing_template.output_format = output_format
            existing_template.template_storage_path = object_name
            existing_template.original_file_path = object_name
            existing_template.jinja2_variables = jinja2_vars
            existing_template.rendering_mode = rendering_mode
            existing_template.is_active = True
            if description is not None:
                existing_template.description = description
            await db.flush()
            await db.refresh(existing_template)
            logger.info("기존 템플릿 업데이트 완료: id=%s, name=%s", existing_template.id, existing_template.name)
            return existing_template

        template = DocumentTemplate(
            id=template_id,
            organization_id=org_id,
            name=effective_name,
            description=description,
            template_type=template_type,
            tone=tone,
            output_format=output_format,
            template_storage_path=object_name,
            original_file_path=object_name,
            jinja2_variables=jinja2_vars,
            rendering_mode=rendering_mode,
            is_active=True,
            created_by=user_id,
        )
        db.add(template)
        await db.flush()
        await db.refresh(template)

        logger.info("템플릿 레코드 생성 완료: id=%s, name=%s, mode=%s", template.id, template.name, rendering_mode)
        return template

    # ------------------------------------------------------------------
    # upload_blank_form — 빈 양식 업로드 → 구조 분석 → Jinja2 자동 변환
    # ------------------------------------------------------------------

    @staticmethod
    async def upload_blank_form(
        db: AsyncSession,
        org_id: UUID,
        user_id: UUID,
        file_bytes: bytes,
        filename: str,
        template_type: str,
        tone: str,
        output_format: str,
        name: str | None = None,
        description: str | None = None,
    ) -> DocumentTemplate:
        """빈 양식 파일을 분석하고 Jinja2 변수를 자동 삽입하여 템플릿을 생성한다.

        기존 upload_template()과 달리, 이 메서드는 Jinja2 변수가 이미 들어있는
        파일이 아닌 빈 양식(제목/항목만 있는 문서)을 받아서:
          1. 문서 구조를 분석 (analyze_blank_form)
          2. 빈 섹션에 Jinja2 변수를 자동 삽입 (auto_generate_jinja2_from_structure)
          3. 원본과 변환된 파일 모두 MinIO에 저장
          4. DB에 템플릿 레코드 생성

        Args:
            db: 비동기 DB 세션
            org_id: 조직 ID
            user_id: 업로드한 사용자 ID
            file_bytes: 업로드된 빈 양식 파일의 바이트 데이터
            filename: 원본 파일 이름 (확장자 포함)
            template_type: 템플릿 유형 (자유 텍스트)
            tone: 문서 어조
            output_format: 출력 형식 (docx 또는 pptx)
            name: 템플릿 이름 (미지정 시 파일명 사용)
            description: 템플릿 설명

        Returns:
            생성된 DocumentTemplate ORM 객체

        Raises:
            HTTPException 400: 파일 형식이 지원되지 않거나 분석/변환에 실패한 경우
        """
        # 0. 중복 체크: 같은 이름 + 같은 org_id의 템플릿이 이미 있는지 확인한다
        effective_name = name or filename
        existing_stmt = select(DocumentTemplate).where(
            DocumentTemplate.organization_id == org_id,
            DocumentTemplate.name == effective_name,
        )
        existing_result = await db.execute(existing_stmt)
        existing_template = existing_result.scalar_one_or_none()

        # 1. 기존 템플릿이 있으면 해당 ID 재사용, 없으면 새 UUID 생성
        template_id = existing_template.id if existing_template else uuid_mod.uuid4()

        # 2. 파일 확장자를 확인하여 지원 여부를 검증한다
        ext = PurePosixPath(filename).suffix.lower().lstrip(".")
        if ext not in ("docx", "pptx"):
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=(f"빈 양식 자동 변환은 DOCX/PPTX만 지원합니다. 업로드된 파일 형식: {ext}"),
            )

        # 3. 빈 양식 구조 분석 — 휴리스틱 먼저 시도, 부족하면 AI fallback
        try:
            structure = analyze_blank_form(file_bytes, ext)
            logger.info(
                "빈 양식 휴리스틱 분석 완료: 섹션 %d개 (빈 섹션 %d개)",
                structure["metadata"]["total_sections"],
                structure["metadata"]["empty_sections"],
            )
        except ValueError as e:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"빈 양식 구조 분석 실패: {e}",
            ) from e

        # 3-1. 휴리스틱 결과가 부실한 경우 AI fallback을 시도한다
        # 조건: 섹션이 0개이거나, DOCX에 표가 있는데 섹션이 감지되지 않은 경우
        needs_ai_fallback = False
        if structure["metadata"]["total_sections"] == 0:
            # 섹션이 전혀 감지되지 않음
            needs_ai_fallback = True
        elif ext == "docx":
            # DOCX에 표가 있는데 섹션으로 감지되지 않은 경우를 확인한다
            try:
                from docx import Document as _Doc

                _test_doc = _Doc(io.BytesIO(file_bytes))
                if _test_doc.tables and structure["metadata"]["total_sections"] == 0:
                    needs_ai_fallback = True
            except Exception:
                # 체크 실패 시 휴리스틱 결과를 그대로 사용한다
                pass

        if needs_ai_fallback and ext == "docx":
            # AI 분석에는 OpenAI API 키가 필요하다
            ai_api_key = settings.openai_api_key
            if ai_api_key:
                try:
                    logger.info(
                        "휴리스틱 결과 부족 (섹션 %d개) → AI 양식 분석 시도",
                        structure["metadata"]["total_sections"],
                    )
                    structure = analyze_blank_form_with_ai(
                        file_bytes=file_bytes,
                        format=ext,
                        api_key=ai_api_key,
                        model=settings.llm_model or "gpt-4o",
                    )
                    logger.info(
                        "AI 양식 분석 완료: 섹션 %d개 (빈 섹션 %d개)",
                        structure["metadata"]["total_sections"],
                        structure["metadata"]["empty_sections"],
                    )
                except ValueError as e:
                    # AI도 실패하면 에러
                    logger.error("AI 양식 분석도 실패: %s", e)
                    raise HTTPException(
                        status_code=status.HTTP_400_BAD_REQUEST,
                        detail=f"휴리스틱/AI 양식 분석 모두 실패: {e}",
                    ) from e
            else:
                logger.warning("AI fallback 불가: OPENAI_API_KEY가 설정되지 않음")

        # 4. 섹션이 감지되었으면 Jinja2 변수 자동 삽입, 없으면 슬라이드마스터 양식으로 판단
        if structure["metadata"]["total_sections"] == 0:
            # PPTX 슬라이드마스터 양식 등 — Jinja2 변수 삽입 불가
            # rendering_mode를 "structured"로 설정하여 기존 AI 자유형 경로 사용
            logger.info("섹션 0개 → 슬라이드마스터 양식으로 판단, structured 모드로 저장")
            converted_bytes = file_bytes  # 원본 그대로 사용
            jinja2_vars = {"variables": []}
            rendering_mode = "structured"
        else:
            try:
                converted_bytes, jinja2_vars = auto_generate_jinja2_from_structure(
                    file_bytes,
                    ext,
                    structure,
                )
                logger.info(
                    "Jinja2 변수 자동 삽입 완료: 변수 %d개",
                    len(jinja2_vars.get("variables", [])),
                )
            except ValueError as e:
                raise HTTPException(
                    status_code=status.HTTP_400_BAD_REQUEST,
                    detail=f"Jinja2 자동 변환 실패: {e}",
                ) from e
            rendering_mode = "jinja2"

        # 5. MinIO에 원본 파일 업로드
        content_type = _CONTENT_TYPES.get(ext, "application/octet-stream")
        minio = MinIOService()

        # 원본 파일 경로: templates/{org_id}/{template_id}/original_{filename}
        original_object = f"templates/{org_id}/{template_id}/original_{filename}"
        minio.upload_file(
            bucket=settings.minio_bucket,
            object_name=original_object,
            file_data=file_bytes,
            content_type=content_type,
        )
        logger.info("원본 빈 양식 MinIO 업로드 완료: %s", original_object)

        # 6. MinIO에 변환된 템플릿 파일 업로드
        # 변환된 파일 경로: templates/{org_id}/{template_id}/{filename}
        template_object = f"templates/{org_id}/{template_id}/{filename}"
        minio.upload_file(
            bucket=settings.minio_bucket,
            object_name=template_object,
            file_data=converted_bytes,
            content_type=content_type,
        )
        logger.info("변환된 Jinja2 템플릿 MinIO 업로드 완료: %s", template_object)

        # 7. 기존 템플릿이 있으면 업데이트, 없으면 새로 생성
        if existing_template:
            existing_template.template_type = template_type
            existing_template.tone = tone
            existing_template.output_format = output_format
            existing_template.template_storage_path = template_object
            existing_template.original_file_path = original_object
            existing_template.jinja2_variables = jinja2_vars
            existing_template.rendering_mode = rendering_mode
            existing_template.is_active = True
            if description is not None:
                existing_template.description = description
            await db.flush()
            await db.refresh(existing_template)
            logger.info(
                "기존 빈 양식 템플릿 업데이트 완료: id=%s, name=%s, 변수 %d개",
                existing_template.id,
                existing_template.name,
                len(jinja2_vars.get("variables", [])),
            )
            return existing_template

        template = DocumentTemplate(
            id=template_id,
            organization_id=org_id,
            name=effective_name,
            description=description,
            template_type=template_type,
            tone=tone,
            output_format=output_format,
            template_storage_path=template_object,  # 변환된 파일 경로
            original_file_path=original_object,  # 원본 파일 경로
            jinja2_variables=jinja2_vars,  # 자동 생성된 변수 메타데이터
            rendering_mode=rendering_mode,  # 렌더링 방식
            is_active=True,
            created_by=user_id,
        )
        db.add(template)
        await db.flush()
        await db.refresh(template)

        logger.info(
            "빈 양식 → 템플릿 변환 및 저장 완료: id=%s, name=%s, 변수 %d개",
            template.id,
            template.name,
            len(jinja2_vars.get("variables", [])),
        )
        return template

    # ------------------------------------------------------------------
    # upload_smart — 스마트 업로드 (파일 분석 후 자동 처리 경로 결정)
    # ------------------------------------------------------------------

    @staticmethod
    async def upload_smart(
        db: AsyncSession,
        org_id: UUID,
        user_id: UUID,
        file_bytes: bytes,
        filename: str,
        name: str | None = None,
        description: str | None = None,
        template_type: str | None = None,
        tone: str = "formal",
    ) -> DocumentTemplate:
        """파일 내용을 자동 분석하여 적절한 업로드 처리를 수행한다.

        기존 upload_template(Jinja2 변수가 있는 파일)과 upload_blank_form(빈 양식)을
        하나로 통합한 메서드이다. 파일 내용에 {{ }} 패턴이 있는지 검사하여
        자동으로 적절한 처리 경로를 결정한다.

        자동 판단 로직:
          1. 파일 확장자로 output_format 결정 (docx/pptx)
          2. 파일 내용에서 {{ }} 패턴 검색
             - 있으면 → upload_template() 호출 (Jinja2 변수 파싱)
             - 없으면 → upload_blank_form() 호출 (구조 분석 + AI + 변수 자동 삽입)
          3. name이 없으면 파일명(확장자 제외)을 기본 이름으로 사용
          4. template_type이 없으면 파일 내용에서 키워드 매칭으로 추측

        Args:
            db: 비동기 DB 세션
            org_id: 조직 ID
            user_id: 업로드한 사용자 ID
            file_bytes: 업로드된 파일의 바이트 데이터
            filename: 원본 파일 이름 (확장자 포함)
            name: 템플릿 이름 (미지정 시 파일명에서 확장자를 제거한 값 사용)
            description: 템플릿 설명
            template_type: 템플릿 유형 (미지정 시 파일 내용으로 자동 추측)
            tone: 문서 어조 (기본: "formal")

        Returns:
            생성된 DocumentTemplate ORM 객체

        Raises:
            HTTPException 400: 지원하지 않는 파일 형식인 경우
        """
        # ---------------------------------------------------------------
        # 1단계: 파일 확장자에서 output_format 결정
        # ---------------------------------------------------------------
        # 파일 이름에서 확장자를 추출한다 (예: "회의록.docx" → "docx")
        ext = PurePosixPath(filename).suffix.lower().lstrip(".")

        # docx와 pptx만 지원한다
        if ext not in ("docx", "pptx"):
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=(f"스마트 업로드는 DOCX/PPTX만 지원합니다. 업로드된 파일 형식: {ext}"),
            )

        # 확장자가 곧 출력 형식이 된다
        output_format = ext

        # ---------------------------------------------------------------
        # 2단계: name이 없으면 파일명(확장자 제외)을 기본 이름으로 사용
        # ---------------------------------------------------------------
        # 예: "회의록_양식.docx" → "회의록_양식"
        if not name:
            name = PurePosixPath(filename).stem

        # ---------------------------------------------------------------
        # 3단계: 파일 내용에서 {{ }} 패턴 존재 여부를 검사한다
        # ---------------------------------------------------------------
        has_jinja2_pattern = TemplateService._detect_jinja2_pattern(
            file_bytes,
            ext,
        )

        if has_jinja2_pattern:
            logger.info(
                "스마트 업로드 — {{ }} 패턴 감지됨 → upload_template 경로 사용 (파일: %s)",
                filename,
            )
        else:
            logger.info(
                "스마트 업로드 — {{ }} 패턴 없음 → upload_blank_form 경로 사용 (파일: %s)",
                filename,
            )

        # ---------------------------------------------------------------
        # 4단계: template_type이 없으면 파일 내용/이름에서 자동 추측
        # ---------------------------------------------------------------
        if not template_type:
            template_type = TemplateService._guess_template_type(
                file_bytes,
                ext,
                filename,
            )
            logger.info(
                "스마트 업로드 — template_type 자동 추측 결과: '%s'",
                template_type,
            )

        # ---------------------------------------------------------------
        # 5단계: 감지 결과에 따라 기존 메서드를 호출한다
        # ---------------------------------------------------------------
        if has_jinja2_pattern:
            # {{ }} 패턴이 있음 → 기존 upload_template 로직 (Jinja2 변수 파싱)
            return await TemplateService.upload_template(
                db=db,
                org_id=org_id,
                user_id=user_id,
                file_bytes=file_bytes,
                filename=filename,
                template_type=template_type,
                tone=tone,
                output_format=output_format,
                name=name,
                description=description,
            )
        else:
            # {{ }} 패턴이 없음 → 기존 upload_blank_form 로직 (구조 분석 + 변수 삽입)
            return await TemplateService.upload_blank_form(
                db=db,
                org_id=org_id,
                user_id=user_id,
                file_bytes=file_bytes,
                filename=filename,
                template_type=template_type,
                tone=tone,
                output_format=output_format,
                name=name,
                description=description,
            )

    # ------------------------------------------------------------------
    # _detect_jinja2_pattern — 파일 내용에서 {{ }} 패턴 존재 여부 검사
    # ------------------------------------------------------------------

    @staticmethod
    def _detect_jinja2_pattern(file_bytes: bytes, ext: str) -> bool:
        """파일 내용에서 Jinja2 {{ }} 패턴이 있는지 검사한다.

        DOCX와 PPTX는 ZIP 기반의 XML 파일이므로, 단순히 바이트에서
        검색하면 안 되고 각 파일 형식에 맞는 라이브러리로 텍스트를
        추출한 후 정규식으로 검사해야 한다.

        Args:
            file_bytes: 파일의 바이트 데이터
            ext: 파일 확장자 ("docx" 또는 "pptx")

        Returns:
            True이면 {{ }} 패턴이 존재, False이면 없음
        """
        import re

        # Jinja2 변수 패턴: {{ 변수명 }}
        # 양쪽 중괄호 2개와 그 사이의 텍스트를 찾는다
        jinja2_pattern = re.compile(r"\{\{.*?\}\}")

        try:
            if ext == "docx":
                # python-docx로 DOCX 파일을 열어서 텍스트를 추출한다
                from docx import Document as DocxDocument

                doc = DocxDocument(io.BytesIO(file_bytes))

                # 모든 문단(paragraph)의 텍스트를 검사한다
                for para in doc.paragraphs:
                    if jinja2_pattern.search(para.text):
                        return True

                # 모든 표(table)의 셀 텍스트를 검사한다
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            if jinja2_pattern.search(cell.text):
                                return True

            elif ext == "pptx":
                # python-pptx로 PPTX 파일을 열어서 텍스트를 추출한다
                from pptx import Presentation

                prs = Presentation(io.BytesIO(file_bytes))

                # 모든 슬라이드의 모든 도형(shape)에서 텍스트를 검사한다
                for slide in prs.slides:
                    for shape in slide.shapes:
                        # 텍스트가 있는 도형만 검사한다
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                # 문단의 전체 텍스트를 조합한다
                                full_text = "".join(run.text for run in paragraph.runs)
                                if jinja2_pattern.search(full_text):
                                    return True

                        # 표(table) 도형인 경우 셀 텍스트도 검사한다
                        if shape.has_table:
                            for row in shape.table.rows:
                                for cell in row.cells:
                                    if jinja2_pattern.search(cell.text):
                                        return True

        except Exception as e:
            # 파일 파싱에 실패하면 패턴이 없는 것으로 간주한다
            logger.warning(
                "Jinja2 패턴 감지 중 오류 발생 (패턴 없음으로 처리): %s",
                e,
            )
            return False

        # 어디에서도 패턴을 찾지 못했다
        return False

    # ------------------------------------------------------------------
    # _guess_template_type — 파일 내용/이름에서 템플릿 유형 자동 추측
    # ------------------------------------------------------------------

    @staticmethod
    def _guess_template_type(
        file_bytes: bytes,
        ext: str,
        filename: str,
    ) -> str:
        """파일명과 내용에서 키워드를 매칭하여 템플릿 유형을 추측한다.

        아래 키워드 매핑 테이블을 순서대로 검사하여, 가장 먼저 일치하는
        유형을 반환한다. 아무것도 일치하지 않으면 "기타"를 반환한다.

        키워드 매핑:
          - "회의록", "minutes" → "회의록"
          - "보고서", "report"  → "보고서"
          - "제안서", "proposal" → "제안서"

        Args:
            file_bytes: 파일의 바이트 데이터
            ext: 파일 확장자
            filename: 원본 파일 이름

        Returns:
            추측된 템플릿 유형 문자열
        """
        # 키워드 → 유형 매핑 테이블 (순서대로 검사한다)
        keyword_map: list[tuple[list[str], str]] = [
            (["회의록", "minutes"], "회의록"),
            (["보고서", "report"], "보고서"),
            (["제안서", "proposal"], "제안서"),
        ]

        # 먼저 파일명에서 키워드를 검색한다 (빠르고 정확도 높음)
        filename_lower = filename.lower()
        for keywords, result_type in keyword_map:
            for kw in keywords:
                if kw in filename_lower:
                    return result_type

        # 파일명에서 못 찾았으면, 파일 내용에서 텍스트를 추출하여 검사한다
        content_text = ""
        try:
            if ext == "docx":
                from docx import Document as DocxDocument

                doc = DocxDocument(io.BytesIO(file_bytes))
                # 처음 20개 문단만 검사한다 (전체 파일을 검사하면 느릴 수 있으므로)
                content_text = " ".join(p.text for p in doc.paragraphs[:20])

            elif ext == "pptx":
                from pptx import Presentation

                prs = Presentation(io.BytesIO(file_bytes))
                # 처음 5개 슬라이드만 검사한다
                texts: list[str] = []
                for slide in prs.slides[:5]:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            texts.append(shape.text_frame.text)
                content_text = " ".join(texts)

        except Exception as e:
            # 텍스트 추출에 실패하면 파일명 기반 결과만 사용한다
            logger.warning(
                "template_type 추측을 위한 텍스트 추출 실패: %s",
                e,
            )

        # 추출된 텍스트에서 키워드 검색
        content_lower = content_text.lower()
        for keywords, result_type in keyword_map:
            for kw in keywords:
                if kw in content_lower:
                    return result_type

        # 아무 키워드도 매칭되지 않으면 "기타"를 반환한다
        return "기타"

    # ------------------------------------------------------------------
    # convert_to_template — 일반 문서를 Jinja2 템플릿으로 변환
    # ------------------------------------------------------------------

    @staticmethod
    async def convert_to_template(
        db: AsyncSession,
        template_id: UUID,
        org_id: UUID,
        ai_analysis: dict,
    ) -> DocumentTemplate:
        """기존 템플릿의 원본 파일을 AI 분석 결과를 적용하여 Jinja2 템플릿으로 변환한다.

        처리 순서:
          1. 기존 템플릿 조회 → original_file_path 확인
          2. MinIO에서 원본 파일 다운로드
          3. jinja2_engine.convert_example_to_template()으로 변환
          4. 변환된 파일을 MinIO에 새 경로로 저장
          5. 변수 재추출
          6. DB 업데이트

        Args:
            db: 비동기 DB 세션
            template_id: 변환할 템플릿 ID
            org_id: 조직 ID
            ai_analysis: AI가 분석한 텍스트 → 변수 매핑 정보
                예: {"replacements": [{"original": "홍길동", "variable": "author"}]}

        Returns:
            업데이트된 DocumentTemplate ORM 객체
        """
        # 1. 기존 템플릿 조회
        template = await TemplateService.get_template(db, template_id, org_id)

        # 원본 파일 경로가 없으면 변환 불가
        if not template.original_file_path:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="원본 파일이 없어 변환할 수 없습니다. 먼저 파일을 업로드하세요.",
            )

        # 2. MinIO에서 원본 파일 다운로드
        minio = MinIOService()
        original_bytes = minio.download_file(
            bucket=settings.minio_bucket,
            object_name=template.original_file_path,
        )
        logger.info("원본 파일 다운로드 완료: %s (%d bytes)", template.original_file_path, len(original_bytes))

        # 3. jinja2_engine으로 변환 (텍스트 → {{ variable }} 치환)
        converted_bytes = convert_example_to_template(
            file_bytes=original_bytes,
            format=template.output_format,
            ai_analysis=ai_analysis,
        )

        # 4. 변환된 파일을 MinIO에 저장 (경로에 '_template' 접미사 추가)
        original_path = PurePosixPath(template.original_file_path)
        converted_name = f"{original_path.stem}_template{original_path.suffix}"
        converted_object_name = str(original_path.parent / converted_name)
        ext = original_path.suffix.lower().lstrip(".")
        content_type = _CONTENT_TYPES.get(ext, "application/octet-stream")

        minio.upload_file(
            bucket=settings.minio_bucket,
            object_name=converted_object_name,
            file_data=converted_bytes,
            content_type=content_type,
        )
        logger.info("변환된 템플릿 MinIO 저장 완료: %s", converted_object_name)

        # 5. 변환된 파일에서 Jinja2 변수 재추출
        jinja2_vars: dict | None = None
        try:
            if ext == "docx":
                jinja2_vars = extract_docx_variables(converted_bytes)
            elif ext == "pptx":
                jinja2_vars = extract_pptx_variables(converted_bytes)
        except ValueError as e:
            logger.warning("변환 후 변수 재추출 실패 (계속 진행): %s", e)

        # 6. DB 업데이트
        template.template_storage_path = converted_object_name
        template.jinja2_variables = jinja2_vars
        template.rendering_mode = "jinja2"

        await db.flush()
        await db.refresh(template)

        logger.info("템플릿 변환 완료: id=%s", template.id)
        return template

    # ------------------------------------------------------------------
    # update_variables — 변수 메타데이터 수동 편집
    # ------------------------------------------------------------------

    @staticmethod
    async def update_variables(
        db: AsyncSession,
        template_id: UUID,
        org_id: UUID,
        variables_data: list[dict],
    ) -> DocumentTemplate:
        """템플릿의 Jinja2 변수 메타데이터(라벨, 타입, 설명 등)를 업데이트한다.

        프론트엔드에서 사용자가 변수의 label, var_type, description, required 등을
        수동으로 편집한 결과를 저장한다.

        Args:
            db: 비동기 DB 세션
            template_id: 대상 템플릿 ID
            org_id: 조직 ID
            variables_data: 변수 메타데이터 리스트
                예: [{"name": "title", "var_type": "string", "label": "제목", ...}]

        Returns:
            업데이트된 DocumentTemplate ORM 객체
        """
        template = await TemplateService.get_template(db, template_id, org_id)

        # 기존 jinja2_variables에 메타데이터를 병합한다
        # variables 키 아래에 전체 변수 목록을 덮어쓴다
        template.jinja2_variables = {"variables": variables_data}

        await db.flush()
        await db.refresh(template)

        logger.info(
            "템플릿 변수 메타데이터 업데이트 완료: id=%s, 변수 %d개",
            template.id,
            len(variables_data),
        )
        return template

    # ------------------------------------------------------------------
    # get_template_file — 템플릿 파일 다운로드
    # ------------------------------------------------------------------

    @staticmethod
    async def get_template_file(
        db: AsyncSession,
        template_id: UUID,
        org_id: UUID,
    ) -> tuple[bytes, str]:
        """MinIO에서 템플릿 파일을 다운로드하여 바이트와 파일명을 반환한다.

        template_storage_path가 설정된 경우 변환된 템플릿 파일을,
        없으면 original_file_path의 원본 파일을 반환한다.

        Args:
            db: 비동기 DB 세션
            template_id: 대상 템플릿 ID
            org_id: 조직 ID

        Returns:
            (file_bytes, filename) 튜플

        Raises:
            HTTPException 404: 파일 경로가 없는 경우
        """
        template = await TemplateService.get_template(db, template_id, org_id)

        # 템플릿 파일 경로 결정 (변환된 파일 우선, 없으면 원본)
        storage_path = template.template_storage_path or template.original_file_path
        if not storage_path:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail="이 템플릿에는 다운로드 가능한 파일이 없습니다.",
            )

        # MinIO에서 파일 다운로드
        minio = MinIOService()
        file_bytes = minio.download_file(
            bucket=settings.minio_bucket,
            object_name=storage_path,
        )

        # 파일명은 스토리지 경로의 마지막 부분
        filename = PurePosixPath(storage_path).name

        logger.info(
            "템플릿 파일 다운로드: id=%s, path=%s, size=%d bytes",
            template_id,
            storage_path,
            len(file_bytes),
        )
        return file_bytes, filename

    # ------------------------------------------------------------------
    # auto_fill_variables — AI 자동 채우기 (소스 문서 기반 변수값 생성)
    # ------------------------------------------------------------------

    @staticmethod
    async def auto_fill_variables(
        db: AsyncSession,
        template_id: UUID,
        org_id: UUID,
        source_document_ids: list[UUID],
        ai_prompt: str | None = None,
    ) -> dict:
        """소스 문서 내용을 참고하여 AI가 템플릿 변수값을 자동으로 생성한다.

        처리 순서:
          1. 템플릿 조회 → jinja2_variables에서 변수 목록 추출
          2. source_document_ids로 tb_document_chunks에서 청크 텍스트 조회
          3. 변수 스키마를 JSON Schema properties로 변환
          4. OpenAI GPT-4o에 Structured Outputs로 호출
          5. 생성된 변수값을 dict로 반환

        Args:
            db: 비동기 DB 세션
            template_id: 대상 템플릿 ID
            org_id: 조직 ID
            source_document_ids: 참고할 소스 문서 ID 목록
            ai_prompt: AI에게 전달할 추가 지시사항 (선택)

        Returns:
            AI가 생성한 변수값 딕셔너리
            예: {"title": "보고서 제목", "items": ["항목1", "항목2"]}

        Raises:
            HTTPException 404: 템플릿이 존재하지 않는 경우
            HTTPException 400: 변수가 정의되지 않았거나 소스 문서가 없는 경우
        """
        # 1. 템플릿을 조회한다
        template = await TemplateService.get_template(db, template_id, org_id)

        # 2. jinja2_variables에서 변수 목록을 추출한다
        # jinja2_variables 구조: {"variables": [{"name": "...", "type": "...", ...}, ...]}
        if not template.jinja2_variables or "variables" not in template.jinja2_variables:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="이 템플릿에는 정의된 변수가 없습니다. 먼저 변수를 설정하세요.",
            )

        variables = template.jinja2_variables["variables"]
        if not variables:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="이 템플릿에는 정의된 변수가 없습니다. 먼저 변수를 설정하세요.",
            )

        # 3. 소스 문서의 청크 텍스트를 DB에서 조회한다
        # 최대 30개 청크만 가져와서 토큰 제한을 방지한다
        chunk_stmt = select(DocumentChunk.content).where(DocumentChunk.document_id.in_(source_document_ids)).limit(30)
        chunk_result = await db.execute(chunk_stmt)
        chunks = chunk_result.scalars().all()

        if not chunks:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="소스 문서에서 텍스트를 찾을 수 없습니다. 문서가 처리 완료되었는지 확인하세요.",
            )

        # 청크 텍스트를 하나의 문자열로 합친다
        context_text = "\n\n".join(chunks)
        logger.info(
            "자동 채우기 — 청크 %d개 로드 완료 (총 %d자)",
            len(chunks),
            len(context_text),
        )

        # 4. 변수 스키마를 OpenAI JSON Schema properties 형식으로 변환한다
        # 각 변수의 type에 따라 적절한 JSON Schema 속성을 생성한다
        properties: dict[str, dict] = {}
        for var in variables:
            var_name = var.get("name", "")
            # var_type 또는 type 필드에서 타입을 읽는다 (호환성 고려)
            var_type = var.get("type", var.get("var_type", "string"))

            if var_type == "array":
                # 배열 타입: 문자열 배열로 생성
                properties[var_name] = {
                    "type": "array",
                    "items": {"type": "string"},
                }
            elif var_type == "boolean":
                # 불리언 타입
                properties[var_name] = {"type": "boolean"}
            else:
                # 기본: 문자열 타입 (string, image 등 모두 문자열로 처리)
                properties[var_name] = {"type": "string"}

        # 5. LLMClient 추상화를 통해 Structured Outputs로 호출한다
        try:
            # 템플릿 자동 채우기용 프로바이더를 설정에서 가져온다
            provider = get_provider_for_task("template")
            llm_client = create_llm_client(provider)

            # 시스템 프롬프트: AI에게 역할과 규칙을 알려준다
            system_message = (
                "주어진 문서 내용을 참고하여 각 변수에 맞는 적절한 값을 생성하세요. "
                "한국어로 작성하세요. "
                "문서에 명시된 정보를 우선적으로 사용하고, "
                "문서에 없는 정보는 맥락에 맞게 합리적으로 추론하세요."
            )

            # 사용자 메시지: 문서 내용 + 변수 목록 + 추가 지시사항
            user_message = (
                f"## 참고 문서 내용\n{context_text}\n\n## 생성할 변수 목록\n{json.dumps(variables, ensure_ascii=False)}"
            )
            # 추가 프롬프트가 있으면 사용자 메시지에 덧붙인다
            if ai_prompt:
                user_message += f"\n\n## 추가 지시사항\n{ai_prompt}"

            # JSON 스키마 정의 — AI가 반드시 이 형식에 맞춰 응답하도록 강제한다
            json_schema = {
                "name": "auto_fill",
                "strict": True,
                "schema": {
                    "type": "object",
                    "properties": properties,
                    "required": list(properties.keys()),
                    "additionalProperties": False,
                },
            }

            messages = [
                {"role": "system", "content": system_message},
                {"role": "user", "content": user_message},
            ]

            # generate_structured()는 dict를 직접 반환하므로 JSON 파싱이 불필요하다
            result = await llm_client.generate_structured(messages, json_schema, temperature=0.3, max_tokens=4096)
            logger.info(
                "자동 채우기 완료: template_id=%s, 변수 %d개 생성",
                template_id,
                len(result),
            )
            return result

        except Exception as e:
            # LLM API 오류 (네트워크, 인증, 할당량 초과 등) — 프로바이더 독립적 처리
            logger.error("LLM API 오류 (자동 채우기): %s", e, exc_info=True)
            return {}

    # ------------------------------------------------------------------
    # get_template_structure — 에디터용 문서 구조 추출
    # ------------------------------------------------------------------

    @staticmethod
    async def get_template_structure(
        db: AsyncSession,
        template_id: UUID,
        org_id: UUID,
    ) -> dict:
        """에디터에서 렌더링할 수 있도록 템플릿 원본의 문서 구조를 반환한다.

        MinIO에서 원본 DOCX 파일을 다운로드하고,
        extract_docx_structure_for_editor()를 호출하여
        표/문단의 상세 구조와 기존 변수 매핑 정보를 함께 반환한다.

        Args:
            db: 비동기 DB 세션
            template_id: 대상 템플릿 ID
            org_id: 조직 ID

        Returns:
            에디터용 구조 딕셔너리 (paragraphs, tables, existing_variables)

        Raises:
            HTTPException 404: 템플릿이 존재하지 않는 경우
            HTTPException 400: 원본 파일이 없거나 DOCX가 아닌 경우
        """
        # 1. 템플릿 조회
        template = await TemplateService.get_template(db, template_id, org_id)

        # 2. 원본 파일 경로 결정 (original_file_path 우선, 없으면 template_storage_path)
        file_path = template.original_file_path or template.template_storage_path
        if not file_path:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="이 템플릿에는 업로드된 파일이 없습니다.",
            )

        # 3. DOCX만 에디터 구조 추출을 지원한다
        if not file_path.lower().endswith(".docx"):
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="에디터 구조 추출은 DOCX 파일만 지원합니다.",
            )

        # 4. MinIO에서 원본 파일 다운로드
        minio = MinIOService()
        try:
            file_bytes = minio.download_file(
                bucket=settings.minio_bucket,
                object_name=file_path,
            )
        except Exception as e:
            logger.error("에디터용 파일 다운로드 실패: %s", e)
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"파일 다운로드에 실패했습니다: {e}",
            ) from e

        logger.info(
            "에디터용 파일 다운로드 완료: %s (%d bytes)",
            file_path,
            len(file_bytes),
        )

        # 5. 기존 jinja2_variables를 existing_variables로 전달한다
        existing_variables = template.jinja2_variables

        # 6. 문서 구조를 추출한다
        try:
            structure = extract_docx_structure_for_editor(
                file_bytes=file_bytes,
                existing_variables=existing_variables,
            )
        except Exception as e:
            logger.error("에디터용 구조 추출 실패: %s", e)
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"문서 구조 추출에 실패했습니다: {e}",
            ) from e

        return structure

    # ------------------------------------------------------------------
    # apply_variable_mapping — 에디터에서 설정한 변수 매핑 적용
    # ------------------------------------------------------------------

    @staticmethod
    async def apply_variable_mapping(
        db: AsyncSession,
        template_id: UUID,
        org_id: UUID,
        mappings: list[dict],
    ) -> DocumentTemplate:
        """에디터에서 사용자가 설정한 변수 매핑을 원본 DOCX에 적용한다.

        처리 흐름:
          1. MinIO에서 원본 DOCX 파일 다운로드
          2. python-docx로 열기
          3. 각 매핑에 대해 해당 위치에 {{ 변수명 }} 삽입
          4. 변환된 파일을 MinIO에 저장 (template_storage_path)
          5. jinja2_variables 업데이트
          6. DB 업데이트

        Args:
            db: 비동기 DB 세션
            template_id: 대상 템플릿 ID
            org_id: 조직 ID
            mappings: 변수 매핑 정보 리스트
                각 항목: {"location_type": "table_cell", "table_index": 0,
                          "row": 2, "col": 1, "variable_name": "장소", ...}

        Returns:
            업데이트된 DocumentTemplate ORM 객체

        Raises:
            HTTPException 404: 템플릿이 존재하지 않는 경우
            HTTPException 400: 원본 파일이 없거나 매핑 적용에 실패한 경우
        """
        import io as _io

        from docx import Document

        # 1. 템플릿 조회
        template = await TemplateService.get_template(db, template_id, org_id)

        # 2. 원본 파일 경로 결정
        file_path = template.original_file_path or template.template_storage_path
        if not file_path:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="이 템플릿에는 업로드된 파일이 없습니다.",
            )

        # 3. MinIO에서 원본 파일 다운로드
        minio = MinIOService()
        try:
            file_bytes = minio.download_file(
                bucket=settings.minio_bucket,
                object_name=file_path,
            )
        except Exception as e:
            logger.error("매핑 적용용 파일 다운로드 실패: %s", e)
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"파일 다운로드에 실패했습니다: {e}",
            ) from e

        # 4. python-docx로 문서를 연다
        try:
            doc = Document(_io.BytesIO(file_bytes))
        except Exception as e:
            logger.error("DOCX 파일 열기 실패: %s", e)
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"DOCX 파일을 열 수 없습니다: {e}",
            ) from e

        # 5. 각 매핑에 대해 {{ 변수명 }}을 삽입한다
        variables_meta: list[dict] = []
        for mapping in mappings:
            location_type = mapping.get("location_type")
            variable_name = mapping.get("variable_name", "")
            jinja2_text = "{{ " + variable_name + " }}"

            if location_type == "table_cell":
                # --- 표 셀에 변수 삽입 ---
                t_idx = mapping.get("table_index", 0)
                row_idx = mapping.get("row", 0)
                col_idx = mapping.get("col", 0)

                # 표/행/열 범위를 검증한다
                if t_idx >= len(doc.tables):
                    logger.warning(
                        "매핑 건너뜀: table_index %d 범위 초과 (표 %d개)",
                        t_idx,
                        len(doc.tables),
                    )
                    continue

                table = doc.tables[t_idx]
                if row_idx >= len(table.rows):
                    logger.warning(
                        "매핑 건너뜀: row %d 범위 초과 (행 %d개)",
                        row_idx,
                        len(table.rows),
                    )
                    continue

                row_cells = table.rows[row_idx].cells
                if col_idx >= len(row_cells):
                    logger.warning(
                        "매핑 건너뜀: col %d 범위 초과 (셀 %d개)",
                        col_idx,
                        len(row_cells),
                    )
                    continue

                # 대상 셀에 Jinja2 변수를 삽입한다
                target_cell = row_cells[col_idx]
                if target_cell.paragraphs:
                    para = target_cell.paragraphs[0]
                    if para.runs:
                        # 기존 텍스트를 덮어쓴다
                        para.runs[0].text = jinja2_text
                        for run in para.runs[1:]:
                            run.text = ""
                    else:
                        para.add_run(jinja2_text)
                else:
                    target_cell.text = jinja2_text

                logger.debug(
                    "표 셀에 변수 삽입: table=%d, row=%d, col=%d → %s",
                    t_idx,
                    row_idx,
                    col_idx,
                    jinja2_text,
                )

            elif location_type == "paragraph":
                # --- 문단에 변수 삽입 ---
                p_idx = mapping.get("paragraph_index", 0)

                if p_idx >= len(doc.paragraphs):
                    logger.warning(
                        "매핑 건너뜀: paragraph_index %d 범위 초과 (문단 %d개)",
                        p_idx,
                        len(doc.paragraphs),
                    )
                    continue

                para = doc.paragraphs[p_idx]
                if para.runs:
                    para.runs[0].text = jinja2_text
                    for run in para.runs[1:]:
                        run.text = ""
                else:
                    para.add_run(jinja2_text)

                logger.debug(
                    "문단에 변수 삽입: paragraph=%d → %s",
                    p_idx,
                    jinja2_text,
                )
            else:
                logger.warning(
                    "매핑 건너뜀: 알 수 없는 location_type '%s'",
                    location_type,
                )
                continue

            # 변수 메타데이터를 수집한다
            variables_meta.append(
                {
                    "name": variable_name,
                    "type": mapping.get("var_type", "string"),
                    "label": mapping.get("label", variable_name),
                    "description": f"에디터에서 매핑된 변수: {variable_name}",
                    "required": True,
                    "category": mapping.get("category", "ai_generated"),
                    "field_type": mapping.get("field_type", "short"),
                }
            )

        # 6. 변환된 파일을 바이트로 저장한다
        output = _io.BytesIO()
        doc.save(output)
        output.seek(0)
        converted_bytes = output.read()

        # 7. MinIO에 변환된 파일 업로드 (template_storage_path 경로)
        # 기존 template_storage_path가 있으면 같은 경로에 덮어쓰기,
        # 없으면 새 경로 생성
        from pathlib import PurePosixPath as _PurePosixPath

        if template.template_storage_path:
            target_path = template.template_storage_path
        else:
            # 원본 파일 경로에서 파일명을 변형하여 템플릿 경로 생성
            orig_path = _PurePosixPath(file_path)
            target_path = str(orig_path.parent / f"{orig_path.stem}_mapped{orig_path.suffix}")

        ext = _PurePosixPath(target_path).suffix.lower().lstrip(".")
        content_type = _CONTENT_TYPES.get(ext, "application/octet-stream")

        minio.upload_file(
            bucket=settings.minio_bucket,
            object_name=target_path,
            file_data=converted_bytes,
            content_type=content_type,
        )
        logger.info(
            "변수 매핑 적용된 템플릿 MinIO 업로드 완료: %s",
            target_path,
        )

        # 8. DB 업데이트 — jinja2_variables와 template_storage_path 갱신
        template.template_storage_path = target_path
        template.jinja2_variables = {"variables": variables_meta}
        template.rendering_mode = "jinja2"

        await db.flush()
        await db.refresh(template)

        logger.info(
            "변수 매핑 적용 완료: id=%s, 변수 %d개",
            template.id,
            len(variables_meta),
        )
        return template
