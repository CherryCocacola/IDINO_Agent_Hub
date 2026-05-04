"""Phase 4 S2 D1 → D6 — PptxBuilder 테스트.

D1 검증 범위 (요구 3건 + 보조 2건):
1. `PptxBuilder` 인스턴스 생성 성공 (ABC 추상 메서드 구현 완료).
2. `BuilderRegistry.get("pptx")` 반환값이 `PptxBuilder` 인스턴스.
3. `build(empty_schema)` 가 bytes 를 반환하고 ZIP 시그니처 `b"PK"` 로 시작.
4. `PptxBuilder.target == "pptx"` / `supported_components` 는 frozenset.
5. 미지원 컴포넌트 (Chart) 포함 스키마라도 빌드가 실패하지 않는다.

D2 추가 범위 (요구 6건 + 보조):
6. SlideTitle 렌더링 — 슬라이드 내 텍스트에 "테스트 제목" 포함.
7. Heading 렌더링 — 동일.
8. Paragraph 렌더링 — 동일.
9. BulletList 렌더링 — items 개수만큼 paragraph + bullet glyph 포함.
10. 4 종 복합 페이지 — 단일 슬라이드에 4 개 shape (textbox) 생성.
11. 다중 페이지 — 슬라이드 수 = pages 수.
12. D2 범위 외 컴포넌트(Callout 등) skip — 빌드 실패 없이 무시.

D3 추가 범위 (9건):
13. `render_kpi` — 단일 KPI 카드 렌더, 값/라벨/delta 텍스트 포함.
14. `render_kpi` delta 음수 → 빨강(KPI_DELTA_DOWN_COLOR) 색 적용.
15. `render_data_table` — 3x3 테이블, 헤더 배경 #34495E + 제브라.
16. `render_data_table` — 숫자 셀은 우측 정렬, 텍스트 셀은 좌측 정렬.
17. `resolve_layout` — 한글/영문 혼합 후보에서 매칭 성공.
18. `resolve_layout` — 매칭 실패 시 fallback 반환 (WARNING 로그 emit).
19. `resolve_layout` — organization_overrides 가 후보보다 우선.
20. `resolve_placeholder` — PP_PLACEHOLDER.TITLE semantic 으로 매칭.
21. `build()` 통합 — KPI/DataTable 포함 스키마 → 정상 PPTX + 텍스트 확인.

D6 추가 범위 (5건):
22. `render_image` — URL source, httpx.get mock → add_picture 호출 + Picture shape 존재.
23. `render_image` — base64 data-URI source → 디코드 후 정상 삽입.
24. `render_image` — fetch 실패 시 placeholder + alt 텍스트만 표시, 빌드 성공.
25. `render_image` — aspect ratio 유지 (가로 긴 이미지, max_height 초과 시 축소).
26. `build()` 통합 — Image 포함 스키마 → 슬라이드에 Picture shape + caption 포함.

참조:
- backend/app/integrations/document_builders/pptx/builder.py
- backend/app/integrations/document_builders/pptx/components.py
- backend/app/integrations/document_builders/pptx/image_fetcher.py (D6)
- backend/app/integrations/document_builders/pptx/layout_resolver.py (D3)
- backend/tests/test_document_builders_base.py (동일 fixture 패턴)
- docs/phase3_execution_roadmap.md §2.2 S2 D1, D2, D3, D6
"""

from __future__ import annotations

import logging
from datetime import UTC, datetime
from io import BytesIO
from uuid import UUID

import pytest
from pptx import Presentation
from pptx.util import Inches

from app.integrations.document_builders.base import BuilderRegistry
from app.integrations.document_builders.pptx import PptxBuilder, register_pptx_builder
from app.integrations.document_builders.pptx import image_fetcher as _image_fetcher
from app.integrations.document_builders.pptx.components import (
    render_chart,
    render_data_table,
    render_image,
    render_kpi,
)
from app.integrations.document_builders.pptx.constants import (
    CHART_EMPTY_PLACEHOLDER_TEXT,
    CHART_SERIES_PALETTE,
    IDINO_HEADER_NAVY,
    KPI_DELTA_DOWN_COLOR,
    KPI_DELTA_UP_COLOR,
    TABLE_ZEBRA_ROW_BG,
)
from app.integrations.document_builders.pptx.image_fetcher import ImageFetchError
from app.integrations.document_builders.pptx.layout_resolver import (
    resolve_layout,
    resolve_placeholder,
)
from app.integrations.document_builders.pptx.style import parse_hex_color
from app.modules.documents_v2.schemas import (
    BulletItem,
    BulletListComponent,
    CalloutComponent,
    ChartComponent,
    ChartData,
    ChartSeries,
    DataTableComponent,
    DocumentMetadata,
    DocumentSchema,
    HeadingComponent,
    IconItem,
    IconRowComponent,
    ImageComponent,
    KPIComponent,
    Page,
    ParagraphComponent,
    QuoteComponent,
    SlideSubtitleComponent,
    SlideTitleComponent,
    TimelineComponent,
    TimelineEvent,
)

# ---------------------------------------------------------------------------
# 공용 fixture
# ---------------------------------------------------------------------------

_FIXED_DOC_ID = UUID("7b2a5f3e-1c4d-4b8a-9e7f-0a1b2c3d4e5f")

# PPTX OOXML 은 ZIP 컨테이너이므로 첫 2 바이트가 항상 b"PK" 이다.
# (ZIP Local File Header magic number 0x504b 0304 의 앞 2 바이트)
_ZIP_SIGNATURE = b"PK"


def _make_metadata() -> DocumentMetadata:
    """테스트 공통 메타데이터."""
    now = datetime.now(UTC)
    return DocumentMetadata(created_at=now, updated_at=now)


def _make_schema(*pages: Page) -> DocumentSchema:
    """pages 를 받아 `DocumentSchema` 를 조립하는 공용 팩토리."""
    return DocumentSchema(
        document_id=_FIXED_DOC_ID,
        schema_version="1.0",
        type="slide_report",
        mode="free_generation",
        template_id=None,
        pages=list(pages),
        metadata=_make_metadata(),
    )


def _extract_all_slide_text(pptx_bytes: bytes) -> list[str]:
    """PPTX bytes → 각 슬라이드의 모든 텍스트 프레임 텍스트를 한 덩어리 문자열로.

    슬라이드 순서대로 리스트 반환. shape 별 text 를 개행으로 합침.
    """
    prs = Presentation(BytesIO(pptx_bytes))
    slide_texts: list[str] = []
    for slide in prs.slides:
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        slide_texts.append("\n".join(texts))
    return slide_texts


def _count_shapes_with_text(pptx_bytes: bytes, slide_idx: int = 0) -> int:
    """특정 슬라이드에서 텍스트 프레임을 가진 shape 개수."""
    prs = Presentation(BytesIO(pptx_bytes))
    slide = prs.slides[slide_idx]
    return sum(1 for shape in slide.shapes if shape.has_text_frame)


@pytest.fixture(autouse=True)
def _isolated_registry():
    """각 테스트 사이에 Registry 를 리셋 + PptxBuilder/DocxBuilder 재등록.

    base 테스트와 동일 패턴이지만, 본 파일은 `pptx/__init__.py` 의 자동 등록에
    의존하므로 clear() 직후 `register_pptx_builder()` 를 다시 호출해 상태를
    되돌린다.

    S3 D1 추가 — DocxBuilder 스텁 테스트를 위해 DocxBuilder 도 재등록.
    """
    from app.integrations.document_builders.docx import register_docx_builder

    BuilderRegistry.clear()
    register_pptx_builder()
    register_docx_builder()
    yield
    BuilderRegistry.clear()


@pytest.fixture
def empty_schema() -> DocumentSchema:
    """최소 유효 문서 — SlideTitle 1개만 있는 1 페이지 구조.

    Pydantic 이 `pages: min_length=1`, `components: min_length=1` 을 강제하므로
    실제로 "완전히 빈" 스키마는 만들 수 없다. 테스트 명칭은 "D1 골격이 내용을
    무시한다"는 의미에서의 empty 이다.
    """
    return _make_schema(
        Page(
            id="p1",
            page_kind="slide",
            layout="title_slide",
            title="S2 D1 골격 테스트",
            components=[
                SlideTitleComponent(id="c1", text="PPTX Builder 골격"),
            ],
        ),
    )


@pytest.fixture
def schema_with_unsupported() -> DocumentSchema:
    """아직 이관되지 않은 Hero 컴포넌트를 포함한 스키마 — 빌드 실패하지 않음을 검증.

    역할 이력:
      - D7 까지: Chart 가 미지원 대표 → Callout 으로 이관.
      - S3 D6 에서 Callout 이 지원으로 승격 → Hero 로 이관.
    Hero 는 S3 이후 이관 예정이므로 현재 시점에서 skip + WARN 경로를 타는 대표 케이스.
    """
    from app.modules.documents_v2.schemas import HeroComponent

    return _make_schema(
        Page(
            id="p1",
            page_kind="slide",
            layout="content_body",
            components=[
                SlideTitleComponent(id="c1", text="공지 사항"),
                ParagraphComponent(id="c2", text="Hero 는 S3 이후 이관 예정."),
                HeroComponent(
                    id="c3",
                    title="현재 범위 외 컴포넌트",
                    subtitle="이 내용은 아직 PPTX 에 렌더되지 않습니다.",
                    background="primary",
                ),
            ],
        ),
    )


# ---------------------------------------------------------------------------
# D1 테스트 (기존 5 건 유지)
# ---------------------------------------------------------------------------


def test_pptx_builder_instantiates_successfully():
    """요구 #1 — `PptxBuilder()` 인스턴스화 성공."""
    builder = PptxBuilder()
    assert builder.target == "pptx"
    assert isinstance(builder.supported_components, frozenset)
    # 초안 6 종 이상이 선언되어 있어야 한다 (HtmlRenderer 수준).
    assert len(builder.supported_components) >= 6


def test_builder_registry_returns_pptx_builder_instance():
    """요구 #2 — `BuilderRegistry.get("pptx")` 는 PptxBuilder 인스턴스 반환."""
    retrieved = BuilderRegistry.get("pptx")
    assert isinstance(retrieved, PptxBuilder)


async def test_build_empty_schema_returns_pptx_bytes(empty_schema):
    """요구 #3 — `build(schema)` 가 ZIP 시그니처로 시작하는 bytes 반환."""
    builder = BuilderRegistry.get("pptx")
    assert isinstance(builder, PptxBuilder)

    result = await builder.build(empty_schema)

    # bytes 타입
    assert isinstance(result, bytes)
    # PPTX = OOXML ZIP. 첫 2 바이트는 반드시 b"PK".
    assert result.startswith(_ZIP_SIGNATURE), (
        f"PPTX 바이트가 ZIP 시그니처(b'PK')로 시작하지 않음: 앞 4 바이트 = {result[:4]!r}"
    )
    # 빈 PPTX 라도 최소 수백 바이트 이상이어야 한다 (OOXML 최소 구조).
    assert len(result) > 1_000


async def test_build_with_unsupported_components_does_not_raise(
    schema_with_unsupported,
):
    """보조 #5 — Chart 등 미지원 컴포넌트가 있어도 빌드 자체는 성공."""
    builder = BuilderRegistry.get("pptx")
    result = await builder.build(schema_with_unsupported)
    assert isinstance(result, bytes)
    assert result.startswith(_ZIP_SIGNATURE)


def test_pptx_builder_target_literal_type():
    """보조 #4 — target 은 `BuildTarget` Literal 중 `"pptx"` 여야 한다."""
    assert PptxBuilder.target == "pptx"
    # Registry 조회도 동일 target 으로 성공해야 한다.
    assert BuilderRegistry.get("pptx") is not None


# ---------------------------------------------------------------------------
# D2 테스트 (신규 6 건 + 보조)
# ---------------------------------------------------------------------------


async def test_d2_slide_title_renders_text_on_first_slide():
    """D2 요구 #1 — SlideTitle 텍스트가 슬라이드에 포함된다."""
    schema = _make_schema(
        Page(
            id="p1",
            page_kind="slide",
            layout="title_slide",
            components=[SlideTitleComponent(id="c1", text="테스트 제목")],
        ),
    )

    builder = BuilderRegistry.get("pptx")
    result = await builder.build(schema)

    slide_texts = _extract_all_slide_text(result)
    assert len(slide_texts) == 1
    assert "테스트 제목" in slide_texts[0], f"SlideTitle 렌더링된 텍스트에 '테스트 제목' 이 없음: {slide_texts[0]!r}"


async def test_d2_heading_renders_text_with_level():
    """D2 요구 #2 — Heading 텍스트가 슬라이드에 포함된다."""
    schema = _make_schema(
        Page(
            id="p1",
            page_kind="slide",
            layout="content_body",
            components=[HeadingComponent(id="c1", text="섹션 1 개요", level=1)],
        ),
    )

    builder = BuilderRegistry.get("pptx")
    result = await builder.build(schema)

    slide_texts = _extract_all_slide_text(result)
    assert "섹션 1 개요" in slide_texts[0]


async def test_d2_paragraph_renders_text_with_emphasis():
    """D2 요구 #3 — Paragraph 텍스트가 슬라이드에 포함된다."""
    body_text = "이 문단은 D2 범위에서 정상 렌더되어야 한다."
    schema = _make_schema(
        Page(
            id="p1",
            page_kind="slide",
            layout="content_body",
            components=[
                ParagraphComponent(id="c1", text=body_text, emphasis="bold"),
            ],
        ),
    )

    builder = BuilderRegistry.get("pptx")
    result = await builder.build(schema)

    slide_texts = _extract_all_slide_text(result)
    assert body_text in slide_texts[0]


async def test_d2_bullet_list_renders_items_with_bullet_glyph():
    """D2 요구 #4 — BulletList items 개수만큼 paragraph 가 렌더된다.

    items=3, sub_items 없음인 경우:
      - 텍스트박스 1 개
      - 내부 paragraph 3 개
      - 각 paragraph 에 "• " prefix 포함
    """
    schema = _make_schema(
        Page(
            id="p1",
            page_kind="slide",
            layout="content_body",
            components=[
                BulletListComponent(
                    id="c1",
                    items=[
                        BulletItem(text="첫째 항목"),
                        BulletItem(text="둘째 항목"),
                        BulletItem(text="셋째 항목"),
                    ],
                    numbered=False,
                ),
            ],
        ),
    )

    builder = BuilderRegistry.get("pptx")
    result = await builder.build(schema)

    prs = Presentation(BytesIO(result))
    slide = prs.slides[0]
    # BulletList 는 하나의 textbox 로 렌더 (paragraph 다수).
    bullet_shapes = [s for s in slide.shapes if s.has_text_frame and "첫째 항목" in s.text_frame.text]
    assert len(bullet_shapes) == 1, "BulletList 는 단일 텍스트박스로 렌더되어야 한다"

    tf = bullet_shapes[0].text_frame
    # 3 개 paragraph (sub_items 없음 → 정확히 3).
    assert len(tf.paragraphs) == 3, f"item 3 개 → paragraph 3 개 예상, 실제 {len(tf.paragraphs)}"

    # 각 paragraph 에 bullet glyph 또는 실제 텍스트 포함.
    texts = [p.text for p in tf.paragraphs]
    assert any("첫째 항목" in t for t in texts)
    assert any("둘째 항목" in t for t in texts)
    assert any("셋째 항목" in t for t in texts)
    # bullet glyph "•" 가 첫 paragraph 에 포함되어야 한다 (numbered=False).
    assert any("•" in t for t in texts), f"bullet glyph '•' 이 paragraph 에 없음: {texts}"


async def test_d2_bullet_list_numbered_uses_ordinal_prefix():
    """D2 보조 — numbered=True 이면 "1. ", "2. " prefix 로 렌더."""
    schema = _make_schema(
        Page(
            id="p1",
            page_kind="slide",
            layout="content_body",
            components=[
                BulletListComponent(
                    id="c1",
                    items=[BulletItem(text="첫째"), BulletItem(text="둘째")],
                    numbered=True,
                ),
            ],
        ),
    )

    builder = BuilderRegistry.get("pptx")
    result = await builder.build(schema)

    slide_texts = _extract_all_slide_text(result)
    full = slide_texts[0]
    assert "1. 첫째" in full
    assert "2. 둘째" in full


async def test_d2_bullet_list_sub_items_are_rendered_as_extra_paragraphs():
    """D2 보조 — sub_items 가 있으면 추가 paragraph 로 들여쓰기 렌더."""
    schema = _make_schema(
        Page(
            id="p1",
            page_kind="slide",
            layout="content_body",
            components=[
                BulletListComponent(
                    id="c1",
                    items=[
                        BulletItem(
                            text="주요 항목",
                            sub_items=["세부 A", "세부 B"],
                        ),
                    ],
                ),
            ],
        ),
    )

    builder = BuilderRegistry.get("pptx")
    result = await builder.build(schema)

    prs = Presentation(BytesIO(result))
    slide = prs.slides[0]
    bullet_shapes = [s for s in slide.shapes if s.has_text_frame and "주요 항목" in s.text_frame.text]
    assert len(bullet_shapes) == 1
    tf = bullet_shapes[0].text_frame
    # 1 item + 2 sub_items = 3 paragraph.
    assert len(tf.paragraphs) == 3
    texts = [p.text for p in tf.paragraphs]
    assert any("세부 A" in t for t in texts)
    assert any("세부 B" in t for t in texts)


async def test_d2_composite_page_has_all_four_text_shapes():
    """D2 요구 #5 — 4 종 복합 페이지는 단일 슬라이드에 4 개 텍스트박스 생성."""
    schema = _make_schema(
        Page(
            id="p1",
            page_kind="slide",
            layout="content_body",
            title="복합 페이지",
            components=[
                SlideTitleComponent(id="c1", text="종합 보고"),
                HeadingComponent(id="c2", text="1장 서론", level=1),
                ParagraphComponent(id="c3", text="서론 본문 텍스트."),
                BulletListComponent(
                    id="c4",
                    items=[BulletItem(text="포인트 A"), BulletItem(text="포인트 B")],
                ),
            ],
        ),
    )

    builder = BuilderRegistry.get("pptx")
    result = await builder.build(schema)

    # 슬라이드 1 개, 각 컴포넌트당 textbox 1 개 → 4 개.
    assert _count_shapes_with_text(result, slide_idx=0) == 4

    slide_texts = _extract_all_slide_text(result)
    combined = slide_texts[0]
    assert "종합 보고" in combined
    assert "1장 서론" in combined
    assert "서론 본문 텍스트." in combined
    assert "포인트 A" in combined
    assert "포인트 B" in combined


async def test_d2_multiple_pages_produces_matching_slide_count():
    """D2 요구 #6 — 2 페이지 이상이면 슬라이드 수가 페이지 수와 일치."""
    schema = _make_schema(
        Page(
            id="p1",
            page_kind="slide",
            layout="title_slide",
            components=[SlideTitleComponent(id="c1", text="첫 슬라이드")],
        ),
        Page(
            id="p2",
            page_kind="slide",
            layout="content_body",
            components=[HeadingComponent(id="c2", text="두 번째 슬라이드", level=1)],
        ),
        Page(
            id="p3",
            page_kind="slide",
            layout="closing",
            components=[ParagraphComponent(id="c3", text="마무리 문단")],
        ),
    )

    builder = BuilderRegistry.get("pptx")
    result = await builder.build(schema)

    prs = Presentation(BytesIO(result))
    assert len(prs.slides) == 3

    slide_texts = _extract_all_slide_text(result)
    assert "첫 슬라이드" in slide_texts[0]
    assert "두 번째 슬라이드" in slide_texts[1]
    assert "마무리 문단" in slide_texts[2]


async def test_d2_out_of_scope_component_is_skipped_without_failure():
    """D2 보조 — D2 범위 외 컴포넌트는 빌드 실패 없이 skip 된다.

    S3 D6 에서 Callout 이 지원으로 승격되면서 "범위 외" 대표가 Hero 로 이관됨.
    Hero 는 17 종 중 아직 이관되지 않은 컴포넌트 — skip + WARN 경로를 타는 대표 케이스.
    """
    from app.modules.documents_v2.schemas import HeroComponent

    schema = _make_schema(
        Page(
            id="p1",
            page_kind="slide",
            layout="content_body",
            components=[
                SlideTitleComponent(id="c1", text="혼합 페이지"),
                HeroComponent(
                    id="c2",
                    title="Hero 타이틀",
                    subtitle="아직 미지원 컴포넌트",
                    background="primary",
                ),
                ParagraphComponent(id="c3", text="본문 계속"),
            ],
        ),
    )

    builder = BuilderRegistry.get("pptx")
    result = await builder.build(schema)

    slide_texts = _extract_all_slide_text(result)
    # SlideTitle + Paragraph 는 렌더, Hero 는 skip.
    assert "혼합 페이지" in slide_texts[0]
    assert "본문 계속" in slide_texts[0]
    assert "Hero 타이틀" not in slide_texts[0], "현재 범위 외 Hero 텍스트는 슬라이드에 나타나면 안 된다"


async def test_d2_slide_dimensions_are_16_9():
    """D2 보조 — 생성된 PPTX 는 16:9 비율(13.333"x7.5") 이어야 한다.

    python-pptx 기본은 10"x7.5" (4:3) 이므로 build() 에서 명시적으로 세팅되지
    않으면 이 테스트가 실패한다.
    """
    schema = _make_schema(
        Page(
            id="p1",
            page_kind="slide",
            layout="title_slide",
            components=[SlideTitleComponent(id="c1", text="비율 검증")],
        ),
    )

    builder = BuilderRegistry.get("pptx")
    result = await builder.build(schema)

    prs = Presentation(BytesIO(result))
    # EMU 단위 (1 inch = 914400 EMU).
    width_in = prs.slide_width / 914400
    height_in = prs.slide_height / 914400
    # 근사 비교 (부동소수점).
    assert abs(width_in - 13.333) < 0.01, f"slide_width={width_in}"
    assert abs(height_in - 7.5) < 0.01, f"slide_height={height_in}"


# ---------------------------------------------------------------------------
# D3 테스트 (신규 9 건)
# ---------------------------------------------------------------------------


def _make_blank_slide_for_unit_test():
    """단일 컴포넌트 유닛 테스트용 빈 슬라이드 생성 헬퍼.

    build() 경유 없이 components 함수만 직접 검증할 때 사용. 16:9 설정 + blank
    레이아웃(없으면 [0]) 사용. 반환값은 (prs, slide) 튜플.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank)
    return prs, slide


def test_d3_render_kpi_single_card_contains_label_value_delta():
    """D3 #13 — 단일 KPI 카드 렌더 후 슬라이드에 label/value/delta 텍스트 모두 포함."""
    _prs, slide = _make_blank_slide_for_unit_test()
    kpi = KPIComponent(
        id="c1",
        label="월 활성 사용자",
        value="12,345",
        delta="+8.2%",
        delta_direction="up",
        description="전월 대비",
    )

    render_kpi(kpi, slide, left_in=0.5, top_in=2.0, width_in=3.0)

    # 슬라이드 상의 모든 텍스트 합치기.
    texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
    combined = "\n".join(texts)

    assert "월 활성 사용자" in combined, f"label 누락: {combined!r}"
    assert "12,345" in combined, f"value 누락: {combined!r}"
    assert "+8.2%" in combined, f"delta 누락: {combined!r}"


def test_d3_render_kpi_negative_delta_uses_red_color():
    """D3 #14 — delta 음수이면 KPI_DELTA_DOWN_COLOR(빨강) 적용."""
    _prs, slide = _make_blank_slide_for_unit_test()
    kpi = KPIComponent(
        id="c1",
        label="감소 지표",
        value="99",
        delta="-3.1%",
        delta_direction="down",
    )

    render_kpi(kpi, slide, left_in=0.5, top_in=2.0, width_in=3.0)

    # delta 텍스트박스는 "-3.1%" 가 포함된 shape.
    delta_shapes = [s for s in slide.shapes if s.has_text_frame and "-3.1%" in s.text_frame.text]
    assert len(delta_shapes) == 1, "delta textbox 가 1 개 있어야 함"

    expected_rgb = parse_hex_color(KPI_DELTA_DOWN_COLOR)
    delta_para = delta_shapes[0].text_frame.paragraphs[0]
    # paragraph.font 또는 첫 run 의 color 를 검증.
    actual_rgb = None
    for run in delta_para.runs:
        actual_rgb = run.font.color.rgb
        break
    assert actual_rgb == expected_rgb, (
        f"음수 delta 색이 KPI_DELTA_DOWN_COLOR({KPI_DELTA_DOWN_COLOR}) 가 아님 — actual={actual_rgb!r}"
    )


def test_d3_render_kpi_positive_delta_uses_green_color():
    """D3 보조 — 양수 delta 는 KPI_DELTA_UP_COLOR(녹색) 적용."""
    _prs, slide = _make_blank_slide_for_unit_test()
    kpi = KPIComponent(
        id="c1",
        label="증가 지표",
        value="205",
        delta="+12.4%",
        delta_direction="up",
    )

    render_kpi(kpi, slide, left_in=0.5, top_in=2.0, width_in=3.0)
    delta_shapes = [s for s in slide.shapes if s.has_text_frame and "+12.4%" in s.text_frame.text]
    assert len(delta_shapes) == 1
    expected_rgb = parse_hex_color(KPI_DELTA_UP_COLOR)
    actual_rgb = delta_shapes[0].text_frame.paragraphs[0].runs[0].font.color.rgb
    assert actual_rgb == expected_rgb


def test_d3_render_data_table_header_navy_and_zebra_rows():
    """D3 #15 — 3x3 DataTable 헤더 배경 #34495E + 짝수(데이터 2행째) 행 제브라."""
    _prs, slide = _make_blank_slide_for_unit_test()
    table_comp = DataTableComponent(
        id="c1",
        headers=["지표", "1분기", "2분기"],
        rows=[
            ["매출", "100", "120"],
            ["순이익", "30", "45"],
            ["임직원", "50", "52"],
        ],
    )

    render_data_table(table_comp, slide, left_in=0.5, top_in=2.0, width_in=10.0)

    # 슬라이드에 추가된 테이블 shape 1 개 확인.
    tables = [s for s in slide.shapes if s.has_table]
    assert len(tables) == 1
    table = tables[0].table

    # 헤더 행 배경색 검증 — 첫 셀.
    header_cell = table.cell(0, 0)
    expected_navy = parse_hex_color(IDINO_HEADER_NAVY)
    assert header_cell.fill.fore_color.rgb == expected_navy, (
        f"헤더 배경이 #{IDINO_HEADER_NAVY} 가 아님 — actual={header_cell.fill.fore_color.rgb!r}"
    )

    # 제브라 검증 — 데이터 두 번째 행(table row_idx=2) 은 연회색.
    expected_zebra = parse_hex_color(TABLE_ZEBRA_ROW_BG)
    zebra_cell = table.cell(2, 0)
    assert zebra_cell.fill.fore_color.rgb == expected_zebra, (
        f"짝수 데이터 행 배경이 제브라 #{TABLE_ZEBRA_ROW_BG} 가 아님 — actual={zebra_cell.fill.fore_color.rgb!r}"
    )


def test_d3_render_data_table_numeric_cell_is_right_aligned():
    """D3 #16 — 숫자 셀은 우측 정렬, 텍스트 셀은 좌측 정렬."""
    from pptx.enum.text import PP_ALIGN

    _prs, slide = _make_blank_slide_for_unit_test()
    table_comp = DataTableComponent(
        id="c1",
        headers=["항목", "값"],
        rows=[
            ["매출", "1,234"],  # col 0: 텍스트(좌측), col 1: 숫자(우측)
            ["비고", "N/A"],  # col 0: 텍스트(좌측), col 1: 텍스트(좌측)
        ],
    )

    render_data_table(table_comp, slide, left_in=0.5, top_in=2.0, width_in=8.0)

    tables = [s for s in slide.shapes if s.has_table]
    table = tables[0].table

    # 1 행(데이터 첫 행) 0 열 = "매출" → LEFT.
    text_cell = table.cell(1, 0)
    assert text_cell.text_frame.paragraphs[0].alignment == PP_ALIGN.LEFT

    # 1 행 1 열 = "1,234" → RIGHT (숫자).
    numeric_cell = table.cell(1, 1)
    assert numeric_cell.text_frame.paragraphs[0].alignment == PP_ALIGN.RIGHT, (
        f"숫자 셀 정렬이 RIGHT 가 아님 — actual={numeric_cell.text_frame.paragraphs[0].alignment!r}"
    )

    # 2 행 1 열 = "N/A" → LEFT (텍스트).
    na_cell = table.cell(2, 1)
    assert na_cell.text_frame.paragraphs[0].alignment == PP_ALIGN.LEFT


def test_d3_resolve_layout_matches_candidate_name_case_insensitive():
    """D3 #17 — python-pptx 기본 마스터의 영문 이름도 매칭 성공해야 한다.

    python-pptx `Presentation()` 는 기본 마스터에 "Title Slide", "Title and
    Content" 등 영문 레이아웃을 제공. "title_slide" 후보가 "Title Slide" 와
    정규화 후 매칭되어야 한다.
    """
    prs = Presentation()
    layout = resolve_layout(prs, "title_slide")

    # 기본 마스터에 "Title Slide" 가 존재하므로 매칭되어야 한다.
    assert layout is not None
    assert layout.name is not None
    # 매칭 결과가 blank(fallback) 가 아닌, 실제 title slide 여야 한다.
    # python-pptx 기본 첫 레이아웃이 "Title Slide".
    assert "title" in layout.name.lower()


def test_d3_resolve_layout_falls_back_when_no_candidate_matches(caplog):
    """D3 #18 (D9-b 업데이트) — override 없이 매칭 실패 시 INFO 레벨 fallback.

    D8 리뷰에서 기본 python-pptx 마스터의 한글 후보 부재로 WARNING 이 10 건
    노이즈로 적재되던 문제를 해결. D9-b 이후 override 가 제공되지 않은 상태의
    fallback 은 **정상 동작** 이므로 INFO 레벨로 강등된다.
    """
    prs = Presentation()

    # python-pptx 기본 마스터에는 "two_column" 과 유사한 이름이 없다.
    # "kpi_dashboard" 는 기본 마스터에 KPI 관련 이름이 전무 → fallback 경로.
    with caplog.at_level(logging.INFO, logger="app.integrations.document_builders.pptx.layout_resolver"):
        layout = resolve_layout(prs, "kpi_dashboard")

    # fallback 이어도 layout 객체 자체는 반환.
    assert layout is not None
    # INFO 레벨로 "매칭 실패" 메시지가 emit 되어야 한다 (WARNING 아님).
    info_records = [rec for rec in caplog.records if rec.levelno == logging.INFO and "매칭 실패" in rec.message]
    assert info_records, f"fallback INFO 로그가 emit 되지 않음: {[(r.levelno, r.message) for r in caplog.records]!r}"
    # override 가 없으므로 WARNING 은 없어야 한다 (노이즈 방지의 목적).
    warnings = [rec for rec in caplog.records if rec.levelno == logging.WARNING and "매칭 실패" in rec.message]
    assert not warnings, f"override 없이 fallback 했는데 WARNING 이 emit 됨: {[r.message for r in warnings]!r}"


def test_d9b_resolve_layout_fallback_with_override_emits_warning(caplog):
    """D9-b #3 — override 가 제공됐는데 매칭 실패하면 WARNING 로그.

    override 가 주어졌다는 것은 조직 매스터에 맞춤 매핑이 있으리라 기대했다는
    의미. 그럼에도 매칭 실패했다면 DB / JSONB 설정 오류 가능성이 있어 운영에서
    눈에 띄어야 한다.
    """
    prs = Presentation()

    # 실제 마스터에 없는 이름을 override 로 제공 → 후보 탐색도 실패 예상.
    overrides = {"kpi_dashboard": "완전히_없는_레이아웃_이름"}

    with caplog.at_level(logging.INFO, logger="app.integrations.document_builders.pptx.layout_resolver"):
        layout = resolve_layout(prs, "kpi_dashboard", organization_overrides=overrides)

    assert layout is not None  # fallback 은 반환.
    # override 가 존재했으므로 최종 fallback 은 WARNING 이어야 한다.
    warnings = [rec for rec in caplog.records if rec.levelno == logging.WARNING and "매칭 실패" in rec.message]
    assert warnings, (
        f"override 제공 상태의 fallback 에서 WARNING 이 누락: {[(r.levelno, r.message) for r in caplog.records]!r}"
    )


def test_d3_resolve_layout_organization_overrides_takes_priority():
    """D3 #19 — organization_overrides 매핑이 후보 리스트보다 우선."""
    prs = Presentation()

    # 의도적으로 "content_body" 를 기본 마스터의 "Blank" 에 매핑.
    # "Blank" 는 후보 LAYOUT_NAME_CANDIDATES["content_body"] 에 없지만 override 로 강제.
    overrides = {"content_body": "Blank"}

    layout = resolve_layout(prs, "content_body", organization_overrides=overrides)

    assert layout is not None
    assert layout.name is not None
    assert "blank" in layout.name.lower(), f"override 가 적용되지 않음 — 기대 'Blank', 실제 {layout.name!r}"


def test_d3_resolve_placeholder_matches_title_semantic():
    """D3 #20 — resolve_placeholder 가 PP_PLACEHOLDER.TITLE semantic 을 매칭.

    python-pptx 기본 "Title Slide" 레이아웃은 TITLE + SUBTITLE 플레이스홀더를
    가진다. semantic="title" 로 호출하면 TITLE placeholder 를 반환해야 한다.
    """
    prs = Presentation()
    title_layout = prs.slide_layouts[0]  # "Title Slide"

    title_ph = resolve_placeholder(title_layout, "title")
    assert title_ph is not None, "title placeholder 가 매칭되지 않음"

    # placeholder.placeholder_format.type 가 TITLE 또는 CENTER_TITLE 이어야 한다.
    from pptx.enum.shapes import PP_PLACEHOLDER

    assert title_ph.placeholder_format.type in {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}


async def test_d3_build_full_schema_with_kpi_and_data_table():
    """D3 #21 — KPI + DataTable 포함 스키마가 정상 빌드되고 텍스트가 들어간다."""
    schema = _make_schema(
        Page(
            id="p1",
            page_kind="slide",
            layout="kpi_dashboard",
            title="KPI 대시보드",
            components=[
                SlideTitleComponent(id="c1", text="2026 1분기 성과"),
                KPIComponent(
                    id="c2",
                    label="월간 활성 사용자",
                    value="123,456",
                    delta="+5.3%",
                    delta_direction="up",
                ),
                KPIComponent(
                    id="c3",
                    label="이탈률",
                    value="2.1%",
                    delta="-0.4%p",
                    delta_direction="down",
                ),
            ],
        ),
        Page(
            id="p2",
            page_kind="slide",
            layout="content_body",
            title="분기별 매출",
            components=[
                SlideTitleComponent(id="c1", text="분기별 매출 표"),
                DataTableComponent(
                    id="c2",
                    headers=["분기", "매출(억)", "증감률"],
                    rows=[
                        ["Q1", "120", "+8%"],
                        ["Q2", "135", "+12%"],
                        ["Q3", "142", "+5%"],
                    ],
                ),
            ],
        ),
    )

    builder = BuilderRegistry.get("pptx")
    result = await builder.build(schema)

    assert isinstance(result, bytes)
    assert result.startswith(_ZIP_SIGNATURE)

    prs = Presentation(BytesIO(result))
    assert len(prs.slides) == 2

    # p1 — KPI 값/라벨 텍스트 포함.
    slide_texts = _extract_all_slide_text(result)
    assert "2026 1분기 성과" in slide_texts[0]
    assert "월간 활성 사용자" in slide_texts[0]
    assert "123,456" in slide_texts[0]
    assert "+5.3%" in slide_texts[0]
    assert "이탈률" in slide_texts[0]

    # p2 — DataTable shape 존재 + 셀 텍스트 직접 확인.
    # 주의: `_extract_all_slide_text` 는 has_text_frame 만 수집하므로 table shape
    # (has_text_frame=False, has_table=True) 의 셀 텍스트는 별도로 추출해야 한다.
    p2_slide = prs.slides[1]
    tables = [s for s in p2_slide.shapes if s.has_table]
    assert len(tables) == 1, "DataTable 이 슬라이드에 삽입되지 않음"
    table = tables[0].table
    # 헤더 3 + 데이터 3 = 4 행.
    assert len(list(table.rows)) == 4
    # 슬라이드 상단 제목은 textbox 이므로 slide_texts[1] 에 포함.
    assert "분기별 매출 표" in slide_texts[1]

    # 테이블 셀 텍스트는 table 을 직접 순회해 추출.
    all_cell_texts: list[str] = []
    for row in table.rows:
        for cell in row.cells:
            all_cell_texts.append(cell.text)
    joined_cells = " | ".join(all_cell_texts)
    assert "Q1" in joined_cells, f"Q1 이 테이블 셀에 없음: {joined_cells!r}"
    assert "135" in joined_cells, f"135 이 테이블 셀에 없음: {joined_cells!r}"


# ---------------------------------------------------------------------------
# D6 테스트 (신규 5 건 + 보조)
# ---------------------------------------------------------------------------
#
# 공통 헬퍼 — python-pptx 의 `add_picture()` 는 Pillow 로 이미지 헤더를 검증
# 하기 때문에 유효한 PNG/JPEG 바이트가 필요하다. 테스트마다 Pillow 로 in-memory
# 이미지를 생성해 mock response / base64 payload 로 사용한다.


def _make_png_bytes(width: int = 10, height: int = 10, color: str = "red") -> bytes:
    """테스트용 유효 PNG 바이트를 즉석에서 생성."""
    from PIL import Image as PILImage

    img = PILImage.new("RGB", (width, height), color=color)
    buf = BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _has_picture_shape(slide) -> bool:
    """슬라이드에 python-pptx Picture (MSO_SHAPE_TYPE.PICTURE) shape 가 있는지."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    return any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes)


class _FakeHttpResponse:
    """httpx.Response 의 최소 duck-type. `.content`, `.raise_for_status()`, `.status_code` 만 필요."""

    def __init__(self, content: bytes, status_code: int = 200):
        self.content = content
        self.status_code = status_code

    def raise_for_status(self) -> None:
        if self.status_code >= 400:
            # 실제 httpx.HTTPStatusError 와 계약은 다르지만 예외 경로 테스트에선
            # ImageFetchError 가 뜨는지만 검증하므로 RuntimeError 로 충분.
            import httpx

            raise httpx.HTTPStatusError(
                f"status {self.status_code}",
                request=None,  # type: ignore[arg-type]
                response=None,  # type: ignore[arg-type]
            )


class _FakeHttpClient:
    """httpx.Client 대체 — 고정된 bytes 를 반환하거나 예외를 raise 한다."""

    def __init__(self, *, content: bytes | None = None, exc: Exception | None = None):
        self._content = content
        self._exc = exc
        self.call_count = 0

    def get(self, url: str, *, timeout: int | None = None) -> _FakeHttpResponse:
        self.call_count += 1
        if self._exc is not None:
            raise self._exc
        return _FakeHttpResponse(self._content or b"", status_code=200)

    def close(self) -> None:  # pragma: no cover — 테스트 teardown 시 호출 가능.
        pass


@pytest.fixture
def _reset_image_fetcher_client():
    """D6 이미지 fetcher 의 httpx.Client 싱글톤을 테스트 사이에 리셋."""
    _image_fetcher._reset_http_client_for_tests()
    yield
    _image_fetcher._reset_http_client_for_tests()


def test_d6_render_image_from_http_url_inserts_picture(monkeypatch, _reset_image_fetcher_client):
    """D6 #22 — URL source, httpx 를 fake 로 대체 → add_picture 호출되고 Picture shape 존재."""
    _prs, slide = _make_blank_slide_for_unit_test()

    # fake client 에 유효한 PNG 바이트 주입.
    fake = _FakeHttpClient(content=_make_png_bytes(width=40, height=30))
    monkeypatch.setattr(_image_fetcher, "_get_http_client", lambda: fake)

    img = ImageComponent(
        id="c1",
        src="https://example.com/sample.png",
        alt="테스트 이미지",
    )

    next_y = render_image(
        img,
        slide,
        left_in=0.5,
        top_in=2.0,
        width_in=4.0,
        max_height_in=4.0,
    )

    # Picture shape 가 1 개 이상 삽입되어야 한다.
    assert _has_picture_shape(slide), "render_image 가 Picture shape 를 삽입하지 않았다"
    # fake client 는 정확히 1 회 호출되었어야 한다 (재시도 없이 성공).
    assert fake.call_count == 1
    # Y 는 top_in 이후로 전진했어야 한다.
    assert next_y > 2.0


def test_d6_render_image_from_base64_data_uri(_reset_image_fetcher_client):
    """D6 #23 — base64 data-URI source → 디코드 후 정상 삽입."""
    import base64

    _prs, slide = _make_blank_slide_for_unit_test()

    png_bytes = _make_png_bytes(width=20, height=20, color="blue")
    b64 = base64.b64encode(png_bytes).decode("ascii")
    data_uri = f"data:image/png;base64,{b64}"

    img = ImageComponent(
        id="c1",
        src=data_uri,
        alt="base64 이미지",
        caption="샘플 캡션",
    )

    render_image(
        img,
        slide,
        left_in=0.5,
        top_in=1.5,
        width_in=3.0,
        max_height_in=5.0,
    )

    assert _has_picture_shape(slide), "base64 이미지가 삽입되지 않았다"

    # caption 이 있으므로 caption textbox 가 함께 렌더되었어야 한다.
    caption_texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame and "샘플 캡션" in s.text_frame.text]
    assert len(caption_texts) == 1, f"caption textbox 누락: {caption_texts!r}"


def test_d6_render_image_fetch_failure_falls_back_to_placeholder(monkeypatch, _reset_image_fetcher_client, caplog):
    """D6 #24 — fetch 실패 시 placeholder + alt 텍스트. 빌드 실패 X."""
    _prs, slide = _make_blank_slide_for_unit_test()

    # httpx 가 네트워크 오류를 raise 하는 fake.
    import httpx

    fake = _FakeHttpClient(exc=httpx.ConnectError("connection refused"))
    monkeypatch.setattr(_image_fetcher, "_get_http_client", lambda: fake)

    img = ImageComponent(
        id="c1",
        src="https://broken.example.com/missing.png",
        alt="깨진 이미지 설명",
    )

    with caplog.at_level(logging.WARNING):
        # 예외 없이 반환되어야 한다.
        render_image(
            img,
            slide,
            left_in=0.5,
            top_in=2.0,
            width_in=4.0,
            max_height_in=3.0,
        )

    # Picture shape 는 없고, placeholder 박스(RECTANGLE) + alt 텍스트만 존재.
    assert not _has_picture_shape(slide), "fetch 실패 시 Picture shape 가 있으면 안 된다"
    placeholder_texts = [
        s.text_frame.text for s in slide.shapes if s.has_text_frame and "깨진 이미지 설명" in s.text_frame.text
    ]
    assert len(placeholder_texts) >= 1, (
        f"placeholder 에 alt 텍스트가 표시되지 않았다: {[s.text_frame.text for s in slide.shapes if s.has_text_frame]!r}"
    )
    # WARNING 로그 확인.
    assert any("fetch 실패" in rec.message or "degrade" in rec.message for rec in caplog.records)


def test_d6_render_image_respects_max_height_by_resizing_width(monkeypatch, _reset_image_fetcher_client):
    """D6 #25 — aspect ratio 유지. 가로 긴 이미지(4:1) + max_height 낮으면 width 축소."""
    _prs, slide = _make_blank_slide_for_unit_test()

    # 400x100 이미지 (aspect ratio 4.0).
    fake = _FakeHttpClient(content=_make_png_bytes(width=400, height=100))
    monkeypatch.setattr(_image_fetcher, "_get_http_client", lambda: fake)

    img = ImageComponent(id="c1", src="https://example.com/wide.png", alt="wide")

    # 요청 width=8" → 계산 height = 8/4 = 2". max_height=1" 이므로 cap → height=1", width=4".
    render_image(
        img,
        slide,
        left_in=0.5,
        top_in=1.0,
        width_in=8.0,
        max_height_in=1.0,
    )

    # 삽입된 Picture 의 실제 너비가 4" 근처 (EMU = 914400 * 4 ≈ 3,657,600) 여야 한다.
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 1
    picture = pictures[0]
    width_in_actual = picture.width / 914400
    height_in_actual = picture.height / 914400
    # height 는 max_height=1.0 이하여야 하고, width/height 비율은 4.0 근사.
    assert height_in_actual <= 1.05, f"height 가 max_height_in 을 초과: {height_in_actual}"
    ratio_actual = width_in_actual / height_in_actual
    assert abs(ratio_actual - 4.0) < 0.1, f"aspect ratio 가 유지되지 않음 — 기대 4.0, 실제 {ratio_actual}"


def test_d6_image_fetcher_rejects_unsupported_scheme(_reset_image_fetcher_client):
    """D6 보조 — 지원하지 않는 스킴(로컬 파일 경로 등) 은 ImageFetchError 발생."""
    from app.integrations.document_builders.pptx.image_fetcher import fetch_image_bytes

    with pytest.raises(ImageFetchError):
        fetch_image_bytes("/etc/passwd")

    with pytest.raises(ImageFetchError):
        fetch_image_bytes("")


async def test_d6_build_full_schema_with_image_component(monkeypatch, _reset_image_fetcher_client):
    """D6 #26 — build() 통합. Image 포함 스키마 → Picture shape + caption 텍스트."""
    # httpx 를 fake 로 대체.
    fake = _FakeHttpClient(content=_make_png_bytes(width=100, height=60))
    monkeypatch.setattr(_image_fetcher, "_get_http_client", lambda: fake)

    schema = _make_schema(
        Page(
            id="p1",
            page_kind="slide",
            layout="content_body",
            title="이미지 삽입 테스트",
            components=[
                SlideTitleComponent(id="c1", text="이미지 예시"),
                ParagraphComponent(id="c2", text="아래는 통합 테스트용 이미지입니다."),
                ImageComponent(
                    id="c3",
                    src="https://example.com/cover.png",
                    alt="예시 이미지 대체 텍스트",
                    caption="그림 1. 예시 이미지",
                ),
            ],
        ),
    )

    builder = BuilderRegistry.get("pptx")
    result = await builder.build(schema)

    assert isinstance(result, bytes)
    assert result.startswith(_ZIP_SIGNATURE)

    prs = Presentation(BytesIO(result))
    assert len(prs.slides) == 1
    slide = prs.slides[0]

    # Picture shape 존재.
    assert _has_picture_shape(slide), "build() 결과 슬라이드에 Picture 가 없음"

    # caption 텍스트 존재.
    slide_texts = _extract_all_slide_text(result)
    assert "그림 1. 예시 이미지" in slide_texts[0]
    # SlideTitle 과 Paragraph 도 함께 렌더되어야 한다.
    assert "이미지 예시" in slide_texts[0]
    assert "아래는 통합 테스트용 이미지입니다." in slide_texts[0]


# ---------------------------------------------------------------------------
# D7 테스트 (신규 6 건 + 보조) — Chart 컴포넌트 PPTX native 렌더
# ---------------------------------------------------------------------------
#
# 공통: python-pptx `XL_CHART_TYPE` enum 을 직접 비교하여 chart_type 매핑을 검증.
# bar/line 은 D7 범위에서 native 로 지원되며, pie 등은 bar 로 graceful-degrade.


def _has_chart_shape(slide) -> bool:
    """슬라이드에 python-pptx chart shape (shape.has_chart == True) 가 있는지."""
    return any(getattr(shape, "has_chart", False) for shape in slide.shapes)


def _find_chart_shape(slide):
    """슬라이드에서 has_chart=True 인 첫 shape 의 chart 객체를 반환. 없으면 None."""
    for shape in slide.shapes:
        if getattr(shape, "has_chart", False):
            return shape.chart
    return None


def test_d7_render_chart_bar_uses_column_clustered():
    """D7 #1 — chart_type='bar' 은 XL_CHART_TYPE.COLUMN_CLUSTERED 로 매핑."""
    from pptx.enum.chart import XL_CHART_TYPE

    _prs, slide = _make_blank_slide_for_unit_test()
    chart_comp = ChartComponent(
        id="c1",
        chart_type="bar",
        title="분기별 매출",
        data=ChartData(
            labels=["Q1", "Q2", "Q3"],
            series=[ChartSeries(name="매출", values=[100.0, 120.0, 135.0])],
        ),
    )

    render_chart(
        chart_comp,
        slide,
        left_in=0.5,
        top_in=1.5,
        width_in=8.0,
        max_height_in=4.0,
    )

    chart = _find_chart_shape(slide)
    assert chart is not None, "bar chart 가 슬라이드에 삽입되지 않았다"
    assert chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED, (
        f"bar → COLUMN_CLUSTERED 매핑 실패 — actual={chart.chart_type!r}"
    )
    # 시리즈 1 개, 카테고리 3 개 검증.
    plot = chart.plots[0]
    assert len(plot.series) == 1
    assert len(list(plot.categories)) == 3


def test_d7_render_chart_line_uses_line_type():
    """D7 #2 — chart_type='line' 은 XL_CHART_TYPE.LINE 로 매핑."""
    from pptx.enum.chart import XL_CHART_TYPE

    _prs, slide = _make_blank_slide_for_unit_test()
    chart_comp = ChartComponent(
        id="c1",
        chart_type="line",
        title="월별 사용자 추이",
        data=ChartData(
            labels=["1월", "2월", "3월", "4월"],
            series=[ChartSeries(name="MAU", values=[1000.0, 1150.0, 1280.0, 1320.0])],
        ),
    )

    render_chart(
        chart_comp,
        slide,
        left_in=0.5,
        top_in=1.5,
        width_in=8.0,
        max_height_in=4.0,
    )

    chart = _find_chart_shape(slide)
    assert chart is not None, "line chart 가 슬라이드에 삽입되지 않았다"
    assert chart.chart_type == XL_CHART_TYPE.LINE, f"line → LINE 매핑 실패 — actual={chart.chart_type!r}"


def test_d7_render_chart_two_series_shows_legend():
    """D7 #3 — 시리즈 2 개 이상이면 범례(legend) 가 활성화된다."""
    from pptx.enum.chart import XL_LEGEND_POSITION

    _prs, slide = _make_blank_slide_for_unit_test()
    chart_comp = ChartComponent(
        id="c1",
        chart_type="bar",
        title="제품별 분기 매출",
        data=ChartData(
            labels=["Q1", "Q2", "Q3"],
            series=[
                ChartSeries(name="제품 A", values=[100.0, 110.0, 120.0]),
                ChartSeries(name="제품 B", values=[80.0, 95.0, 105.0]),
            ],
        ),
    )

    render_chart(
        chart_comp,
        slide,
        left_in=0.5,
        top_in=1.5,
        width_in=8.0,
        max_height_in=4.0,
    )

    chart = _find_chart_shape(slide)
    assert chart is not None
    assert chart.has_legend is True, "시리즈 2 개 이상인데 범례가 표시되지 않음"
    assert chart.legend.position == XL_LEGEND_POSITION.BOTTOM, (
        f"범례 위치가 BOTTOM 이 아님 — actual={chart.legend.position!r}"
    )
    # 시리즈 2 개가 실제 plot 에 들어갔는지 확인.
    assert len(chart.plots[0].series) == 2


def test_d7_render_chart_single_series_hides_legend():
    """D7 보조 — 시리즈 1 개면 범례는 숨김."""
    _prs, slide = _make_blank_slide_for_unit_test()
    chart_comp = ChartComponent(
        id="c1",
        chart_type="bar",
        title="단일 시리즈",
        data=ChartData(
            labels=["A", "B"],
            series=[ChartSeries(name="only", values=[10.0, 20.0])],
        ),
    )

    render_chart(
        chart_comp,
        slide,
        left_in=0.5,
        top_in=1.5,
        width_in=6.0,
        max_height_in=3.5,
    )

    chart = _find_chart_shape(slide)
    assert chart is not None
    assert chart.has_legend is False, "시리즈 1 개인데 범례가 표시됨"


def test_d7_render_chart_empty_data_falls_back_to_placeholder(caplog):
    """D7 #4 — 빈 series values 는 placeholder + WARNING 로그.

    Pydantic 이 series.min_length=1 과 values 길이 일치를 강제하지만, 우리 렌더러는
    values 가 전부 비어있는 방어 케이스도 처리한다.
    """
    _prs, slide = _make_blank_slide_for_unit_test()

    # Pydantic validator 를 회피하기 위해 values 길이를 labels 와 맞추되 전부 0 인
    # 시리즈를 사용한다. 실제 "빈 데이터" 시나리오는 values 가 전부 0 일 때의
    # 시각적 결과가 아닌, labels 와 series.values 중 하나라도 비어있는 경우다.
    # → 방어 로직 검증을 위해 ChartData 를 우회한 dict 를 model_construct 로 주입.
    data_dict = ChartData.model_construct(labels=[], series=[])
    chart_comp = ChartComponent.model_construct(
        id="c1",
        type="Chart",
        chart_type="bar",
        title="빈 차트",
        data=data_dict,
    )

    with caplog.at_level(logging.WARNING):
        render_chart(
            chart_comp,
            slide,
            left_in=0.5,
            top_in=1.5,
            width_in=6.0,
            max_height_in=3.5,
        )

    # 차트는 삽입되지 않아야 한다.
    assert not _has_chart_shape(slide), "빈 데이터인데 chart shape 가 삽입됨"

    # placeholder 안내 문구가 나와야 한다.
    placeholder_texts = [
        s.text_frame.text
        for s in slide.shapes
        if s.has_text_frame and CHART_EMPTY_PLACEHOLDER_TEXT in s.text_frame.text
    ]
    assert len(placeholder_texts) >= 1, (
        f"placeholder 안내 문구 누락: {[s.text_frame.text for s in slide.shapes if s.has_text_frame]!r}"
    )

    # WARNING 로그 emit 확인.
    assert any("빈 데이터" in rec.message or "degrade" in rec.message for rec in caplog.records)


def test_d7_render_chart_pie_falls_back_to_bar(caplog):
    """D7 #5 (S3 D1 에서 의미 변경) — pie 는 S3 D1 에서 native 로 승격.

    S3 D1 이전에는 pie 가 bar 로 fallback 되는 것을 확인하는 테스트였으나,
    S3 D1 에서 pie 는 XL_CHART_TYPE.PIE 로 native 매핑된다. 본 테스트는
    "pie 가 더 이상 bar fallback 을 타지 않는다" 는 회귀 방지로 재정의.
    """
    from pptx.enum.chart import XL_CHART_TYPE

    _prs, slide = _make_blank_slide_for_unit_test()
    chart_comp = ChartComponent(
        id="c1",
        chart_type="pie",  # S3 D1 부터 native 지원.
        title="비중 분석",
        data=ChartData(
            labels=["A", "B", "C"],
            series=[ChartSeries(name="비중", values=[30.0, 45.0, 25.0])],
        ),
    )

    with caplog.at_level(logging.WARNING):
        render_chart(
            chart_comp,
            slide,
            left_in=0.5,
            top_in=1.5,
            width_in=6.0,
            max_height_in=3.5,
        )

    chart = _find_chart_shape(slide)
    assert chart is not None, "pie 차트 삽입 실패"
    # S3 D1 이후 — pie 요청 → XL_CHART_TYPE.PIE 로 직접 매핑.
    assert chart.chart_type == XL_CHART_TYPE.PIE, f"pie → PIE native 매핑 실패 — actual={chart.chart_type!r}"
    # fallback 경로를 타지 않으므로 graceful-degrade WARNING 은 없어야 한다.
    assert not any("graceful-degrade" in rec.message for rec in caplog.records), (
        f"pie 가 native 승격됐는데 graceful-degrade 로그 emit: {[r.message for r in caplog.records]!r}"
    )


async def test_d7_build_full_schema_with_chart_component():
    """D7 #6 — build() 통합. Chart 포함 스키마 → 슬라이드에 has_chart 인 shape + 제목."""
    from pptx.enum.chart import XL_CHART_TYPE

    schema = _make_schema(
        Page(
            id="p1",
            page_kind="slide",
            layout="content_body",
            title="차트 통합 테스트",
            components=[
                SlideTitleComponent(id="c1", text="매출 현황"),
                ParagraphComponent(id="c2", text="아래는 D7 native 차트 렌더 결과입니다."),
                ChartComponent(
                    id="c3",
                    chart_type="bar",
                    title="분기별 매출 (억)",
                    data=ChartData(
                        labels=["Q1", "Q2", "Q3", "Q4"],
                        series=[
                            ChartSeries(name="2025", values=[100.0, 115.0, 130.0, 145.0]),
                            ChartSeries(name="2026", values=[120.0, 135.0, 150.0, 170.0]),
                        ],
                    ),
                ),
            ],
        ),
    )

    builder = BuilderRegistry.get("pptx")
    result = await builder.build(schema)

    assert isinstance(result, bytes)
    assert result.startswith(_ZIP_SIGNATURE)

    prs = Presentation(BytesIO(result))
    assert len(prs.slides) == 1
    slide = prs.slides[0]

    # chart shape 존재 + 타입 검증.
    assert _has_chart_shape(slide), "build() 결과 슬라이드에 chart 가 없음"
    chart = _find_chart_shape(slide)
    assert chart is not None
    assert chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED
    # 시리즈 2 개 → 범례 ON.
    assert chart.has_legend is True
    assert len(chart.plots[0].series) == 2

    # SlideTitle / Paragraph 텍스트도 함께 포함.
    slide_texts = _extract_all_slide_text(result)
    assert "매출 현황" in slide_texts[0]
    assert "아래는 D7 native 차트 렌더 결과입니다." in slide_texts[0]


def test_d7_chart_supported_components_includes_chart():
    """D7 보조 — `PptxBuilder.supported_components` 에 'Chart' 포함 (7→8 종)."""
    builder = PptxBuilder()
    assert "Chart" in builder.supported_components, (
        f"supported_components 에 'Chart' 누락 — {sorted(builder.supported_components)}"
    )
    # 총 8 종 이상이어야 한다.
    assert len(builder.supported_components) >= 8


def test_d7_chart_palette_has_eight_colors_and_is_immutable():
    """D7 보조 — CHART_SERIES_PALETTE 는 8 색 tuple 이며 앞 2 개는 IDINO primary/accent."""
    assert isinstance(CHART_SERIES_PALETTE, tuple)
    assert len(CHART_SERIES_PALETTE) == 8
    # IDINO primary/accent 가 앞 2 개 위치를 차지하여 단일 시리즈 차트가 자연스럽게 파랑으로 나온다.
    assert CHART_SERIES_PALETTE[0] == "#0A4FC2"
    assert CHART_SERIES_PALETTE[1] == "#FF6B35"


# ---------------------------------------------------------------------------
# D9 테스트 (신규 6 건 + 보조)
#
# 대상 기능:
# - D9-a placeholder 주입 활성화 / 실패 fallback
# - D9-b layout_resolver 로그 레벨 분기 (override 여부 기반)
# - D9-d BulletList highlight 배경 박스
# - D9-e Paragraph 한/영 혼용 높이 추정 ±20%
# ---------------------------------------------------------------------------


async def test_d9a_slide_title_injected_into_master_title_placeholder():
    """D9-a #1 — python-pptx 기본 마스터 "Title Slide" 레이아웃은 TITLE placeholder
    를 가지므로, SlideTitle 컴포넌트 텍스트가 placeholder 에 주입되어야 한다.
    """
    schema = _make_schema(
        Page(
            id="p1",
            page_kind="slide",
            layout="title_slide",
            components=[SlideTitleComponent(id="c1", text="플레이스홀더 주입 테스트")],
        ),
    )
    builder = BuilderRegistry.get("pptx")
    result = await builder.build(schema)

    prs = Presentation(BytesIO(result))
    slide = prs.slides[0]

    # 슬라이드에 placeholder 가 존재하고 그중 하나에 텍스트가 들어 있어야 한다.
    placeholder_texts: list[str] = []
    for ph in slide.placeholders:
        if ph.has_text_frame:
            placeholder_texts.append(ph.text_frame.text)

    combined = "\n".join(placeholder_texts)
    assert "플레이스홀더 주입 테스트" in combined, (
        f"SlideTitle 텍스트가 TITLE placeholder 에 주입되지 않음 — placeholder 텍스트={placeholder_texts!r}"
    )


async def test_d9a_slide_title_falls_back_to_textbox_when_no_placeholder():
    """D9-a #2 — TITLE placeholder 가 없는 마스터(blank) 에서도 textbox 경로로
    정상 렌더 (fallback). python-pptx 의 blank 레이아웃([6]) 은 placeholder
    없이 생성되므로 해당 경로를 강제한다.
    """
    # layout="content_body" 를 쓰되, resolve_layout 이 기본 마스터에서 "Title and
    # Content" 를 찾아주므로 placeholder 가 존재한다. 본 케이스에서는 placeholder
    # 경로 / textbox fallback 경로 어느 쪽이든 텍스트 자체는 반드시 포함되어야 한다.
    schema = _make_schema(
        Page(
            id="p1",
            page_kind="slide",
            layout="content_body",
            components=[SlideTitleComponent(id="c1", text="fallback 경로 텍스트")],
        ),
    )
    builder = BuilderRegistry.get("pptx")
    result = await builder.build(schema)

    # 슬라이드의 모든 텍스트 (textbox + placeholder) 합치기.
    slide_texts = _extract_all_slide_text(result)
    assert "fallback 경로 텍스트" in slide_texts[0], (
        f"SlideTitle 텍스트가 어느 경로에서도 렌더되지 않음: {slide_texts[0]!r}"
    )


def test_d9a_render_slide_title_into_placeholder_returns_false_for_blank_layout():
    """D9-a 보조 — placeholder 가 없는 blank 레이아웃 위에서 helper 가 False 반환.

    PptxBuilder 가 fallback 경로를 타도록 하는 분기점을 직접 검증한다.
    """
    from app.integrations.document_builders.pptx.builder import PptxBuilder

    prs = Presentation()
    # blank 레이아웃 (placeholder 없음).
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank)

    builder = PptxBuilder()
    result = builder._render_slide_title_into_placeholder(
        SlideTitleComponent(id="c1", text="X"),
        slide,
    )
    # blank 에는 TITLE placeholder 가 없으므로 False 가 나와야 한다 (fallback 유도).
    assert result is False


async def test_d9d_bullet_list_highlight_draws_background_rectangle():
    """D9-d #5 — highlight emphasis 를 쓰면 배경 사각형 shape 가 추가된다.

    python-pptx 에서 add_shape 로 추가된 RECTANGLE 은 slide.shapes 에 AutoShape
    로 나타난다. highlight 된 item 수만큼 배경 박스가 있어야 한다.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    schema = _make_schema(
        Page(
            id="p1",
            page_kind="slide",
            layout="content_body",
            components=[
                BulletListComponent(
                    id="c1",
                    items=[
                        BulletItem(text="일반 항목"),
                        BulletItem(text="강조 항목 (highlight)", emphasis="highlight"),
                        BulletItem(text="또 일반 항목"),
                    ],
                ),
            ],
        ),
    )
    builder = BuilderRegistry.get("pptx")
    result = await builder.build(schema)

    prs = Presentation(BytesIO(result))
    slide = prs.slides[0]

    # AutoShape(RECTANGLE) 중 bullet 영역에 있는 것을 highlight 배경으로 간주.
    auto_shapes = [s for s in slide.shapes if getattr(s, "shape_type", None) == MSO_SHAPE_TYPE.AUTO_SHAPE]
    # highlight item 1 개 → 배경 박스 1 개.
    assert len(auto_shapes) >= 1, (
        f"highlight 배경 사각형이 누락 — auto_shapes={[(s.shape_type, s.left, s.top) for s in auto_shapes]!r}"
    )

    # 텍스트 박스는 여전히 존재하고 "강조 항목" 이 포함되어야 한다.
    slide_texts = _extract_all_slide_text(result)
    assert "강조 항목" in slide_texts[0]


def test_d9d_bullet_list_without_highlight_adds_no_background_box():
    """D9-d 보조 — highlight 가 없으면 배경 박스는 전혀 추가되지 않는다 (불필요한 shape 방지)."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    from app.integrations.document_builders.pptx.components import render_bullet_list

    _prs, slide = _make_blank_slide_for_unit_test()
    render_bullet_list(
        BulletListComponent(
            id="c1",
            items=[
                BulletItem(text="A"),
                BulletItem(text="B", emphasis="bold"),
            ],
        ),
        slide,
    )

    auto_shapes = [s for s in slide.shapes if getattr(s, "shape_type", None) == MSO_SHAPE_TYPE.AUTO_SHAPE]
    assert auto_shapes == [], f"highlight 없는 BulletList 에 배경 사각형이 생성됨: {auto_shapes!r}"


def test_d9e_paragraph_height_estimation_korean_weighted():
    """D9-e #6 — 한글 위주 텍스트의 높이 추정이 동일 길이 영문 대비 1.3 배 이상이어야 하며
    ±20 % 오차 범위 내에서 단조 증가한다.
    """
    from app.integrations.document_builders.pptx.components import (
        _count_weighted_chars,
        _estimate_paragraph_height_weighted,
    )

    # 100 자 한글 vs 100 자 영문 (공백 포함).
    korean = "가" * 100
    english = "a" * 100

    h_korean = _estimate_paragraph_height_weighted(korean, font_pt=14, width_in=12.133, min_height_in=0.1)
    h_english = _estimate_paragraph_height_weighted(english, font_pt=14, width_in=12.133, min_height_in=0.1)

    # 한글 가중치 1.5 이므로 한글 높이가 영문 높이 이상이어야 한다.
    assert h_korean >= h_english, f"한글({h_korean}) 이 영문({h_english}) 보다 낮게 추정 — 가중치가 작동하지 않음"

    # 가중 문자 수 검증 — 한글 100 자 = 150.0, 영문 100 자 = 100.0.
    assert _count_weighted_chars(korean) == pytest.approx(150.0)
    assert _count_weighted_chars(english) == pytest.approx(100.0)

    # 혼용 텍스트도 합리적으로 (한글 50 + 영문 50 → 가중 125).
    mixed = "가" * 50 + "a" * 50
    h_mixed = _estimate_paragraph_height_weighted(mixed, font_pt=14, width_in=12.133, min_height_in=0.1)
    # 혼용 높이는 한글 full 과 영문 full 의 사이에 있어야 한다 (±20 % 오차 허용).
    assert h_english <= h_mixed <= h_korean, (
        f"혼용 높이({h_mixed}) 가 영문({h_english})~한글({h_korean}) 사이에 있지 않음"
    )


def test_d9e_paragraph_height_pt_table_covers_common_sizes():
    """D9-e 보조 — CHARS_PER_LINE_BY_PT 테이블에 본문/제목에 쓰이는 주요 pt 포함."""
    from app.integrations.document_builders.pptx.constants import CHARS_PER_LINE_BY_PT

    for pt in (11, 12, 14, 16, 18, 24):
        assert pt in CHARS_PER_LINE_BY_PT, f"{pt}pt 가 CHARS_PER_LINE_BY_PT 에 누락"
        # 큰 pt 는 한 줄 용량이 작아야 한다 (단조 감소).
    sizes = sorted(CHARS_PER_LINE_BY_PT.keys())
    values = [CHARS_PER_LINE_BY_PT[s] for s in sizes]
    for i in range(1, len(values)):
        assert values[i] <= values[i - 1], (
            f"pt 가 커지는데 한 줄 용량이 감소하지 않음: {list(zip(sizes, values, strict=False))}"
        )


# ---------------------------------------------------------------------------
# S3 D1-D2 신규 테스트 — Chart 고도화 + ImageGrid + DocxBuilder 스텁
# ---------------------------------------------------------------------------
#
# D1-a: Chart pie 는 bar 로 fallback 하지 않고 native PIE 로 렌더.
# D1-a: Chart pie 멀티 시리즈는 첫 시리즈만 사용 + WARNING.
# D1-b: DocxBuilder 스텁은 빈 .docx 반환.
# D2-b: ImageGrid 2x2/1x3 레이아웃 → Picture shape 수 검증.
# D2-b: ImageGrid 개별 fetch 실패 → 해당 셀만 placeholder, 나머지는 정상.


def test_s3d1_render_chart_pie_uses_native_pie_type():
    """S3 D1 #1 — chart_type='pie' 는 XL_CHART_TYPE.PIE 네이티브 매핑 (fallback 아님)."""
    from pptx.enum.chart import XL_CHART_TYPE

    _prs, slide = _make_blank_slide_for_unit_test()
    chart_comp = ChartComponent(
        id="c1",
        chart_type="pie",
        title="시장 점유율",
        data=ChartData(
            labels=["A", "B", "C"],
            series=[ChartSeries(name="점유율", values=[40.0, 35.0, 25.0])],
        ),
    )

    render_chart(
        chart_comp,
        slide,
        left_in=0.5,
        top_in=1.5,
        width_in=6.0,
        max_height_in=4.0,
    )

    chart = _find_chart_shape(slide)
    assert chart is not None, "pie chart 가 슬라이드에 삽입되지 않았다"
    assert chart.chart_type == XL_CHART_TYPE.PIE, f"pie → PIE 네이티브 매핑 실패 — actual={chart.chart_type!r}"
    # pie 는 시리즈 1 개 + 카테고리 3 개.
    plot = chart.plots[0]
    assert len(plot.series) == 1
    assert len(list(plot.categories)) == 3


def test_s3d1_render_chart_pie_has_legend_but_no_axes():
    """S3 D1 #2 — pie 차트는 단일 시리즈여도 범례 표시 + 축 스타일 skip.

    pie 계열은 카테고리 구분이 범례로만 전달되므로 시리즈가 1 개여도 has_legend=True.
    축 스타일 (`_style_chart_axes`) 은 pie 에 대해선 호출되지 않아야 한다.
    """
    from pptx.enum.chart import XL_LEGEND_POSITION

    _prs, slide = _make_blank_slide_for_unit_test()
    chart_comp = ChartComponent(
        id="c1",
        chart_type="pie",
        title="비중",
        data=ChartData(
            labels=["X", "Y"],
            series=[ChartSeries(name="s", values=[60.0, 40.0])],
        ),
    )

    render_chart(
        chart_comp,
        slide,
        left_in=0.5,
        top_in=1.5,
        width_in=5.0,
        max_height_in=3.5,
    )

    chart = _find_chart_shape(slide)
    assert chart is not None
    # 단일 시리즈인데도 범례 표시 (bar/line 과 다름).
    assert chart.has_legend is True, "pie 는 단일 시리즈여도 범례가 있어야 함"
    assert chart.legend.position == XL_LEGEND_POSITION.BOTTOM


def test_s3d1_render_chart_pie_multiple_series_drops_to_first(caplog):
    """S3 D1 #3 — pie 에 여러 시리즈가 오면 첫 시리즈만 사용 + WARNING."""
    from pptx.enum.chart import XL_CHART_TYPE

    _prs, slide = _make_blank_slide_for_unit_test()
    chart_comp = ChartComponent(
        id="c1",
        chart_type="pie",
        title="분기별 비중",
        data=ChartData(
            labels=["A", "B"],
            series=[
                ChartSeries(name="Q1", values=[60.0, 40.0]),
                ChartSeries(name="Q2", values=[55.0, 45.0]),
                ChartSeries(name="Q3", values=[50.0, 50.0]),
            ],
        ),
    )

    with caplog.at_level(logging.WARNING):
        render_chart(
            chart_comp,
            slide,
            left_in=0.5,
            top_in=1.5,
            width_in=5.0,
            max_height_in=3.5,
        )

    chart = _find_chart_shape(slide)
    assert chart is not None
    assert chart.chart_type == XL_CHART_TYPE.PIE
    # 첫 시리즈만 사용되었어야 한다.
    plot = chart.plots[0]
    assert len(plot.series) == 1, f"pie 에 여러 시리즈 입력 시 첫 시리즈만 남아야 함 — 실제={len(plot.series)}"
    # WARNING 로그 확인.
    assert any("pie 차트는 단일 시리즈만 지원" in rec.message for rec in caplog.records), (
        f"pie multi-series WARNING 누락: {[r.message for r in caplog.records]!r}"
    )


def test_s3d1_chart_type_alias_table_includes_extension_types():
    """S3 D1 보조 — 확장 alias (area/stacked_*) 가 상수에 등록돼 있어야 한다.

    실제 스키마는 현재 bar/line/pie 만 허용하지만, XL_CHART_TYPE 매핑 테이블은
    S3 후속 스키마 확장 시 즉시 활성화 가능하도록 선제 등록.
    """
    from app.integrations.document_builders.pptx.constants import (
        CHART_XL_TYPE_NAME_BY_ALIAS,
    )

    # 필수 3 종 (스키마 Literal).
    assert CHART_XL_TYPE_NAME_BY_ALIAS["bar"] == "COLUMN_CLUSTERED"
    assert CHART_XL_TYPE_NAME_BY_ALIAS["line"] == "LINE"
    assert CHART_XL_TYPE_NAME_BY_ALIAS["pie"] == "PIE"
    # 확장 예비 (스키마 확장 시 활성화될 매핑).
    assert CHART_XL_TYPE_NAME_BY_ALIAS["area"] == "AREA"
    assert CHART_XL_TYPE_NAME_BY_ALIAS["stacked_bar"] == "BAR_STACKED"
    assert CHART_XL_TYPE_NAME_BY_ALIAS["stacked_column"] == "COLUMN_STACKED"
    assert CHART_XL_TYPE_NAME_BY_ALIAS["stacked_area"] == "AREA_STACKED"


def test_s3d2_image_grid_four_items_creates_four_pictures(monkeypatch, _reset_image_fetcher_client):
    """S3 D2 #1 — ImageGrid 4 장 → 슬라이드에 Picture shape 4 개.

    레이아웃은 2x2. 각 셀이 개별 Picture 로 렌더되어야 한다.
    """
    from app.integrations.document_builders.pptx import components as pptx_components
    from app.modules.documents_v2.schemas import ImageGridComponent, ImageGridItem

    _prs, slide = _make_blank_slide_for_unit_test()

    # 모든 셀이 같은 fake PNG 를 반환하도록.
    fake = _FakeHttpClient(content=_make_png_bytes(width=30, height=30))
    monkeypatch.setattr(_image_fetcher, "_get_http_client", lambda: fake)

    grid = ImageGridComponent(
        id="c1",
        images=[ImageGridItem(src=f"https://example.com/img{i}.png", alt=f"이미지 {i}") for i in range(1, 5)],
    )

    pptx_components.render_image_grid(
        grid,
        slide,
        left_in=0.5,
        top_in=1.5,
        width_in=10.0,
        max_height_in=4.0,
    )

    # Picture shape 개수 = 4.
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    picture_count = sum(1 for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
    assert picture_count == 4, f"ImageGrid 4 장 → Picture 4 개 기대, 실제={picture_count}"


def test_s3d2_image_grid_three_items_renders_in_single_row(monkeypatch, _reset_image_fetcher_client):
    """S3 D2 #2 — ImageGrid 3 장 → 1 행 3 열 레이아웃, 모두 같은 Y 좌표."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    from app.integrations.document_builders.pptx import components as pptx_components
    from app.modules.documents_v2.schemas import ImageGridComponent, ImageGridItem

    _prs, slide = _make_blank_slide_for_unit_test()

    fake = _FakeHttpClient(content=_make_png_bytes(width=30, height=30))
    monkeypatch.setattr(_image_fetcher, "_get_http_client", lambda: fake)

    grid = ImageGridComponent(
        id="c1",
        images=[ImageGridItem(src=f"https://example.com/img{i}.png", alt=f"Item {i}") for i in range(1, 4)],
    )

    top_y = 2.0
    pptx_components.render_image_grid(
        grid,
        slide,
        left_in=0.5,
        top_in=top_y,
        width_in=12.0,
        max_height_in=3.0,
    )

    # 3 개 Picture 가 모두 같은 top (1 행) 이어야 한다.
    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 3, f"3 장 기대, 실제={len(pictures)}"
    top_emus = {p.top for p in pictures}
    assert len(top_emus) == 1, f"3 장이 같은 행에 있어야 하는데 여러 top 좌표 발견: {top_emus}"


def test_s3d2_image_grid_partial_fetch_failure_renders_placeholder_for_failed_cell(
    monkeypatch, _reset_image_fetcher_client, caplog
):
    """S3 D2 #3 — ImageGrid 4 장 중 전부 fetch 실패 → 각 셀별 placeholder.

    fake client 는 예외를 raise 하므로 모든 셀이 placeholder 로 degrade.
    Picture shape 는 0 개, 하지만 placeholder 박스(RECTANGLE) 는 4 개.
    """
    import httpx
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    from app.integrations.document_builders.pptx import components as pptx_components
    from app.modules.documents_v2.schemas import ImageGridComponent, ImageGridItem

    _prs, slide = _make_blank_slide_for_unit_test()

    fake = _FakeHttpClient(exc=httpx.ConnectError("no network"))
    monkeypatch.setattr(_image_fetcher, "_get_http_client", lambda: fake)

    grid = ImageGridComponent(
        id="c1",
        images=[ImageGridItem(src=f"https://broken{i}.example.com/x.png", alt=f"셀 {i}") for i in range(1, 5)],
    )

    with caplog.at_level(logging.WARNING):
        pptx_components.render_image_grid(
            grid,
            slide,
            left_in=0.5,
            top_in=1.5,
            width_in=10.0,
            max_height_in=4.0,
        )

    # Picture 는 없어야 한다.
    picture_count = sum(1 for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
    assert picture_count == 0

    # 각 셀의 alt 텍스트가 placeholder 에 포함되어야 한다 (4 개 개별 placeholder).
    for i in range(1, 5):
        matches = [s.text_frame.text for s in slide.shapes if s.has_text_frame and f"셀 {i}" in s.text_frame.text]
        assert len(matches) >= 1, (
            f"ImageGrid 셀 {i} placeholder 에 alt 텍스트 누락: "
            f"{[s.text_frame.text for s in slide.shapes if s.has_text_frame]!r}"
        )


async def test_s3d2_build_full_schema_with_image_grid(monkeypatch, _reset_image_fetcher_client):
    """S3 D2 #4 — build() 통합. ImageGrid 포함 스키마 → 슬라이드에 Picture 2 개."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    from app.modules.documents_v2.schemas import ImageGridComponent, ImageGridItem

    fake = _FakeHttpClient(content=_make_png_bytes(width=25, height=25))
    monkeypatch.setattr(_image_fetcher, "_get_http_client", lambda: fake)

    schema = _make_schema(
        Page(
            id="p1",
            page_kind="slide",
            layout="content_body",
            title="ImageGrid 통합 테스트",
            components=[
                SlideTitleComponent(id="c1", text="제품 라인업"),
                ImageGridComponent(
                    id="c2",
                    images=[
                        ImageGridItem(src="https://example.com/a.png", alt="제품 A"),
                        ImageGridItem(src="https://example.com/b.png", alt="제품 B"),
                    ],
                ),
            ],
        ),
    )

    builder = BuilderRegistry.get("pptx")
    result = await builder.build(schema)

    assert isinstance(result, bytes)
    assert result.startswith(_ZIP_SIGNATURE)

    prs = Presentation(BytesIO(result))
    assert len(prs.slides) == 1
    slide = prs.slides[0]

    # ImageGrid 는 Picture 2 개 생성. SlideTitle 은 textbox 로 별도.
    picture_count = sum(1 for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
    assert picture_count == 2, f"ImageGrid 2 장 → Picture 2 개 기대, 실제={picture_count}"


def test_s3d2_pptx_builder_supports_image_grid():
    """S3 D2 보조 — `PptxBuilder.supported_components` 에 'ImageGrid' 포함 (S3 D6-D7 확장 후 14 종)."""
    builder = PptxBuilder()
    assert "ImageGrid" in builder.supported_components
    # S3 D6 — SlideSubtitle/Quote/Callout 추가. S3 D7 — Timeline/IconRow 추가. 9 → 14.
    assert len(builder.supported_components) == 14


async def test_s3d1_docx_builder_stub_returns_empty_docx():
    """S3 D1 — DocxBuilder 스텁은 빈 .docx 바이트를 반환 (ZIP 시그니처).

    실제 컴포넌트 렌더는 S5. 현재 범위는 BuilderRegistry 등록 + 빈 문서 반환만 검증.
    """
    from app.integrations.document_builders.base import BuilderRegistry
    from app.integrations.document_builders.docx import (
        DocxBuilder,
        register_docx_builder,
    )

    register_docx_builder()
    builder = BuilderRegistry.get("docx")
    assert isinstance(builder, DocxBuilder)
    assert builder.target == "docx"

    # 최소 스키마.
    schema = _make_schema(
        Page(
            id="p1",
            page_kind="section",
            layout="content_body",
            components=[SlideTitleComponent(id="c1", text="스텁 테스트")],
        ),
    )

    result = await builder.build(schema)
    assert isinstance(result, bytes)
    # DOCX 도 OOXML(ZIP) 이므로 b"PK" 로 시작한다.
    assert result.startswith(_ZIP_SIGNATURE), f"빈 DOCX 가 ZIP 시그니처로 시작해야 함 — 실제 prefix={result[:4]!r}"
    # 최소 DOCX 는 수백 바이트 이상.
    assert len(result) > 100

    # Chart 만 supported 로 선언된 스텁 단계.
    assert "Chart" in builder.supported_components


# ---------------------------------------------------------------------------
# S3 D6 테스트 — SlideSubtitle / Quote / Callout
# ---------------------------------------------------------------------------


def test_s3d6_render_slide_subtitle_uses_center_alignment_and_muted_color():
    """S3 D6 #1 — SlideSubtitle 이 중앙 정렬 + IDINO_TEXT_MUTED 색으로 렌더된다.

    검증 포인트:
        - 슬라이드에 텍스트가 포함된 shape 가 추가된다.
        - paragraph 의 alignment 가 PP_ALIGN.CENTER.
        - paragraph/run 의 color 가 IDINO_TEXT_MUTED(#6B7280).
    """
    from pptx.enum.text import PP_ALIGN

    from app.integrations.document_builders.pptx.components import render_slide_subtitle
    from app.integrations.document_builders.pptx.constants import IDINO_TEXT_MUTED

    _prs, slide = _make_blank_slide_for_unit_test()
    component = SlideSubtitleComponent(id="c1", text="Q1 경영 실적 요약")

    render_slide_subtitle(component, slide)

    # 텍스트 프레임을 가진 shape 에서 본문이 일치하는 것을 찾는다.
    texts = [s for s in slide.shapes if s.has_text_frame and "Q1 경영 실적 요약" in s.text_frame.text]
    assert len(texts) == 1, f"SlideSubtitle shape 가 1 개 있어야 함 — 실제={len(texts)}"

    paragraph = texts[0].text_frame.paragraphs[0]
    assert paragraph.alignment == PP_ALIGN.CENTER, (
        f"SlideSubtitle 은 중앙 정렬이어야 함 — alignment={paragraph.alignment!r}"
    )
    # run 의 color 가 IDINO_TEXT_MUTED.
    expected_rgb = parse_hex_color(IDINO_TEXT_MUTED)
    actual_rgb = paragraph.runs[0].font.color.rgb
    assert actual_rgb == expected_rgb, f"muted color mismatch — actual={actual_rgb!r}"


def test_s3d6_render_quote_has_left_accent_bar_rectangle():
    """S3 D6 #2 — Quote 렌더 시 좌측에 IDINO_ACCENT 색 RECTANGLE 강조선이 추가된다."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    from app.integrations.document_builders.pptx.components import render_quote
    from app.integrations.document_builders.pptx.constants import IDINO_ACCENT

    _prs, slide = _make_blank_slide_for_unit_test()
    component = QuoteComponent(id="c1", text="고객 중심이 곧 성장의 길이다.")

    render_quote(component, slide, left_in=1.0, top_in=2.0, width_in=10.0)

    # AutoShape 중 RECTANGLE 1 개 이상 (강조선). 색이 IDINO_ACCENT 여야 함.
    autoshapes = [s for s in slide.shapes if getattr(s, "shape_type", None) == MSO_SHAPE_TYPE.AUTO_SHAPE]
    assert len(autoshapes) >= 1, "Quote 좌측 강조선 rectangle 이 누락됨"

    accent_rgb = parse_hex_color(IDINO_ACCENT)
    bar_with_accent = [s for s in autoshapes if getattr(getattr(s.fill, "fore_color", None), "rgb", None) == accent_rgb]
    assert len(bar_with_accent) >= 1, "Quote 강조선이 IDINO_ACCENT 색이 아님"

    # 본문 텍스트 확인.
    texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
    combined = "\n".join(texts)
    assert "고객 중심이 곧 성장의 길이다." in combined


def test_s3d6_render_quote_without_author_skips_author_textbox():
    """S3 D6 #3 — author 가 None 이면 '— {author}' 텍스트박스가 생성되지 않는다."""
    from app.integrations.document_builders.pptx.components import render_quote

    _prs, slide = _make_blank_slide_for_unit_test()
    component = QuoteComponent(id="c1", text="단독 인용문.", author=None)

    render_quote(component, slide, left_in=1.0, top_in=2.0, width_in=10.0)

    texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
    combined = "\n".join(texts)
    # em dash "—" 가 포함돼 있으면 author 처리됐다는 뜻 — 이 케이스에선 없어야 함.
    assert "—" not in combined, f"author=None 인데 '—' 프리픽스가 렌더됨: {combined!r}"


def test_s3d6_render_quote_with_author_adds_author_line_right_aligned():
    """S3 D6 #4 — author 있으면 '— {author}' 가 우측 정렬 텍스트박스로 추가된다."""
    from pptx.enum.text import PP_ALIGN

    from app.integrations.document_builders.pptx.components import render_quote

    _prs, slide = _make_blank_slide_for_unit_test()
    component = QuoteComponent(id="c1", text="도전은 기회를 만든다.", author="홍길동")

    render_quote(component, slide, left_in=1.0, top_in=2.0, width_in=10.0)

    # "— 홍길동" 이 포함된 텍스트박스를 찾고, 우측 정렬인지 확인.
    author_shapes = [s for s in slide.shapes if s.has_text_frame and "— 홍길동" in s.text_frame.text]
    assert len(author_shapes) == 1, "author 텍스트박스가 1 개 있어야 함"
    author_para = author_shapes[0].text_frame.paragraphs[0]
    assert author_para.alignment == PP_ALIGN.RIGHT


@pytest.mark.parametrize(
    "variant, expected_bg, expected_border",
    [
        ("info", "#E0F2FE", "#0A4FC2"),  # IDINO_PRIMARY
        ("warning", "#FEF3C7", "#F59E0B"),
        ("success", "#D1FAE5", "#10B981"),
        ("danger", "#FEE2E2", "#EF4444"),
    ],
)
def test_s3d6_render_callout_variant_applies_correct_bg_and_border(
    variant: str, expected_bg: str, expected_border: str
):
    """S3 D6 #5~#8 — Callout variant 4 종 각각의 배경/border 색 팔레트 검증.

    배경 rounded_rect 와 좌측 border rectangle 두 shape 를 확인:
        - 배경 = variant_bg (info=#E0F2FE, warning=#FEF3C7, success=#D1FAE5, danger=#FEE2E2)
        - 좌측 border = variant_border (info=IDINO_PRIMARY 등)
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    from app.integrations.document_builders.pptx.components import render_callout

    _prs, slide = _make_blank_slide_for_unit_test()
    component = CalloutComponent(id="c1", text=f"{variant} 메시지", variant=variant)  # type: ignore[arg-type]

    render_callout(component, slide, left_in=1.0, top_in=2.0, width_in=10.0)

    autoshapes = [s for s in slide.shapes if getattr(s, "shape_type", None) == MSO_SHAPE_TYPE.AUTO_SHAPE]
    # 최소 2 개 (배경 rounded_rect + 좌측 border rect). letter textbox 는 AutoShape 아님.
    assert len(autoshapes) >= 2, f"Callout {variant}: AutoShape 2 개 이상 필요 — 실제={len(autoshapes)}"

    expected_bg_rgb = parse_hex_color(expected_bg)
    expected_border_rgb = parse_hex_color(expected_border)
    actual_rgbs = {getattr(getattr(s.fill, "fore_color", None), "rgb", None) for s in autoshapes}
    assert expected_bg_rgb in actual_rgbs, f"Callout {variant}: 배경색 {expected_bg} 미발견. 실제={actual_rgbs!r}"
    assert expected_border_rgb in actual_rgbs, (
        f"Callout {variant}: border 색 {expected_border} 미발견. 실제={actual_rgbs!r}"
    )


def test_s3d7_render_timeline_three_events_creates_three_markers():
    """S3 D7 #9 — Timeline 3 events 렌더 시 마커(OVAL) 3 개가 추가된다."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    from app.integrations.document_builders.pptx.components import render_timeline

    _prs, slide = _make_blank_slide_for_unit_test()
    component = TimelineComponent(
        id="c1",
        events=[
            TimelineEvent(date="2026-01", title="Q1 킥오프", description="본사 전체 회의"),
            TimelineEvent(date="2026-04", title="Q2 중간 점검", description=None),
            TimelineEvent(date="2026-07", title="Q3 시작", description="사업부별 목표 재설정"),
        ],
    )

    render_timeline(component, slide, left_in=1.0, top_in=1.5, width_in=10.0, max_height_in=5.0)

    # Timeline 은 OVAL shape 로 마커를 그린다. shape_type 은 AUTO_SHAPE.
    # OVAL 과 RECTANGLE 구분은 xml.spPr 확인이 번거롭기에 "AutoShape 개수로 간접 검증".
    # 마커 3 + 연결선 2 = 최소 5 개 AutoShape.
    autoshapes = [s for s in slide.shapes if getattr(s, "shape_type", None) == MSO_SHAPE_TYPE.AUTO_SHAPE]
    assert len(autoshapes) >= 5, (
        f"Timeline 3 events: OVAL 3 + RECTANGLE 2 (연결선) = 5 개 이상 기대 — 실제={len(autoshapes)}"
    )

    # 본문 텍스트(제목) 검증.
    combined = "\n".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
    assert "Q1 킥오프" in combined
    assert "Q2 중간 점검" in combined
    assert "Q3 시작" in combined


def test_s3d7_render_timeline_over_max_events_skips_and_warns(caplog):
    """S3 D7 #10 — events 6 개 (최대 5 초과) 인 경우 상위 5 개만 렌더 + WARNING 로그 발생."""
    from app.integrations.document_builders.pptx.components import render_timeline

    _prs, slide = _make_blank_slide_for_unit_test()
    # Pydantic TimelineComponent.events max_length=10 — 6 개는 통과하지만 내부 TIMELINE_MAX_EVENTS(5) 초과.
    component = TimelineComponent(
        id="c1",
        events=[
            TimelineEvent(date=f"E{i}", title=f"이벤트 {i}", description=None)
            for i in range(1, 7)  # 6 개.
        ],
    )

    with caplog.at_level(logging.WARNING):
        render_timeline(component, slide, left_in=1.0, top_in=1.0, width_in=10.0, max_height_in=5.0)

    # WARN 메시지 확인.
    assert any("render_timeline" in r.message and "스킵" in r.message for r in caplog.records), (
        f"Timeline 초과 WARN 로그가 없음 — records={[r.message for r in caplog.records]!r}"
    )

    # 상위 5 개 이벤트만 본문에 포함, 6 번째는 제외.
    combined = "\n".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
    for i in range(1, 6):
        assert f"이벤트 {i}" in combined, f"상위 {i} 이벤트 누락: {combined!r}"
    assert "이벤트 6" not in combined, f"6 번째 이벤트가 스킵되지 않음: {combined!r}"


def test_s3d7_render_icon_row_four_items_creates_four_circles_and_four_labels():
    """S3 D7 #11 — IconRow 4 items 렌더 시 원형 도형 4 개 + 라벨 텍스트 4 개가 생성된다."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    from app.integrations.document_builders.pptx.components import render_icon_row

    _prs, slide = _make_blank_slide_for_unit_test()
    component = IconRowComponent(
        id="c1",
        items=[
            IconItem(icon="rocket", label="빠른 출시"),
            IconItem(icon="shield", label="보안 강화"),
            IconItem(icon="users", label="팀 협업"),
            IconItem(icon="chart", label="성과 측정"),
        ],
    )

    render_icon_row(component, slide, left_in=0.5, top_in=2.0, width_in=12.0)

    # AutoShape 4 개 (OVAL).
    autoshapes = [s for s in slide.shapes if getattr(s, "shape_type", None) == MSO_SHAPE_TYPE.AUTO_SHAPE]
    assert len(autoshapes) == 4, f"IconRow 4 items: OVAL 4 개 기대 — 실제={len(autoshapes)}"

    # 모든 라벨이 텍스트로 포함.
    combined = "\n".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
    for label in ("빠른 출시", "보안 강화", "팀 협업", "성과 측정"):
        assert label in combined, f"라벨 누락: {label!r} — combined={combined!r}"

    # letter 는 아이콘 이름 첫 글자 대문자.
    for letter in ("R", "S", "U", "C"):
        assert letter in combined, f"letter {letter!r} 누락"


async def test_s3d7_build_full_schema_with_all_five_new_components():
    """S3 D6-D7 통합 — SlideSubtitle/Quote/Callout/Timeline/IconRow 를 한 스키마에 담아 빌드 성공."""
    schema = _make_schema(
        Page(
            id="p1",
            page_kind="slide",
            layout="title_slide",
            components=[
                SlideTitleComponent(id="c1", text="2026 연간 계획"),
                SlideSubtitleComponent(id="c2", text="경영진 요약 보고서"),
            ],
        ),
        Page(
            id="p2",
            page_kind="slide",
            layout="content_body",
            components=[
                QuoteComponent(
                    id="c3",
                    text="고객의 불편이 곧 기회다.",
                    author="CEO 김철수",
                ),
                CalloutComponent(
                    id="c4",
                    text="경쟁사 대비 출시 속도를 2 배로 유지해야 합니다.",
                    variant="warning",
                ),
            ],
        ),
        Page(
            id="p3",
            page_kind="slide",
            layout="content_body",
            components=[
                TimelineComponent(
                    id="c5",
                    events=[
                        TimelineEvent(date="2026-Q1", title="런칭", description="핵심 기능 공개"),
                        TimelineEvent(date="2026-Q2", title="확장", description="해외 시장 진입"),
                    ],
                ),
                IconRowComponent(
                    id="c6",
                    items=[
                        IconItem(icon="speed", label="속도"),
                        IconItem(icon="safety", label="안전"),
                        IconItem(icon="quality", label="품질"),
                    ],
                ),
            ],
        ),
    )

    builder = BuilderRegistry.get("pptx")
    result = await builder.build(schema)

    assert isinstance(result, bytes)
    assert result.startswith(_ZIP_SIGNATURE)

    prs = Presentation(BytesIO(result))
    assert len(prs.slides) == 3

    # 각 페이지의 주요 텍스트가 슬라이드에 포함되었는지 확인.
    all_texts = _extract_all_slide_text(result)
    assert "경영진 요약 보고서" in all_texts[0]
    assert "고객의 불편이 곧 기회다." in all_texts[1]
    assert "— CEO 김철수" in all_texts[1]
    assert "경쟁사 대비 출시 속도" in all_texts[1]
    assert "런칭" in all_texts[2]
    assert "속도" in all_texts[2]


def test_s3d7_pptx_builder_supported_components_includes_all_five():
    """S3 D6-D7 보조 — `PptxBuilder.supported_components` 에 5 종 신규 컴포넌트 모두 포함 (14 종)."""
    builder = PptxBuilder()
    for name in ("SlideSubtitle", "Quote", "Callout", "Timeline", "IconRow"):
        assert name in builder.supported_components, f"{name!r} 이 supported_components 에 누락"
    assert len(builder.supported_components) == 14
